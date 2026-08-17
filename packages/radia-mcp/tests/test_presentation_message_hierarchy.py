from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from radia_mcp.presentation.tools import (
    presentation_check_slide_message_hierarchy,
    presentation_check_slide_title_specificity,
)


def _add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "磁気モーメント法"
    slide.placeholders[1].text = "研究会発表"


def _add_content_slide(
    prs: Presentation,
    title: str,
    takeaway: str | None,
    footer: str = "近畿大学 2",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(2.2))
    body.text_frame.text = "中央の比較結果"
    if takeaway is not None:
        box = slide.shapes.add_textbox(
            Inches(0.8), Inches(5.35), Inches(11.5), Inches(0.45)
        )
        box.text_frame.text = takeaway
    foot = slide.shapes.add_textbox(
        Inches(10.8), Inches(7.15), Inches(2.0), Inches(0.25)
    )
    foot.text_frame.text = footer


def test_message_hierarchy_accepts_claim_title_and_bottom_takeaway(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_title_slide(prs)
    _add_content_slide(
        prs,
        "ひずみ要素に対するMMPMの精度評価",
        "分かったこと：局所閉包がひずみ感度を支配する",
    )
    pptx_path = tmp_path / "good.pptx"
    prs.save(pptx_path)

    result = presentation_check_slide_message_hierarchy(str(pptx_path))

    assert result["score"] == 10.0
    assert result["slides_checked"] == 1
    assert result["slides_passing"] == 1
    assert result["slides"][0]["skipped"] is True
    assert result["slides"][1]["title_is_message"] is True
    assert result["slides"][1]["has_bottom_takeaway"] is True


def test_message_hierarchy_flags_generic_title_and_missing_takeaway(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_title_slide(prs)
    _add_content_slide(prs, "結果", None)
    pptx_path = tmp_path / "bad.pptx"
    prs.save(pptx_path)

    result = presentation_check_slide_message_hierarchy(str(pptx_path))

    slide = result["slides"][1]
    assert result["score"] == 0.0
    assert slide["issues"] == ["title_not_specific", "bottom_takeaway_missing"]


def test_concrete_target_viewpoint_title_is_valid(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_title_slide(prs)
    _add_content_slide(
        prs,
        "バブル変位によるひずみ感度評価",
        "分かったこと：バブル変位により形状変化を分離できる",
    )
    pptx_path = tmp_path / "method-intent.pptx"
    prs.save(pptx_path)

    result = presentation_check_slide_message_hierarchy(str(pptx_path))

    slide = result["slides"][1]
    assert slide["title_is_message"] is True
    assert slide["has_bottom_takeaway"] is True
    assert result["score"] == 10.0


def test_footer_and_citation_do_not_count_as_takeaway(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_title_slide(prs)
    _add_content_slide(
        prs,
        "BDM1の2次曲面対応",
        "https://example.com/source",
        footer="12",
    )
    pptx_path = tmp_path / "footer-only.pptx"
    prs.save(pptx_path)

    result = presentation_check_slide_message_hierarchy(str(pptx_path))

    slide = result["slides"][1]
    assert slide["title_is_message"] is True
    assert slide["has_bottom_takeaway"] is False
    assert slide["issues"] == ["bottom_takeaway_missing"]


def test_structural_slide_is_skipped(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_title_slide(prs)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "参考文献"
    pptx_path = tmp_path / "structural.pptx"
    prs.save(pptx_path)

    result = presentation_check_slide_message_hierarchy(str(pptx_path))

    assert result["slides_checked"] == 0
    assert result["slides"][1]["skipped"] is True


def test_cued_last_line_in_bottom_spanning_body_counts(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_title_slide(prs)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "コロケーション型MMPMのHACApK高速化"
    body = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.2)
    )
    body.text_frame.text = (
        "中央の比較結果\n"
        "分かったこと：コロケーション法を維持できる"
    )
    pptx_path = tmp_path / "cued-body.pptx"
    prs.save(pptx_path)

    result = presentation_check_slide_message_hierarchy(str(pptx_path))

    assert result["slides"][1]["has_bottom_takeaway"] is True
    assert result["slides"][1]["bottom_takeaway"].startswith("分かったこと：")


def test_title_specificity_distinguishes_generic_specific_and_result(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_title_slide(prs)
    _add_content_slide(prs, "数値計算結果", "分かったこと：差は0.28%以内")
    _add_content_slide(prs, "モデル1の計算精度評価", "分かったこと：差は0.28%以内")
    _add_content_slide(
        prs,
        "モデル1はFEMと0.28%以内で一致した",
        "分かったこと：差は0.28%以内",
    )
    pptx_path = tmp_path / "title-specificity.pptx"
    prs.save(pptx_path)

    result = presentation_check_slide_title_specificity(str(pptx_path))

    assert result["slides"][1]["issues"] == [
        "title_not_specific",
        "viewpoint_not_explicit",
    ]
    assert result["slides"][2]["passed"] is True
    assert result["slides"][2]["score"] == 10
    assert all(result["slides"][2]["criteria"].values())
    assert "title_is_result_sentence" in result["slides"][3]["issues"]


def test_title_specificity_flags_duplicate_titles(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_title_slide(prs)
    _add_content_slide(prs, "モデル1の計算精度評価", "分かったこと：差は0.28%以内")
    _add_content_slide(prs, "モデル1の計算精度評価", "分かったこと：差は0.10%以内")
    pptx_path = tmp_path / "duplicate-title.pptx"
    prs.save(pptx_path)

    result = presentation_check_slide_title_specificity(str(pptx_path))

    assert result["passed"] is False
    assert result["slides"][1]["criteria"]["unique_in_deck"] is False
    assert result["slides"][1]["score"] == 8
    assert result["slides"][1]["issues"] == ["duplicate_title"]


def test_title_specificity_separates_model_size_from_numeric_result(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_title_slide(prs)
    _add_content_slide(
        prs,
        "16.2万自由度モデルの計算性能評価",
        "分かったこと：求解時間は60.7秒",
    )
    _add_content_slide(
        prs,
        "16.2万自由度を26.1 GiB・60.7秒で解析",
        "分かったこと：求解時間は60.7秒",
    )
    pptx_path = tmp_path / "numeric-title.pptx"
    prs.save(pptx_path)

    result = presentation_check_slide_title_specificity(str(pptx_path))

    assert result["slides"][1]["passed"] is True
    assert "title_is_result_sentence" in result["slides"][2]["issues"]


def test_title_specificity_accepts_loop_free_as_viewpoint(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_title_slide(prs)
    _add_content_slide(
        prs,
        "HDiv-MMMによる完全loop-free化",
        "分かったこと：大域的なloop成分を除去できる",
    )
    pptx_path = tmp_path / "loop-free-title.pptx"
    prs.save(pptx_path)

    result = presentation_check_slide_title_specificity(str(pptx_path))

    assert result["slides"][1]["passed"] is True
    assert result["slides"][1]["criteria"]["viewpoint_is_explicit"] is True
