"""Two conventions for a slide title, and neither is wrong.

- ``topic``   the title names the target and the viewpoint; the finding goes
              in the bottom line.
- ``message`` the title IS the finding ("MMPM reduces the deviation to 1/145").

The check used to enforce ``topic`` on every deck, and scored a title that
states its result as a defect. Measured on a 16-slide deck written in the
message style (2026-08-21): 8 of 16 titles failed, and every one of them was a
correct title. A check that fails half a correct deck is not a check.

So the convention is read off the deck and reported, and can be pinned. What
must NOT change is that a bare label -- 結果, まとめ, 提案手法の概要 -- fails
under either convention.
"""

from __future__ import annotations

import pytest

pytest.importorskip("pptx")

from pptx import Presentation  # noqa: E402
from pptx.util import Cm, Pt  # noqa: E402

from radia_mcp.presentation.tools import (  # noqa: E402
    presentation_check_slide_message_hierarchy as check,
)

MESSAGE_TITLES = [
    "MMPMは要素ひずみによる偏差を最大1/145に低減",
    "HACApKにより16.2万自由度を60.7秒で求解",
    "従来法は補償点配置により解が変化",
    "低次モーメント保存で6面磁荷を局所閉包",
]
TOPIC_TITLES = [
    "ひずみ要素での閉包条件の評価",
    "H行列圧縮率の条件依存性",
    "粗密メッシュ間の整合の比較",
    "非線形材料モデルの検証条件",
]


def _deck(tmp_path, titles, name="deck.pptx", bottom="下部に置いた知見の一文"):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Cm(33.87), Cm(19.05)
    prs.slides.add_slide(prs.slide_layouts[6])          # cover: always skipped
    for title in titles:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        box = s.shapes.add_textbox(Cm(0), Cm(0), Cm(33.87), Cm(2.19))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = title
        run.font.size = Pt(34)
        if bottom:
            bar = s.shapes.add_textbox(Cm(1), Cm(16.7), Cm(31), Cm(1.3))
            run = bar.text_frame.paragraphs[0].add_run()
            run.text = bottom
            run.font.size = Pt(24)
    out = tmp_path / name
    prs.save(str(out))
    return str(out)


def test_a_message_written_deck_is_read_as_message(tmp_path):
    r = check(_deck(tmp_path, MESSAGE_TITLES))
    assert r["title_style_detected"] == "message"
    assert r["title_style_used"] == "message"
    assert r["title_message_failures"] == 0, [s["issues"] for s in r["slides"]]


def test_the_same_deck_under_the_other_convention_fails_as_it_used_to(tmp_path):
    """The old behaviour is still reachable, and still means what it meant."""
    deck = _deck(tmp_path, MESSAGE_TITLES)
    r = check(deck, title_style="topic")
    assert r["title_message_failures"] > 0
    flagged = [i for s in r["slides"] if not s.get("skipped") for i in s["issues"]]
    assert "title_is_result_sentence" in flagged


def test_a_topic_written_deck_is_read_as_topic_and_passes(tmp_path):
    r = check(_deck(tmp_path, TOPIC_TITLES))
    assert r["title_style_detected"] == "topic"
    assert r["title_message_failures"] == 0, [s["issues"] for s in r["slides"]]


def test_a_bare_label_fails_under_either_convention(tmp_path):
    labels = ["結果", "まとめ", "実験結果", "考察"]
    for style in ("message", "topic"):
        r = check(_deck(tmp_path, labels, name="l_%s.pptx" % style),
                  title_style=style)
        assert r["title_message_failures"] == len(labels), (style, r["slides"])


def test_the_conventions_disagree_exactly_where_they_should(tmp_path):
    """「提案手法の概要」 is a target and a viewpoint -- and states nothing.

    That is the whole difference between the two conventions, so it is the one
    case that must come out differently: fine as a topic title, not a message.
    """
    deck = _deck(tmp_path, ["提案手法の概要"] * 1 + TOPIC_TITLES[:3],
                 name="overview.pptx")
    assert check(deck, title_style="topic")["title_message_failures"] == 0
    msg = check(deck, title_style="message")
    failed = [s["title"] for s in msg["slides"]
              if not s.get("skipped") and "title_states_no_message" in s["issues"]]
    assert "提案手法の概要" in failed


def test_a_noun_phrase_is_not_a_message_just_because_it_holds_a_particle(tmp_path):
    """「非線形計算に用いた実測B-H曲線」 is what the figure IS, not a finding.

    Its に belongs to a clause modifying the noun, so a naive "has a particle"
    rule would call it a message. What separates it from 「…6面磁荷を局所閉包」
    is that the action noun there ENDS the title, right after the particle.
    """
    r = check(_deck(tmp_path, ["非線形計算に用いた実測B–H曲線"] + MESSAGE_TITLES),
              title_style="message")
    by_title = {s["title"]: s for s in r["slides"] if not s.get("skipped")}
    assert by_title["非線形計算に用いた実測B–H曲線"]["title_asserts_message"] is False
    assert by_title["低次モーメント保存で6面磁荷を局所閉包"]["title_asserts_message"] is True


def test_a_short_deck_keeps_the_established_default(tmp_path):
    """Two slides are not evidence of a convention; do not guess from them."""
    r = check(_deck(tmp_path, MESSAGE_TITLES[:2]))
    assert r["title_style_detected"] == "topic"


def test_the_bottom_line_check_is_unaffected_by_the_convention(tmp_path):
    r = check(_deck(tmp_path, MESSAGE_TITLES, bottom=""))
    assert r["bottom_takeaway_failures"] == len(MESSAGE_TITLES)


def test_an_unknown_style_is_refused(tmp_path):
    r = check(_deck(tmp_path, MESSAGE_TITLES), title_style="assertive")
    assert "error" in r
