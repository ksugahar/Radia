from __future__ import annotations

import json
import hashlib
import xml.etree.ElementTree as ET
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from radia_mcp.presentation.tools import (
    _picture_source,
    _svg_dimensions,
    presentation_check_embedded_figure_text_size,
)


def _deck_with_picture(tmp_path: Path) -> tuple[Path, str]:
    image_path = tmp_path / "figure.png"
    Image.new("RGB", (1000, 500), "white").save(image_path)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    picture = slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(1), width=Inches(10)
    )
    path = tmp_path / "deck.pptx"
    prs.save(path)
    return path, picture.name


def _write_manifest(
        tmp_path: Path, shape_name: str, bbox_height_px: int) -> Path:
    path = tmp_path / "ocr.json"
    payload = {
        "pictures": [
            {
                "slide": 1,
                "shape": shape_name,
                "words": [
                    {
                        "text": "HACApK",
                        "confidence": 0.99,
                        "bbox": [100, 100, 300, 100 + bbox_height_px],
                    }
                ],
            }
        ]
    }
    path.write_text(json.dumps(payload), encoding="utf-8")
    return path


def test_embedded_figure_text_at_24pt_passes(tmp_path: Path) -> None:
    deck, shape_name = _deck_with_picture(tmp_path)
    manifest = _write_manifest(tmp_path, shape_name, bbox_height_px=24)

    result = presentation_check_embedded_figure_text_size(
        str(deck), ocr_backend="manifest", ocr_manifest_path=str(manifest)
    )

    assert result["passed"] is True
    assert result["violation_count"] == 0
    assert result["pictures"][0]["minimum_estimated_font_pt"] == 24.0


def test_embedded_figure_text_below_20pt_is_a_violation(
        tmp_path: Path) -> None:
    deck, shape_name = _deck_with_picture(tmp_path)
    manifest = _write_manifest(tmp_path, shape_name, bbox_height_px=14)

    result = presentation_check_embedded_figure_text_size(
        str(deck), ocr_backend="manifest", ocr_manifest_path=str(manifest)
    )

    assert result["passed"] is False
    assert result["violation_count"] == 1
    assert result["pictures"][0]["minimum_estimated_font_pt"] == 14.0


def test_source_size_evidence_uses_actual_paste_width(tmp_path: Path) -> None:
    deck, shape_name = _deck_with_picture(tmp_path)
    manifest = tmp_path / "source-evidence.json"
    manifest.write_text(json.dumps({
        "pictures": [{
            "slide": 1,
            "shape": shape_name,
            "source_evidence": {
                "minimum_source_font_pt": 24,
                "source_width_in": 12,
            },
        }]
    }), encoding="utf-8")

    result = presentation_check_embedded_figure_text_size(
        str(deck), ocr_backend="manifest", ocr_manifest_path=str(manifest)
    )

    assert result["passed"] is True
    assert result["pictures"][0]["evidence_type"] == "source_size"
    assert result["pictures"][0]["embed_scale"] == 0.8333
    assert result["pictures"][0]["minimum_estimated_font_pt"] == 20.0


def test_source_size_evidence_below_floor_fails(tmp_path: Path) -> None:
    deck, shape_name = _deck_with_picture(tmp_path)
    manifest = tmp_path / "source-evidence-small.json"
    manifest.write_text(json.dumps({
        "pictures": [{
            "slide": 1,
            "shape": shape_name,
            "source_evidence": {
                "minimum_source_font_pt": 18,
                "source_width_in": 10,
            },
        }]
    }), encoding="utf-8")

    result = presentation_check_embedded_figure_text_size(
        str(deck), ocr_backend="manifest", ocr_manifest_path=str(manifest)
    )

    assert result["passed"] is False
    assert result["pictures"][0]["minimum_estimated_font_pt"] == 18.0


def test_unverified_picture_does_not_pass(tmp_path: Path) -> None:
    deck, _ = _deck_with_picture(tmp_path)

    result = presentation_check_embedded_figure_text_size(str(deck))

    assert result["passed"] is False
    assert result["unresolved_count"] == 1


def test_visually_confirmed_textless_picture_is_exempt(tmp_path: Path) -> None:
    deck, shape_name = _deck_with_picture(tmp_path)

    result = presentation_check_embedded_figure_text_size(
        str(deck), confirmed_textless_shapes=[f"1:{shape_name}"]
    )

    assert result["passed"] is True
    assert result["pictures"][0]["status"] == "confirmed_textless"


def test_source_evidence_can_follow_picture_by_asset_sha1(
        tmp_path: Path) -> None:
    deck, _ = _deck_with_picture(tmp_path)
    prs = Presentation(str(deck))
    picture = next(
        shape for shape in prs.slides[0].shapes
        if getattr(shape, "image", None) is not None
    )
    asset_sha1 = hashlib.sha1(picture.image.blob).hexdigest()
    manifest = tmp_path / "source-evidence-by-hash.json"
    manifest.write_text(json.dumps({
        "pictures": [{
            "slide": 1,
            "shape": "stale PowerPoint shape name",
            "asset_sha1": asset_sha1,
            "source_evidence": {
                "minimum_source_font_pt": 24,
                "source_width_in": 10,
            },
        }]
    }), encoding="utf-8")

    result = presentation_check_embedded_figure_text_size(
        str(deck), ocr_backend="manifest", ocr_manifest_path=str(manifest)
    )

    assert result["passed"] is True
    assert result["unresolved_count"] == 0
    assert result["pictures"][0]["evidence_type"] == "source_size"
    assert result["pictures"][0]["asset_sha1"] == asset_sha1[:10]


def test_svg_dimensions_prefer_viewbox() -> None:
    blob = (
        b'<svg xmlns="http://www.w3.org/2000/svg" width="1" height="1" '
        b'viewBox="0 0 640 360"/>'
    )

    assert _svg_dimensions(blob) == (640.0, 360.0)


def test_picture_source_follows_powerpoint_svg_relationship() -> None:
    svg_blob = (
        b'<svg xmlns="http://www.w3.org/2000/svg" '
        b'viewBox="0 0 800 400"/>'
    )
    element = ET.fromstring(
        '<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main">'
        '<asvg:svgBlip r:embed="rId7"/></p:pic>'
    )

    class _SvgPart:
        blob = svg_blob

    class _SlidePart:
        @staticmethod
        def related_part(relationship_id: str):
            assert relationship_id == "rId7"
            return _SvgPart()

    class _Shape:
        _element = element
        part = _SlidePart()

        @property
        def image(self):
            raise AssertionError("PNG fallback must not be used")

    blob, size, source_type = _picture_source(_Shape())

    assert blob == svg_blob
    assert size == (800.0, 400.0)
    assert source_type == "svg"
