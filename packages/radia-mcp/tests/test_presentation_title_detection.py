"""Which shape is the slide's title, when nothing declares it.

A deck drawn on the Blank layout has no title placeholder, so the title has to
be inferred. Taking the first text box in shape order -- what this did before --
reads a COVER slide backwards: the venue and paper number sit above the title,
so the 24 pt metadata line was called the title and held to the 32 pt title
floor, while the actual 54 pt title was checked as body text. Two false
failures on a correctly typeset cover, measured on the MMPM deck 2026-08-20.

The two readings that replace it are ordered, and the order matters: a
full-width bar at the very top is a title band even when something larger sits
below it (slide 9 of that deck puts a 48 pt number under a 34 pt title bar),
and only when there is no such bar does size decide.
"""

from __future__ import annotations

import pytest

pytest.importorskip("pptx")

from pptx import Presentation  # noqa: E402
from pptx.util import Cm, Pt  # noqa: E402

from radia_mcp.presentation.tools import (  # noqa: E402
    _slide_title,
    presentation_check_pptx_font_size,
)

SLIDE_W_CM = 33.87
SLIDE_H_CM = 19.05


def _deck():
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W_CM)
    prs.slide_height = Cm(SLIDE_H_CM)
    return prs


def _box(slide, text, left, top, width, height, size_pt):
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    return box


def test_a_title_placeholder_still_wins():
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "the declared title"
    _box(slide, "much larger text lower down", 2.0, 6.0, 20.0, 3.0, 60)
    assert _slide_title(slide) == "the declared title"


def test_a_full_width_top_bar_is_the_title_even_with_bigger_text_below():
    """Slide 9 of the MMPM deck: a 34 pt title bar over a 48 pt result number."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _box(slide, "HACApK solves 162,000 DOF in 60.7 s",
         0.0, 0.0, SLIDE_W_CM, 2.19, 34)
    _box(slide, "162,000", 22.05, 3.39, 10.23, 2.40, 48)
    assert _slide_title(slide).startswith("HACApK solves")


def test_a_cover_slides_title_is_the_large_text_not_the_line_above_it():
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _box(slide, "IEEJ SA-26-069 / RM-26-075", 2.05, 1.20, 29.63, 1.20, 24)
    _box(slide, "Magnetic Moment Method", 2.05, 6.60, 29.46, 2.54, 54)
    assert _slide_title(slide) == "Magnetic Moment Method"


def test_the_cover_metadata_line_is_held_to_the_body_floor(tmp_path):
    """The incident: two false title-floor failures on a correct cover.

    24 pt for a venue line is right; it only failed because it was mistaken
    for the title and measured against the 32 pt title floor.
    """
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _box(slide, "IEEJ SA-26-069 / RM-26-075", 2.05, 1.20, 29.63, 1.20, 24)
    _box(slide, "Magnetic Moment Method", 2.05, 6.60, 29.46, 2.54, 54)
    out = tmp_path / "cover.pptx"
    prs.save(str(out))

    report = presentation_check_pptx_font_size(str(out))
    assert report["total_violations"] == 0, report["violations"]


def test_a_body_line_under_the_title_floor_is_still_caught(tmp_path):
    """The check must not have been softened into silence."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _box(slide, "A title bar", 0.0, 0.0, SLIDE_W_CM, 2.19, 34)
    _box(slide, "far too small to read from the back", 2.0, 6.0, 20.0, 1.2, 12)
    out = tmp_path / "small.pptx"
    prs.save(str(out))

    report = presentation_check_pptx_font_size(str(out))
    assert report["total_violations"] == 1
    assert report["violations"][0]["size_pt"] == 12.0


def test_a_reconstructed_figure_label_is_never_the_title():
    """FIGURE_TEXT:: shapes are labels lifted out of a picture.

    They sit wherever the label sat inside the figure, so one can easily be the
    only text on a slide -- and being called the title put a correct 20 pt
    figure label against the 32 pt title floor.
    """
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    label = _box(slide, "HACApK", 6.0, 5.0, 6.0, 1.2, 20)
    label.name = "FIGURE_TEXT::picture 1::0"
    _box(slide, "the actual title", 0.0, 0.0, SLIDE_W_CM, 2.19, 34)
    assert _slide_title(slide) == "the actual title"


def test_shape_order_is_the_last_resort():
    """No sizes anywhere: fall back to the old behaviour rather than guess."""
    prs = _deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for text, top in (("first in order", 1.0), ("second", 6.0)):
        box = slide.shapes.add_textbox(Cm(2.0), Cm(top), Cm(20.0), Cm(1.2))
        box.text_frame.paragraphs[0].add_run().text = text
    assert _slide_title(slide) == "first in order"
