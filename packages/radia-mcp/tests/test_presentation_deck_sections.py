from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from radia_mcp.presentation.plans.T28 import (
    presentation_speaking_pace_estimate,
)
from radia_mcp.presentation.tools import _slide_title
from radia_mcp.presentation._deck_sections import classify_deck_sections


def _add_slide_with_notes(
    prs: Presentation,
    title: str,
    note_text: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.notes_slide.notes_text_frame.text = note_text


def test_backup_slides_are_excluded_from_speaking_time_by_default(
        tmp_path: Path) -> None:
    prs = Presentation()
    _add_slide_with_notes(prs, "手法の概要", "本編の説明です。" * 10)
    _add_slide_with_notes(prs, "謝辞", "ご清聴ありがとうございました。")
    _add_slide_with_notes(prs, "補足：rank構造", "質疑用の長い説明です。" * 30)
    path = tmp_path / "deck-with-backup.pptx"
    prs.save(path)

    result = presentation_speaking_pace_estimate(str(path))

    assert result["n_slides"] == 2
    assert result["backup_slide_numbers"] == [3]
    assert result["per_slide"][2]["excluded_reason"] == "backup_slide"


def test_sources_block_is_not_counted_as_spoken_notes(tmp_path: Path) -> None:
    prs = Presentation()
    spoken = "ここだけを読み上げます。"
    _add_slide_with_notes(
        prs,
        "精度評価",
        spoken + "\n[Sources]\nhttps://example.com/very-long-source",
    )
    path = tmp_path / "deck-with-sources.pptx"
    prs.save(path)

    result = presentation_speaking_pace_estimate(str(path))

    assert result["total_chars_in_notes"] == len(spoken)
    assert result["source_blocks_excluded"] == 1
    assert result["per_slide"][0]["source_block_excluded"] is True


def test_section_classifier_keeps_thanks_and_marks_appendix() -> None:
    prs = Presentation()
    for title in ("手法", "謝辞", "補足：計算条件"):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title

    sections = classify_deck_sections(list(prs.slides), _slide_title)

    assert sections["main_slide_numbers"] == [1, 2]
    assert sections["closing_slide_numbers"] == [2]
    assert sections["backup_slide_numbers"] == [3]
