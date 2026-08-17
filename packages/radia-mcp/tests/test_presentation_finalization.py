from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from radia_mcp.presentation.tools import presentation_check_final_deck_directory


def _save_deck(path: Path, slide_specs: list[tuple[str, str]]) -> None:
    prs = Presentation()
    for title, body in slide_specs:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(10), Inches(2)
        )
        textbox.text_frame.text = body
    prs.save(path)


def test_one_final_deck_with_same_file_backup_slides_passes(
        tmp_path: Path) -> None:
    _save_deck(
        tmp_path / "SA-26-069_MMPM_発表資料.pptx",
        [
            ("MMPMの計算精度評価", "非線形C型鉄心でFEMとの差を評価した。"),
            ("補足：近零モード", "質問時に示す追加の固有値分布である。"),
        ],
    )

    result = presentation_check_final_deck_directory(str(tmp_path))

    assert result["passed"] is True
    assert result["top_level_pptx_count"] == 1
    assert result["backup_slide_count"] == 1


def test_multiple_root_decks_and_archived_revisions_fail(tmp_path: Path) -> None:
    slides = [("MMPMの計算精度評価", "比較結果を示す本文である。")]
    _save_deck(tmp_path / "final.pptx", slides)
    _save_deck(tmp_path / "final2.pptx", slides)
    archive = tmp_path / "_archive"
    archive.mkdir()
    _save_deck(archive / "old.pptx", slides)

    result = presentation_check_final_deck_directory(str(tmp_path))

    assert result["passed"] is False
    assert "top_level_pptx_count_must_be_one" in result["issues"]
    assert "archived_revision_pptx_must_be_deleted" in result["issues"]


def test_exact_duplicate_slides_must_be_consolidated(tmp_path: Path) -> None:
    duplicate = (
        "MMPMの計算精度評価",
        "同じ主張と同じ証拠を持つ完全重複スライドである。",
    )
    _save_deck(tmp_path / "final.pptx", [duplicate, duplicate])

    result = presentation_check_final_deck_directory(str(tmp_path))

    assert result["passed"] is False
    assert result["duplicate_slide_groups"] == [[1, 2]]
    assert "duplicate_slides_must_be_consolidated" in result["issues"]


def test_unverified_figure_text_blocks_finalization(tmp_path: Path) -> None:
    image_path = tmp_path / "figure.png"
    Image.new("RGB", (400, 200), "white").save(image_path)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "MMPMの計算精度評価"
    slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(1.5), width=Inches(8)
    )
    prs.save(tmp_path / "final.pptx")

    result = presentation_check_final_deck_directory(str(tmp_path))

    assert result["passed"] is False
    assert "embedded_figure_text_must_be_verified_at_20pt" in result["issues"]
    assert result["figure_text_audit"]["unresolved_count"] == 1
