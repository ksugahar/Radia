from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from radia_mcp.presentation.tools import presentation_check_pptx_font_size


def _deck_with_body_and_footer(tmp_path: Path, body_pt: int) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "本文は24ポイント以上とする"
    slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(32)

    body = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5), Inches(11.4), Inches(2.0)
    )
    body.text_frame.text = "聴衆が読む本文"
    body.text_frame.paragraphs[0].runs[0].font.size = Pt(body_pt)

    footer = slide.shapes.add_textbox(
        Inches(10.7), Inches(7.0), Inches(2.0), Inches(0.3)
    )
    footer.text_frame.text = "近畿大学 1"
    footer.text_frame.paragraphs[0].runs[0].font.size = Pt(12)

    path = tmp_path / f"font-{body_pt}.pptx"
    prs.save(path)
    return path


def test_default_body_floor_is_24pt_and_chrome_is_excluded(tmp_path: Path) -> None:
    path = _deck_with_body_and_footer(tmp_path, body_pt=23)

    result = presentation_check_pptx_font_size(str(path))

    assert result["min_body_pt"] == 24
    assert result["min_title_pt"] == 32
    assert result["total_violations"] == 1
    assert result["violations"][0]["text"] == "聴衆が読む本文"
    assert result["excluded_chrome_runs"] == 1


def test_24pt_body_passes(tmp_path: Path) -> None:
    path = _deck_with_body_and_footer(tmp_path, body_pt=24)

    result = presentation_check_pptx_font_size(str(path))

    assert result["total_violations"] == 0


def test_chrome_can_be_included_explicitly(tmp_path: Path) -> None:
    path = _deck_with_body_and_footer(tmp_path, body_pt=24)

    result = presentation_check_pptx_font_size(
        str(path), exclude_chrome=False
    )

    assert result["total_violations"] == 1
    assert result["violations"][0]["text"] == "近畿大学 1"


def test_title_slide_metadata_is_not_misclassified_as_title(
        tmp_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    metadata = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.35)
    )
    metadata.text_frame.text = "静止器・回転機合同研究会 2026年8月27日"
    metadata.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
    title = slide.shapes.add_textbox(
        Inches(0.9), Inches(1.4), Inches(11.5), Inches(1.4)
    )
    title.text_frame.text = "多重極モーメント拘束に基づく磁気モーメント法"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
    path = tmp_path / "title-slide-metadata.pptx"
    prs.save(path)

    result = presentation_check_pptx_font_size(str(path))

    assert result["total_violations"] == 0
