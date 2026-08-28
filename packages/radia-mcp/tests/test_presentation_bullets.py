from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches

from radia_mcp.presentation.tools import (
    presentation_check_bullet_count_per_slide,
    presentation_check_bullet_ending_style,
)


def _mark_bullet(paragraph) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    bullet = OxmlElement("a:buChar")
    bullet.set("char", "•")
    p_pr.append(bullet)


def test_bullet_count_ignores_plain_textbox_paragraphs(tmp_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for index in range(8):
        box = slide.shapes.add_textbox(
            Inches(0.5 + index), Inches(1.0), Inches(0.8), Inches(0.5)
        )
        box.text_frame.text = f"項目{index + 1}"
    pptx_path = tmp_path / "plain-paragraphs.pptx"
    prs.save(pptx_path)

    result = presentation_check_bullet_count_per_slide(str(pptx_path), 5)

    assert result["over_bullet_limit_count"] == 0
    assert result["violations"] == []


def test_bullet_count_detects_explicit_powerpoint_bullets(tmp_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    for index in range(6):
        paragraph = (
            box.text_frame.paragraphs[0]
            if index == 0 else box.text_frame.add_paragraph()
        )
        paragraph.text = f"bullet {index + 1}"
        _mark_bullet(paragraph)
    pptx_path = tmp_path / "explicit-bullets.pptx"
    prs.save(pptx_path)

    result = presentation_check_bullet_count_per_slide(str(pptx_path), 5)

    assert result["over_bullet_limit_count"] == 1
    assert result["violations"] == [{"slide": 1, "bullets": 6}]


def test_bullet_ending_style_ignores_nonbullet_prose(tmp_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    prose = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    prose.text_frame.text = "本文には句点がある。"
    bullets = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    first = bullets.text_frame.paragraphs[0]
    first.text = "bullet without period"
    _mark_bullet(first)
    second = bullets.text_frame.add_paragraph()
    second.text = "bullet with period."
    _mark_bullet(second)
    pptx_path = tmp_path / "bullet-endings.pptx"
    prs.save(pptx_path)

    result = presentation_check_bullet_ending_style(str(pptx_path))

    assert result["bullets_with_period"] == 1
    assert result["bullets_without_period"] == 1
    assert result["mixed"] is True
