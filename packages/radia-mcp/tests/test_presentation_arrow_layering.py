from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches

from radia_mcp.presentation.tools import presentation_check_arrow_layering


def _write_crossing_shapes(
    path: Path,
    *,
    arrow_in_front: bool,
    artifact_style_line: bool = False,
) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def add_arrow() -> None:
        slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(2.0),
            Inches(2.0),
            Inches(2.0),
            Inches(0.8),
        )

    def add_divider() -> None:
        if artifact_style_line:
            divider = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(3.0),
                Inches(1.0),
                0,
                Inches(3.0),
            )
            divider._element.spPr.prstGeom.set("prst", "line")
            return
        slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(3.0),
            Inches(1.0),
            Inches(3.0),
            Inches(4.0),
        )

    if arrow_in_front:
        add_divider()
        add_arrow()
    else:
        add_arrow()
        add_divider()
    prs.save(path)


def test_arrow_in_front_of_intersecting_divider_passes(tmp_path: Path) -> None:
    path = tmp_path / "arrow-front.pptx"
    _write_crossing_shapes(path, arrow_in_front=True)

    result = presentation_check_arrow_layering(str(path))

    assert result["passed"] is True
    assert result["intersection_count"] == 1
    assert result["violation_count"] == 0


def test_arrow_behind_intersecting_divider_is_violation(tmp_path: Path) -> None:
    path = tmp_path / "arrow-behind.pptx"
    _write_crossing_shapes(path, arrow_in_front=False)

    result = presentation_check_arrow_layering(str(path))

    assert result["passed"] is False
    assert result["violation_count"] == 1
    violation = result["violations"][0]
    assert violation["issue"] == "arrow_behind_straight_line"
    assert violation["arrow_z_order"] < violation["line_z_order"]


def test_artifact_tool_preset_line_is_checked(tmp_path: Path) -> None:
    path = tmp_path / "artifact-line.pptx"
    _write_crossing_shapes(
        path,
        arrow_in_front=False,
        artifact_style_line=True,
    )

    result = presentation_check_arrow_layering(str(path))

    assert result["straight_line_count"] == 1
    assert result["intersection_count"] == 1
    assert result["violation_count"] == 1
