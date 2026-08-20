"""The script-vs-slide check has to be able to read Japanese prose.

Slides put an acronym in a box of its own; a script says it in the middle of a
sentence.  The tool used to see the first and not the second, so every acronym
in a deck came back as never spoken -- the worst kind of wrong, because the
report looks specific and is confidently false.

Each test here is one way that went wrong on a real deck (SA-26-069, 2026-08).
"""

from __future__ import annotations

import pytest

pytest.importorskip("pptx")

from radia_mcp.presentation.plans.T30 import (  # noqa: E402
    _extract_tokens,
    presentation_script_vs_slide_coverage,
)


def test_an_acronym_inside_a_japanese_sentence_is_found():
    r"""``\b`` does not work next to CJK: Python's ``\w`` matches Japanese, so
    "のHACApK" has no word boundary before the H."""
    spoken = "既存のHACApK、階層行列の実装をそのまま継承できます。"
    assert "HACApK" in _extract_tokens(spoken)


def test_acronyms_are_found_wherever_they_sit():
    cases = {
        "ブロックJacobiを前処理に置いた": "Jacobi",
        "置いたRichardson反復で": "Richardson",
        "HACApK": "HACApK",
        "（HACApK）を使う": "HACApK",
        "H(div)のGalerkin定式化": "Galerkin",
    }
    for text, token in cases.items():
        assert token in _extract_tokens(text), text


def test_a_longer_latin_word_is_not_split_into_a_false_token():
    assert "ABCdef" in _extract_tokens("これはABCdefです")
    # a token must not be found inside a longer run of ASCII
    assert "ABC" not in _extract_tokens("これはABCdefです")


def build(tmp_path, slides):
    """A deck of (body_text, note_text) pairs.

    Each slide gets a title placeholder AND a separate body box, because the
    tool deliberately leaves the title out of the body it scores -- a deck of
    bare text boxes would be read as title-only and skipped.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    title_only = prs.slide_layouts[5]
    for i, (body, note) in enumerate(slides, 1):
        slide = prs.slides.add_slide(title_only)
        slide.shapes.title.text = "見出し%d" % i
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(3))
        box.text_frame.text = body
        slide.notes_slide.notes_text_frame.text = note
    path = str(tmp_path / "deck.pptx")
    prs.save(path)
    return path


def test_a_spoken_acronym_is_not_reported_as_missing(tmp_path):
    deck = build(tmp_path, [
        ("表紙", "表紙です。"),
        ("HACApK 階層行列 線形解法",
         "既存のHACApKで階層行列を作り、線形解法に渡します。"),
    ])
    got = presentation_script_vs_slide_coverage(deck)
    missing = [t for s in got["low_coverage_slides"] for t in s["missing_tokens"]]
    assert "HACApK" not in missing
    assert got["avg_coverage_pct"] == 100.0


def test_saying_more_than_the_slide_says_is_not_a_gap(tmp_path):
    r"""A maximal kanji run makes "二重積分" and "Gram二重積分" two tokens, so
    the longer phrase used to read as an omission."""
    deck = build(tmp_path, [
        ("表紙", "表紙です。"),
        ("二重積分", "GalerkinのGram二重積分は要りません。"),
    ])
    got = presentation_script_vs_slide_coverage(deck)
    assert got["avg_coverage_pct"] == 100.0


def test_deck_furniture_is_not_scored_against_the_script(tmp_path):
    """The affiliation in the corner of every slide is not spoken."""
    deck = build(tmp_path, [
        ("表紙 近畿大学", "表紙です。"),
        ("面磁荷 近畿大学", "面磁荷の話です。"),
        ("閉包条件 近畿大学", "閉包条件の話です。"),
        ("数値実験 近畿大学", "数値実験の話です。"),
        ("結論 近畿大学", "結論です。"),
    ])
    got = presentation_script_vs_slide_coverage(deck)
    assert "近畿大学" in got["deck_boilerplate_ignored"]
    missing = [t for s in got["low_coverage_slides"] for t in s["missing_tokens"]]
    assert "近畿大学" not in missing


def test_the_cover_is_not_scored(tmp_path):
    """Its English subtitle and author list are shown, not read out."""
    deck = build(tmp_path, [
        ("Magnetic Moment Method Based on Multipole Moment Constraints",
         "近畿大学の菅原です。よろしくお願いいたします。"),
        ("面磁荷", "面磁荷の話です。"),
    ])
    got = presentation_script_vs_slide_coverage(deck)
    cover = got["per_slide"][0]
    assert cover["coverage"] is None
    assert cover["skip_reason"] == "cover_slide"


def _tiny_png(tmp_path):
    """A 1x1 PNG, so the fixture can place a real picture."""
    import base64
    p = tmp_path / "dot.png"
    if not p.exists():
        p.write_bytes(base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk"
            "YPhfDwAChwGA60e6kgAAAABJRU5ErkJggg=="))
    return str(p)


def build_rich(tmp_path, slides, slide_h=6858000):
    """A deck where each slide can carry a body, a figure with a label on it,
    a table, and a take-home banner along the bottom."""
    from pptx import Presentation
    from pptx.util import Emu, Inches

    prs = Presentation()
    title_only = prs.slide_layouts[5]
    for i, spec in enumerate(slides, 1):
        slide = prs.slides.add_slide(title_only)
        slide.shapes.title.text = "見出し%d" % i
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
        box.text_frame.text = spec.get("body", "")

        if spec.get("figure_label"):
            # a real picture, and a label sitting on top of it -- an autoshape
            # would not do: the classifier looks for PICTURE, as a deck has
            for _ in (0,):
                slide.shapes.add_picture(_tiny_png(tmp_path),
                                         Emu(4000000), Emu(2500000),
                                         Emu(2000000), Emu(1500000))
            lbl = slide.shapes.add_textbox(Emu(4500000), Emu(3000000),
                                           Emu(1000000), Emu(300000))
            lbl.text_frame.text = spec["figure_label"]

        if spec.get("banner"):
            b = slide.shapes.add_textbox(Emu(300000), int(0.89 * slide_h),
                                         Emu(8000000), Emu(400000))
            b.text_frame.text = spec["banner"]

        slide.notes_slide.notes_text_frame.text = spec.get("note", "")
    path = str(tmp_path / "rich.pptx")
    prs.save(path)
    return path


def test_an_inflected_form_counts_as_spoken():
    r"""A slide compresses where a speaker inflects.  A maximal kanji run
    makes these different words, and they are the same thing."""
    from radia_mcp.presentation.plans.T30 import _extract_tokens, _is_spoken
    cases = [
        ("独立参照", "独立な参照解との一致ではありません。"),
        ("補償点法", "補償点磁荷法と比べます。"),
        ("内部節点", "内部の節点だけを動かします。"),
    ]
    for token, script in cases:
        assert _is_spoken(token, script, _extract_tokens(script)), token


def test_scattered_pieces_do_not_count_as_spoken():
    """The relaxation must not accept a word whose halves merely appear."""
    from radia_mcp.presentation.plans.T30 import _extract_tokens, _is_spoken
    script = "停止した。条件を書く。"
    assert not _is_spoken("停止条件測定", script, _extract_tokens(script))


def test_a_label_on_a_figure_is_not_scored(tmp_path):
    """"Mesh A" is read by the eye while the speaker talks; it is not a line."""
    deck = build_rich(tmp_path, [
        {"body": "表紙", "note": "表紙です。"},
        {"body": "粗密比較", "figure_label": "Mesh A",
         "note": "粗いメッシュと細かいメッシュを比較します。"},
    ])
    got = presentation_script_vs_slide_coverage(deck)
    missing = [t for s in got["low_coverage_slides"] for t in s["missing_tokens"]]
    assert "Mesh" not in missing


def test_the_take_home_banner_is_checked_on_its_own(tmp_path):
    """Averaged in with the slide it disappears -- a deck can score well while
    the speaker never lands one of its claims out loud."""
    deck = build_rich(tmp_path, [
        {"body": "表紙", "note": "表紙です。"},
        {"body": "面磁荷 閉包条件 数値実験 収束",
         "banner": "局所閉包だけを置換し大規模計算基盤を再利用",
         "note": "面磁荷の閉包条件について、数値実験の収束を示します。"},
    ])
    got = presentation_script_vs_slide_coverage(deck)
    assert got["banner_checked"] == 1
    unspoken = {b["slide_no"] for b in got["banners_not_spoken"]}
    assert unspoken == {2}, got["banner_detail"]


def test_a_spoken_banner_is_not_flagged(tmp_path):
    deck = build_rich(tmp_path, [
        {"body": "表紙", "note": "表紙です。"},
        {"body": "面磁荷",
         "banner": "局所閉包だけを置換し大規模計算基盤を再利用",
         "note": "局所閉包だけを置換して、大規模計算基盤はそのまま再利用できます。"},
    ])
    got = presentation_script_vs_slide_coverage(deck)
    assert got["banners_not_spoken"] == [], got["banner_detail"]


def test_an_unwritten_script_is_still_reported(tmp_path):
    """The point of the check: a slide nobody has scripted must be visible."""
    deck = build(tmp_path, [
        ("表紙", "表紙です。"),
        ("面磁荷 閉包条件", ""),
    ])
    got = presentation_script_vs_slide_coverage(deck)
    assert got["n_low_coverage_slides"] == 1
    assert got["low_coverage_slides"][0]["slide_no"] == 2
