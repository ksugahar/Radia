"""A vector figure has to be audited, not skipped.

Every case here is something the deck did on 2026-08-20, the day its figures
were changed from PNG to SVG. The first one is the reason the rest exist: the
paste audit reported "0 pictures, 0 flagged" for a deck carrying eight
figures, because python-pptx raises on a picture that has no raster and the
audit swallowed the exception. Nothing was wrong with the deck and nothing was
checked either -- the report read exactly like a clean one.

The others are readings that came out wrong before the measurement was right,
and each is pinned to the number that was actually measured on the file:

- font sizes inside a reduction matrix, which put a slide label at "237.7 pt";
- an SVG whose width is unitless, where guessing CSS pixels shrank every label
  by a quarter and reported a legible figure as too small;
- an empty <text> element, which is not a label;
- subscripts, which are smaller than the body by construction and must not be
  what the floor is measured against.
"""

from __future__ import annotations

import pytest

pytest.importorskip("pptx")

from pptx import Presentation  # noqa: E402
from pptx.opc.constants import RELATIONSHIP_TYPE as RT  # noqa: E402
from pptx.opc.package import Part  # noqa: E402
from pptx.opc.packuri import PackURI  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402
from pptx.util import Emu  # noqa: E402

from radia_mcp.figure._pptx_audit import audit_pptx_figures  # noqa: E402
from radia_mcp.figure._svg_pptx import (  # noqa: E402
    svg_font_sizes_units,
    svg_picture_row,
    svg_size_pt,
    svg_viewbox,
)

SVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"

# a 1x1 transparent PNG, so a picture can exist without carrying artwork
_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
    "890000000a49444154789c6360000002000100ffff03000006000557bfabd400"
    "00000049454e44ae426082")

EMU_CM = 360000


def _png(tmp_path):
    p = tmp_path / "dot.png"
    p.write_bytes(_PNG)
    return str(p)


def _deck_with_svg(tmp_path, svg_text, width_cm, height_cm, strip_raster=False):
    """A one-slide deck holding one SVG picture, the way PowerPoint stores one."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = slide.shapes.add_picture(_png(tmp_path), 0, 0,
                                   Emu(int(width_cm * EMU_CM)),
                                   Emu(int(height_cm * EMU_CM)))
    part = Part(PackURI("/ppt/media/figure.svg"), "image/svg+xml",
                prs.part.package, svg_text.encode("utf-8"))
    rid = slide.part.relate_to(part, RT.IMAGE)
    blip = pic._element.blipFill.blip
    ext_lst = blip.makeelement(qn("a:extLst"), {})
    ext = blip.makeelement(qn("a:ext"),
                           {"uri": "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"})
    ext.append(blip.makeelement("{%s}svgBlip" % SVG_NS, {qn("r:embed"): rid}))
    ext_lst.append(ext)
    blip.append(ext_lst)
    if strip_raster:
        # COM's AddPicture leaves the blip with no r:embed at all -- the case
        # python-pptx cannot read and the audit used to drop.
        del blip.attrib[qn("r:embed")]
    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    return str(out)


def _stub_shape(width_cm, height_cm):
    class _S:
        width = int(width_cm * EMU_CM)
        height = int(height_cm * EMU_CM)
    return _S()


def _row(svg_text, width_cm, height_cm, floor=20.0):
    return svg_picture_row(1, "pic", _stub_shape(width_cm, height_cm),
                           svg_text.encode("utf-8"),
                           slide_area_pt=960.0 * 540.0,
                           scale_tol=0.02, aspect_tol=0.01,
                           min_area_fraction=0.01, min_visible_font_pt=floor)


def test_an_svg_picture_is_counted_not_dropped(tmp_path):
    """The incident: eight figures, "0 pictures, 0 flagged"."""
    svg = ('<svg width="200pt" height="100pt" viewBox="0 0 200 100">'
           '<text style="font-size: 20px">label</text></svg>')
    rep = audit_pptx_figures(_deck_with_svg(tmp_path, svg, 20.0, 10.0,
                                            strip_raster=True))
    assert rep["n_pictures"] == 1
    assert rep["n_vector"] == 1
    assert rep["pictures"][0]["kind"] == "svg"


def test_a_picture_with_no_artwork_at_all_is_reported(tmp_path):
    """Linked or OLE artwork cannot be measured -- and must not vanish."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = slide.shapes.add_picture(_png(tmp_path), 0, 0, Emu(10 * EMU_CM),
                                   Emu(5 * EMU_CM))
    del pic._element.blipFill.blip.attrib[qn("r:embed")]
    out = tmp_path / "linked.pptx"
    prs.save(str(out))

    rep = audit_pptx_figures(str(out))
    assert rep["n_pictures"] == 1
    assert rep["n_unmeasurable"] == 1
    assert any("NOT MEASURABLE" in r for r in rep["pictures"][0]["risks"])


def test_the_on_slide_size_comes_from_the_pasted_width(tmp_path):
    """The C-type model figure: 25 units in a 453.48-unit page, pasted 13.05 cm.

    Its width attribute is unitless, so reading it as CSS pixels made every
    label three quarters of its real size and reported the figure as 16.3 pt
    against a 20 pt floor.  The unit cancels when the size is taken against the
    pasted width, and the answer is 25 * (13.05 cm in pt) / 453.48 = 20.4 pt.
    """
    svg = ('<svg width="453.48" height="467.76" viewBox="0 0 453.48 467.76">'
           '<text font-size="25">182.5 mm</text></svg>')
    row = _row(svg, 13.05, 13.47)
    assert row["largest_font_on_slide_pt"] == pytest.approx(20.4, abs=0.15)
    assert row["risks"] == []


def test_a_font_size_inside_a_reduction_is_read_at_its_drawn_size():
    """The concept figure: font-size="173.9" inside matrix(.0733 ...).

    Read as written it is 237.7 pt on the slide, which is not a size anything
    on a slide is set in.  Multiplied by its own transform it is 12.7 units.
    """
    svg = ('<svg width="453.48" height="212.64" viewBox="0 0 453.48 212.64">'
           '<text transform="matrix(.07325287 .0000000363 .000000016 .07314812'
           ' 104.25 21.75)" font-size="173.9"><tspan>2</tspan></text></svg>')
    sizes = svg_font_sizes_units(svg)
    assert sizes == pytest.approx([12.7], abs=0.1)


def test_an_empty_text_element_is_not_a_label():
    svg = ('<svg width="200pt" height="100pt" viewBox="0 0 200 100">'
           '<text font-size="8"> </text>'
           '<text font-size="20">real</text></svg>')
    assert svg_font_sizes_units(svg) == pytest.approx([20.0])


def test_the_floor_is_measured_against_the_largest_label():
    """A 10^4 tick's exponent is smaller than the axis label by construction.

    Checking the smallest -- or the most common -- size would fail a figure
    whose actual labels are fine, and a check that fires on a good figure is
    one nobody runs twice.
    """
    svg = ('<svg width="200pt" height="100pt" viewBox="0 0 200 100">'
           '<text font-size="24">Memory (GiB)</text>'
           '<text font-size="14">4</text><text font-size="14">5</text>'
           '<text font-size="14">6</text></svg>')
    row = _row(svg, 200 / 72 * 2.54, 100 / 72 * 2.54)
    assert row["largest_font_on_slide_pt"] == pytest.approx(24.0, abs=0.1)
    assert row["smallest_font_on_slide_pt"] == pytest.approx(14.0, abs=0.1)
    assert row["risks"] == [], "the subscripts are not the body text"


def test_a_figure_whose_biggest_label_is_under_the_floor_is_flagged():
    svg = ('<svg width="200pt" height="100pt" viewBox="0 0 200 100">'
           '<text font-size="12">everything is small</text></svg>')
    row = _row(svg, 200 / 72 * 2.54, 100 / 72 * 2.54)
    assert any("TOO SMALL" in r for r in row["risks"])


def test_a_stretched_vector_picture_is_flagged():
    svg = ('<svg width="200pt" height="100pt" viewBox="0 0 200 100">'
           '<text font-size="24">label</text></svg>')
    row = _row(svg, 200 / 72 * 2.54, 100 / 72 * 2.54 * 1.5)
    assert any("ASPECT DISTORTED" in r for r in row["risks"])


def test_resizing_vector_artwork_is_not_a_defect_by_itself():
    """A raster pasted at 180% is interpolated; a vector is not.

    Slide 3 carries its figure at 182% and reads 16-30 pt on the slide, which
    is legible -- flagging the scale would have been a false alarm.
    """
    svg = ('<svg width="200pt" height="100pt" viewBox="0 0 200 100">'
           '<text font-size="14">label</text></svg>')
    row = _row(svg, 200 / 72 * 2.54 * 1.8, 100 / 72 * 2.54 * 1.8)
    assert row["scale"] == pytest.approx(1.8, abs=0.01)
    assert row["risks"] == []
    assert row["largest_font_on_slide_pt"] == pytest.approx(25.2, abs=0.2)


def test_a_document_that_states_no_size_is_reported_as_such():
    row = _row("<svg><text font-size='24'>label</text></svg>", 10.0, 5.0)
    assert any("STATES NO SIZE" in r for r in row["risks"])
    assert row["text_check"] == "not-verifiable"


def test_size_and_viewbox_are_read_in_the_documents_own_terms():
    svg = '<svg width="786.667" height="393.333" viewBox="0 0 590 295"></svg>'
    # PowerPoint rewrites matplotlib's "590.4pt" as unitless CSS pixels
    assert svg_size_pt(svg)[0] == pytest.approx(590.0, abs=0.1)
    assert svg_viewbox(svg) == (590.0, 295.0)


def test_a_raster_picture_is_still_audited_the_old_way(tmp_path):
    """The SVG branch must not have changed what a PNG picture reports."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(_png(tmp_path), 0, 0, Emu(10 * EMU_CM),
                             Emu(5 * EMU_CM))
    out = tmp_path / "raster.pptx"
    prs.save(str(out))

    rep = audit_pptx_figures(str(out))
    assert rep["n_pictures"] == 1
    assert rep["pictures"][0]["kind"] == "raster"
    assert rep["n_vector"] == 0
    assert rep["n_raster"] == 1


def test_a_raster_figure_is_flagged_because_a_figure_should_be_vector(tmp_path):
    """The rule is not taste: a raster's text cannot be measured from the file.

    That is how a 16 pt label shipped unnoticed -- the size check on a raster
    returns "OCR required", not a verdict, so the 20 pt floor is unenforceable
    on one.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(_png(tmp_path), 0, 0, Emu(20 * EMU_CM),
                             Emu(10 * EMU_CM))
    out = tmp_path / "raster_big.pptx"
    prs.save(str(out))

    rep = audit_pptx_figures(str(out))
    assert any("RASTER FIGURE" in r for r in rep["pictures"][0]["risks"])


def test_a_photograph_can_declare_itself_raster(tmp_path):
    """A screen capture cannot be vector, and a flag nobody can clear is noise.

    The decision is recorded in the shape name, where a reader of the file can
    see it -- the same convention as FIGURE_TEXT::.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = slide.shapes.add_picture(_png(tmp_path), 0, 0, Emu(20 * EMU_CM),
                                   Emu(10 * EMU_CM))
    pic.name = "RASTER_OK::Cubit screen capture"
    out = tmp_path / "declared.pptx"
    prs.save(str(out))

    row = audit_pptx_figures(str(out))["pictures"][0]
    assert row["declared_raster"] is True
    assert not any("RASTER FIGURE" in r for r in row["risks"])


def test_a_logo_is_not_asked_to_be_vector(tmp_path):
    """Decorations below the area threshold were never in scope."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(_png(tmp_path), 0, 0, Emu(1 * EMU_CM),
                             Emu(1 * EMU_CM))
    out = tmp_path / "logo.pptx"
    prs.save(str(out))

    row = audit_pptx_figures(str(out))["pictures"][0]
    assert row["minor"] is True
    assert not any("RASTER FIGURE" in r for r in row["risks"])
