"""Deck-integrity checks (2026-08-16, MMPM SA-26-069 defects).

Locks the two failure classes every other presentation lint walked past:
unrendered math markup left in slide text, and the same artwork reused on
two slides.
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from radia_mcp.presentation.tools import (
    presentation_apply_math_subscripts,
    presentation_check_duplicate_slide_images,
    presentation_check_raw_math_markup,
)


def _deck_with_text(tmp_path: Path, lines: list[str], name: str = "d.pptx") -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    frame = box.text_frame
    for i, line in enumerate(lines):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(24)
    path = tmp_path / name
    prs.save(path)
    return path


def _image(tmp_path: Path, name: str, color: str = "white") -> Path:
    path = tmp_path / name
    Image.new("RGB", (600, 400), color).save(path, dpi=(300, 300))
    return path


# ---------------------------------------------------------------- math markup
def test_underscore_subscript_is_flagged() -> None:
    from radia_mcp.presentation._deck_integrity import _segment_math

    segs = _segment_math("M_ij = B_fj(x) a_j")
    assert segs is not None
    assert "".join(t for t, _b in segs) == "Mij = Bfj(x) aj"
    assert [t for t, b in segs if b] == ["ij", "fj", "j"]


def test_braced_group_flattens_onto_the_script_level() -> None:
    from radia_mcp.presentation._deck_integrity import _segment_math

    segs = _segment_math("Q_{S_f}")
    assert "".join(t for t, _b in segs) == "QSf"
    # one contiguous subscript chunk: the inner token is flattened, not nested
    assert [t for t, b in segs if b] == ["Sf"]


def test_token_split_across_runs_is_repaired(tmp_path: Path) -> None:
    """PowerPoint splits a formula at every font change, so `Sigma_f` lands as
    a run ending in `_` plus a run starting with `f` (MMPM slide 5)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    frame = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2)).text_frame
    for piece in ("Σ_", "f dS"):                 # one paragraph, two runs
        run = frame.paragraphs[0].add_run()
        run.text = piece
        run.font.size = Pt(24)
    deck = tmp_path / "split.pptx"
    prs.save(deck)

    assert presentation_check_raw_math_markup(str(deck))["n_findings"] == 1
    out = tmp_path / "split_fixed.pptx"
    presentation_apply_math_subscripts(str(deck), out_path=str(out), dry_run=False)

    fixed = Presentation(str(out))
    runs = [r for s in fixed.slides for sh in s.shapes if sh.has_text_frame
            for p in sh.text_frame.paragraphs for r in p.runs]
    assert "".join(r.text for r in runs) == "Σf dS"
    subs = [r.text for r in runs if r._r.find(qn("a:rPr")) is not None
            and r._r.find(qn("a:rPr")).get("baseline") == "-25000"]
    assert subs == ["f"]
    assert presentation_check_raw_math_markup(str(out))["ok"] is True


def test_identifiers_and_report_numbers_are_not_math() -> None:
    from radia_mcp.presentation._deck_integrity import _segment_math

    assert _segment_math("figure_text_source_evidence.json") is None
    assert _segment_math("SA-26-069 / RM-26-075") is None
    assert _segment_math("HACApK") is None


def test_check_reports_slides_and_repairability(tmp_path: Path) -> None:
    deck = _deck_with_text(tmp_path, ["H(x_c) の評価", r"\frac{a}{b} は手作業"])

    result = presentation_check_raw_math_markup(str(deck))

    assert result["ok"] is False
    assert result["n_findings"] == 2
    kinds = {f["kind"]: f["auto_repairable"] for f in result["findings"]}
    assert kinds["subscript"] is True
    assert kinds["latex_macro"] is False


def test_clean_deck_passes(tmp_path: Path) -> None:
    deck = _deck_with_text(tmp_path, ["磁荷中性で6面磁荷を閉じる", "162,000 自由度"])

    assert presentation_check_raw_math_markup(str(deck))["ok"] is True


# ---------------------------------------------------------------- repair pass
def test_apply_subscripts_dry_run_does_not_write(tmp_path: Path) -> None:
    deck = _deck_with_text(tmp_path, ["H(x_c)"])
    before = deck.read_bytes()

    result = presentation_apply_math_subscripts(str(deck))

    assert result["dry_run"] is True
    assert result["n_paragraphs_changed"] == 1
    assert result["changes"][0]["after"] == "H(xc)"
    assert deck.read_bytes() == before


def test_apply_subscripts_splits_runs_and_sets_baseline(tmp_path: Path) -> None:
    deck = _deck_with_text(tmp_path, ["M_ij = 0"])
    out = tmp_path / "fixed.pptx"

    result = presentation_apply_math_subscripts(
        str(deck), out_path=str(out), dry_run=False)
    assert result["n_paragraphs_changed"] == 1

    prs = Presentation(str(out))
    runs = [r for s in prs.slides for sh in s.shapes if sh.has_text_frame
            for p in sh.text_frame.paragraphs for r in p.runs]
    assert "".join(r.text for r in runs) == "Mij = 0"
    subs = [r for r in runs if r._r.find(qn("a:rPr")) is not None
            and r._r.find(qn("a:rPr")).get("baseline") == "-25000"]
    assert [r.text for r in subs] == ["ij"]
    # the original run formatting survives the split
    assert all(r.font.size == Pt(24) for r in runs)
    assert presentation_check_raw_math_markup(str(out))["ok"] is True


def test_apply_subscripts_leaves_latex_macros_to_a_human(tmp_path: Path) -> None:
    deck = _deck_with_text(tmp_path, [r"\frac{a_i}{b}"])

    result = presentation_apply_math_subscripts(str(deck))

    assert result["n_paragraphs_changed"] == 0
    assert result["unhandled"][0]["reason"] == "LaTeX macro"


def test_superscript_uses_the_raised_baseline(tmp_path: Path) -> None:
    deck = _deck_with_text(tmp_path, ["10^{-3} 以下"])
    out = tmp_path / "sup.pptx"

    presentation_apply_math_subscripts(str(deck), out_path=str(out), dry_run=False)

    prs = Presentation(str(out))
    runs = [r for s in prs.slides for sh in s.shapes if sh.has_text_frame
            for p in sh.text_frame.paragraphs for r in p.runs]
    raised = [r.text for r in runs if r._r.find(qn("a:rPr")) is not None
              and r._r.find(qn("a:rPr")).get("baseline") == "30000"]
    assert raised == ["-3"]


# ------------------------------------------------------------ duplicate image
def test_same_picture_on_two_slides_is_reported(tmp_path: Path) -> None:
    img = _image(tmp_path, "cube.png")
    prs = Presentation()
    for width in (Inches(6), Inches(6.5)):          # reused AND at two scales
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(img), Inches(1), Inches(1), width=width)
    deck = tmp_path / "dup.pptx"
    prs.save(deck)

    result = presentation_check_duplicate_slide_images(str(deck))

    assert result["ok"] is False
    group = result["groups"][0]
    assert group["slides"] == [1, 2]
    assert group["adjacent"] is True
    assert group["same_scale"] is False


def test_distinct_pictures_pass(tmp_path: Path) -> None:
    prs = Presentation()
    for name, color in (("a.png", "white"), ("b.png", "black")):
        img = _image(tmp_path, name, color)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(img), Inches(1), Inches(1), width=Inches(6))
    deck = tmp_path / "distinct.pptx"
    prs.save(deck)

    assert presentation_check_duplicate_slide_images(str(deck))["ok"] is True


def test_small_repeated_marks_such_as_logos_are_ignored(tmp_path: Path) -> None:
    img = _image(tmp_path, "logo.png")
    prs = Presentation()
    for _ in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(img), Inches(9), Inches(6.8), width=Inches(0.8))
    deck = tmp_path / "logo.pptx"
    prs.save(deck)

    assert presentation_check_duplicate_slide_images(str(deck))["ok"] is True
