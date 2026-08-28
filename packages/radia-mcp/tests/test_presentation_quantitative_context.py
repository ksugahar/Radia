from pptx import Presentation
from pptx.util import Inches

from radia_mcp.presentation.tools import (
    presentation_check_quantitative_claim_context,
)


def _deck(tmp_path, texts):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for index, text in enumerate(texts):
        box = slide.shapes.add_textbox(
            Inches(1), Inches(1 + index * 0.7), Inches(8), Inches(0.5)
        )
        box.text = text
    path = tmp_path / "quantitative-context.pptx"
    prs.save(path)
    return path


def test_flags_percentage_without_quantity_location_or_baseline(tmp_path):
    path = _deck(tmp_path, ["0.28%", "FEMとの差"])
    result = presentation_check_quantitative_claim_context(str(path))
    assert result["passed"] is False
    finding = result["findings"][0]
    assert finding["issue"] == "ambiguous_percentage"
    assert "quantity" in finding["missing"]
    assert "location_or_aggregation" in finding["missing"]


def test_accepts_self_contained_percentage_claim(tmp_path):
    path = _deck(tmp_path, ["0.28%", "ギャップ中央 Bz：FEM基準との差"])
    result = presentation_check_quantitative_claim_context(str(path))
    assert result["passed"] is True


def test_flags_undefined_internal_mesh_identifier(tmp_path):
    path = _deck(tmp_path, ["q6 / q10 / q20", "ギャップ中央 Bz"])
    result = presentation_check_quantitative_claim_context(str(path))
    assert result["passed"] is False
    assert any(
        finding["issue"] == "undefined_internal_identifier"
        for finding in result["findings"]
    )


def test_accepts_defined_internal_mesh_identifier(tmp_path):
    path = _deck(tmp_path, ["q6 = 7,200 DoF", "q10 = 18,900 DoF"])
    result = presentation_check_quantitative_claim_context(str(path))
    assert result["passed"] is True
