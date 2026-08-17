from __future__ import annotations

import json
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from radia_mcp.presentation.tools import (
    presentation_check_embedded_figure_text_size,
    presentation_check_pptx_font_size,
    presentation_replace_embedded_figure_text,
)


def _deck_with_baked_text(tmp_path: Path) -> tuple[Path, str]:
    image_path = tmp_path / "figure.png"
    image = Image.new("RGB", (1000, 500), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((100, 98, 220, 111), fill="black")
    draw.rectangle((230, 98, 350, 111), fill="black")
    image.save(image_path)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    picture = slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(1), width=Inches(10)
    )
    deck = tmp_path / "input.pptx"
    prs.save(deck)
    return deck, picture.name


def _manifest(
        tmp_path: Path,
        shape_name: str,
        bbox_height_px: int = 14,
        confidence: float = 0.99) -> Path:
    path = tmp_path / "ocr.json"
    path.write_text(
        json.dumps({
            "pictures": [{
                "slide": 1,
                "shape": shape_name,
                "words": [{
                    "text": "HACApK",
                    "confidence": confidence,
                    "bbox": [90, 98, 240, 98 + bbox_height_px],
                }],
            }]
        }),
        encoding="utf-8",
    )
    return path


def test_repair_dry_run_does_not_write_files(tmp_path: Path) -> None:
    deck, shape_name = _deck_with_baked_text(tmp_path)
    manifest = _manifest(tmp_path, shape_name)
    output = tmp_path / "candidate.pptx"

    result = presentation_replace_embedded_figure_text(
        str(deck), str(output), str(manifest), dry_run=True
    )

    assert result["status"] == "dry_run"
    assert result["candidate_count"] == 1
    assert result["written"] is False
    assert not output.exists()


def test_repair_writes_inpainted_picture_and_native_20pt_text(
        tmp_path: Path) -> None:
    deck, shape_name = _deck_with_baked_text(tmp_path)
    manifest = _manifest(tmp_path, shape_name)
    output = tmp_path / "candidate.pptx"
    before = Presentation(str(deck)).slides[0].shapes[0].image.blob

    result = presentation_replace_embedded_figure_text(
        str(deck), str(output), str(manifest), dry_run=False,
        source_unavailable_confirmed=True,
    )

    assert result["status"] == "candidate_ready"
    assert result["replacement_count"] == 1
    assert output.exists()
    repaired = Presentation(str(output))
    picture = next(
        shape for shape in repaired.slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    assert picture.image.blob != before
    text_shapes = [
        shape for shape in repaired.slides[0].shapes
        if getattr(shape, "name", "").startswith("FIGURE_TEXT::")
    ]
    assert len(text_shapes) == 1
    assert text_shapes[0].text_frame.text == "HACApK"
    run = text_shapes[0].text_frame.paragraphs[0].runs[0]
    assert run.font.size.pt == 20.0
    font_audit = presentation_check_pptx_font_size(str(output))
    assert font_audit["total_violations"] == 0

    repaired_manifest = json.loads(
        Path(result["output_manifest"]).read_text(encoding="utf-8")
    )
    assert repaired_manifest["pictures"][0]["confirmed_textless"] is True
    image_text_audit = presentation_check_embedded_figure_text_size(
        str(output),
        ocr_backend="manifest",
        ocr_manifest_path=result["output_manifest"],
    )
    assert image_text_audit["passed"] is True


def test_low_confidence_word_is_not_replaced(tmp_path: Path) -> None:
    deck, shape_name = _deck_with_baked_text(tmp_path)
    manifest = _manifest(tmp_path, shape_name, confidence=0.40)

    result = presentation_replace_embedded_figure_text(
        str(deck), str(tmp_path / "candidate.pptx"), str(manifest),
        dry_run=True,
    )

    assert result["candidate_count"] == 0
    assert result["unresolved_count"] == 1


def test_text_already_above_minimum_is_not_replaced(tmp_path: Path) -> None:
    deck, shape_name = _deck_with_baked_text(tmp_path)
    manifest = _manifest(tmp_path, shape_name, bbox_height_px=24)

    result = presentation_replace_embedded_figure_text(
        str(deck), str(tmp_path / "candidate.pptx"), str(manifest),
        dry_run=True,
    )

    assert result["candidate_count"] == 0


def test_reconstructed_text_overlap_requires_review(tmp_path: Path) -> None:
    deck, shape_name = _deck_with_baked_text(tmp_path)
    manifest = tmp_path / "overlap.json"
    manifest.write_text(
        json.dumps({
            "pictures": [{
                "slide": 1,
                "shape": shape_name,
                "words": [
                    {"text": "Long reconstructed label A", "confidence": 0.99,
                     "bbox": [100, 98, 220, 112]},
                    {"text": "Long reconstructed label B", "confidence": 0.99,
                     "bbox": [230, 98, 350, 112]},
                ],
            }]
        }),
        encoding="utf-8",
    )

    result = presentation_replace_embedded_figure_text(
        str(deck), str(tmp_path / "overlap-candidate.pptx"), str(manifest),
        dry_run=False, source_unavailable_confirmed=True,
    )

    assert result["status"] == "needs_review"
    assert result["overlap_warning_count"] >= 1


def test_repair_write_is_blocked_when_source_may_exist(tmp_path: Path) -> None:
    deck, shape_name = _deck_with_baked_text(tmp_path)
    manifest = _manifest(tmp_path, shape_name)

    result = presentation_replace_embedded_figure_text(
        str(deck), str(tmp_path / "candidate.pptx"), str(manifest),
        dry_run=False,
    )

    assert "fallback-only" in result["error"]
    assert not (tmp_path / "candidate.pptx").exists()
