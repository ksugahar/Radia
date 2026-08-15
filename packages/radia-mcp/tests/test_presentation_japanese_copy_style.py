from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from radia_mcp.presentation.tools import presentation_check_japanese_copy_style


def _presentation_with_body(text: str) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "磁気モーメント法"
    title_slide.placeholders[1].text = "研究会発表"
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "PEEC法は計算時間を10分の1に低減する"
    box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5), Inches(11.4), Inches(3.5)
    )
    box.text_frame.text = text
    return prs


def test_accepts_contextual_breaks_and_nominal_endings(tmp_path: Path) -> None:
    prs = _presentation_with_body("法線連続性の維持\v高次曲面への拡張")
    pptx_path = tmp_path / "good-copy.pptx"
    prs.save(pptx_path)

    result = presentation_check_japanese_copy_style(str(pptx_path))

    assert result["passed"] is True
    assert result["awkward_line_break_count"] == 0
    assert result["nominal_like_count"] == 1
    assert result["verbal_sentence_count"] == 0


def test_flags_break_after_particle(tmp_path: Path) -> None:
    prs = _presentation_with_body(
        "需要家の近くに発電設備をつくって電力を\v"
        "確保し，自律的に需給バランスを図る"
    )
    pptx_path = tmp_path / "bad-break.pptx"
    prs.save(pptx_path)

    result = presentation_check_japanese_copy_style(str(pptx_path))

    assert result["awkward_line_break_count"] == 1
    assert result["awkward_line_breaks"][0]["before"].endswith("を")


def test_flags_verbal_sentence_and_screen_only_lead_in(tmp_path: Path) -> None:
    prs = _presentation_with_body("そこで本研究では非線形解析を行う。")
    pptx_path = tmp_path / "wordy-copy.pptx"
    prs.save(pptx_path)

    result = presentation_check_japanese_copy_style(str(pptx_path))

    assert result["passed"] is False
    assert result["nominal_ratio_passed"] is False
    assert result["verbal_sentence_count"] == 1
    assert result["screen_only_lead_in_count"] == 1


def test_excludes_question_formula_title_and_footer(tmp_path: Path) -> None:
    prs = _presentation_with_body("なぜ誤差が低減するか？")
    slide = prs.slides[1]
    formula = slide.shapes.add_textbox(
        Inches(0.8), Inches(5.1), Inches(5.0), Inches(0.5)
    )
    formula.text_frame.text = "H(div)-Galerkin + Q = 0"
    footer = slide.shapes.add_textbox(
        Inches(9.5), Inches(7.1), Inches(3.0), Inches(0.2)
    )
    footer.text_frame.text = "近畿大学 2"
    pptx_path = tmp_path / "excluded-copy.pptx"
    prs.save(pptx_path)

    result = presentation_check_japanese_copy_style(str(pptx_path))

    assert result["paragraphs_checked"] == 0
    assert result["nominal_like_ratio"] == 1.0
    assert result["verbal_sentence_count"] == 0
