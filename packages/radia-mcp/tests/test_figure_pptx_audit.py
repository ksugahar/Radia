"""PPTX paste-scale audit -- the slide analogue of audit_tex_figures.

A figure authored at W cm carries 24 pt text only while it is PASTED at W cm;
dragging it smaller rescales the text silently (MMPM SA-26-069, 2026-08-16:
16.49 cm artwork pasted at 13.97 cm -> 20.3 pt).
"""
from __future__ import annotations

from pathlib import Path
import xml.etree.ElementTree as ET

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from radia_mcp.figure import audit_pptx_figures
from radia_mcp.figure.tools import figure_audit_pptx_figures
from radia_mcp.common.pptx_svg import svg_length_pt as _svg_length_pt
from radia_mcp.figure._pptx_audit import _svg_asset


def _png(tmp_path: Path, name: str, px: tuple[int, int], dpi: int | None = 300) -> Path:
    path = tmp_path / name
    img = Image.new("RGB", px, "white")
    img.save(path, **({"dpi": (dpi, dpi)} if dpi else {}))
    return path


def _deck(tmp_path: Path, image: Path, width_in: float,
          height_in: float | None = None, name: str = "d.pptx") -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    kw = {"width": Inches(width_in)}
    if height_in is not None:
        kw["height"] = Inches(height_in)
    slide.shapes.add_picture(str(image), Inches(0.5), Inches(0.5), **kw)
    path = tmp_path / name
    prs.save(path)
    return path


def test_paste_at_the_authored_width_is_clean(tmp_path: Path) -> None:
    # 1800 px @ 300 dpi == 6 in authored; pasted at 6 in.
    deck = _deck(tmp_path, _png(tmp_path, "f.png", (1800, 1200)), width_in=6.0)

    rep = audit_pptx_figures(str(deck))

    assert rep["n_pictures"] == 1
    pic = rep["pictures"][0]
    # The paste itself is clean.  The only thing reported is that the figure is
    # a raster at all -- a separate rule (a figure should be vector), added
    # 2026-08-21, which this test is not about.
    assert [r for r in pic["risks"] if "RASTER FIGURE" not in r] == []
    assert pic["scale"] == 1.0
    assert pic["displayed_figure_font_pt"] == 24.0
    assert pic["effective_dpi"] == 300.0


def test_downscaled_paste_reports_the_displayed_font(tmp_path: Path) -> None:
    # authored 6 in, pasted 4 in -> scale 2/3 -> 24 pt text displayed at 16 pt.
    deck = _deck(tmp_path, _png(tmp_path, "f.png", (1800, 1200)), width_in=4.0)

    pic = audit_pptx_figures(str(deck))["pictures"][0]

    assert round(pic["scale"], 3) == 0.667
    assert pic["displayed_figure_font_pt"] == 16.0
    assert any("DOWNSCALED" in r for r in pic["risks"])
    assert any("below the 20 pt slide floor" in r for r in pic["risks"])


def test_mild_downscale_stays_above_the_floor(tmp_path: Path) -> None:
    # 0.9 scale: flagged as rescaled, but 21.6 pt is still above the floor.
    deck = _deck(tmp_path, _png(tmp_path, "f.png", (1800, 1200)), width_in=5.4)

    pic = audit_pptx_figures(str(deck))["pictures"][0]

    assert pic["displayed_figure_font_pt"] == 21.6
    assert any("DOWNSCALED" in r for r in pic["risks"])
    assert not any("below the 20 pt slide floor" in r for r in pic["risks"])


def test_upscaled_paste_reports_low_effective_dpi(tmp_path: Path) -> None:
    deck = _deck(tmp_path, _png(tmp_path, "f.png", (600, 400)), width_in=6.0)

    pic = audit_pptx_figures(str(deck))["pictures"][0]

    assert any("UPSCALED" in r for r in pic["risks"])
    assert any("LOW EFFECTIVE DPI" in r for r in pic["risks"])
    assert pic["effective_dpi"] == 100.0


def test_stretched_picture_is_reported(tmp_path: Path) -> None:
    deck = _deck(tmp_path, _png(tmp_path, "f.png", (1800, 1200)),
                 width_in=6.0, height_in=6.0)

    pic = audit_pptx_figures(str(deck))["pictures"][0]

    assert any("ASPECT DISTORTED" in r for r in pic["risks"])


def test_missing_dpi_metadata_is_unverifiable_not_a_pass(tmp_path: Path) -> None:
    deck = _deck(tmp_path, _png(tmp_path, "f.png", (1800, 1200), dpi=None),
                 width_in=6.0)

    pic = audit_pptx_figures(str(deck))["pictures"][0]

    assert any("NO DPI METADATA" in r for r in pic["risks"])


def test_tiny_marks_are_measured_but_never_flagged(tmp_path: Path) -> None:
    deck = _deck(tmp_path, _png(tmp_path, "logo.png", (600, 400)), width_in=0.8)

    rep = audit_pptx_figures(str(deck))

    assert rep["n_pictures"] == 1
    assert rep["n_flagged"] == 0
    assert rep["pictures"][0]["minor"] is True


def test_report_lists_every_picture_and_the_fix(tmp_path: Path) -> None:
    deck = _deck(tmp_path, _png(tmp_path, "f.png", (1800, 1200)), width_in=4.0)

    text = figure_audit_pptx_figures(str(deck))

    assert "1 pictures, 1 flagged" in text
    assert "DOWNSCALED" in text
    assert "lab_figure(medium='presentation'" in text


def test_missing_file_is_reported(tmp_path: Path) -> None:
    assert "file not found" in figure_audit_pptx_figures(str(tmp_path / "nope.pptx"))


def test_pdf_rasterised_at_the_paste_width_audits_clean(tmp_path: Path) -> None:
    fitz = pytest.importorskip("fitz")
    from radia_mcp.figure import slide_png_from_pdf

    src = tmp_path / "vector.pdf"
    doc = fitz.open()
    doc.new_page(width=453.5, height=212.6)           # a 16 cm manuscript figure
    doc.save(src)
    doc.close()

    info = slide_png_from_pdf(str(src), str(tmp_path / "slide.png"),
                              paste_width_pt=500.0)
    assert info["text_scale"] == round(500.0 / 453.5, 3)

    deck = _deck(tmp_path, Path(info["png"]), width_in=500.0 / 72.0)
    rep = audit_pptx_figures(str(deck))

    assert [r for r in rep["pictures"][0]["risks"]
            if "RASTER FIGURE" not in r] == []
    # integer pixel rounding leaves a sub-0.1% residue, not a paste-scale error
    assert abs(rep["pictures"][0]["scale"] - 1.0) < 0.002
    assert abs(rep["pictures"][0]["effective_dpi"] - 300.0) < 1.0


def test_svg_length_uses_css_pixel_to_point_conversion() -> None:
    assert _svg_length_pt("96") == 72.0
    assert _svg_length_pt("2.54cm") == pytest.approx(72.0)


def test_svg_asset_follows_svgblip_without_png_fallback() -> None:
    svg_blob = (
        b'<svg xmlns="http://www.w3.org/2000/svg" width="960" '
        b'height="480" viewBox="0 0 720 360"/>'
    )
    element = ET.fromstring(
        '<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main">'
        '<asvg:svgBlip r:embed="rId9"/></p:pic>'
    )

    class _SvgPart:
        blob = svg_blob

    class _SlidePart:
        @staticmethod
        def related_part(relationship_id: str):
            assert relationship_id == "rId9"
            return _SvgPart()

    class _Shape:
        _element = element
        part = _SlidePart()

    asset = _svg_asset(_Shape())

    assert asset is not None
    assert asset["source_type"] == "svg"
    assert asset["width_pt"] == 720.0
    assert asset["height_pt"] == 360.0
    assert asset["pixels"] is None
