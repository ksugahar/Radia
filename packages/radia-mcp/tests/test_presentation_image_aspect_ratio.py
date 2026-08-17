from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from radia_mcp.presentation.tools import presentation_check_image_aspect_ratio


def _source_image(tmp_path: Path) -> Path:
    path = tmp_path / "source-2-to-1.png"
    Image.new("RGB", (400, 200), "white").save(path)
    return path


def test_picture_scaled_uniformly_passes(tmp_path: Path) -> None:
    image_path = _source_image(tmp_path)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), width=Inches(6))
    path = tmp_path / "uniform.pptx"
    prs.save(path)

    result = presentation_check_image_aspect_ratio(str(path))

    assert result["passed"] is True
    assert result["images_checked"] == 1
    assert result["violation_count"] == 0


def test_picture_stretched_non_uniformly_is_a_violation(tmp_path: Path) -> None:
    image_path = _source_image(tmp_path)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(1), Inches(6), Inches(6)
    )
    path = tmp_path / "stretched.pptx"
    prs.save(path)

    result = presentation_check_image_aspect_ratio(str(path))

    assert result["passed"] is False
    assert result["violation_count"] == 1
    assert result["violations"][0]["issue"] == "image_aspect_ratio_changed"


def test_crop_that_preserves_visible_source_ratio_passes(tmp_path: Path) -> None:
    image_path = _source_image(tmp_path)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    picture = slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(1), Inches(4), Inches(4)
    )
    picture.crop_left = 0.25
    picture.crop_right = 0.25
    path = tmp_path / "cropped.pptx"
    prs.save(path)

    result = presentation_check_image_aspect_ratio(str(path))

    assert result["passed"] is True
    assert result["violation_count"] == 0
