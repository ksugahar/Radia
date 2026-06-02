"""Tool functions for the doc_convert sub-skill.

Document format conversion via Office / Adobe COM:
    - PowerPoint .ppt/.pptx -> .pdf  (PowerPoint COM)
    - PDF -> .jpg                    (Adobe Acrobat COM)
plus slide extraction + lint. Windows-only for COM-based tools; the
``_load_pptx_com`` / ``_load_adobe_com`` helpers return clean errors on
other platforms. ``doc_convert_extract_slide`` works cross-platform via
python-pptx but the resulting file's chart / SmartArt fidelity depends
on python-pptx's support for those shape types.

Tools (prefix ``doc_convert_``):
    Tier 1 — Conversion (COM):
        doc_convert_pptx_to_pdf(pptx_path, output_dir=None)
        doc_convert_pptx_to_pdf_dir(directory, recursive=False)
        doc_convert_pdf_to_jpg(pdf_path, output_dir=None)
        doc_convert_pdf_to_jpg_dir(directory, recursive=False)

    Tier 2 — Extraction:
        doc_convert_extract_slide(src_pptx, slide_index, out_pptx)
        doc_convert_extract_slides(src_pptx, mapping)
        doc_convert_inventory(pptx_path)
        doc_convert_slide_to_image(pptx_path, slide_index, out_path, fmt="PNG")
        doc_convert_extract_speaker_notes(pptx_path, out_path=None)
        doc_convert_business_card_ocr_dir(directory, recursive=False)
        doc_convert_rename_business_card(jpg_path, company, name, dry_run=False)
        doc_convert_extract_pdf_title_by_cover(pdf_path, render_dpi=144)
        doc_convert_extract_title_from_image(image_path, timeout_s=180)
        doc_convert_rename_pdf_by_title(pdf_path, title, output_dir=None)

    Tier 3 — Lint:
        doc_convert_slide_lint(pptx_path)
        doc_convert_hidden_audit(pptx_path)

    Tier 5 — Intelligence:
        doc_convert_health_report(pptx_path)
"""
from __future__ import annotations

import os
import pathlib
import re
import subprocess
import sys
import time


# ---------------------------------------------------------------------------
# Shared: PDF lock release (mirror of poster/_lock release)
# ---------------------------------------------------------------------------
_PDF_VIEWER_PROCESSES = (
    "Acrobat", "AcroRd32", "SumatraPDF", "FoxitReader", "FoxitPDFReader",
    "msedge", "chrome", "firefox",
)


def _release_pdf_lock(pdf_path: pathlib.Path) -> list[str]:
    """Kill PDF viewers holding ``pdf_path``. Windows-only."""
    killed: list[str] = []
    if not pdf_path.exists() or sys.platform != "win32":
        return killed
    try:
        with open(pdf_path, "r+b"):
            return killed
    except PermissionError:
        pass
    except OSError:
        return killed
    basename = pdf_path.name
    ps_script = (
        "Get-Process | Where-Object { $_.MainWindowTitle -ne '' } | "
        f"Where-Object {{ $_.MainWindowTitle -match '{re.escape(basename)}' "
        "-or $_.ProcessName -in @("
        + ",".join(f"'{p}'" for p in _PDF_VIEWER_PROCESSES)
        + ") } | Select-Object -ExpandProperty Id"
    )
    try:
        out = subprocess.run(
            ["powershell", "-NoProfile", "-Command", ps_script],
            capture_output=True, text=True, timeout=15,
        )
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return killed
    for pid in (out.stdout or "").split():
        pid = pid.strip()
        if not pid.isdigit():
            continue
        try:
            subprocess.run(
                ["powershell", "-NoProfile", "-Command", f"Stop-Process -Id {pid} -Force"],
                capture_output=True, text=True, timeout=10,
            )
            killed.append(pid)
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass
    return killed


# ---------------------------------------------------------------------------
# COM-based conversion (Windows-only)
# ---------------------------------------------------------------------------
def _load_pptx_com():
    """Import win32com bits or return None when unavailable."""
    if sys.platform != "win32":
        return None
    try:
        import win32com.client  # type: ignore
        import winerror  # type: ignore
        from win32com.client.dynamic import ERRORS_BAD_CONTEXT  # type: ignore
        if winerror.E_NOTIMPL not in ERRORS_BAD_CONTEXT:
            ERRORS_BAD_CONTEXT.append(winerror.E_NOTIMPL)
        return win32com.client
    except ImportError:
        return None


# PowerPoint SaveAs format codes; see PpSaveAsFileType enumeration.
_PP_SAVE_AS_PDF = 32


def _convert_one(ppt_app, pptx_path: pathlib.Path,
                 out_path: pathlib.Path) -> tuple[bool, str]:
    """Open ``pptx_path`` in the given Application, SaveAs PDF, close.

    Returns ``(ok, message)``. The caller owns the COM Application lifecycle.
    """
    presentation = None
    try:
        # Match the original doc_convert.py: no extra arguments. Some
        # PowerPoint versions reject WithWindow / ReadOnly with E_FAIL.
        presentation = ppt_app.Presentations.Open(str(pptx_path))
    except Exception as e:
        return False, f"Open failed: {type(e).__name__}: {e}"
    try:
        presentation.SaveAs(str(out_path), _PP_SAVE_AS_PDF)
    except Exception as e:
        return False, f"SaveAs failed: {type(e).__name__}: {e}"
    finally:
        try:
            presentation.Close()
        except Exception:
            pass
    if not out_path.exists():
        return False, f"SaveAs reported success but {out_path} not found"
    return True, f"wrote {out_path} ({out_path.stat().st_size / 1024:.1f} KB)"


def doc_convert_pptx_to_pdf(pptx_path: str,
                    output_dir: str | None = None) -> str:
    """Convert a single .ppt/.pptx to .pdf via PowerPoint COM.

    PowerPoint must be installed and licensed on the machine. The
    Application is launched in invisible mode (``WithWindow=False``).

    Args:
        pptx_path: Path to the .ppt/.pptx source.
        output_dir: Optional directory for the PDF. Defaults to the same
                    directory as the input. The output filename is
                    ``<stem>.pdf``.

    Returns:
        Status string with the output path + size, or an error message.
    """
    com = _load_pptx_com()
    if com is None:
        return ("Error: PowerPoint COM unavailable. doc_convert_pptx_to_pdf requires "
                "Windows + pywin32 + an installed PowerPoint.")
    src = pathlib.Path(pptx_path)
    if not src.is_absolute():
        src = pathlib.Path.cwd() / src
    if not src.exists():
        return f"Error: file not found: {src}"
    if src.suffix.lower() not in (".ppt", ".pptx", ".pptm"):
        return f"Error: not a PowerPoint file: {src}"

    out_dir = pathlib.Path(output_dir) if output_dir else src.parent
    if not out_dir.is_absolute():
        out_dir = pathlib.Path.cwd() / out_dir
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"{src.stem}.pdf"
    _release_pdf_lock(out_path)

    app = com.Dispatch("PowerPoint.Application")
    try:
        ok, msg = _convert_one(app, src, out_path)
    finally:
        try:
            app.Quit()
        except Exception:
            pass
    return f"doc_convert_pptx_to_pdf: {src} -> {msg}"


def doc_convert_pptx_to_pdf_dir(directory: str,
                        recursive: bool = False,
                        output_dir: str | None = None) -> str:
    """Convert every .ppt/.pptx in ``directory`` to PDF.

    Mirrors the original ``doc_convert.py`` behaviour but with explicit
    directory targeting (no implicit CWD scan).

    Args:
        directory: Folder to scan.
        recursive: If True, descend into subdirectories.
        output_dir: Optional unified output directory. Defaults to
                    each PPT file's own directory.

    Returns:
        Per-file status + summary count.
    """
    com = _load_pptx_com()
    if com is None:
        return ("Error: PowerPoint COM unavailable. doc_convert_pptx_to_pdf_dir "
                "requires Windows + pywin32 + an installed PowerPoint.")
    root = pathlib.Path(directory)
    if not root.is_absolute():
        root = pathlib.Path.cwd() / root
    if not root.is_dir():
        return f"Error: not a directory: {root}"

    pattern = "**/*" if recursive else "*"
    targets = [p for p in root.glob(pattern)
               if p.is_file() and p.suffix.lower() in (".ppt", ".pptx", ".pptm")]
    if not targets:
        return f"doc_convert_pptx_to_pdf_dir: no PowerPoint files in {root}"

    app = com.Dispatch("PowerPoint.Application")
    out_root = pathlib.Path(output_dir) if output_dir else None
    if out_root and not out_root.is_absolute():
        out_root = pathlib.Path.cwd() / out_root
    if out_root:
        out_root.mkdir(parents=True, exist_ok=True)

    results: list[str] = []
    ok_count = 0
    try:
        for src in targets:
            out_dir = out_root or src.parent
            out_path = out_dir / f"{src.stem}.pdf"
            _release_pdf_lock(out_path)
            ok, msg = _convert_one(app, src, out_path)
            if ok:
                ok_count += 1
            results.append(f"  - {src.name}: {msg}")
    finally:
        try:
            app.Quit()
        except Exception:
            pass

    head = f"doc_convert_pptx_to_pdf_dir: {ok_count}/{len(targets)} converted in {root}"
    return "\n".join([head, *results])


# ---------------------------------------------------------------------------
# Slide extraction (python-pptx, cross-platform)
# ---------------------------------------------------------------------------
# PowerPoint SaveAs file format codes.
_PP_SAVE_AS_PPTX = 24  # ppSaveAsOpenXMLPresentation


def _copy_for_single_slide_com(app, src: pathlib.Path, dst: pathlib.Path,
                               keep_index: int) -> tuple[bool, str]:
    """COM-based single-slide extraction (Windows + PowerPoint).

    Open source, delete every slide except ``keep_index`` (1-based,
    iterating in reverse so indices remain stable), SaveAs pptx.
    """
    presentation = None
    try:
        presentation = app.Presentations.Open(str(src))
    except Exception as e:
        return False, f"Open failed: {type(e).__name__}: {e}"
    try:
        n = presentation.Slides.Count
        if not (1 <= keep_index <= n):
            return False, f"keep_index {keep_index} out of range [1, {n}]"
        # Delete in reverse to keep indices stable.
        for i in range(n, 0, -1):
            if i != keep_index:
                presentation.Slides(i).Delete()
        presentation.SaveAs(str(dst), _PP_SAVE_AS_PPTX)
        return True, f"kept slide {keep_index} of {n}"
    except Exception as e:
        return False, f"Extract failed: {type(e).__name__}: {e}"
    finally:
        try:
            presentation.Close()
        except Exception:
            pass


def _copy_for_single_slide_pptx(src: pathlib.Path, dst: pathlib.Path,
                                keep_index: int) -> tuple[bool, str]:
    """Cross-platform python-pptx fallback. Brittle for fancy slides.

    Note: PowerPoint COM is preferred (used when available). This
    fallback relies on python-pptx semi-private APIs and may produce
    pptx files that PowerPoint refuses to re-open if the source uses
    SmartArt, complex relationships, or media.
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return False, ("python-pptx not installed. Install with "
                       "`pip install -e .[pptx]`.")
    if not (1 <= keep_index <= 999):
        return False, f"keep_index out of range: {keep_index}"
    import shutil
    shutil.copy2(src, dst)
    prs = Presentation(dst)
    n = len(prs.slides)
    if keep_index > n:
        return False, f"keep_index {keep_index} exceeds slide count {n}"
    keep_zero = keep_index - 1
    sldIdLst = prs.slides._sldIdLst  # noqa: SLF001
    id_elements = list(sldIdLst)
    for i, sldId in enumerate(id_elements):
        if i == keep_zero:
            continue
        rId = sldId.attrib.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(sldId)
    prs.save(dst)
    return True, f"kept slide {keep_index} of {n} (python-pptx fallback)"


def _extract_with_best_backend(src: pathlib.Path,
                               mapping: dict[int, pathlib.Path]) -> list[tuple[int, pathlib.Path, bool, str]]:
    """Try COM extraction first, fall back to python-pptx per-file.

    Returns a list of ``(idx, dst, ok, msg)`` tuples.
    """
    results: list[tuple[int, pathlib.Path, bool, str]] = []
    com = _load_pptx_com()
    if com is not None:
        app = com.Dispatch("PowerPoint.Application")
        try:
            for idx, dst in mapping.items():
                ok, msg = _copy_for_single_slide_com(app, src, dst, idx)
                results.append((idx, dst, ok, msg))
        finally:
            try:
                app.Quit()
            except Exception:
                pass
        return results
    for idx, dst in mapping.items():
        ok, msg = _copy_for_single_slide_pptx(src, dst, idx)
        results.append((idx, dst, ok, msg))
    return results


def doc_convert_extract_slide(src_pptx: str,
                          slide_index: int,
                          out_pptx: str) -> str:
    """Write a single slide of ``src_pptx`` as its own .pptx file.

    Uses PowerPoint COM when available (round-trip-safe), otherwise
    falls back to python-pptx (best-effort).

    Args:
        src_pptx: Source PowerPoint file.
        slide_index: 1-based slide index to keep.
        out_pptx: Destination .pptx path. Parent dir is created if needed.

    Returns:
        Status string.
    """
    src = pathlib.Path(src_pptx)
    if not src.is_absolute():
        src = pathlib.Path.cwd() / src
    if not src.exists():
        return f"Error: file not found: {src}"
    out = pathlib.Path(out_pptx)
    if not out.is_absolute():
        out = pathlib.Path.cwd() / out
    out.parent.mkdir(parents=True, exist_ok=True)

    results = _extract_with_best_backend(src, {slide_index: out})
    _, dst, ok, msg = results[0]
    if not ok:
        if dst.exists():
            try:
                dst.unlink()
            except OSError:
                pass
        return f"doc_convert_extract_slide: FAIL — {msg}"
    return f"doc_convert_extract_slide: wrote {dst} — {msg}"


def doc_convert_extract_slides(src_pptx: str,
                           mapping: dict[int, str]) -> str:
    """Bulk variant: write multiple single-slide files in one call.

    Args:
        src_pptx: Source PowerPoint file.
        mapping: ``{slide_index: out_pptx_path, ...}`` (1-based indices).

    Returns:
        Per-mapping status + summary count.
    """
    src = pathlib.Path(src_pptx)
    if not src.is_absolute():
        src = pathlib.Path.cwd() / src
    if not src.exists():
        return f"Error: file not found: {src}"

    resolved: dict[int, pathlib.Path] = {}
    for idx, out_str in mapping.items():
        out = pathlib.Path(out_str)
        if not out.is_absolute():
            out = pathlib.Path.cwd() / out
        out.parent.mkdir(parents=True, exist_ok=True)
        resolved[idx] = out

    results = _extract_with_best_backend(src, resolved)

    lines = [f"doc_convert_extract_slides: {src}"]
    ok_count = 0
    for idx, dst, ok, msg in results:
        if ok:
            ok_count += 1
            lines.append(f"  [OK] slide {idx} -> {dst.name} ({msg})")
        else:
            if dst.exists():
                try:
                    dst.unlink()
                except OSError:
                    pass
            lines.append(f"  [FAIL] slide {idx}: {msg}")
    lines.append(f"summary: {ok_count}/{len(mapping)} extracted")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Tier 2 expansion — Inventory + slide-to-image + speaker notes
# ---------------------------------------------------------------------------
from .plans.T5_inventory import doc_convert_inventory  # noqa: E402,F401
from .plans.T8_slide_to_image import doc_convert_slide_to_image  # noqa: E402,F401
from .plans.T9_extract_speaker_notes import doc_convert_extract_speaker_notes  # noqa: E402,F401

# ---------------------------------------------------------------------------
# Tier 3 — Lint
# ---------------------------------------------------------------------------
from .plans.T6_slide_lint import doc_convert_slide_lint  # noqa: E402,F401
from .plans.T7_hidden_audit import doc_convert_hidden_audit  # noqa: E402,F401

# ---------------------------------------------------------------------------
# Tier 5 — Intelligence
# ---------------------------------------------------------------------------
from .plans.T10_health_report import doc_convert_health_report  # noqa: E402,F401
from .plans.T11_figure_readability import doc_convert_figure_readability_check  # noqa: E402,F401
from .plans.T12_classify import doc_convert_classify  # noqa: E402,F401
from .plans.T13_resize_for_print import doc_convert_resize_for_print  # noqa: E402,F401

# ---------------------------------------------------------------------------
# Tier 1 — Adobe Acrobat COM: PDF -> JPG
# ---------------------------------------------------------------------------
from .plans.T14_pdf_to_jpg import (  # noqa: E402,F401
    doc_convert_pdf_to_jpg,
    doc_convert_pdf_to_jpg_dir,
)

# ---------------------------------------------------------------------------
# Tier 2 — Business card OCR + rename (calls EasyOCR via ocr.core)
# ---------------------------------------------------------------------------
from .plans.T15_business_card_rename import (  # noqa: E402,F401
    doc_convert_business_card_ocr_dir,
    doc_convert_rename_business_card,
)

# ---------------------------------------------------------------------------
# Tier 2 — Book-cover title extract + PDF rename (Claude CLI + Win-safe FS op)
# ---------------------------------------------------------------------------
from .plans.T16_book_pdf_rename import (  # noqa: E402,F401
    doc_convert_extract_pdf_title_by_cover,
    doc_convert_extract_title_from_image,
    doc_convert_rename_pdf_by_title,
)
