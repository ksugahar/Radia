"""Tier 2 — doc_convert_inventory.

List every slide with its title, body word/character count, hidden
flag, animation count, and notes presence. Useful for picking which
slides to extract.
"""
from __future__ import annotations

import pathlib


def _load_pptx():
    try:
        from pptx import Presentation  # type: ignore
        return Presentation
    except ImportError:
        return None


def doc_convert_inventory(pptx_path: str) -> str:
    """Return a per-slide inventory of a .pptx.

    Args:
        pptx_path: Path to .ppt/.pptx.

    Returns:
        Tabular per-slide summary: ``# | hidden | title | body_chars | notes_chars``
    """
    Pres = _load_pptx()
    if Pres is None:
        return "Error: python-pptx not installed. `pip install -e .[pptx]`."
    p = pathlib.Path(pptx_path)
    if not p.is_absolute():
        p = pathlib.Path.cwd() / p
    if not p.exists():
        return f"Error: file not found: {p}"
    prs = Pres(p)
    n = len(prs.slides)
    lines = [f"doc_convert_inventory: {p}",
             f"  slides: {n}",
             "  -" * 30]
    for i, slide in enumerate(prs.slides, 1):
        hidden = "H" if (slide.element.attrib.get("show", "1") == "0") else " "
        # Title
        title = ""
        body_chars = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                txt = "".join(r.text for r in para.runs)
                if not title and txt.strip():
                    title = txt.strip()
                else:
                    body_chars += len(txt)
        notes_chars = 0
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            notes_chars = len(notes)
        lines.append(f"  {i:3d} [{hidden}] {title[:55]:55s}  body={body_chars:4d}  notes={notes_chars:4d}")
    return "\n".join(lines)
