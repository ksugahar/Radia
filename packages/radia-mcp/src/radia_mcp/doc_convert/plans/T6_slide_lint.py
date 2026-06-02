"""Tier 3 — doc_convert_slide_lint.

Per-slide quality lint for *presentation* pptx (as opposed to figure
pptx). Rules:
    - HIGH: slide has no title (screen-reader breaks)
    - HIGH: body text exceeds 100 characters (CSU/WCAG threshold)
    - HIGH: font size on any text run < 20pt (lab rule: presentations
            should use the largest font possible, never < 20pt)
    - MEDIUM: image without alt text (accessibility)
    - LOW: more than 2 fonts in use across the slide

For *figure* pptx (to be embedded in a paper/poster), use
``doc_convert_figure_readability_check`` instead — that tool scales source
pt by the embed width.
"""
from __future__ import annotations

import pathlib


def _load_pptx():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Pt  # type: ignore
        return Presentation, Pt
    except ImportError:
        return None, None


def doc_convert_slide_lint(pptx_path: str) -> str:
    """Run accessibility + density lints on every slide.

    Severity thresholds align with WCAG 2.1 AA + Texas Tech / Harvard
    accessibility guidelines.

    Args:
        pptx_path: Path to .ppt/.pptx.

    Returns:
        Per-slide issue list + summary count.
    """
    Pres, Pt = _load_pptx()
    if Pres is None:
        return "Error: python-pptx not installed."
    p = pathlib.Path(pptx_path)
    if not p.is_absolute():
        p = pathlib.Path.cwd() / p
    if not p.exists():
        return f"Error: file not found: {p}"

    prs = Pres(p)
    issues: list[tuple[int, str, str]] = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        body_chars = 0
        fonts: set[str] = set()
        min_pt = 999
        images_no_alt = 0
        images_with_alt = 0
        for shape in slide.shapes:
            try:
                # Image check
                if shape.shape_type == 13:  # PICTURE
                    name = shape.name or ""
                    alt = (getattr(shape, "alt_text", None)
                           or shape.element.nvSpPr.cNvPr.get("descr", "") if hasattr(shape.element, "nvSpPr") else "")
                    if not alt:
                        images_no_alt += 1
                    else:
                        images_with_alt += 1
            except Exception:
                pass

            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    txt = run.text
                    if run.font.name:
                        fonts.add(run.font.name)
                    if run.font.size is not None:
                        try:
                            pt_val = run.font.size.pt
                            if pt_val < min_pt:
                                min_pt = pt_val
                        except Exception:
                            pass
                    body_chars += len(txt)
                if not title and "".join(r.text for r in para.runs).strip():
                    title = "".join(r.text for r in para.runs).strip()

        if not title:
            issues.append((i, "HIGH", "no slide title (screen reader skips, outline empty)"))
        if body_chars > 100:
            issues.append((i, "HIGH", f"body text {body_chars} chars > 100 (cognitive overload + WCAG)"))
        if min_pt < 20 and min_pt != 999:
            issues.append((i, "HIGH", f"smallest font {min_pt}pt < 20pt (lab rule: presentations never use <20pt)"))
        if images_no_alt > 0:
            issues.append((i, "MEDIUM", f"{images_no_alt} image(s) without alt text (accessibility)"))
        if len(fonts) > 2:
            issues.append((i, "LOW", f"{len(fonts)} fonts on one slide: {sorted(fonts)}"))

    lines = [f"doc_convert_slide_lint: {p}", f"  slides: {len(prs.slides)}"]
    if not issues:
        lines.append("PASS — no slide-level issues detected")
        return "\n".join(lines)
    sev_order = {"HIGH": 0, "MEDIUM": 1, "LOW": 2}
    issues.sort(key=lambda x: (sev_order[x[1]], x[0]))
    lines.append(f"FAIL — {len(issues)} issue(s):")
    for idx, sev, msg in issues:
        lines.append(f"  [{sev}] slide {idx}: {msg}")
    return "\n".join(lines)
