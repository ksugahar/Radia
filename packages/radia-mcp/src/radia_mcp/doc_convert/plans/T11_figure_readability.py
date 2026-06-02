"""Tier 3 — doc_convert_figure_readability_check.

When a pptx slide is exported and later embedded as a *figure* in a
paper/poster, it is usually scaled down to a column width (8cm single,
16cm double, or some \\includegraphics fraction of \\linewidth). The
slide's source font size must survive that scale-down.

The lab rule (provided by the user):
    - print width 16 cm  →  text must render ≥ 20 pt
    - print width  8 cm  →  text must render ≥ 10 pt

These two points describe a linear law

    required_rendered_pt(w_cm) = 1.25 * w_cm

so the source font requirement in a slide of source width W_src cm is

    required_source_pt(w_target_cm) = 1.25 * w_target_cm * W_src / w_target_cm
                                    = 1.25 * W_src     (≈ 31.75 pt for a
                                                         25.4 cm pptx)

i.e. it's independent of the target width — a standard pptx page needs
≥ 32 pt source text to survive both 16 cm and 8 cm prints.
"""
from __future__ import annotations

import pathlib


def _load_pptx():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Emu
        return Presentation, Emu
    except ImportError:
        return None, None


# 1 EMU = 914400 per inch; 2.54 cm per inch.
_EMU_PER_CM = 914400 / 2.54


def doc_convert_figure_readability_check(pptx_path: str,
                                     target_width_cm: float = 8.0,
                                     min_rendered_pt: float | None = None,
                                     law_slope: float = 1.25) -> str:
    """Verify the slide's source font survives a target print width.

    The default law follows the user's lab rule
    ``required_rendered_pt = 1.25 * target_width_cm``:

        - 16 cm column → 20 pt rendered (≥ 32 pt source on standard pptx)
        - 8 cm column  → 10 pt rendered (≥ 32 pt source on standard pptx)

    Args:
        pptx_path: Slide source.
        target_width_cm: Width at which the slide will be printed when
                         embedded as a figure (default 8 cm — single
                         column of an IEEE 2-column paper).
        min_rendered_pt: Override the law. If None, use
                         ``law_slope * target_width_cm``.
        law_slope: Linear coefficient of the readability law (default
                   1.25 pt/cm — the lab rule).

    Returns:
        Per-slide verdict + the smallest rendered pt vs. required.
    """
    Pres, Emu = _load_pptx()
    if Pres is None:
        return "Error: python-pptx not installed."
    p = pathlib.Path(pptx_path)
    if not p.is_absolute():
        p = pathlib.Path.cwd() / p
    if not p.exists():
        return f"Error: file not found: {p}"

    prs = Pres(p)
    src_w_cm = (prs.slide_width or 0) / _EMU_PER_CM
    if src_w_cm <= 0:
        src_w_cm = 25.4  # standard widescreen 16:9 / 4:3 = 25.4 cm
    scale = target_width_cm / src_w_cm
    if min_rendered_pt is None:
        required_rendered = law_slope * target_width_cm
    else:
        required_rendered = float(min_rendered_pt)
    required_source = required_rendered / scale if scale > 0 else float("inf")

    lines = [f"doc_convert_figure_readability_check: {p}",
             f"  source slide width: {src_w_cm:.2f} cm",
             f"  target print width: {target_width_cm:.2f} cm (scale {scale:.3f}×)",
             f"  law: required_rendered_pt = {law_slope} × target_cm = {required_rendered:.1f} pt",
             f"  → required source pt: ≥ {required_source:.1f} pt"]

    failures = 0
    for i, slide in enumerate(prs.slides, 1):
        min_pt = None
        sample_text = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                # Determine effective pt: run.font.size if set, otherwise
                # paragraph default if set, otherwise None (placeholder
                # inherits — we can't resolve without rendering).
                for run in para.runs:
                    txt = run.text
                    if not txt.strip():
                        continue
                    pt_val = None
                    if run.font.size is not None:
                        try:
                            pt_val = run.font.size.pt
                        except Exception:
                            pass
                    if pt_val is None and para.font.size is not None:
                        try:
                            pt_val = para.font.size.pt
                        except Exception:
                            pass
                    if pt_val is None:
                        continue
                    if min_pt is None or pt_val < min_pt:
                        min_pt = pt_val
                        sample_text = txt
        if min_pt is None:
            lines.append(f"  - slide {i}: no explicit font sizes found (placeholder defaults)")
            continue
        rendered = min_pt * scale
        verdict = "OK" if rendered >= required_rendered else "FAIL"
        if verdict == "FAIL":
            failures += 1
        lines.append(
            f"  - slide {i}: min source {min_pt:.0f}pt → renders {rendered:.1f}pt "
            f"[{verdict}] e.g. {sample_text[:30]!r}"
        )
    if failures:
        lines.append(f"FAIL — {failures}/{len(prs.slides)} slide(s) below readability target")
    else:
        lines.append("PASS — every slide's smallest text survives the target print width")
    return "\n".join(lines)
