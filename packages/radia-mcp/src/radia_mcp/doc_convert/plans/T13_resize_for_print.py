"""Tier 2 — doc_convert_resize_for_print.

Resize a figure pptx to a target embed width and enforce a minimum
font size. The intended workflow:

    1. Author the figure at a comfortable canvas size (default pptx
       25.4 × 19.05 cm).
    2. Run this tool to shrink the canvas to the *target embed width*
       (e.g. 16 cm — the typical 2-column journal full-width figure).
    3. All shapes scale proportionally (positions, sizes, font sizes).
    4. Any font that drops below ``min_font_pt`` after scaling is
       bumped back up.
    5. Convert with ``doc_convert_pptx_to_pdf``; the resulting PDF is "natively
       sized" — embed at the same width with no further scaling.

The benefit: ``rendered_pt == source_pt`` at the target width, so the
lab readability rule (≥20pt at 16cm, ≥10pt at 8cm) reduces to
``min_font_pt ≥ 20``.
"""
from __future__ import annotations

import pathlib


def _load_pptx():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Cm, Pt  # type: ignore
        return Presentation, Cm, Pt
    except ImportError:
        return None, None, None


def doc_convert_resize_for_print(pptx_path: str,
                             target_width_cm: float = 16.0,
                             target_font_pt: float = 20.0,
                             out_path: str | None = None) -> str:
    """Shrink a pptx to ``target_width_cm`` and set ALL fonts to ``target_font_pt``.

    The lab convention for figure pptx:
        - Canvas width = target embed width (so source pt == rendered pt)
        - Every text run = the same target_font_pt (uniform, no hierarchy)

    This matches the "20pt straight across" rule: a figure printed at
    16cm with all-20pt text renders at exactly 20pt — readable per the
    lab rule, and clean per visual consistency.

    Args:
        pptx_path: Source pptx.
        target_width_cm: New canvas width in cm (height scales proportionally).
        target_font_pt: All runs with an explicit font size are set to this.
        out_path: Optional output path. If None, overwrite the source.

    Returns:
        Status string with the before/after dimensions and font count.
    """
    Pres, Cm, Pt = _load_pptx()
    if Pres is None:
        return "Error: python-pptx not installed."
    src = pathlib.Path(pptx_path)
    if not src.is_absolute():
        src = pathlib.Path.cwd() / src
    if not src.exists():
        return f"Error: file not found: {src}"

    prs = Pres(src)
    old_w = prs.slide_width
    old_h = prs.slide_height
    new_w = Cm(target_width_cm)
    scale = new_w / old_w
    new_h = int(old_h * scale)

    prs.slide_width = new_w
    prs.slide_height = new_h

    n_shapes_scaled = 0
    n_fonts_set = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            # Scale shape geometry. Some shapes (placeholders inheriting
            # from the layout) don't have left/top — protect with try.
            try:
                if shape.left is not None:
                    shape.left = int(shape.left * scale)
                if shape.top is not None:
                    shape.top = int(shape.top * scale)
                if shape.width is not None:
                    shape.width = int(shape.width * scale)
                if shape.height is not None:
                    shape.height = int(shape.height * scale)
                n_shapes_scaled += 1
            except Exception:
                pass

            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is None:
                        # No explicit size set — inherits from layout.
                        # Set explicitly so the rendered output is uniform.
                        run.font.size = Pt(target_font_pt)
                        n_fonts_set += 1
                        continue
                    run.font.size = Pt(target_font_pt)
                    n_fonts_set += 1

    out = pathlib.Path(out_path) if out_path else src
    if not out.is_absolute():
        out = pathlib.Path.cwd() / out
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)

    return (f"doc_convert_resize_for_print: {src}\n"
            f"  canvas: {old_w/360000:.1f}×{old_h/360000:.1f} cm "
            f"→ {target_width_cm:.1f}×{new_h/360000:.1f} cm (scale {scale:.3f}×)\n"
            f"  shapes scaled: {n_shapes_scaled}\n"
            f"  fonts set to {target_font_pt}pt: {n_fonts_set}\n"
            f"  saved: {out}")
