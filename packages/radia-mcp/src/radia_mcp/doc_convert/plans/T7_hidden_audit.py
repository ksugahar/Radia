"""Tier 3 — doc_convert_hidden_audit.

List hidden slides + presenter notes. These are normally invisible
during a presentation but DO appear in PDF SaveAs unless excluded.
Helpful before publishing to a public repository.
"""
from __future__ import annotations

import pathlib


def _load_pptx():
    try:
        from pptx import Presentation  # type: ignore
        return Presentation
    except ImportError:
        return None


def doc_convert_hidden_audit(pptx_path: str) -> str:
    """Audit hidden slides and presenter notes for accidental leak risk.

    Args:
        pptx_path: Path to .ppt/.pptx.

    Returns:
        List of hidden slide indices + slides with non-empty notes.
    """
    Pres = _load_pptx()
    if Pres is None:
        return "Error: python-pptx not installed."
    p = pathlib.Path(pptx_path)
    if not p.is_absolute():
        p = pathlib.Path.cwd() / p
    if not p.exists():
        return f"Error: file not found: {p}"
    prs = Pres(p)
    hidden: list[tuple[int, str]] = []
    has_notes: list[tuple[int, str, int]] = []
    for i, slide in enumerate(prs.slides, 1):
        if slide.element.attrib.get("show", "1") == "0":
            title = next(
                (("".join(r.text for r in para.runs).strip())
                 for shape in slide.shapes if shape.has_text_frame
                 for para in shape.text_frame.paragraphs
                 if "".join(r.text for r in para.runs).strip()),
                "(no title)",
            )
            hidden.append((i, title))
        if slide.has_notes_slide:
            txt = slide.notes_slide.notes_text_frame.text.strip()
            if len(txt) > 20:  # ignore trivial whitespace
                title = next(
                    (("".join(r.text for r in para.runs).strip())
                     for shape in slide.shapes if shape.has_text_frame
                     for para in shape.text_frame.paragraphs
                     if "".join(r.text for r in para.runs).strip()),
                    "(no title)",
                )
                has_notes.append((i, title, len(txt)))

    lines = [f"doc_convert_hidden_audit: {p}",
             f"  slides: {len(prs.slides)}"]
    if hidden:
        lines.append(f"[MEDIUM] {len(hidden)} hidden slide(s):")
        for idx, title in hidden:
            lines.append(f"  - slide {idx}: {title[:60]}")
    if has_notes:
        lines.append(f"[LOW] {len(has_notes)} slide(s) with presenter notes:")
        for idx, title, n in has_notes:
            lines.append(f"  - slide {idx} ({n} chars): {title[:60]}")
    if not hidden and not has_notes:
        lines.append("PASS — no hidden slides, no presenter notes")
    return "\n".join(lines)
