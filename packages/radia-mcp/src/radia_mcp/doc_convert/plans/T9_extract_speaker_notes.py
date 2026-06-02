"""Tier 2 — doc_convert_extract_speaker_notes.

Extract presenter notes from every slide as a plain-text file. Useful
for converting a slide deck into a speaking script.
"""
from __future__ import annotations

import pathlib


def _load_pptx():
    try:
        from pptx import Presentation  # type: ignore
        return Presentation
    except ImportError:
        return None


def doc_convert_extract_speaker_notes(pptx_path: str,
                                  out_path: str | None = None) -> str:
    """Extract presenter notes per slide.

    Args:
        pptx_path: Source pptx.
        out_path: Optional output .txt file. If None, returns the text.

    Returns:
        The extracted notes text, or a status line when written to disk.
    """
    Pres = _load_pptx()
    if Pres is None:
        return "Error: python-pptx not installed."
    p = pathlib.Path(pptx_path)
    if not p.is_absolute():
        p = pathlib.Path.cwd() / p
    if not p.exists():
        return f"Error: file not found: {p}"
    prs = Pres(p)
    out_lines: list[str] = []
    n_with_notes = 0
    for i, slide in enumerate(prs.slides, 1):
        # Slide title for context
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    title = t.splitlines()[0]
                    break
        out_lines.append(f"=== Slide {i}: {title} ===")
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                n_with_notes += 1
                out_lines.append(notes)
        out_lines.append("")
    body = "\n".join(out_lines)
    if out_path is None:
        return body
    o = pathlib.Path(out_path)
    if not o.is_absolute():
        o = pathlib.Path.cwd() / o
    o.parent.mkdir(parents=True, exist_ok=True)
    o.write_text(body, encoding="utf-8")
    return f"doc_convert_extract_speaker_notes: wrote {o}  ({n_with_notes}/{len(prs.slides)} slide(s) had notes)"
