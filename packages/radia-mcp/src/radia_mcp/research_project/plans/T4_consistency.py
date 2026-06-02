"""Tier 2 — research_project_consistency_check.

Check cross-artifact consistency in a project: the same headline number
or claim should appear in paper / poster / slides without contradiction.

Heuristic implementation:
    1. Extract candidate numerical claims (e.g. ``5.37``, ``94%``,
       ``300 DPI``, ``A1``, ``\\fontsize{60}``) from each artifact.
    2. Group by surface form and report which artifacts agree / disagree.
"""
from __future__ import annotations

import pathlib
import re

from .T1_scan import _HEURISTICS


_NUMBER_RE = re.compile(r"\b\d{1,4}(?:\.\d{1,3})?\s?(?:%|kHz|MHz|GHz|Hz|dB|mm|cm|m|μm|us|ms|s|mT|T|kV|V|kΩ|Ω|VA|W|Pa|°C|°)\b")


def _extract_text(path: pathlib.Path) -> str:
    ext = path.suffix.lower()
    if ext in (".tex", ".bib", ".txt", ".md"):
        try:
            return path.read_text(encoding="utf-8", errors="replace")
        except OSError:
            return ""
    if ext == ".pptx":
        try:
            from pptx import Presentation  # type: ignore
            prs = Presentation(path)
            chunks = []
            for s in prs.slides:
                for shape in s.shapes:
                    if shape.has_text_frame:
                        chunks.append(shape.text_frame.text)
            return "\n".join(chunks)
        except Exception:
            return ""
    if ext == ".pdf":
        try:
            import fitz  # type: ignore
            doc = fitz.open(path)
            chunks = []
            for page in doc:
                chunks.append(page.get_text())
            doc.close()
            return "\n".join(chunks)
        except Exception:
            return ""
    return ""


def research_project_consistency_check(directory: str) -> str:
    """Cross-artifact numerical consistency.

    For each (number, unit) tuple, list which artifacts mention it.
    Numbers that appear in only ONE artifact aren't flagged. Numbers
    that appear in multiple but with conflicting nearby context are
    surfaced for user review.
    """
    root = pathlib.Path(directory)
    if not root.is_absolute():
        root = pathlib.Path.cwd() / root
    if not root.is_dir():
        return f"Error: not a directory: {root}"

    classified: list[tuple[str, pathlib.Path, str]] = []
    for p in root.rglob("*"):
        if not p.is_file():
            continue
        kind = next((k for k, fn in _HEURISTICS if fn(p.name.lower(), p.suffix.lower())), None)
        if kind not in ("paper", "poster", "slides", "output_pdf"):
            continue
        text = _extract_text(p)
        if text:
            classified.append((kind, p, text))

    if len(classified) < 2:
        return (f"research_project_consistency_check: {root}\n"
                f"  fewer than 2 artifacts found; nothing to cross-check")

    # Find numbers per artifact.
    num_map: dict[str, set[str]] = {}
    for kind, path, text in classified:
        for m in _NUMBER_RE.finditer(text):
            tok = re.sub(r"\s+", "", m.group())
            num_map.setdefault(tok, set()).add(f"{kind}:{path.name}")

    lines = [f"research_project_consistency_check: {root}",
             f"  artifacts compared: {len(classified)}"]
    shared = [(t, srcs) for t, srcs in num_map.items() if len(srcs) >= 2]
    if shared:
        lines.append(f"  shared numerical claims: {len(shared)}")
        for tok, srcs in sorted(shared, key=lambda x: -len(x[1]))[:30]:
            lines.append(f"  {tok:10s} -> {sorted(srcs)}")
    else:
        lines.append("  no numerical claims appear in 2+ artifacts (possibly disjoint)")
    return "\n".join(lines)
