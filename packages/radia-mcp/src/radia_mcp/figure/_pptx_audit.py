"""radia_mcp.figure._pptx_audit -- PPTX side of the "author AT the embed width" rule.

:func:`audit_tex_figures` (``_lab_api``) guards the LaTeX embed: a figure authored
at 8 cm must be pasted with ``\\includegraphics[width=8cm]`` or its on-page font
silently changes.  A slide deck has exactly the same failure mode with no LaTeX to
audit -- a picture dragged to a "nice looking" size on the slide is scaled away from
the width it was authored for, so the 24 pt authored labels land at 24*scale pt and
the 20 pt slide floor (``PRESENTATION_MIN_VISIBLE_FONT_PT``) is breached silently.

This module reconstructs the authored width from the embedded image itself
(pixels / DPI for raster images, SVG width / viewBox for vector images, minus
any PowerPoint crop) and compares it against the width the shape actually
occupies on the slide.  It needs no manifest and no OCR: it is the
automatic, file-only complement to
``presentation_check_embedded_figure_text_size`` (which converts *known* source
font sizes through the paste width).

Incident that motivated it (MMPM SA-26-069 deck, 2026-08-16): three figures
regenerated at their exact paste widths were pasted at scale 1.000, but the mode
figure -- authored 16.49 cm -- was pasted at 13.97 cm on one slide and 14.64 cm on
the next, i.e. the same artwork ran at two different scales and its 24 pt text was
displayed at ~20 pt without anything reporting it.
"""
from __future__ import annotations

import hashlib
import os

from radia_mcp.common.pptx_svg import picture_svg_blob, svg_geometry

# A figure authored for a slide carries 24 pt text (lab_figure(medium=
# 'presentation')); the displayed floor is 20 pt.  Scaling below this ratio puts
# the authored text under the floor.
_SLIDE_AUTHORED_FONT_PT = 24.0
_SLIDE_MIN_VISIBLE_FONT_PT = 20.0
_MIN_SAFE_SCALE = _SLIDE_MIN_VISIBLE_FONT_PT / _SLIDE_AUTHORED_FONT_PT  # 0.8333


def _pptx():
    try:
        import pptx  # noqa: F401
    except ImportError as exc:                       # pragma: no cover - env dependent
        raise RuntimeError(
            "python-pptx is required for the PPTX figure audit "
            "(pip install python-pptx)."
        ) from exc
    import pptx as _p
    return _p


def _svg_asset(shape) -> dict | None:
    blob = picture_svg_blob(shape)
    if blob is None:
        return None
    geometry = svg_geometry(blob)
    if geometry["width_pt"] <= 0 or geometry["height_pt"] <= 0:
        # A viewBox-only SVG has an aspect but no authored physical width, and
        # this audit's whole verdict is authored-vs-pasted width.
        raise ValueError("SVG authored width/height is unavailable")
    view_width = geometry["view_width"]
    view_height = geometry["view_height"]
    return {
        "blob": blob,
        "source_type": "svg",
        "width_pt": geometry["width_pt"],
        "height_pt": geometry["height_pt"],
        "aspect_width": view_width if view_width > 0 else geometry["width_pt"],
        "aspect_height": (
            view_height if view_height > 0 else geometry["height_pt"]
        ),
        "pixels": None,
        "dpi": None,
    }


def _picture_asset(shape) -> dict:
    """Resolve a picture's authored asset, preferring its vector original.

    A broken or dimensionless SVG must not remove the picture from the audit:
    PowerPoint stores a raster fallback beside every ``svgBlip``, so the audit
    continues against that fallback and reports the SVG failure instead of
    silently dropping the figure (the same contract as ``NO DPI METADATA``).
    """
    svg_error = ""
    try:
        svg = _svg_asset(shape)
    except Exception as exc:
        svg = None
        svg_error = f"{type(exc).__name__}: {exc}"
    if svg is not None:
        return svg
    image = shape.image
    px_w, px_h = image.size
    dpi_w, dpi_h = image.dpi
    return {
        "blob": bytes(image.blob),
        "source_type": image.ext,
        "width_pt": px_w / float(dpi_w) * 72.0,
        "height_pt": px_h / float(dpi_h) * 72.0,
        "aspect_width": float(px_w),
        "aspect_height": float(px_h),
        "pixels": (px_w, px_h),
        "dpi": (dpi_w, dpi_h),
        "svg_error": svg_error,
    }


def _iter_pictures(shapes, slide_no, out, prefix=""):
    """Collect picture shapes, descending into groups.

    Each entry is ``(slide, name, shape, asset)``. An unreadable embedded,
    linked, or OLE picture is still collected because silently dropping it can
    make a deck report zero findings while carrying unchecked figures.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _iter_pictures(shape.shapes, slide_no, out,
                           prefix=f"{prefix}{shape.name}/")
            continue
        if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            asset = _picture_asset(shape)
        except Exception as exc:
            # Linked / OLE picture, or an SVG with no usable raster fallback.
            # Report it as unreadable; never drop a figure from the audit.
            asset = {"unreadable": f"{type(exc).__name__}: {exc}"}
        out.append((slide_no, prefix + shape.name, shape, asset))
    return out


def audit_pptx_figures(pptx_path: str,
                       scale_tol: float = 0.02,
                       min_effective_dpi: float = 150.0,
                       aspect_tol: float = 0.01,
                       min_area_fraction: float = 0.01) -> dict:
    """Audit every picture on every slide for paste-scale / resolution defects.

    For each picture the authored width is ``pixels / DPI`` (with the PowerPoint
    crop removed) and the displayed width comes from the shape geometry, so
    ``scale = displayed / authored``.

    Risks reported per picture:
      * ``DOWNSCALED`` -- pasted narrower than authored: figure text is displayed
        at ``authored_pt * scale``; below ``scale=0.833`` a 24 pt label breaches
        the 20 pt slide floor.
      * ``UPSCALED`` -- a raster image pasted wider than authored
        (soft / pixelated artwork). SVG enlargement is not a quality defect.
      * ``ASPECT DISTORTED`` -- displayed aspect differs from the image aspect,
        i.e. the picture was stretched (fonts distorted, circles become ovals).
      * ``LOW EFFECTIVE DPI`` -- raster pixels per displayed inch below
        ``min_effective_dpi``. SVG has no DPI floor.
      * ``NO DPI METADATA`` -- the image carries no resolution, so the authored
        width is UNVERIFIABLE (python-pptx then reports the 72 dpi default).
        Reported, never silently passed.

    Args:
        pptx_path: deck to audit.
        scale_tol: allowed |scale - 1| before a paste-scale risk is raised.
        min_effective_dpi: floor for pixels per displayed inch.
        aspect_tol: allowed relative aspect mismatch before "stretched".
        min_area_fraction: pictures covering less than this fraction of the slide
            (logos, decorations) are measured but never flagged.

    Returns ``{"pptx", "n_pictures", "n_flagged", "pictures": [...]}``.
    """
    _p = _pptx()
    from pptx.util import Emu
    if not os.path.isfile(pptx_path):
        return {"error": f"file not found: {pptx_path}"}

    prs = _p.Presentation(pptx_path)
    slide_area_pt = float(Emu(prs.slide_width).pt) * float(Emu(prs.slide_height).pt)

    found = []
    for i, slide in enumerate(prs.slides, 1):
        _iter_pictures(slide.shapes, i, found)

    rows = []
    for slide_no, name, shape, asset in found:
        if "unreadable" in asset:
            disp_pt_w = float(Emu(shape.width).pt)
            disp_pt_h = float(Emu(shape.height).pt)
            rows.append({
                "slide": slide_no,
                "shape": name,
                "kind": "unmeasurable",
                "source_type": None,
                "pixels": None,
                "dpi": None,
                "authored_cm": float("nan"),
                "authored_cm_height": float("nan"),
                "displayed_cm": round(float(Emu(shape.width).cm), 2),
                "displayed_pt": round(disp_pt_w, 1),
                "displayed_cm_height": round(float(Emu(shape.height).cm), 2),
                "scale": float("nan"),
                "displayed_figure_font_pt": float("nan"),
                "effective_dpi": None,
                "aspect_error": 0.0,
                "area_fraction": round(
                    disp_pt_w * disp_pt_h / slide_area_pt, 4
                ) if slide_area_pt else 0.0,
                "cropped": False,
                "sha1": None,
                "minor": False,
                "risks": [
                    "UNREADABLE PICTURE -- authored width UNVERIFIABLE "
                    f"({asset['unreadable']}); the image is linked/OLE or has "
                    "no embeddable asset. Re-insert it as an embedded file."
                ],
            })
            continue
        if asset["source_type"] == "svg":
            # Measure the actual vector text instead of estimating every SVG
            # from the generic 24 pt authoring profile.
            from ._svg_pptx import svg_picture_row
            rows.append(svg_picture_row(
                slide_no, name, shape, asset["blob"], slide_area_pt,
                scale_tol, aspect_tol, min_area_fraction,
                _SLIDE_MIN_VISIBLE_FONT_PT))
            continue
        source_type = asset["source_type"]
        pixels = asset["pixels"]
        dpi = asset["dpi"]
        px_w, px_h = pixels if pixels is not None else (None, None)
        dpi_w, dpi_h = dpi if dpi is not None else (None, None)
        # python-pptx substitutes 72 dpi when the file carries no resolution.
        no_dpi = (
            source_type != "svg"
            and (int(dpi_w), int(dpi_h)) == (72, 72)
        )

        crop_x = float(shape.crop_left or 0.0) + float(shape.crop_right or 0.0)
        crop_y = float(shape.crop_top or 0.0) + float(shape.crop_bottom or 0.0)
        visible_x = max(0.0, 1.0 - crop_x)
        visible_y = max(0.0, 1.0 - crop_y)
        authored_pt_w = asset["width_pt"] * visible_x
        authored_pt_h = asset["height_pt"] * visible_y
        authored_cm_w = authored_pt_w / 72.0 * 2.54
        authored_cm_h = authored_pt_h / 72.0 * 2.54
        disp_pt_w = float(Emu(shape.width).pt)
        disp_pt_h = float(Emu(shape.height).pt)
        disp_cm_w = float(Emu(shape.width).cm)
        disp_cm_h = float(Emu(shape.height).cm)

        scale = disp_cm_w / authored_cm_w if authored_cm_w > 0 else float("nan")
        used_px_w = px_w * visible_x if px_w is not None else None
        eff_dpi = (
            used_px_w / (disp_pt_w / 72.0)
            if used_px_w is not None and disp_pt_w > 0 else None
        )
        native_aspect = (
            asset["aspect_height"] * visible_y
            / (asset["aspect_width"] * visible_x)
            if asset["aspect_width"] * visible_x else 0.0
        )
        disp_aspect = (disp_pt_h / disp_pt_w) if disp_pt_w else 0.0
        aspect_err = (disp_aspect / native_aspect - 1.0) if native_aspect else 0.0
        area_fraction = (disp_pt_w * disp_pt_h / slide_area_pt) if slide_area_pt else 0.0

        risks = []
        minor = area_fraction < min_area_fraction
        declared_raster = name.split("/")[-1].startswith("RASTER_OK::")
        if asset.get("svg_error"):
            risks.append(
                "SVG UNREADABLE -- the vector original could not be measured "
                f"({asset['svg_error']}); the authored width below comes from "
                "the raster fallback and may not match the vector artwork.")
        elif not minor and not declared_raster:
            # The lab rule is that a figure is vector.  It is not a matter of
            # taste: the text inside a raster cannot be measured from the file
            # (that needs OCR), so the 20 pt floor is unenforceable on one --
            # which is exactly how a 16 pt label shipped unnoticed.
            risks.append(
                "RASTER FIGURE -- re-export as SVG (matplotlib: savefig(..., "
                "format='svg'); a PDF converts with PyMuPDF get_svg_image). "
                "Its on-slide text size cannot be verified from the file. If "
                "raster is genuinely right (a photograph or a screen capture), "
                "rename the shape RASTER_OK::<why> to record that decision.")
        if no_dpi:
            risks.append(
                "NO DPI METADATA -- authored width UNVERIFIABLE (python-pptx "
                "reports the 72 dpi default); save with lab_savefig(dpi=300).")
        if not minor and not no_dpi:
            if source_type == "svg" and scale < _MIN_SAFE_SCALE:
                displayed_pt = _SLIDE_AUTHORED_FONT_PT * scale
                risks.append(
                    f"DOWNSCALED to {scale * 100:.1f}% of the authored "
                    f"{authored_cm_w:.2f} cm -- {_SLIDE_AUTHORED_FONT_PT:.0f} pt "
                    f"figure text is displayed at {displayed_pt:.1f} pt "
                    f"(below the {_SLIDE_MIN_VISIBLE_FONT_PT:.0f} pt slide floor).")
            elif source_type != "svg" and scale < 1.0 - scale_tol:
                displayed_pt = _SLIDE_AUTHORED_FONT_PT * scale
                risks.append(
                    f"DOWNSCALED to {scale * 100:.1f}% of the authored "
                    f"{authored_cm_w:.2f} cm -- {_SLIDE_AUTHORED_FONT_PT:.0f} pt "
                    f"figure text is displayed at {displayed_pt:.1f} pt"
                    + (f" (below the {_SLIDE_MIN_VISIBLE_FONT_PT:.0f} pt slide floor)"
                       if scale < _MIN_SAFE_SCALE else "") + ".")
            elif source_type != "svg" and scale > 1.0 + scale_tol:
                risks.append(
                    f"UPSCALED to {scale * 100:.1f}% of the authored "
                    f"{authored_cm_w:.2f} cm -- the artwork is interpolated; "
                    "re-author at the paste width instead.")
        if not minor and abs(aspect_err) > aspect_tol:
            risks.append(
                f"ASPECT DISTORTED by {aspect_err * 100:+.1f}% -- the picture was "
                "stretched; hold the aspect or re-author the figure.")
        if (not minor and eff_dpi is not None
                and eff_dpi < min_effective_dpi):
            risks.append(
                f"LOW EFFECTIVE DPI {eff_dpi:.0f} (< {min_effective_dpi:.0f}) -- "
                f"{int(used_px_w)} px over {disp_pt_w:.0f} pt.")

        rows.append({
            "slide": slide_no,
            "shape": name,
            "kind": "raster",
            "declared_raster": bool(declared_raster),
            "source_type": source_type,
            "pixels": (
                [int(px_w), int(px_h)] if pixels is not None else None
            ),
            "dpi": (
                [int(dpi_w), int(dpi_h)] if dpi is not None else None
            ),
            "authored_cm": round(authored_cm_w, 2),
            "authored_cm_height": round(authored_cm_h, 2),
            "displayed_cm": round(disp_cm_w, 2),
            "displayed_pt": round(disp_pt_w, 1),
            "displayed_cm_height": round(disp_cm_h, 2),
            "scale": round(scale, 4),
            "displayed_figure_font_pt": round(_SLIDE_AUTHORED_FONT_PT * scale, 1),
            "effective_dpi": (
                round(eff_dpi, 1) if eff_dpi is not None else None
            ),
            "aspect_error": round(aspect_err, 4),
            "area_fraction": round(area_fraction, 4),
            "cropped": bool(crop_x or crop_y),
            "sha1": hashlib.sha1(asset["blob"]).hexdigest()[:10],
            "minor": bool(minor),
            "risks": risks,
        })

    return {
        "pptx": os.path.abspath(pptx_path),
        "n_pictures": len(rows),
        "n_vector": sum(1 for r in rows if r.get("kind") == "svg"),
        "n_raster": sum(1 for r in rows if r.get("kind") == "raster"),
        "n_unmeasurable": sum(1 for r in rows if r.get("kind") == "unmeasurable"),
        "n_flagged": sum(1 for r in rows if r["risks"]),
        "authored_font_pt": _SLIDE_AUTHORED_FONT_PT,
        "min_visible_font_pt": _SLIDE_MIN_VISIBLE_FONT_PT,
        "pictures": rows,
    }


def slide_png_from_pdf(pdf_path: str,
                       out_png: str,
                       paste_width_pt: float,
                       page: int = 0,
                       dpi: int = 300) -> dict:
    """Rasterise a vector figure so that its AUTHORED width equals the width it
    will be pasted at -- the paste-scale-1.000 path for artwork that is a PDF
    (a CAD export, a drawing inherited from a manuscript) rather than a
    matplotlib figure.

    Two things have to line up and both are easy to get wrong by hand: the
    render scale (so the pixel count matches ``paste_width_pt`` at ``dpi``) and
    the PNG's own DPI metadata (so :func:`audit_pptx_figures`, and PowerPoint's
    "reset picture size", recover that authored width instead of assuming 72).

    Args:
        pdf_path: source vector figure.
        out_png: PNG to write.
        paste_width_pt: width the picture will occupy on the slide, in points
            (PowerPoint's own unit for shape width).
        page: page index to render.
        dpi: output resolution; also written into the PNG metadata.

    Returns ``{"png", "pixels", "authored_cm", "paste_width_pt",
    "text_scale"}``; ``text_scale`` is how much the drawing's own text grows or
    shrinks versus the PDF page (paste width / page width), so a 20 pt in-page
    label lands at ``20 * text_scale`` on the slide.
    """
    try:
        import fitz                                   # PyMuPDF
    except ImportError as exc:                        # pragma: no cover - env dependent
        raise RuntimeError(
            "PyMuPDF (fitz) is required to rasterise a vector figure at the "
            "paste width (pip install pymupdf)."
        ) from exc
    if paste_width_pt <= 0:
        raise ValueError("paste_width_pt must be > 0.")

    doc = fitz.open(pdf_path)
    try:
        pg = doc[page]
        page_width_pt = float(pg.rect.width)
        if page_width_pt <= 0:
            raise ValueError(f"{pdf_path}: page {page} has zero width.")
        # Render so that (pixels / dpi) inches == paste_width_pt / 72 inches.
        # `get_pixmap(dpi=...)` takes an int, so drive the fractional scale
        # through the zoom matrix (pixels per PDF point) instead.
        zoom = (paste_width_pt / 72.0 * dpi) / page_width_pt
        pix = pg.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
        pix.set_dpi(int(round(dpi)), int(round(dpi)))  # authored width recoverable
        pix.save(out_png)
        px_w, px_h = pix.width, pix.height
    finally:
        doc.close()

    return {
        "png": os.path.abspath(out_png),
        "pixels": [px_w, px_h],
        "authored_cm": round(px_w / dpi * 2.54, 2),
        "paste_width_pt": round(float(paste_width_pt), 1),
        "text_scale": round(paste_width_pt / page_width_pt, 3),
    }


def figure_audit_pptx_figures(pptx_path: str,
                              scale_tol: float = 0.02,
                              min_effective_dpi: float = 150.0) -> str:
    """Lint every picture in a PPTX deck for paste-scale defects -- the slide
    analogue of ``figure_audit_embeds`` (which audits ``\\includegraphics``).

    A slide figure authored with ``lab_figure(medium='presentation',
    embed_width_cm=W)`` carries 24 pt text and must be pasted AT ``W``.  Dragging
    it to another size rescales the text silently: at 83% a 24 pt label is
    displayed at 20 pt, the slide floor.  This tool recovers the authored width
    from the embedded image (pixels / DPI, crop removed), compares it with the
    width the shape occupies, and reports the scale, the displayed figure-font
    size, the effective DPI, and any aspect stretch.

    Args:
        pptx_path: path to the .pptx deck.
        scale_tol: allowed |scale - 1| before flagging (default 2%).
        min_effective_dpi: pixels-per-displayed-inch floor (default 150).

    Returns a multi-line report; every picture is listed with its scale so a
    clean deck is auditable, not just the flagged ones.
    """
    rep = audit_pptx_figures(pptx_path, scale_tol=scale_tol,
                             min_effective_dpi=min_effective_dpi)
    if "error" in rep:
        return f"figure-paste audit: {rep['error']}"
    unchecked = sum(1 for r in rep["pictures"]
                    if r.get("text_check") == "not-verifiable")
    lines = [f"figure-paste audit: {rep['pptx']}",
             f"  {rep['n_pictures']} pictures, {rep['n_flagged']} flagged "
             f"(authored {rep['authored_font_pt']:.0f} pt, floor "
             f"{rep['min_visible_font_pt']:.0f} pt)"]
    if unchecked:
        lines.append(
            f"  {unchecked} vector picture(s) do not state their text size "
            f"(PDF-converted artwork); those are NOT text-checked here.")
    lines.append("")
    for r in rep["pictures"]:
        tag = "FLAG" if r["risks"] else " ok "
        head = f"  [{tag}] slide {r['slide']:>2} {r['shape']}: "
        if r.get("kind") == "svg":
            font = r.get("body_font_on_slide_pt")
            lines.append(
                head + (f"authored {r['authored_cm']:.2f} cm -> pasted "
                        f"{r['displayed_cm']:.2f} cm (scale "
                        f"{r['scale']:.3f}, " if r.get("scale") else
                        f"pasted {r['displayed_cm']:.2f} cm (")
                + "vector, text "
                + (f"{r['smallest_font_on_slide_pt']:.0f}-"
                   f"{r['largest_font_on_slide_pt']:.0f} pt on the slide)"
                   if font else "size not readable from the file)"))
        elif r.get("kind") == "unmeasurable":
            lines.append(head + f"pasted {r['displayed_cm']:.2f} cm "
                                f"(no embedded artwork to measure)")
        else:
            lines.append(
                head + f"authored "
                f"{r['authored_cm']:.2f} cm -> pasted {r['displayed_cm']:.2f} cm "
                f"(scale {r['scale']:.3f}, figure text "
                f"{r['displayed_figure_font_pt']:.1f} pt, "
                f"{r['effective_dpi']:.0f} dpi)")
        for risk in r["risks"]:
            lines.append(f"         - {risk}")
    if rep["n_pictures"] == 0:
        lines.append("  (no pictures found)")
    elif rep["n_flagged"] == 0:
        lines.append("\nEvery picture is pasted at its authored width -- clean.")
    else:
        lines.append(
            "\nFix: re-author the figure at the paste width with "
            "lab_figure(medium='presentation', embed_width_cm=W) + lab_savefig, "
            "then paste at 100% (PowerPoint: size the picture to W cm).")
    return "\n".join(lines)
