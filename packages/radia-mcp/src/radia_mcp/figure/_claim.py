"""Draw the figure the claim needs, and check the claim is actually visible.

A slide asserts something specific -- "reduced to 1/145", "within 0.18%",
"162,000 DOF in 60.7 s" -- and then shows a figure. The figure is usually a
correct plot of the data and still fails the slide, because plotting the data
is not the same as making the claim *visible*: the audience is left to find
1/145 in two curves, or to locate 162,000 on an axis, while the speaker has
already moved on.

Two things here:

- `figure_plan_for_claim` says what form makes a given claim visible, and what
  to annotate. The form follows the KIND of claim, which is what the wording
  tells you: a ratio wants the two things side by side with the ratio drawn
  between them; a bound wants a band; an operating point wants a marked point.
- `figure_claim_visibility` checks a finished deck: for every slide that makes
  a numeric claim, is that number anywhere the eye can find it?

Measured on the deck this was written for, three claims out of four were
visible and one -- an operating point on a scaling curve -- was asserted in
the title and shown nowhere.
"""

from __future__ import annotations

import pathlib
import re

# Numbers as a claim states them.  "1/145" first: a bare 145 would otherwise
# swallow it.
NUMERIC = re.compile(
    r"1\s*/\s*\d+"
    r"|\d+(?:[.,]\d+)?\s*分の\s*\d+"
    r"|\d+(?:[.,]\d+)?\s*(?:%|％|倍|万自由度|自由度|秒|ms|s\b|GiB|MiB|T\b|kW|kA/m)"
)

KINDS = (
    ("ratio", re.compile(r"1\s*/\s*\d+|分の|倍|低減|削減|に抑制|改善")),
    ("bound", re.compile(r"以内|以下|未満|within|抑制|差は")),
    ("operating_point", re.compile(r"自由度.*(?:秒|s\b)|(?:秒|s\b).*自由度|求解|で解析|スループット")),
    ("agreement", re.compile(r"一致|整合|参照|FEM|実測|理論|agree|match")),
    ("trend", re.compile(r"傾き|スケール|収束|次数|slope|order|scaling")),
)

FORM = {
    "ratio": {
        "form": "the two things side by side, with the ratio drawn between them",
        "annotate": "the ratio itself, on the gap it measures -- not in a legend",
        "why": ("A ratio is a comparison, so the eye must be able to make the "
                "comparison without arithmetic. Two curves on one axis leave "
                "the audience to divide."),
        "avoid": "two separate panels the reader has to hold in memory",
    },
    "bound": {
        "form": "the difference itself, with the bound as a shaded band",
        "annotate": "the bound value on the band edge",
        "why": ("A claim of 'within X' is about the difference, not about the "
                "two quantities. Plotting both and letting the reader subtract "
                "hides exactly the thing being claimed."),
        "avoid": "plotting the two absolute curves and asserting they are close",
    },
    "operating_point": {
        "form": "the curve, with that point marked and labelled",
        "annotate": "both coordinates at the point, e.g. '162,000 DOF, 60.7 s'",
        "why": ("A specific achievement on a curve is invisible unless it is "
                "marked. The audience cannot find it on a log axis while the "
                "speaker is talking."),
        "avoid": "a curve whose claimed point is left for the reader to locate",
    },
    "agreement": {
        "form": "the two overlaid, plus a residual panel underneath",
        "annotate": "the worst-case difference, where it occurs",
        "why": ("Overlaid curves that agree look like one curve; the residual "
                "is where the claim lives and it needs its own axis."),
        "avoid": "an overlay alone, where agreement and over-plotting look the same",
    },
    "trend": {
        "form": "log-log, with the fitted slope drawn on the fitted range",
        "annotate": "the slope, and the range it was fitted over",
        "why": ("A slope claim is about a range. Drawing the fit over the "
                "range it came from keeps the claim honest about where it "
                "holds."),
        "avoid": "a slope quoted in a caption with no line on the plot",
    },
}


def classify_claim(claim: str) -> dict:
    """What kind of claim is this, and what numbers does it assert?"""
    numbers = [m.group(0).strip() for m in NUMERIC.finditer(claim or "")]
    kind = None
    for name, pattern in KINDS:
        if pattern.search(claim or ""):
            kind = name
            break
    return {"claim": claim, "kind": kind, "numbers": numbers,
            "is_quantitative": bool(numbers)}


def figure_plan_for_claim(claim: str, embed_width_cm: float = 8.0) -> dict:
    """What figure makes this claim visible, and what to annotate on it."""
    got = classify_claim(claim)
    kind = got["kind"]
    plan = dict(FORM.get(kind, {
        "form": "whatever shows the quantity the claim is about",
        "annotate": "the claimed quantity itself",
        "why": ("The claim's kind could not be read from its wording. State "
                "the claim as a comparison, a bound, a point, an agreement or "
                "a trend, and the form follows."),
        "avoid": "a plot that shows the data but not the claim",
    }))
    plan.update(got)
    # the lab's own rules, restated where they bite
    plan["rules"] = [
        "No title inside the figure -- it goes in the slide title or caption.",
        "Annotating the claimed NUMBER is not a title; it is what makes the "
        "claim visible, and it belongs on the feature it measures.",
        "%.0f pt on-page text at %.1f cm embed width (1.25 pt per cm)."
        % (1.25 * embed_width_cm, embed_width_cm),
    ]
    return plan


def figure_claim_visibility(pptx_path: str) -> dict:
    """For every slide that makes a numeric claim, is the number visible?

    Looks in the figure's own labels, the slide body and any callout -- that
    is, anywhere the audience's eye can land. A claim asserted only in the
    title, over a figure that does not carry it, is the case this finds.
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}

    from ..presentation.plans.T30 import classify_slide_text, _is_backup

    prs = _pptx.Presentation(str(p))
    rows: list[dict] = []
    invisible: list[dict] = []

    for i, slide in enumerate(prs.slides, 1):
        if i == 1 or _is_backup(slide):
            continue
        part = classify_slide_text(slide)
        boxes = sorted(
            [(s.top, (s.text_frame.text or "").replace("\n", ""))
             for s in slide.shapes
             if s.has_text_frame and (s.text_frame.text or "").strip()
             and s.top is not None],
            key=lambda b: b[0])
        title = boxes[0][1] if boxes else ""
        claim = title if NUMERIC.search(title) else part["banner"]
        got = classify_claim(claim)
        if not got["numbers"]:
            continue

        pool = (part["figure_labels"] + "\n" + part["body"] + "\n"
                + part["table_text"]).replace(" ", "")
        shown = [n for n in got["numbers"] if n.replace(" ", "") in pool]
        row = {
            "slide_no": i,
            "claim": claim,
            "kind": got["kind"],
            "numbers": got["numbers"],
            "shown": shown,
            "visible": len(shown) == len(got["numbers"]),
        }
        rows.append(row)
        if not row["visible"]:
            row = dict(row)
            row["missing"] = [n for n in got["numbers"] if n not in shown]
            row["plan"] = figure_plan_for_claim(claim)
            invisible.append(row)

    n = len(rows)
    ratio = (n - len(invisible)) / n if n else 1.0
    return {
        "score": round(ratio * 10, 1),
        "score_max": 10,
        "claims_checked": n,
        "n_not_visible": len(invisible),
        "claims_not_visible": invisible,
        "per_claim": rows,
        "hint": (
            "A claim stated in the title and shown nowhere leaves the audience "
            "to find it. Mark it on the figure -- an annotated number is not "
            "an in-figure title, it is the claim made visible. "
            "figure_plan_for_claim says which form each kind of claim wants."
        ),
        "source": "presentation script-first support (2026-08-20).",
    }
