"""Is each figure earned -- does the speaker ever point at it?

`figure_audit_pptx_figures` answers "is this figure pasted correctly and big
enough to read". It cannot answer the question that decides whether a figure
belongs on the slide at all: **does the script send the audience to it.**

A figure nobody points at is left for the audience to decode alone while the
speaker talks about something else. It is not a rendering fault, so no amount
of paste-scale checking finds it -- and it is common, because figures and
script are usually written at different times by different hands. On the deck
this was built for, five of eight figures were never pointed at, one of them
because a sentence was cut while trimming the talk to fit its slot.

The check is deliberately generous: any deictic reference counts, and so does
naming something the figure itself is labelled with. What it reports is the
slide where the script says nothing at all that would make a listener look up.
"""

from __future__ import annotations

import pathlib
import re

# Words that send a listener to the screen.  Kept wide on purpose: the cost of
# missing a real gap is a figure nobody looks at; the cost of a false alarm is
# one glance at a slide that turns out fine.
DEIXIS = re.compile(
    r"図|グラフ|プロット|横軸|縦軸|軸|凡例|左|右|上段|下段|こちら|ご覧|"
    r"見え|示し|表[はにをのが]|"
    r"figure|Fig\.|graph|plot|axis|legend|left|right|shown|see ",
    re.IGNORECASE,
)


def _pictures(shapes, out):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _pictures(shape.shapes, out)
        elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            out.append(shape)
    return out


def _labels_on(slide, pictures) -> set[str]:
    """Text sitting on top of a figure -- its own labels."""
    rects = []
    for p in pictures:
        if None not in (p.left, p.top, p.width, p.height):
            rects.append((p.left, p.top, p.left + p.width, p.top + p.height))
    words: set[str] = set()
    for shape in slide.shapes:
        try:
            if not shape.has_text_frame or None in (shape.left, shape.top):
                continue
            cx = shape.left + (shape.width or 0) // 2
            cy = shape.top + (shape.height or 0) // 2
            if not any(x0 <= cx <= x1 and y0 <= cy <= y1 for x0, y0, x1, y1 in rects):
                continue
            for m in re.finditer(r"[一-鿿]{2,}|[A-Za-z]{3,}", shape.text_frame.text or ""):
                words.add(m.group(0))
        except Exception:
            continue
    return words


def figure_audit_script_reference(pptx_path: str) -> dict:
    """Report figures the script never points at.

    Returns per-slide findings plus `figures_not_referenced`, the list to act
    on: either add a sentence that sends the audience to the figure, or drop
    the figure.
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}

    prs = _pptx.Presentation(str(p))
    per_slide: list[dict] = []
    unreferenced: list[dict] = []
    n_pictures = 0

    for i, slide in enumerate(prs.slides, 1):
        pics = _pictures(slide.shapes, [])
        if not pics:
            continue
        n_pictures += len(pics)
        note = ""
        try:
            if slide.has_notes_slide:
                note = (slide.notes_slide.notes_text_frame.text or "").split("---")[0]
        except Exception:
            note = ""
        note = note.strip()

        deictic = bool(DEIXIS.search(note))
        labels = _labels_on(slide, pics)
        named = sorted(w for w in labels if w in note)
        referenced = deictic or bool(named)

        row = {
            "slide_no": i,
            "n_pictures": len(pics),
            "script_chars": len(note.replace("\n", "")),
            "deictic_reference": deictic,
            "figure_labels_named": named[:5],
            "referenced": referenced,
        }
        per_slide.append(row)
        if not referenced:
            unreferenced.append({
                "slide_no": i,
                "n_pictures": len(pics),
                "reason": ("script never sends the audience to the figure "
                           "(no deictic word, and none of its labels named)"),
            })

    n = len(per_slide)
    ratio = (n - len(unreferenced)) / n if n else 1.0
    return {
        "score": round(ratio * 10, 1),
        "score_max": 10,
        "slides_with_figures": n,
        "n_pictures": n_pictures,
        "n_not_referenced": len(unreferenced),
        "figures_not_referenced": unreferenced,
        "per_slide": per_slide,
        "hint": (
            "A figure the script never points at is left for the audience to "
            "decode alone. Either add a sentence that sends them to it "
            "(\"左の図は…\", \"横軸は…\") or drop the figure. Note that trimming "
            "a talk to fit its slot is a common way to create these: the "
            "pointing sentence is short and looks expendable."
        ),
        "source": "presentation script-first support (2026-08-20).",
    }
