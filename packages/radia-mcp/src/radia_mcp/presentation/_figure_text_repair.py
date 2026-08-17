"""Repair small text baked into raster figures in a PowerPoint deck."""

from __future__ import annotations

import copy
import json
import math
import pathlib
import tempfile


def _bbox_points(bbox) -> list[tuple[float, float]]:
    if (isinstance(bbox, list) and len(bbox) == 4
            and all(isinstance(value, (int, float)) for value in bbox)):
        x0, y0, x1, y1 = (float(value) for value in bbox)
        return [(x0, y0), (x1, y0), (x1, y1), (x0, y1)]
    if isinstance(bbox, list) and len(bbox) >= 4:
        return [
            (float(point[0]), float(point[1])) for point in bbox[:4]
        ]
    return []


def _bbox_height(points: list[tuple[float, float]]) -> float:
    if len(points) != 4:
        return 0.0
    left = math.dist(points[0], points[3])
    right = math.dist(points[1], points[2])
    return max(left, right)


def _bbox_angle(points: list[tuple[float, float]]) -> float:
    if len(points) != 4:
        return 0.0
    return math.degrees(math.atan2(
        points[1][1] - points[0][1],
        points[1][0] - points[0][0],
    ))


def _text_width_pt(text: str, font_pt: float) -> float:
    width_em = 0.0
    for char in text:
        if char.isspace():
            width_em += 0.32
        elif ord(char) < 128:
            width_em += 0.58
        else:
            width_em += 1.0
    return max(font_pt, width_em * font_pt)


def _overlap_ratio(a, b) -> float:
    ax0, ay0, ax1, ay1 = a
    bx0, by0, bx1, by1 = b
    width = max(0, min(ax1, bx1) - max(ax0, bx0))
    height = max(0, min(ay1, by1) - max(ay0, by0))
    intersection = width * height
    area_a = max(1, (ax1 - ax0) * (ay1 - ay0))
    area_b = max(1, (bx1 - bx0) * (by1 - by0))
    return intersection / min(area_a, area_b)


def presentation_replace_embedded_figure_text(
        pptx_path: str,
        output_pptx_path: str,
        ocr_manifest_path: str,
        min_font_pt: float = 20.0,
        font_name: str = "Yu Gothic",
        min_confidence: float = 0.70,
        replace_scope: str = "below_minimum",
        dry_run: bool = True,
        source_unavailable_confirmed: bool = False,
        overwrite: bool = False,
        inpaint_padding_px: int = 1,
        inpaint_radius_px: int = 3,
        ink_threshold: int = 28,
        output_manifest_path: str = "") -> dict:
    """Replace OCR text in raster figures with native PowerPoint text.

    The input deck is never overwritten. OCR word boxes are removed from the
    embedded bitmap with local inpainting and reinserted as editable native
    text boxes named ``FIGURE_TEXT::...``. The output is a candidate for visual
    review, not an automatic publication artifact: mathematical notation,
    subscripts, coloured labels, and text crossing plot lines remain manual
    review points.

    ``replace_scope`` is ``below_minimum`` (default) or ``all``. ``dry_run``
    reports candidates without creating files.
    """
    try:
        import pptx as _pptx
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Pt
    except ImportError:
        return {"error": "python-pptx not installed."}

    source = pathlib.Path(pptx_path)
    output = pathlib.Path(output_pptx_path)
    manifest_path = pathlib.Path(ocr_manifest_path)
    if not source.exists():
        return {"error": f"file not found: {pptx_path}"}
    if not manifest_path.exists():
        return {"error": f"OCR manifest not found: {ocr_manifest_path}"}
    if source.resolve() == output.resolve():
        return {"error": "output_pptx_path must differ from the input deck."}
    if output.exists() and not overwrite and not dry_run:
        return {"error": f"output already exists: {output_pptx_path}"}
    manifest_output = pathlib.Path(output_manifest_path) if output_manifest_path else (
        output.with_suffix(".figure-text-ocr.json")
    )
    if manifest_output.exists() and not overwrite and not dry_run:
        return {"error": f"output manifest already exists: {manifest_output}"}
    if min_font_pt < 20.0:
        return {"error": "min_font_pt must be >= 20."}
    if not dry_run and not source_unavailable_confirmed:
        return {
            "error": (
                "OCR replacement is fallback-only. Regenerate the source "
                "figure when possible; otherwise set "
                "source_unavailable_confirmed=true."
            )
        }
    if not 0.0 <= min_confidence <= 1.0:
        return {"error": "min_confidence must be between 0 and 1."}
    if replace_scope not in {"below_minimum", "all"}:
        return {"error": "replace_scope must be below_minimum or all."}
    if inpaint_padding_px < 0 or inpaint_radius_px < 1:
        return {"error": "invalid inpaint padding or radius."}
    if not 0 <= ink_threshold <= 255:
        return {"error": "ink_threshold must be between 0 and 255."}

    try:
        manifest_payload = json.loads(manifest_path.read_text(encoding="utf-8"))
        manifest_items = manifest_payload.get("pictures", [])
        manifest = {
            (int(item["slide"]), str(item["shape"])): item
            for item in manifest_items
        }
    except Exception as exc:
        return {"error": f"invalid OCR manifest: {exc}"}

    try:
        prs = _pptx.Presentation(str(source))
    except Exception as exc:
        return {"error": f"failed to open input deck: {exc}"}

    emu_per_point = 12700.0
    candidates = []
    unresolved = []
    per_picture = {}
    for slide_index, slide in enumerate(prs.slides, 1):
        for shape in list(slide.shapes):
            try:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
            except Exception:
                continue
            shape_name = getattr(shape, "name", "")
            item = manifest.get((slide_index, shape_name))
            if item is None:
                unresolved.append({
                    "slide": slide_index,
                    "shape": shape_name,
                    "reason": "picture is missing from OCR manifest",
                })
                continue
            if item.get("confirmed_textless"):
                continue
            try:
                image_width, image_height = shape.image.size
                crop_left = float(shape.crop_left or 0.0)
                crop_right = float(shape.crop_right or 0.0)
                crop_top = float(shape.crop_top or 0.0)
                crop_bottom = float(shape.crop_bottom or 0.0)
                visible_width = image_width * (1.0 - crop_left - crop_right)
                visible_height = image_height * (1.0 - crop_top - crop_bottom)
                if visible_width <= 0 or visible_height <= 0:
                    raise ValueError("invalid crop dimensions")
                if abs(float(shape.rotation or 0.0)) > 1.0e-6:
                    raise ValueError("rotated picture shape is not supported")
                displayed_height_pt = float(shape.height) / emu_per_point
            except Exception as exc:
                unresolved.append({
                    "slide": slide_index,
                    "shape": shape_name,
                    "reason": str(exc),
                })
                continue

            selected = []
            for word_index, word in enumerate(item.get("words", [])):
                confidence = float(word.get("confidence", 1.0))
                if confidence < min_confidence:
                    unresolved.append({
                        "slide": slide_index,
                        "shape": shape_name,
                        "text": str(word.get("text", ""))[:40],
                        "reason": "OCR confidence below replacement threshold",
                    })
                    continue
                points = _bbox_points(word.get("bbox", []))
                glyph_height_px = _bbox_height(points)
                if not points or glyph_height_px <= 0:
                    unresolved.append({
                        "slide": slide_index,
                        "shape": shape_name,
                        "text": str(word.get("text", ""))[:40],
                        "reason": "invalid OCR bounding box",
                    })
                    continue
                estimated_font_pt = (
                    glyph_height_px * displayed_height_pt
                    / visible_height / 0.72
                )
                if (replace_scope == "all"
                        or estimated_font_pt < min_font_pt):
                    candidate = {
                        "slide": slide_index,
                        "shape": shape_name,
                        "word_index": word_index,
                        "text": str(word.get("text", "")),
                        "confidence": round(confidence, 3),
                        "estimated_font_pt": round(estimated_font_pt, 2),
                        "target_font_pt": round(
                            max(min_font_pt, estimated_font_pt), 2
                        ),
                        "bbox": word.get("bbox", []),
                        "points": points,
                    }
                    selected.append(candidate)
                    candidates.append({
                        key: value for key, value in candidate.items()
                        if key != "points"
                    })
            if selected:
                per_picture[(slide_index, shape_name)] = {
                    "shape": shape,
                    "manifest_item": item,
                    "selected": selected,
                    "image_size": (image_width, image_height),
                    "crop": (crop_left, crop_right, crop_top, crop_bottom),
                }

    base_result = {
        "dry_run": dry_run,
        "input_pptx": str(source),
        "output_pptx": str(output),
        "candidate_count": len(candidates),
        "candidates": candidates,
        "unresolved_count": len(unresolved),
        "unresolved": unresolved,
        "warning": (
            "Generated decks require visual review. OCR surgery is a fallback; "
            "regenerating the source figure remains preferred."
        ),
    }
    if dry_run or not candidates:
        return {
            **base_result,
            "written": False,
            "status": "dry_run" if dry_run else "no_candidates",
        }

    try:
        import cv2
        import numpy as np
    except ImportError:
        return {**base_result, "error": "opencv-python is not installed."}

    repaired_manifest = copy.deepcopy(manifest_payload)
    repaired_lookup = {
        (int(item["slide"]), str(item["shape"])): item
        for item in repaired_manifest.get("pictures", [])
    }
    replacements = []
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_root = pathlib.Path(temp_dir)
        for (slide_index, shape_name), data in per_picture.items():
            shape = data["shape"]
            slide = prs.slides[slide_index - 1]
            image_width, image_height = data["image_size"]
            crop_left, crop_right, crop_top, crop_bottom = data["crop"]
            raw = np.frombuffer(shape.image.blob, dtype=np.uint8)
            image = cv2.imdecode(raw, cv2.IMREAD_UNCHANGED)
            if image is None:
                unresolved.append({
                    "slide": slide_index,
                    "shape": shape_name,
                    "reason": "embedded picture format cannot be rasterized",
                })
                continue
            if image.ndim == 2:
                color = cv2.cvtColor(image, cv2.COLOR_GRAY2BGR)
                alpha = None
            elif image.shape[2] == 4:
                color = image[:, :, :3].copy()
                alpha = image[:, :, 3].copy()
            else:
                color = image[:, :, :3].copy()
                alpha = None
            gray = cv2.cvtColor(color, cv2.COLOR_BGR2GRAY)
            total_mask = np.zeros(gray.shape, dtype=np.uint8)
            text_colors = {}
            successful = []
            for candidate in data["selected"]:
                polygon = np.array(candidate["points"], dtype=np.int32)
                polygon[:, 0] = np.clip(polygon[:, 0], 0, image_width - 1)
                polygon[:, 1] = np.clip(polygon[:, 1], 0, image_height - 1)
                region = np.zeros(gray.shape, dtype=np.uint8)
                cv2.fillConvexPoly(region, polygon, 255)
                border_kernel = np.ones((5, 5), np.uint8)
                expanded = cv2.dilate(region, border_kernel, iterations=1)
                border = (expanded > 0) & (region == 0)
                if not np.any(border):
                    unresolved.append({
                        "slide": slide_index,
                        "shape": shape_name,
                        "text": candidate["text"],
                        "reason": "cannot estimate local background",
                    })
                    continue
                background = float(np.median(gray[border]))
                ink = (np.abs(gray.astype(np.float32) - background)
                       >= ink_threshold) & (region > 0)
                if np.count_nonzero(ink) < 3:
                    unresolved.append({
                        "slide": slide_index,
                        "shape": shape_name,
                        "text": candidate["text"],
                        "reason": "no removable text pixels found",
                    })
                    continue
                ink_mask = (ink.astype(np.uint8) * 255)
                if inpaint_padding_px:
                    kernel_size = 2 * inpaint_padding_px + 1
                    kernel = np.ones((kernel_size, kernel_size), np.uint8)
                    ink_mask = cv2.dilate(ink_mask, kernel, iterations=1)
                total_mask = cv2.bitwise_or(total_mask, ink_mask)
                pixels = color[ink]
                median_bgr = np.median(pixels, axis=0).astype(int)
                text_colors[candidate["word_index"]] = (
                    int(median_bgr[2]), int(median_bgr[1]), int(median_bgr[0])
                )
                successful.append(candidate)
            if not successful:
                continue

            repaired_color = cv2.inpaint(
                color, total_mask, inpaint_radius_px, cv2.INPAINT_TELEA
            )
            repaired_image = (
                np.dstack((repaired_color, alpha)) if alpha is not None
                else repaired_color
            )
            safe_name = "".join(
                char if char.isalnum() else "_" for char in shape_name
            )[:40]
            image_path = temp_root / f"slide-{slide_index}-{safe_name}.png"
            ok, encoded = cv2.imencode(".png", repaired_image)
            if not ok:
                unresolved.append({
                    "slide": slide_index,
                    "shape": shape_name,
                    "reason": "failed to encode repaired image",
                })
                continue
            image_path.write_bytes(encoded.tobytes())

            old_element = shape._element
            parent = old_element.getparent()
            old_index = parent.index(old_element)
            new_picture = slide.shapes.add_picture(
                str(image_path), shape.left, shape.top, shape.width, shape.height
            )
            new_picture.crop_left = crop_left
            new_picture.crop_right = crop_right
            new_picture.crop_top = crop_top
            new_picture.crop_bottom = crop_bottom
            new_picture.rotation = shape.rotation
            new_picture.name = shape_name
            new_element = new_picture._element
            parent.remove(old_element)
            new_element.getparent().remove(new_element)
            parent.insert(old_index, new_element)

            visible_x0 = crop_left * image_width
            visible_y0 = crop_top * image_height
            visible_width = image_width * (1.0 - crop_left - crop_right)
            visible_height = image_height * (1.0 - crop_top - crop_bottom)
            for candidate in successful:
                points = candidate["points"]
                xs = [point[0] for point in points]
                ys = [point[1] for point in points]
                x0, x1 = min(xs), max(xs)
                y0, y1 = min(ys), max(ys)
                left = shape.left + int((x0 - visible_x0) / visible_width * shape.width)
                top = shape.top + int((y0 - visible_y0) / visible_height * shape.height)
                box_width = int((x1 - x0) / visible_width * shape.width)
                box_height = int((y1 - y0) / visible_height * shape.height)
                font_pt = float(candidate["target_font_pt"])
                required_width = int(_text_width_pt(
                    candidate["text"], font_pt
                ) * emu_per_point)
                required_height = int(font_pt * 1.25 * emu_per_point)
                width = max(box_width, required_width)
                height = max(box_height, required_height)
                center_x = left + box_width // 2
                center_y = top + box_height // 2
                left = max(0, min(center_x - width // 2, prs.slide_width - width))
                top = max(0, min(center_y - height // 2, prs.slide_height - height))
                textbox = slide.shapes.add_textbox(left, top, width, height)
                textbox.name = (
                    f"FIGURE_TEXT::s{slide_index}:{shape_name}:"
                    f"{candidate['word_index']}"
                )
                textbox.rotation = _bbox_angle(points)
                textbox.fill.background()
                textbox.line.fill.background()
                frame = textbox.text_frame
                frame.clear()
                frame.margin_left = 0
                frame.margin_right = 0
                frame.margin_top = 0
                frame.margin_bottom = 0
                frame.word_wrap = False
                frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                paragraph = frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER
                run = paragraph.add_run()
                run.text = candidate["text"]
                run.font.name = font_name
                run.font.size = Pt(font_pt)
                rgb = text_colors[candidate["word_index"]]
                run.font.color.rgb = RGBColor(*rgb)
                replacements.append({
                    "slide": slide_index,
                    "picture": shape_name,
                    "textbox": textbox.name,
                    "text": candidate["text"],
                    "font_pt": round(font_pt, 2),
                    "rotation": round(float(textbox.rotation or 0.0), 2),
                    "manual_review": bool(
                        set(candidate["text"]) & set("_^{}\\")
                    ),
                })

            repaired_item = repaired_lookup[(slide_index, shape_name)]
            replaced_indices = {
                candidate["word_index"] for candidate in successful
            }
            repaired_item["words"] = [
                word for index, word in enumerate(repaired_item.get("words", []))
                if index not in replaced_indices
            ]
            repaired_item["confirmed_textless"] = not repaired_item["words"]

        output.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output)

    overlap_warnings = []
    for slide_index, slide in enumerate(prs.slides, 1):
        figure_text_shapes = [
            shape for shape in slide.shapes
            if getattr(shape, "name", "").startswith("FIGURE_TEXT::")
        ]
        other_text_shapes = [
            shape for shape in slide.shapes
            if (getattr(shape, "has_text_frame", False)
                and shape.text_frame.text.strip()
                and not getattr(shape, "name", "").startswith("FIGURE_TEXT::"))
        ]
        for index, shape in enumerate(figure_text_shapes):
            shape_box = (
                int(shape.left), int(shape.top),
                int(shape.left + shape.width), int(shape.top + shape.height),
            )
            for other in figure_text_shapes[index + 1:]:
                other_box = (
                    int(other.left), int(other.top),
                    int(other.left + other.width),
                    int(other.top + other.height),
                )
                ratio = _overlap_ratio(shape_box, other_box)
                if ratio > 0.05:
                    overlap_warnings.append({
                        "slide": slide_index,
                        "shape": shape.name,
                        "other": other.name,
                        "overlap_ratio": round(ratio, 3),
                        "reason": "reconstructed figure texts overlap",
                    })
            for other in other_text_shapes:
                other_box = (
                    int(other.left), int(other.top),
                    int(other.left + other.width),
                    int(other.top + other.height),
                )
                ratio = _overlap_ratio(shape_box, other_box)
                if ratio > 0.05:
                    overlap_warnings.append({
                        "slide": slide_index,
                        "shape": shape.name,
                        "other": getattr(other, "name", ""),
                        "overlap_ratio": round(ratio, 3),
                        "reason": "reconstructed figure text overlaps slide text",
                    })

    manifest_output.write_text(
        json.dumps(repaired_manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return {
        **base_result,
        "written": True,
        "status": (
            "needs_review"
            if (unresolved or overlap_warnings
                or any(item["manual_review"] for item in replacements))
            else "candidate_ready"
        ),
        "replacement_count": len(replacements),
        "replacements": replacements,
        "output_manifest": str(manifest_output),
        "overlap_warning_count": len(overlap_warnings),
        "overlap_warnings": overlap_warnings,
        "unresolved_count": len(unresolved),
        "unresolved": unresolved,
    }
