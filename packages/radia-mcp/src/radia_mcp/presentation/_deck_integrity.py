"""Deck-integrity checks: unrendered math markup and reused artwork.

Two failure classes that every existing presentation_* lint walked straight
past, both caught on the MMPM SA-26-069 deck (2026-08-16):

1. **Raw math markup on the slide.**  Equations pasted from a manuscript or
   written by an agent keep their source form -- ``sigma_f(x) = Sum_j B_fj(x)
   a_j``, ``M_ij = Sum_f int_{S_f} phi_i(d) B_fj(x) dS``, ``H(x_c)``,
   ``r_nl`` -- so the audience reads underscores and braces instead of
   subscripts.  Text-extraction lints score such a slide as ordinary prose,
   and the PDF renders "correctly", so nothing fails.
   :func:`presentation_check_raw_math_markup` finds it and
   :func:`presentation_apply_math_subscripts` repairs the mechanical part by
   splitting the run into baseline-shifted runs (PowerPoint's own subscript),
   leaving genuine LaTeX macros for a human.

2. **The same picture on two slides.**  Slides 3 and 4 of that deck embedded
   byte-identical artwork (the mode-decomposition cube) under two different
   titles, i.e. the "before" slide was illustrated with the "after" figure --
   and at two different paste scales.  :func:`presentation_check_duplicate_slide_images`
   reports artwork reused across slides, ignoring the small repeated marks
   (logos, rules) that are supposed to repeat.
"""
from __future__ import annotations

import hashlib
import pathlib
import re
from copy import deepcopy

# Base: a single symbol (Latin / Greek / operator) or a number, not glued to a
# longer identifier -- so `x_c`, `S_f`, `Sigma_j`, `10^{-3}` match while
# `snake_case_name` and `SA-26-069` do not.
_MATH_TOKEN = re.compile(
    r"(?<![A-Za-z0-9_])"
    r"(?P<base>[A-Za-zΑ-ω∂∇∫∑]|\d+)"
    r"(?P<op>[_^])"
    r"(?:\{(?P<braced>[^{}]{1,16})\}|(?P<plain>[-+]?[A-Za-z0-9]{1,3}))"
)
# A LaTeX control sequence that no run-splitting can repair (\frac, \mathrm...).
_LATEX_MACRO = re.compile(r"\\[A-Za-z]{2,}")
_MATH_DELIM = re.compile(r"\$[^$]{1,80}\$")

_SUB_BASELINE = "-25000"
_SUP_BASELINE = "30000"


def _scan_math(text: str, base_offset: int = 0, forced: str | None = None,
               drop: set | None = None, level: dict | None = None):
    """Mark, per character index, what the rendered form must do with it.

    ``drop`` collects the indices of the markup characters that disappear
    (``_``, ``^``, ``{``, ``}``); ``level`` maps a character index to the
    baseline it must carry.  Braced groups are scanned recursively and their
    inner tokens are flattened onto the outer level, because PowerPoint has no
    nested baseline: ``\\int_{S_f}`` renders as an integral with the subscript
    "Sf".
    """
    drop = set() if drop is None else drop
    level = {} if level is None else level
    for m in _MATH_TOKEN.finditer(text):
        script = m.group("braced") if m.group("braced") is not None else m.group("plain")
        if script is None:
            continue
        drop.add(base_offset + m.start("op"))
        lvl = forced or (_SUB_BASELINE if m.group("op") == "_" else _SUP_BASELINE)
        if forced:                       # inside an outer script: keep the base too
            level[base_offset + m.start("base")] = forced
        if m.group("braced") is not None:
            start, end = m.start("braced"), m.end("braced")
            drop.add(base_offset + start - 1)      # '{'
            drop.add(base_offset + end)            # '}'
            for i in range(start, end):
                level[base_offset + i] = lvl
            _scan_math(script, base_offset + start, forced=lvl, drop=drop, level=level)
        else:
            for i in range(m.start("plain"), m.end("plain")):
                level[base_offset + i] = lvl
    return drop, level


def _segment_math(text: str):
    """Split ``text`` into ``[(chunk, baseline_or_None), ...]``.

    ``baseline`` is ``None`` for ordinary text, ``_SUB_BASELINE`` for a
    subscript chunk and ``_SUP_BASELINE`` for a superscript chunk.  Returns
    ``None`` when the text carries no math markup (caller leaves it untouched).
    """
    drop, level = _scan_math(text)
    if not drop and not level:
        return None
    out = []
    for i, ch in enumerate(text):
        if i in drop:
            continue
        bl = level.get(i)
        if out and out[-1][1] == bl:
            out[-1][0] += ch
        else:
            out.append([ch, bl])
    return [(t, b) for t, b in out if t]


def _paragraph_runs(prs):
    """Yield ``(slide_no, shape_name, paragraph, runs, joined_text)``.

    Math markup is analysed per PARAGRAPH, not per run: PowerPoint splits a
    formula across runs at every font change, so ``Sigma_f`` routinely lands as
    a run ending in ``_`` followed by a run starting with ``f``.  Run-local
    analysis silently misses exactly those (observed on MMPM slide 5).
    """
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                runs = list(para.runs)
                if not runs:
                    continue
                yield i, shape.name, para, runs, "".join(r.text for r in runs)


def presentation_check_raw_math_markup(pptx_path: str,
                                       max_findings: int = 40) -> dict:
    """スライド本文に残った未整形の数式マークアップを検出する。

    原稿からコピーした式やAIが書いた式は ``sigma_f(x)``, ``M_ij``,
    ``int_{S_f}``, ``x_c``, ``r_nl`` のように下付き記号がソース表記のまま
    残りやすい。文字数・行数・可読性の既存lintは通ってしまい、PDFも
    「正しく」描画されるため、発表当日まで気付かない。

    検出する3種類:

    - ``subscript`` / ``superscript``: ``X_y`` ``X^{2}`` 形式。単一記号
      (ラテン/ギリシャ/演算子) または数値が基底のものだけを対象とし、
      ``snake_case`` の識別子や ``SA-26-069`` は除外する。
    - ``latex_macro``: ``\\frac`` ``\\mathrm`` などの制御綴。run分割では
      直せないため人手で書き換える。
    - ``math_delimiter``: ``$...$`` がそのまま残っている。

    Args:
        pptx_path: 点検する .pptx。
        max_findings: 返す指摘の上限。

    Returns: ``{"ok", "n_findings", "n_slides_affected", "findings",
    "repair_hint"}``。``findings`` は slide / shape / kind / text /
    ``rendered_preview`` を持つ。
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    path = pathlib.Path(pptx_path)
    if not path.exists():
        return {"error": f"file not found: {pptx_path}"}

    prs = _pptx.Presentation(str(path))
    findings = []
    slides_hit = set()
    for slide_no, shape_name, _para, _runs, text in _paragraph_runs(prs):
        if not text.strip():
            continue
        kinds = []
        segs = _segment_math(text)
        if segs:
            if any(b == _SUB_BASELINE for _t, b in segs):
                kinds.append("subscript")
            if any(b == _SUP_BASELINE for _t, b in segs):
                kinds.append("superscript")
        if _LATEX_MACRO.search(text):
            kinds.append("latex_macro")
        if _MATH_DELIM.search(text):
            kinds.append("math_delimiter")
        if not kinds:
            continue
        slides_hit.add(slide_no)
        if len(findings) < max_findings:
            findings.append({
                "slide": slide_no,
                "shape": shape_name,
                "kind": "+".join(sorted(set(kinds))),
                "text": text[:120],
                "rendered_preview": "".join(t for t, _b in segs) if segs else text[:120],
                "auto_repairable": bool(segs) and "latex_macro" not in kinds,
            })

    n = len(findings)
    return {
        "ok": n == 0,
        "n_findings": n,
        "n_slides_affected": len(slides_hit),
        "slides": sorted(slides_hit),
        "findings": findings,
        "repair_hint": (
            "auto_repairable の指摘は presentation_apply_math_subscripts で "
            "PowerPoint の下付き/上付き run へ変換できる。latex_macro は "
            "人手で書き換える（式が複雑なら図として貼る）。"
        ),
    }


def presentation_apply_math_subscripts(pptx_path: str,
                                       out_path: str = "",
                                       dry_run: bool = True,
                                       slides: list[int] | None = None) -> dict:
    """``X_y`` / ``X^{2}`` 形式のテキストを PowerPoint の下付き・上付き run へ変換する。

    run を分割し、下付き/上付き部分の ``a:rPr@baseline`` を設定する。元の
    run の書式（フォント・サイズ・色）は複製されるため見た目は変わらず、
    アンダースコアと波括弧だけが消える。

    ``\\frac`` のような LaTeX 制御綴は変換対象外で、``unhandled`` として
    報告する（run分割では表現できない）。

    Args:
        pptx_path: 入力 .pptx。
        out_path: 出力先。空文字なら入力を上書きする。
        dry_run: True（既定）なら書き込まず、変換予定だけを返す。
        slides: 対象スライド番号（1始まり）。None なら全スライド。

    Returns: ``{"ok", "n_runs_changed", "changes", "unhandled", "written"}``。
    """
    try:
        import pptx as _pptx
        from pptx.oxml.ns import qn
    except ImportError:
        return {"error": "python-pptx not installed."}
    path = pathlib.Path(pptx_path)
    if not path.exists():
        return {"error": f"file not found: {pptx_path}"}

    prs = _pptx.Presentation(str(path))
    targets = set(slides) if slides else None
    changes = []
    unhandled = []

    for slide_no, shape_name, _para, runs, text in list(_paragraph_runs(prs)):
        if targets and slide_no not in targets:
            continue
        if not text.strip():
            continue
        if _LATEX_MACRO.search(text):
            unhandled.append({"slide": slide_no, "shape": shape_name,
                              "text": text[:120], "reason": "LaTeX macro"})
            continue
        drop, level = _scan_math(text)
        if not drop and not level:
            continue
        changes.append({
            "slide": slide_no,
            "shape": shape_name,
            "before": text[:120],
            "after": "".join(t for t, _b in _segment_math(text))[:120],
            "n_segments": len(_segment_math(text)),
        })
        if dry_run:
            continue

        # Rewrite run by run using the PARAGRAPH-level plan, so a token split
        # across a font change (`Sigma_` + `f`) is repaired as one token.
        pos = 0
        for run in runs:
            r_text = run.text
            start, pos = pos, pos + len(r_text)
            if not r_text:
                continue
            segs = []
            for i, ch in enumerate(r_text, start=start):
                if i in drop:
                    continue
                bl = level.get(i)
                if segs and segs[-1][1] == bl:
                    segs[-1][0] += ch
                else:
                    segs.append([ch, bl])
            segs = [(t, b) for t, b in segs if t]
            if segs == [(r_text, None)]:
                continue                      # untouched run keeps its element
            r_el = run._r
            parent = r_el.getparent()
            idx = list(parent).index(r_el)
            for offset, (chunk, baseline) in enumerate(segs):
                new = deepcopy(r_el)
                t_el = new.find(qn("a:t"))
                t_el.text = chunk
                rPr = new.find(qn("a:rPr"))
                if baseline is None:
                    if rPr is not None and "baseline" in rPr.attrib:
                        del rPr.attrib["baseline"]
                else:
                    if rPr is None:
                        rPr = new.makeelement(qn("a:rPr"), {})
                        new.insert(0, rPr)
                    rPr.set("baseline", baseline)
                parent.insert(idx + offset, new)
            parent.remove(r_el)

    written = ""
    if not dry_run and changes:
        written = str(pathlib.Path(out_path) if out_path else path)
        prs.save(written)

    return {
        "ok": True,
        "dry_run": bool(dry_run),
        "n_paragraphs_changed": len(changes),
        "changes": changes,
        "unhandled": unhandled,
        "written": written,
    }


def presentation_check_duplicate_slide_images(pptx_path: str,
                                              min_area_fraction: float = 0.02
                                              ) -> dict:
    """同一画像が複数スライドで使い回されていないかを検査する。

    連続するスライドに同じ図が出ると、タイトルが違っても聴衆には「話が
    進んでいない」と映る。実例（MMPM SA-26-069, 2026-08-16）では
    「従来法の弱点」スライドと「モード分解」スライドが同一の立方体図を
    共有し、しかも貼り付け倍率が 0.847 と 0.888 で食い違っていた。

    ロゴや罫線のように繰り返して当然の小さな画像は
    ``min_area_fraction`` 未満として除外する。

    Args:
        pptx_path: 点検する .pptx。
        min_area_fraction: スライド面積に対するこの割合未満の画像は無視。

    Returns: ``{"ok", "n_duplicate_groups", "groups", "recommendation"}``。
    ``groups`` は sha1 / slides / displayed_cm / ``same_scale`` を持つ。
    """
    try:
        import pptx as _pptx
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Emu
    except ImportError:
        return {"error": "python-pptx not installed."}
    path = pathlib.Path(pptx_path)
    if not path.exists():
        return {"error": f"file not found: {pptx_path}"}

    prs = _pptx.Presentation(str(path))
    slide_area = float(Emu(prs.slide_width).pt) * float(Emu(prs.slide_height).pt)
    by_hash: dict[str, list[dict]] = {}

    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                blob = shape.image.blob
            except Exception:
                continue
            w_pt = float(Emu(shape.width).pt)
            h_pt = float(Emu(shape.height).pt)
            if slide_area and (w_pt * h_pt / slide_area) < min_area_fraction:
                continue
            key = hashlib.sha1(blob).hexdigest()[:12]
            by_hash.setdefault(key, []).append({
                "slide": i,
                "shape": shape.name,
                "displayed_cm": round(float(Emu(shape.width).cm), 2),
            })

    groups = []
    for key, uses in by_hash.items():
        slides_used = sorted({u["slide"] for u in uses})
        if len(slides_used) < 2:
            continue
        widths = {u["displayed_cm"] for u in uses}
        groups.append({
            "sha1": key,
            "slides": slides_used,
            "n_uses": len(uses),
            "displayed_cm": sorted(widths),
            "same_scale": len(widths) == 1,
            "adjacent": any(b - a == 1 for a, b in zip(slides_used, slides_used[1:])),
        })
    groups.sort(key=lambda g: g["slides"])

    return {
        "ok": not groups,
        "n_duplicate_groups": len(groups),
        "groups": groups,
        "recommendation": (
            "同一図の再掲は、そのスライド固有の主張を持つ図に差し替えるか、"
            "再掲であることを明示する。貼り付け倍率が食い違う場合は "
            "figure_audit_pptx_figures で実寸を確認する。"
        ),
    }
