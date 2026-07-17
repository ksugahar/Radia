"""presentation/_citations.py -- cite references in research talks.

Added 2026-06-02 on user request ("I like to cite references in my
PowerPoint talks -- add it to radia_mcp.presentation").

Four MCP tools (auto-registered because tools.py re-imports the
``presentation_*`` names from here):

    presentation_cite_format(...)        format ONE reference into the
                                         talk styles (footnote / bracket /
                                         numeric / full).  Accepts a raw
                                         BibTeX entry OR explicit fields.
    presentation_references_slide(...)   build a "References" frame
                                         (beamer thebibliography) AND a
                                         plain numbered list (for a PPTX
                                         text box) from a list of cites.
    presentation_add_citation_footer(..) EXECUTING: drop a small
                                         Times-New-Roman footnote textbox
                                         along the bottom of one slide of
                                         an existing .pptx.
    presentation_citation_audit(...)     lint: every ``[N]`` used on the
                                         slides has a matching reference,
                                         and flag references never cited.

Design notes
------------
* No heavy imports at module load (python-pptx is lazy-imported only in
  ``presentation_add_citation_footer``) so ``mcp-server-paper-writing
  --selftest`` (which serves the merged presentation_* tools since
  2026-07-17) runs on a machine without a PPTX stack.
* Author-string contract (predictable, no silent guessing): separate
  multiple authors with ``" and "`` (BibTeX convention) or ``";"``.  A
  single ``", "`` is read as one ``"Last, First"`` BibTeX name.  Pass a
  pre-formatted display string if you want exact control.
* Times New Roman is the lab figure/slide font, so the PPTX footer uses
  it to match (consistent with radia_mcp.figure and pptx_drawing).
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Optional, Union


# ------------------------------------------------------------------
# Author-string helpers (documented contract above)
# ------------------------------------------------------------------
def _split_authors(raw: str) -> list[str]:
    raw = (raw or "").strip()
    if not raw:
        return []
    if " and " in raw:
        return [t.strip() for t in raw.split(" and ") if t.strip()]
    if ";" in raw:
        return [t.strip() for t in raw.split(";") if t.strip()]
    # No multi-author separator: treat the whole thing as ONE name
    # (possibly "Last, First").
    return [raw]


def _disp_name(tok: str) -> str:
    """Normalize one author token to "F. M. Last" display form.

    "Sugahara, Kengo"  -> "K. Sugahara"
    "K. Sugahara"      -> "K. Sugahara"   (already display form, kept)
    "Kengo Sugahara"   -> "Kengo Sugahara" (kept; we do not force initials
                          on an already-spelled-out display name)
    """
    tok = tok.strip()
    if "," in tok:
        last, first = (x.strip() for x in tok.split(",", 1))
        initials = " ".join(
            p[0] + "." for p in first.replace(".", " ").split() if p
        )
        return f"{initials} {last}".strip()
    return tok


def _first_last(tok: str) -> str:
    """Last name of the first author (for "[Last Year]" bracket style)."""
    tok = tok.strip()
    if "," in tok:
        return tok.split(",", 1)[0].strip()
    parts = tok.split()
    return parts[-1] if parts else tok


def _authors_short(raw: str) -> str:
    """First-author-et-al. form for a talk footnote."""
    toks = _split_authors(raw)
    if not toks:
        return ""
    if len(toks) == 1:
        return _disp_name(toks[0])
    if len(toks) == 2:
        return f"{_disp_name(toks[0])} and {_disp_name(toks[1])}"
    return f"{_disp_name(toks[0])} et al."


def _authors_full(raw: str) -> str:
    """All authors, Oxford-comma joined, for the full reference line."""
    toks = [_disp_name(t) for t in _split_authors(raw)]
    if not toks:
        return ""
    if len(toks) == 1:
        return toks[0]
    if len(toks) == 2:
        return f"{toks[0]} and {toks[1]}"
    return ", ".join(toks[:-1]) + f", and {toks[-1]}"


_BIBTEX_FIELD_RE = re.compile(r'(\w+)\s*=\s*[{"]([^{}"]*)[}"]')


def _extract_bibtex_fields(bibtex: str) -> dict:
    """Light BibTeX-entry field extractor (no nested-brace support; good
    enough for a one-line slide citation)."""
    d = {}
    for m in _BIBTEX_FIELD_RE.finditer(bibtex or ""):
        d[m.group(1).lower()] = m.group(2).strip()
    return d


# ------------------------------------------------------------------
# Tool 1: format one reference
# ------------------------------------------------------------------
def presentation_cite_format(
    bibtex: str = "",
    authors: str = "",
    year: str = "",
    venue: str = "",
    title: str = "",
    number: int = 0,
) -> dict:
    """Format ONE reference into the styles used on talk slides.

    Provide EITHER a raw ``bibtex`` entry OR the explicit fields
    (``authors`` / ``year`` / ``venue`` / ``title``).  Explicit fields
    win over the parsed BibTeX when both are given.

    Args:
        bibtex: a single BibTeX entry, e.g.
            ``@article{key, author={Sugahara, Kengo and Ida, Akihiro},
            title={...}, journal={IEEE Trans. Magn.}, year={2026}}``.
        authors: author string; separate authors with " and " or ";"
            (see module contract).  A lone "Last, First" is one author.
        year: publication year.
        venue: journal or conference name (e.g. "IEEE Trans. Magn.").
        title: paper title (only used by the 'full' style).
        number: if > 0, also emit the "[N]" numeric marker for a
            numbered reference list.

    Returns a dict of formatted strings:
        footnote / talk : "K. Sugahara et al., IEEE Trans. Magn., 2026."
        bracket         : "[Sugahara 2026]"
        numeric         : "[3]"  (only if number > 0)
        full            : 'K. Sugahara, A. Ida, and H. Yano, "Title,"
                           IEEE Trans. Magn., 2026.'
        recommended     : same as 'footnote' (the slide-bottom default)
    """
    f = _extract_bibtex_fields(bibtex) if bibtex else {}
    authors = authors or f.get("author", "")
    year = year or f.get("year", "")
    venue = venue or f.get("journal", "") or f.get("booktitle", "") \
        or f.get("publisher", "")
    title = title or f.get("title", "")

    short = _authors_short(authors)
    full_auth = _authors_full(authors)
    first_last = _first_last(_split_authors(authors)[0]) if authors else ""

    # footnote / talk: authors-short, venue, year
    parts = [p for p in (short, venue, year) if p]
    footnote = ", ".join(parts)
    if footnote and not footnote.endswith("."):
        footnote += "."

    # full (IEEE): A. Author, B. Writer, and C. Co, "Title," Venue, year.
    # The title carries its trailing comma INSIDE the closing quote and is
    # followed by a SPACE (not a comma) before the venue.
    full = full_auth or ""
    if title:
        t_clean = title.rstrip(". ,")
        full = (full + ", " if full else "") + f'"{t_clean},"'
        following = [x for x in (venue, year) if x]
        if following:
            full += " " + ", ".join(following)
    else:
        rest = [x for x in (venue, year) if x]
        if rest:
            full = (full + ", " if full else "") + ", ".join(rest)
    full = full.strip()
    if full and not full.endswith("."):
        full += "."

    bracket = f"[{first_last} {year}]".replace(" ]", "]") if first_last else ""

    out = {
        "footnote": footnote,
        "talk": footnote,
        "bracket": bracket,
        "full": full,
        "recommended": footnote,
        "parsed_fields": {
            "authors": authors, "year": year,
            "venue": venue, "title": title,
        },
    }
    if number and number > 0:
        out["numeric"] = f"[{int(number)}]"
    return out


# ------------------------------------------------------------------
# Tool 2: build a References frame / list
# ------------------------------------------------------------------
def presentation_references_slide(
    references: list,
    fmt: str = "both",
    frame_title: str = "References",
    font_size: str = "footnotesize",
) -> dict:
    """Build a "References" slide from a list of full citation lines.

    Args:
        references: list of full citation strings (one per reference),
            in the order you want them numbered [1], [2], ...  You can
            produce each line with presentation_cite_format(... )['full'].
        fmt: 'beamer' | 'pptx' | 'both' (default).
        frame_title: frame title (beamer).
        font_size: beamer font command WITHOUT the backslash
            ('footnotesize' | 'scriptsize' | 'tiny' | 'small').

    Returns dict with:
        beamer : a complete ``\\begin{frame}[allowframebreaks]...`` block
                 using ``thebibliography`` so [N] matches numeric cites.
        plain  : a newline-joined "[1] ...\\n[2] ..." list to paste into
                 a PPTX text box.
        n      : number of references.
    """
    refs = [str(r).strip() for r in (references or []) if str(r).strip()]
    out: dict = {"n": len(refs)}

    if fmt in ("beamer", "both"):
        fs = (font_size or "footnotesize").lstrip("\\")
        lines = [
            f"\\begin{{frame}}[allowframebreaks]{{{frame_title}}}",
            f"\\{fs}",
            "\\begin{thebibliography}{99}",
        ]
        for i, r in enumerate(refs, 1):
            lines.append(f"  \\bibitem{{ref{i}}} {r}")
        lines += ["\\end{thebibliography}", "\\end{frame}"]
        out["beamer"] = "\n".join(lines)

    if fmt in ("pptx", "plain", "both"):
        out["plain"] = "\n".join(f"[{i}] {r}" for i, r in enumerate(refs, 1))

    return out


# ------------------------------------------------------------------
# Tool 3: drop a citation footnote on a PPTX slide (EXECUTING)
# ------------------------------------------------------------------
def presentation_add_citation_footer(
    pptx_path: str,
    text: str,
    slide_index: int = 0,
    out_path: str = "",
    font_pt: float = 9.0,
    italic: bool = True,
    color_rgb: Optional[list] = None,
) -> dict:
    """Add a small citation footnote textbox along the BOTTOM of one
    slide of an existing .pptx (Times New Roman, to match lab figures).

    Args:
        pptx_path: input .pptx.
        text: the citation text (e.g. the 'footnote' from
            presentation_cite_format).
        slide_index: 0-based slide to annotate.
        out_path: where to write; defaults to overwriting ``pptx_path``.
        font_pt: footnote font size (default 9 pt).
        italic: italicize (default True).
        color_rgb: [r,g,b] 0-255; default a muted grey (90,90,90).

    Returns dict {ok, output, slide_index, text} or {ok: False, error}.
    """
    try:
        from pptx import Presentation
        from pptx.util import Mm, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        return {"ok": False,
                "error": "python-pptx is required: pip install python-pptx"}

    if not (text or "").strip():
        return {"ok": False, "error": "empty citation text"}
    if not Path(pptx_path).is_file():
        return {"ok": False, "error": f"pptx not found: {pptx_path}"}

    color = tuple(color_rgb) if color_rgb else (90, 90, 90)
    prs = Presentation(pptx_path)
    n = len(prs.slides)
    if slide_index < 0 or slide_index >= n:
        return {"ok": False,
                "error": f"slide_index {slide_index} out of range (0..{n - 1})"}

    slide = prs.slides[slide_index]
    sw_mm = prs.slide_width / 36000.0     # EMU -> mm (1 mm = 36000 EMU)
    sh_mm = prs.slide_height / 36000.0
    margin = 6.0
    box_h = max(5.0, font_pt / 72.0 * 25.4 * 1.7)

    tb = slide.shapes.add_textbox(
        Mm(margin), Mm(sh_mm - box_h - 2.0),
        Mm(sw_mm - 2 * margin), Mm(box_h),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Times New Roman"
            r.font.size = Pt(font_pt)
            r.font.italic = bool(italic)
            r.font.color.rgb = RGBColor(*color)

    out = out_path or pptx_path
    Path(out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return {"ok": True, "output": str(Path(out).resolve()),
            "slide_index": slide_index, "text": text}


# ------------------------------------------------------------------
# Tool 4: citation-consistency lint
# ------------------------------------------------------------------
_CITE_MARKER_RE = re.compile(r"\[(\d{1,3})\]")


def presentation_citation_audit(
    slides_text: str,
    references: Union[int, list],
) -> dict:
    """Check numeric ``[N]`` citations on the slides against the
    references list.

    Args:
        slides_text: all slide text concatenated (e.g. from
            presentation_extract_pptx_text, or the .tex source).
        references: the references, as a list (its length = the highest
            valid number) OR an integer count.

    Returns dict:
        cited                  : sorted unique [N] used on slides.
        n_references           : number of reference entries.
        cited_without_entry    : [N] used but with no matching entry
                                 (DANGLING -- fix before the talk).
        references_never_cited : entry numbers never referenced on a
                                 slide (candidate to drop).
        ok                     : True iff no dangling citation.
    """
    cited = sorted({int(m) for m in _CITE_MARKER_RE.findall(slides_text or "")})
    if isinstance(references, int):
        n_refs = max(0, references)
    else:
        n_refs = len(references or [])
    ref_nums = list(range(1, n_refs + 1))
    missing = [n for n in cited if n not in ref_nums]
    uncited = [n for n in ref_nums if n not in cited]
    return {
        "cited": cited,
        "n_references": n_refs,
        "cited_without_entry": missing,
        "references_never_cited": uncited,
        "ok": not missing,
    }
