"""presentation/pptx_drawing.py -- python-pptx wrappers for the lab
PowerPoint-as-figure-editor workflow.

Promoted on 2026-05-26 from `public-safe curated corpus`
where the original 6 scripts lived (`ppt2png.py`, `作図.py`, `期末.py`,
`正弦波.py`, `色をいろいろ.py`, `ppt2pdf.py`).  The originals used
`win32com.client` (Windows-only, requires PowerPoint installed); the
cross-platform equivalents below use `python-pptx` so the same
patterns run in CI / Linux / macOS.

Lab workflow:
  1. Generate .pptx programmatically here (or via the presentation_pptx_*
     MCP tools that wrap these helpers).
  2. Open in PowerPoint, tweak interactively (snap-to-grid, alignment,
     rotation, grouping).
  3. Export the finished slide to PNG / PDF.

Conventions:

  - All lengths in **millimeters** (lab natural unit).  python-pptx
    needs EMU internally; the helpers convert.
  - RGB colors as (r, g, b) tuples 0-255, NOT the COM packed-int format
    (the original scripts used `r + g*256 + b*65536`).
  - Slide size defaults to A4 landscape (297 x 210 mm), but every
    helper takes width_mm / height_mm overrides.

References:
  - python-pptx docs: https://python-pptx.readthedocs.io
  - PowerPoint AutoShape index: see `MSO_SHAPE` enum (140+ shapes).
    Lab scripts used numeric shape codes 1-99 (OOXML enum); the
    `_SHAPE_BY_NAME` map below covers the 20 most common.
"""
from __future__ import annotations

from pathlib import Path
from typing import Iterable, Optional


# Lazy import of python-pptx -- the rest of radia_mcp does
# not require it, so import inside each function call (consistent with
# tools.py pattern at line 583).
def _get_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Mm, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
        return Presentation, Mm, Pt, Emu, RGBColor, MSO_SHAPE, MSO_CONNECTOR
    except ImportError as e:
        raise RuntimeError(
            "python-pptx is required for presentation_pptx_* tools. "
            "Install via: pip install python-pptx"
        ) from e


# ------------------------------------------------------------------
# Common AutoShape lookup -- 20 most-used by the lab
# ------------------------------------------------------------------
_SHAPE_ALIASES = {
    # MSO_SHAPE enum names; user can pass either the lab-friendly
    # short name or the python-pptx enum member name.
    "rectangle":            "RECTANGLE",
    "rect":                 "RECTANGLE",
    "rounded_rectangle":    "ROUNDED_RECTANGLE",
    "rounded_rect":         "ROUNDED_RECTANGLE",
    "oval":                 "OVAL",
    "ellipse":              "OVAL",
    "circle":               "OVAL",
    "diamond":              "DIAMOND",
    "triangle":             "ISOCELES_TRIANGLE",
    "right_triangle":       "RIGHT_TRIANGLE",
    "trapezoid":            "TRAPEZOID",
    "parallelogram":        "PARALLELOGRAM",
    "pentagon":             "PENTAGON",
    "hexagon":              "HEXAGON",
    "star_5":               "STAR_5_POINT",
    "star_6":               "STAR_6_POINT",
    "arrow_right":          "RIGHT_ARROW",
    "arrow_left":           "LEFT_ARROW",
    "arrow_up":             "UP_ARROW",
    "arrow_down":           "DOWN_ARROW",
    "block_arc":            "BLOCK_ARC",
    "donut":                "DONUT",
    "callout":              "RECTANGULAR_CALLOUT",
    "flowchart_process":    "FLOWCHART_PROCESS",
    "flowchart_decision":   "FLOWCHART_DECISION",
    "line":                 "LINE_INVERSE",        # placeholder; use add_line
}


def _resolve_shape(name: str):
    """Resolve a user-friendly name OR python-pptx enum member name to
    an MSO_SHAPE enum value.  Case-insensitive."""
    _, _, _, _, _, MSO_SHAPE, _ = _get_pptx()
    key = name.strip().lower()
    if key in _SHAPE_ALIASES:
        member = _SHAPE_ALIASES[key]
    else:
        # try direct enum member match (case-insensitive)
        for m in MSO_SHAPE.__members__:
            if m.lower() == key:
                return MSO_SHAPE.__members__[m]
        member = key.upper()
    return MSO_SHAPE.__members__[member]


# ==================================================================
# CANVAS
# ==================================================================


def create_canvas(width_mm: float = 297, height_mm: float = 210):
    """Create a blank Presentation with one slide at the given physical
    size (defaults to A4 landscape).

    Returns the Presentation object; caller adds shapes via the
    `add_shape` / `add_line` / `add_text_box` / `add_freeform` helpers
    using the returned `prs.slides[0]` reference.

    Note: PowerPoint's canvas-size scaling caps width / height at 142 cm
    each (max EMU value); for wider scientific layouts (e.g. the lab's
    2000-pt = 70-cm-wide `作図.py` canvas), reduce or use a multi-slide
    composition.
    """
    Presentation, Mm, *_ = _get_pptx()
    prs = Presentation()
    prs.slide_width = Mm(width_mm)
    prs.slide_height = Mm(height_mm)
    # Use the BLANK layout (no placeholders).
    # Layouts vary by template; pick index 6 (BLANK) for the default
    # Office theme; fall back to last layout if missing.
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 \
        else prs.slide_layouts[-1]
    prs.slides.add_slide(blank_layout)
    return prs


# ==================================================================
# ADD SHAPES
# ==================================================================


def add_shape(slide, shape_name: str,
              x_mm: float, y_mm: float,
              w_mm: float, h_mm: float,
              fill_rgb=(255, 255, 255),
              line_rgb=(0, 0, 0),
              line_weight_pt: float = 0.75,
              text: Optional[str] = None,
              text_rgb=(0, 0, 0),
              text_size_pt: float = 10.0):
    """Add an AutoShape to a slide.

    Args:
        slide: prs.slides[i] target.
        shape_name: friendly name ('rect', 'oval', 'arrow_right',
            ...) OR python-pptx MSO_SHAPE enum member name.
        x_mm, y_mm: top-left position in mm from slide top-left.
        w_mm, h_mm: dimensions in mm.
        fill_rgb: (r, g, b) 0-255.  Pass None to skip fill.
        line_rgb: (r, g, b) for outline.  None to skip outline.
        line_weight_pt: outline width in points.
        text: optional text inside the shape (most AutoShapes accept).
        text_rgb: text color.
        text_size_pt: text size in points.

    Returns the python-pptx shape object.
    """
    _, Mm, Pt, _, RGBColor, _, _ = _get_pptx()
    shp = slide.shapes.add_shape(
        _resolve_shape(shape_name),
        Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm),
    )
    if fill_rgb is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(*fill_rgb)
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = RGBColor(*line_rgb)
        shp.line.width = Pt(line_weight_pt)
    if text is not None:
        tf = shp.text_frame
        tf.text = str(text)
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(text_size_pt)
                r.font.color.rgb = RGBColor(*text_rgb)
    return shp


def add_line(slide,
             x1_mm: float, y1_mm: float,
             x2_mm: float, y2_mm: float,
             line_rgb=(0, 0, 0),
             line_weight_pt: float = 0.75,
             arrow_end: bool = False,
             arrow_start: bool = False):
    """Add a straight line connector between two points.

    Args:
        slide: prs.slides[i] target.
        (x1,y1) - (x2,y2): endpoints in mm.
        line_rgb, line_weight_pt: as in add_shape.
        arrow_end: draw an arrowhead at the (x2, y2) end.
        arrow_start: draw an arrowhead at the (x1, y1) end.

    Returns the connector shape.
    """
    _, Mm, Pt, _, RGBColor, _, MSO_CONNECTOR = _get_pptx()
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Mm(x1_mm), Mm(y1_mm), Mm(x2_mm), Mm(y2_mm),
    )
    conn.line.color.rgb = RGBColor(*line_rgb)
    conn.line.width = Pt(line_weight_pt)
    # python-pptx exposes arrows via XML directly.
    if arrow_end or arrow_start:
        from pptx.oxml.ns import qn
        ln = conn.line._get_or_add_ln()
        if arrow_end:
            tail = ln.makeelement(qn("a:tailEnd"),
                                   {"type": "triangle"}, nsmap={})
            ln.append(tail)
        if arrow_start:
            head = ln.makeelement(qn("a:headEnd"),
                                   {"type": "triangle"}, nsmap={})
            ln.append(head)
    return conn


def add_text_box(slide,
                 text: str,
                 x_mm: float, y_mm: float,
                 w_mm: float, h_mm: float,
                 text_rgb=(0, 0, 0),
                 text_size_pt: float = 10.0,
                 bold: bool = False,
                 italic: bool = False,
                 font_name: str = "Times New Roman"):
    """Add a stand-alone text box (no shape outline).

    Defaults to Times New Roman 10 pt -- matches lab paper-figure
    convention (`radia_mcp.figure` profile rule).
    """
    _, Mm, Pt, _, RGBColor, _, _ = _get_pptx()
    tb = slide.shapes.add_textbox(Mm(x_mm), Mm(y_mm),
                                    Mm(w_mm), Mm(h_mm))
    tf = tb.text_frame
    tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = font_name
            r.font.size = Pt(text_size_pt)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = RGBColor(*text_rgb)
    return tb


# ==================================================================
# FREEFORM PATH
# ==================================================================


def add_freeform(slide,
                 points_mm: Iterable[tuple[float, float]],
                 line_rgb=(0, 0, 0),
                 line_weight_pt: float = 0.75,
                 fill_rgb=None,
                 close: bool = False):
    """Draw a freeform (poly-line / polygon) through a list of (x, y) mm.

    Matches the lab's `BuildFreeform + AddNodes` pattern (sine wave in
    `正弦波.py`, multi-segment curve in `作図.py`).  The first point
    is the start anchor; subsequent points are line segments.

    Args:
        slide: prs.slides[i] target.
        points_mm: iterable of (x_mm, y_mm) tuples.
        line_rgb / line_weight_pt: outline.
        fill_rgb: pass an (r, g, b) to fill the closed polygon; None
            leaves it stroke-only.
        close: True closes the path back to the first point (true
            polygon); False leaves open polyline.

    Returns the freeform shape.
    """
    _, Mm, Pt, _, RGBColor, _, _ = _get_pptx()
    pts = list(points_mm)
    if len(pts) < 2:
        raise ValueError("add_freeform: need at least 2 points")
    x0, y0 = pts[0]
    builder = slide.shapes.build_freeform(Mm(x0), Mm(y0))
    # add_line_segments takes a list of (x, y) EMUs OR Length, plus a
    # `close` flag.  We pass mm-as-EMU via Mm().
    rest = [(Mm(x), Mm(y)) for x, y in pts[1:]]
    builder.add_line_segments(rest, close=close)
    shp = builder.convert_to_shape()
    if fill_rgb is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shp.line.color.rgb = RGBColor(*line_rgb)
    shp.line.width = Pt(line_weight_pt)
    return shp


# ==================================================================
# SAVE
# ==================================================================


def save_pptx(prs, path: str) -> str:
    """Save the Presentation to disk; returns the absolute path."""
    p = Path(path).resolve()
    p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(p))
    return str(p)


# ==================================================================
# EXPORT TO PNG / PDF (requires LibreOffice or PowerPoint COM)
# ==================================================================


def export_slide_to_png(pptx_path: str,
                         out_path: str,
                         slide_index: int = 0,
                         width_px: int = 1920,
                         height_px: int = 1080,
                         backend: str = "auto") -> dict:
    """Export one slide of an existing .pptx to PNG.

    python-pptx itself cannot render to PNG -- it only writes .pptx.
    This function tries (in priority order):

      1. PowerPoint COM (Windows + Office installed) -- highest fidelity
      2. LibreOffice headless (`soffice --convert-to png ...`) -- cross-
         platform, lower fidelity (no fancy fonts / kerning)

    Args:
        pptx_path: input .pptx
        out_path: output .png
        slide_index: 0-based slide number to export.
        width_px / height_px: target raster size (LibreOffice ignores
            these; PowerPoint COM honours them).
        backend: 'auto' | 'com' | 'libreoffice'

    Returns dict {'output': path, 'backend': str, 'ok': bool, 'message': str}.
    """
    import os
    import subprocess
    import tempfile

    pptx_abs = str(Path(pptx_path).resolve())
    out_abs = str(Path(out_path).resolve())

    def _try_com():
        try:
            import win32com.client  # type: ignore
        except ImportError:
            return None
        try:
            ppt = win32com.client.Dispatch("PowerPoint.Application")
            ppt.Visible = 1
            ppt.Presentations.Open(pptx_abs)
            ppt.ActivePresentation.Slides(slide_index + 1).Export(
                out_abs, "png", width_px, height_px)
            ppt.ActivePresentation.Close()
            ppt.Quit()
            return {"output": out_abs, "backend": "com",
                    "ok": True, "message": "Exported via PowerPoint COM."}
        except Exception as e:
            return {"ok": False, "backend": "com",
                    "message": f"PowerPoint COM failed: {e}"}

    def _try_libreoffice():
        soffice = None
        for cand in ("soffice", "libreoffice",
                      r"C:\Program Files\LibreOffice\program\soffice.exe",
                      r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"):
            from shutil import which
            if which(cand) or os.path.isfile(cand):
                soffice = cand
                break
        if soffice is None:
            return {"ok": False, "backend": "libreoffice",
                    "message": "LibreOffice (soffice) not found on PATH."}
        with tempfile.TemporaryDirectory() as td:
            try:
                subprocess.run(
                    [soffice, "--headless", "--convert-to", "png",
                      "--outdir", td, pptx_abs],
                    check=True, capture_output=True, timeout=60,
                )
            except subprocess.CalledProcessError as e:
                return {"ok": False, "backend": "libreoffice",
                        "message": f"LibreOffice conversion failed: "
                                    f"{e.stderr.decode(errors='replace')[:300]}"}
            except subprocess.TimeoutExpired:
                return {"ok": False, "backend": "libreoffice",
                        "message": "LibreOffice timed out."}
            # LibreOffice produces <basename>.png next to source.
            produced = Path(td) / (Path(pptx_abs).stem + ".png")
            if not produced.exists():
                return {"ok": False, "backend": "libreoffice",
                        "message": f"Expected {produced.name} not produced."}
            # Move to requested location.
            Path(out_abs).parent.mkdir(parents=True, exist_ok=True)
            produced.replace(out_abs)
            note = ("LibreOffice converted only the FIRST slide; "
                    "slide_index argument was ignored."
                    if slide_index != 0 else "")
            return {"output": out_abs, "backend": "libreoffice",
                    "ok": True,
                    "message": f"Exported via LibreOffice headless. {note}"}

    if backend == "com":
        r = _try_com()
        return r or {"ok": False, "backend": "com",
                      "message": "pywin32 not installed (Windows only)."}
    if backend == "libreoffice":
        return _try_libreoffice()
    # 'auto': COM first (best fidelity), then LibreOffice
    r = _try_com()
    if r and r["ok"]:
        return r
    return _try_libreoffice()


# ==================================================================
# UNIT CONVERSION HELPERS (the lab uses mm)
# ==================================================================

def mm_to_emu(mm: float) -> int:
    """mm -> python-pptx Emu (1 EMU = 1/360000 cm = 1/914400 in)."""
    return int(mm * 360000.0)


def mm_to_pt(mm: float) -> float:
    """mm -> points (1 pt = 1/72 in)."""
    return mm / 25.4 * 72.0


def pt_to_mm(pt: float) -> float:
    """points -> mm."""
    return pt / 72.0 * 25.4
