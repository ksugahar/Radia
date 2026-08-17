"""Tool functions for the presentation domain of mcp-server-document.

All functions are plain callables; they get wrapped by @mcp.tool()
in the top-level server.py. No FastMCP import here (that would
accidentally register a second MCP server).
"""

from __future__ import annotations

import json
import pathlib
import re

# Plan B Tier 1 (v0.12.0) — composite score + human-advisor comments
from .plans.T1 import presentation_opening_hook_strength  # noqa: F401
from .plans.T2 import presentation_takehome_strength  # noqa: F401
from .plans.T3 import presentation_check_pie_3d_charts  # noqa: F401
from .plans.T4 import presentation_check_logo_on_every_slide  # noqa: F401
from .plans.T5 import presentation_check_progress_indicator  # noqa: F401

# Plan B Tier 2 (v0.13.0)
from .plans.T6 import presentation_visual_text_ratio_score  # noqa: F401
from .plans.T7 import presentation_speaker_note_ratio  # noqa: F401
from .plans.T8 import presentation_font_consistency  # noqa: F401
from .plans.T9 import presentation_arrow_usage  # noqa: F401
from .plans.T10 import presentation_check_underline_in_pptx  # noqa: F401
from .plans.T11 import presentation_slide_density_balance  # noqa: F401

# Plan B Meta (v0.14.0)
from .plans.T12 import presentation_health_report  # noqa: F401

# Plan B Tier 3 (v0.21.0) — outline coherence / title-body / single-message /
# 欧米 text-heavy / mini-IMRAD / 理系 minimalism / chart simplification
from .plans.T13 import presentation_slide_titles_outline_coherence  # noqa: F401
from .plans.T14 import presentation_title_body_alignment_check  # noqa: F401
from .plans.T15 import presentation_single_message_per_slide_semantic  # noqa: F401
from .plans.T16 import presentation_text_density_per_slide_western_style  # noqa: F401
from .plans.T17 import presentation_mini_imrad_structure_check  # noqa: F401
from .plans.T18 import presentation_rikei_minimalism_score  # noqa: F401
from .plans.T19 import presentation_chart_simplification_check  # noqa: F401

# Plan B Tier 4 (v0.21.0) — 理系プレゼン特化 (equation / figure / results)
from .plans.T20 import presentation_equation_slide_compliance  # noqa: F401
from .plans.T21 import presentation_figure_slide_compliance  # noqa: F401
from .plans.T22 import presentation_results_slide_statistical_evidence  # noqa: F401

# Plan B Tier 5 / Intelligence Layer (v0.25.0) — synthesis / orchestration
from .plans.T23 import presentation_root_cause_diagnosis  # noqa: F401
from .plans.T24 import presentation_next_5_actions  # noqa: F401
from .plans.T25 import presentation_run_full_workflow  # noqa: F401
from .plans.T26 import presentation_rewrite_suggest  # noqa: F401
from .plans.T27 import presentation_adaptive_health_report  # noqa: F401

# Plan B Tier 6 (v0.26.0) — verbal side 強化
from .plans.T28 import presentation_speaking_pace_estimate  # noqa: F401
from .plans.T29 import presentation_qa_anticipation_list  # noqa: F401
from .plans.T30 import presentation_script_vs_slide_coverage  # noqa: F401
from .plans.T31 import presentation_embed_tts_audio_in_pptx  # noqa: F401

# Talk feedback loop (CEFC/Compumag field notes) — learned Q&A catalog,
# the presentation analog of meta/bug_patterns.py (2026-06)
from .talk_feedback import (  # noqa: F401
    presentation_talk_feedback_lookup,
    presentation_qa_from_history,
    presentation_talk_feedback_stats,
)

# Cross-module JA-lint (v0.13.0) — 台本・スライドテキストに grant の和文 lint
from .cross_lint import (  # noqa: F401
    presentation_check_notation_variants,
    presentation_find_undefined_acronyms,
    presentation_acronym_usage_audit,
    presentation_check_kanji_ratio,
    presentation_lint_bedrock,
    presentation_check_misuse_japanese,
    presentation_suggest_redundancy_fixes,
)

# Citing references on talk slides (2026-06-02) — format / references-frame /
# PPTX footnote insertion / citation-consistency lint.
from ._citations import (  # noqa: F401
    presentation_cite_format,
    presentation_references_slide,
    presentation_add_citation_footer,
    presentation_citation_audit,
)
from ._figure_text_repair import (  # noqa: F401
    presentation_replace_embedded_figure_text,
)

# Deck integrity (2026-08-16, MMPM SA-26-069): unrendered math markup left on
# the slide (`X_y`, `int_{S_f}`) and the same artwork reused across slides --
# two classes every text/layout lint above walked past.
from ._deck_integrity import (  # noqa: F401
    presentation_check_raw_math_markup,
    presentation_apply_math_subscripts,
    presentation_check_duplicate_slide_images,
)


_HERE = pathlib.Path(__file__).resolve().parent
KNOWLEDGE = _HERE


def _load_skill() -> str:
    return (KNOWLEDGE / "skill.md").read_text(encoding="utf-8")


# ------------------------------------------------------------------
# Shared pptx helpers (module-level to keep DRY across checks)
# ------------------------------------------------------------------
def _slide_title(slide) -> str:
    """Return the slide title robustly.

    Bug H3 fix: the previous ad-hoc implementation picked the first
    non-empty ``shape.text_frame`` in shape order, which fails when the
    first shape happens to be a body textbox. We now prefer the title
    placeholder (``placeholder_format.idx == 0``) and only fall back to
    shape-order iteration if no title placeholder exists.

    Returns an empty string if nothing text-like is found.
    """
    # 1) Prefer the real title placeholder (idx == 0).
    try:
        for shape in slide.shapes:
            try:
                pf = shape.placeholder_format
            except Exception:
                pf = None
            if pf is None:
                continue
            try:
                if pf.idx == 0 and shape.has_text_frame:
                    txt = shape.text_frame.text
                    if txt.strip():
                        return txt.splitlines()[0]
            except Exception:
                continue
    except Exception:
        pass
    # 2) Fall back to the first non-empty text frame in shape order.
    try:
        for shape in slide.shapes:
            try:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    return shape.text_frame.text.splitlines()[0]
            except Exception:
                continue
    except Exception:
        pass
    return ""


def _walk_shapes(shape_collection):
    """Yield all non-group shapes recursively from a shape collection.

    Bug H7 fix: ``slide.shapes`` does NOT descend into group shapes, so
    any text nested inside a GROUP was silently skipped. We recurse into
    ``MSO_SHAPE_TYPE.GROUP`` collections.
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except ImportError:  # pragma: no cover - caller guards the import
        MSO_SHAPE_TYPE = None  # type: ignore

    for s in shape_collection:
        try:
            is_group = (MSO_SHAPE_TYPE is not None
                        and s.shape_type == MSO_SHAPE_TYPE.GROUP)
        except Exception:
            is_group = False
        if is_group:
            try:
                yield from _walk_shapes(s.shapes)
            except Exception:
                yield s
        else:
            yield s


def _shape_is_slide_title(shape, slide) -> bool:
    """Return whether *shape* is the audience-facing slide title."""
    try:
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            return True
    except Exception:
        pass
    try:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            return False
        first_line = shape.text_frame.text.splitlines()[0].strip()
        if first_line != _slide_title(slide).strip():
            return False
        slide_h = float(slide.part.package.presentation_part.presentation.slide_height)
        return float(shape.top) < 0.20 * slide_h
    except Exception:
        return False


def _is_footer_or_page_chrome(shape, slide,
                              footer_start_fraction: float = 0.90) -> bool:
    """Exclude footer/page/source chrome from audience-content font limits."""
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        chrome_types = {
            getattr(PP_PLACEHOLDER, "FOOTER", None),
            getattr(PP_PLACEHOLDER, "DATE", None),
            getattr(PP_PLACEHOLDER, "SLIDE_NUMBER", None),
            getattr(PP_PLACEHOLDER, "HEADER", None),
        }
        if shape.is_placeholder and shape.placeholder_format.type in chrome_types:
            return True
    except Exception:
        pass
    try:
        if not shape.has_text_frame:
            return False
        text_value = " ".join(shape.text_frame.text.split()).strip()
        if not text_value:
            return False
        slide_h = float(slide.part.package.presentation_part.presentation.slide_height)
        if float(shape.top) < footer_start_fraction * slide_h:
            return False
        return bool(
            re.fullmatch(r"[\d\s/.-]+", text_value)
            or re.search(
                r"(?:大学|研究所|株式会社|University|Institute|©|"
                r"https?://|www\.|doi\s*:|^\[?sources?\]?|^出典)",
                text_value,
                re.IGNORECASE,
            )
        )
    except Exception:
        return False


# A namespace map used when scanning SmartArt / diagram XML for ``a:t``
# text nodes. (Bug H7)
_PPTX_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _extract_shape_text_items(shape) -> list[str]:
    """Return every text fragment hosted by ``shape``.

    Covers text frames, tables, and SmartArt/diagram XML (as ``a:t``
    elements). Grouped shapes are handled by the caller via
    ``_walk_shapes``; this helper only inspects the leaf shape.
    """
    out: list[str] = []

    # 1) Normal text frames.
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        out.append(run.text)
    except Exception:
        pass

    # 2) Tables.
    try:
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    try:
                        tf = cell.text_frame
                    except Exception:
                        continue
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                out.append(run.text)
    except Exception:
        pass

    # 3) SmartArt / diagram XML fallback — python-pptx doesn't expose a
    # high-level API, so we fish out every ``<a:t>`` text node via lxml.
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
        try:
            st = shape.shape_type
        except Exception:
            st = None
        is_smartart = st in (
            getattr(MSO_SHAPE_TYPE, "DIAGRAM", None),
            getattr(MSO_SHAPE_TYPE, "IGX_GRAPHIC", None),
        )
        if is_smartart:
            elt = getattr(shape, "element", None)
            if elt is not None:
                try:
                    # ``findall`` with a qualified name; lxml and
                    # stdlib ElementTree both accept this syntax.
                    for t in elt.findall(f".//{{{_PPTX_A_NS}}}t"):
                        if t.text and t.text.strip():
                            out.append(t.text)
                except Exception:
                    # Graceful fallback — some environments may not
                    # support the element tree walk.
                    pass
    except Exception:
        pass

    return out


def _resolve_font_size(run, paragraph, shape, slide):
    """Resolve the effective font size in points for ``run``.

    Bug H4 fix: previous code skipped every run whose ``run.font.size``
    was ``None``, which is the common case — most pptx files inherit
    size from the slide layout or slide master. We walk:

    1. ``run.font.size``
    2. ``paragraph.font.size``
    3. If the owning shape is a placeholder, the same (idx, paragraph
       index) slot on the slide layout's matching placeholder.
    4. The slide master's ``text_styles`` (title / body / other).

    Returns the resolved size in points (float), or ``None`` if we
    genuinely cannot determine it.
    """
    # 1) run
    try:
        if run.font.size is not None:
            return run.font.size.pt
    except Exception:
        pass
    # 2) paragraph
    try:
        if paragraph.font.size is not None:
            return paragraph.font.size.pt
    except Exception:
        pass
    # 3) slide layout placeholder inheritance
    try:
        pf = shape.placeholder_format
        if pf is not None and pf.idx is not None:
            layout = slide.slide_layout
            # paragraph index inside the shape
            try:
                para_idx = list(shape.text_frame.paragraphs).index(paragraph)
            except Exception:
                para_idx = 0
            for ph in layout.placeholders:
                try:
                    if ph.placeholder_format.idx != pf.idx:
                        continue
                    paras = list(ph.text_frame.paragraphs)
                    if 0 <= para_idx < len(paras):
                        sz = paras[para_idx].font.size
                        if sz is not None:
                            return sz.pt
                    # Fallback to first paragraph of the placeholder
                    if paras and paras[0].font.size is not None:
                        return paras[0].font.size.pt
                except Exception:
                    continue
    except Exception:
        pass
    # 4) slide master text styles (title / body / other)
    try:
        layout = slide.slide_layout
        master = layout.slide_master
        # Determine which text-style key to inspect.
        is_title = False
        try:
            pf = shape.placeholder_format
            if pf is not None and pf.idx == 0:
                is_title = True
        except Exception:
            pass
        key = "title" if is_title else "body"
        try:
            master_elt = master.element
            # Look for the master text style element for the chosen key.
            # <p:titleStyle> or <p:bodyStyle> contains <a:lvlNpPr>.
            ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
            tag = f"{{{ns_p}}}{'titleStyle' if is_title else 'bodyStyle'}"
            style = master_elt.find(f".//{tag}")
            if style is None and is_title:
                tag = f"{{{ns_p}}}bodyStyle"
                style = master_elt.find(f".//{tag}")
            if style is not None:
                # First-level paragraph properties hold the default sz.
                lvl = style.find(f"{{{_PPTX_A_NS}}}lvl1pPr")
                if lvl is not None:
                    defRPr = lvl.find(f"{{{_PPTX_A_NS}}}defRPr")
                    if defRPr is not None:
                        sz = defRPr.get("sz")
                        if sz:
                            # pptx stores size in hundredths of a point.
                            return int(sz) / 100.0
        except Exception:
            pass
        # Python-pptx convenience: some versions expose `text_styles`
        # with a .level_paragraphs -> font structure — try that too.
        try:
            ts = getattr(master, "text_styles", None)
            if ts is not None:
                attr = "title" if is_title else "body"
                style_obj = getattr(ts, attr, None)
                if style_obj is not None:
                    lvl_paras = getattr(style_obj, "level_paragraphs", None)
                    if lvl_paras:
                        lvl1 = lvl_paras[0]
                        sz = getattr(lvl1.font, "size", None)
                        if sz is not None:
                            return sz.pt
        except Exception:
            pass
    except Exception:
        pass
    return None


# ------------------------------------------------------------------
# Knowledge loader
# ------------------------------------------------------------------
def presentation_usage() -> str:
    """学会発表スライド (IEEJ SA / IEEE conference / セミナー) の作文技術ガイド全体。

    高橋メソッド, Garr Reynolds 流, 横徹流の統合。時間配分、Hook,
    key message 1 枚 1 主張, PDF 印刷耐性, Q&A 準備, 動画収録対応等。

    呼び出しタイミング:
    - 学会発表 (IEEJ SA / IEEE / APS) のドラフト review
    - 「スライドが読めない」「情報が多すぎ」feedback 対応
    - 時間配分 / storyboarding の相談
    - beamer / pptx の overfull / 字数チェック
    - Q&A 想定問答の整理

    活動ツール (実測、v0.14.0 時点):

    ## 基本 lint (existing)
    - presentation_count_underlines / validate_pdf_pages / check_overfull_hbox
    - presentation_analyze_sentences / count_weak_expressions
    - presentation_count_slides / estimate_speaking_time
    - presentation_check_slide_density / check_slide_line_count
    - presentation_check_time_13_rule / check_time_14_rule
    - presentation_check_script_paragraph_length / check_takehome_slide
    - presentation_check_slide_title_specificity / check_slide_message_hierarchy
      (check_slide_title_verb は後方互換 alias)
    - presentation_check_pptx_font_size
    - presentation_check_bullet_count_per_slide / check_bullet_ending_style
    - presentation_check_japanese_copy_style
    - presentation_check_qa_backup_slides / check_image_text_ratio
    - presentation_check_image_aspect_ratio
    - presentation_check_embedded_figure_text_size
    - presentation_replace_embedded_figure_text
    - presentation_check_final_deck_directory
    - presentation_check_color_count_per_slide / check_color_accessibility
    - presentation_estimate_per_slide_time / check_over_politeness
    - presentation_check_hedge_on_key_slides / extract_pptx_text

    ## Plan B Tier 1 (v0.12.0) — score + human-advisor comments
    - presentation_opening_hook_strength (T1)
    - presentation_takehome_strength (T2)
    - presentation_check_pie_3d_charts (T3, 宮野 S8)
    - presentation_check_logo_on_every_slide (T4, 宮野 S13)
    - presentation_check_progress_indicator (T5, 宮野 S15)

    ## Plan B Tier 2 (v0.13.0)
    - presentation_visual_text_ratio_score (T6, Reynolds Zen)
    - presentation_speaker_note_ratio (T7)
    - presentation_font_consistency (T8, 宮野 S11)
    - presentation_arrow_usage (T9, 宮野 S10)
    - presentation_check_underline_in_pptx (T10, 宮野 S14)
    - presentation_slide_density_balance (T11)

    ## Meta (v0.14.0) — 1 コールで総点検
    - presentation_health_report (T12) — 全 Plan B T1-T11 を束ね優先度付け

    ## Cross-module 和文 lint (v0.13.0)
    - presentation_check_notation_variants / find_undefined_acronyms
    - presentation_acronym_usage_audit / check_kanji_ratio
    - presentation_lint_bedrock / check_misuse_japanese
    - presentation_suggest_redundancy_fixes
    """
    return _load_skill()


# ------------------------------------------------------------------
# Shared diagnostic tools (copy from grant-writing pattern)
# ------------------------------------------------------------------
def presentation_count_underlines(tex_path: str) -> dict:
    """beamer ソース内の下線コマンドを実測。

    目標: 1-3 per slide (key number のみ下線)
    警告: スライド全体で >=20 (頻用しすぎ)
    """
    p = pathlib.Path(tex_path)
    if not p.exists():
        return {"error": f"file not found: {tex_path}"}
    text = p.read_text(encoding="utf-8", errors="replace")
    macros = ["uline", "underline", "uwave", "uuline",
              "sout", "dashuline", "dotuline", "alert"]
    by_macro: dict[str, int] = {}
    total = 0
    for m in macros:
        n = len(re.findall(r"\\" + m + r"\{", text))
        if n:
            by_macro[m] = n
            total += n
    return {
        "file": str(p),
        "total_emphasis": total,
        "by_macro": by_macro,
        "target_range": "1-3 per slide (key number only)",
        "warning_threshold": ">=20 total",
    }


def presentation_validate_pdf_pages(pdf_path: str,
                                     max_pages: int) -> dict:
    """スライド PDF のページ数を実測。発表時間 / slot との整合を検証。

    参考 (1 page = 1 slide):
        20 min talk: 10-15 slides (タイトル + 結論含む)
        30 min talk: 15-25 slides
        60 min talk: 30-45 slides
        高橋メソッド: 大字 1 word per slide は 60 slides/30min 可
    """
    try:
        import pymupdf  # type: ignore
    except ImportError:
        return {"error": "pymupdf not installed. `pip install pymupdf`"}
    p = pathlib.Path(pdf_path)
    if not p.exists():
        return {"error": f"file not found: {pdf_path}"}
    doc = pymupdf.open(pdf_path)
    n = doc.page_count
    doc.close()
    return {
        "file": str(p),
        "page_count": n,
        "max_pages": max_pages,
        "within_limit": n <= max_pages,
        "pages_over": max(0, n - max_pages),
        "estimated_time_min": {
            "conservative (2 min/slide)": round(n * 2, 1),
            "moderate (1.5 min/slide)": round(n * 1.5, 1),
            "fast (1 min/slide)": round(n * 1.0, 1),
            "takahashi (0.5 min/slide)": round(n * 0.5, 1),
        },
    }


def presentation_check_overfull_hbox(log_path: str) -> dict:
    """beamer ログ中の Overfull \\hbox をカウント。スライドでは致命的。

    目標: 0。overflow があると図が切れる / 文字がはみ出す。
    """
    p = pathlib.Path(log_path)
    if not p.exists():
        return {"error": f"file not found: {log_path}"}
    text = p.read_text(encoding="utf-8", errors="replace")
    matches = re.findall(
        r"Overfull \\hbox \(([^)]+)\) in paragraph at lines (\d+)(?:--(\d+))?",
        text,
    )
    details = [
        {"severity": m[0], "lines": m[1] + (f"-{m[2]}" if m[2] else "")}
        for m in matches[:20]
    ]
    return {
        "file": str(p),
        "overfull_count": len(matches),
        "overfull_details": details,
        "target": "0 (beamer は overflow で致命的)",
    }


def presentation_analyze_sentences(text: str, max_len: int = 30) -> dict:
    """文長分析。スライドは短文指向。

    目標: 平均 <=20 字 / 文 (高橋メソッド系), <=30 字 (standard)
    警告: >=50 字 / 文 (読まれない)

    引数:
        text: 解析対象 (スライド 1 枚ぶんのテキスト)
        max_len: 1 文上限 (既定 30)
    """
    sentences = [s.strip() for s in re.split(r"[。．\n]", text) if s.strip()]
    if not sentences:
        return {"error": "no sentences found (no 。 or ．)"}
    lengths = [len(s) for s in sentences]
    avg = sum(lengths) / len(lengths)
    long_ones = [
        {"index": i, "length": ln, "head": s[:30] + ("..." if len(s) > 30 else "")}
        for i, (s, ln) in enumerate(zip(sentences, lengths))
        if ln > max_len
    ]
    return {
        "total_sentences": len(sentences),
        "avg_length": round(avg, 1),
        "max_length": max(lengths),
        "threshold": max_len,
        "over_threshold_count": len(long_ones),
        "over_threshold_examples": long_ones[:5],
        "target": f"avg <= {max_len} (slide short-form)",
        "warning": "any sentence >= 50 字",
    }


def presentation_count_weak_expressions(text: str) -> dict:
    """弱気修飾語の出現。presentation では key slide 上で使うと信頼感低下。

    目標: key slide (contribution / result) は 0
    警告: >=2 / スライド
    """
    # v0.13.0: JA hedges は _shared.hedges 経由。presentation 固有の
    # 「だろう/でしょう」「ではないか」は別追加。
    from .._shared.hedges import HEDGE_PATTERNS as _JA_HEDGES
    patterns = dict(_JA_HEDGES)
    patterns.update({
        "だろう": r"だろう|でしょう",
        "ではないか": r"ではないか",
    })
    by_pat: dict[str, int] = {}
    total = 0
    for name, pat in patterns.items():
        n = len(re.findall(pat, text))
        if n:
            by_pat[name] = n
            total += n
    return {
        "total_weak_expressions": total,
        "by_pattern": by_pat,
        "target": "0 on key slides (contribution / result)",
        "warning": ">=2 per slide",
    }


# ------------------------------------------------------------------
# Slide-specific diagnostic tools
# ------------------------------------------------------------------
def presentation_count_slides(input_path: str) -> dict:
    """スライド数を count。beamer (.tex) の \\begin{frame} か、
    pptx の slide count を自動判定。

    引数:
        input_path: .tex (beamer) または .pptx のパス
    """
    p = pathlib.Path(input_path)
    if not p.exists():
        return {"error": f"file not found: {input_path}"}
    suffix = p.suffix.lower()

    if suffix == ".tex":
        text = p.read_text(encoding="utf-8", errors="replace")
        # Bug H2 fix: strip LaTeX comments (unescaped `%...\n`) before
        # counting, otherwise commented-out `% \frame{old}` lines are
        # mistakenly counted as slides.
        stripped_lines = [
            re.sub(r"(?<!\\)%[^\n]*", "", line)
            for line in text.splitlines()
        ]
        stripped = "\n".join(stripped_lines)
        # \begin{frame} ... or \frame{...} both supported
        frames = len(re.findall(r"\\begin\{frame\}", stripped))
        short_frames = len(re.findall(r"\\frame\b", stripped)) - frames
        total = frames + max(0, short_frames)
        return {
            "file": str(p),
            "format": "beamer",
            "frame_environments": frames,
            "short_frame_calls": max(0, short_frames),
            "total_slides": total,
        }

    if suffix == ".pptx":
        try:
            import pptx  # type: ignore
        except ImportError:
            return {"error": "python-pptx not installed. `pip install python-pptx`"}
        prs = pptx.Presentation(str(p))
        return {
            "file": str(p),
            "format": "pptx",
            "total_slides": len(prs.slides),
            "layouts": [s.slide_layout.name for s in prs.slides][:30],
        }

    return {"error": f"unsupported format: {suffix} (expected .tex or .pptx)"}


def presentation_estimate_speaking_time(text: str,
                                         wpm_ja: int = 300,
                                         wpm_en: int = 130) -> dict:
    """原稿テキストから発表時間を推定。

    日本語: 300 字/分 (学会のあっさり目発表) ~ 400 字/分 (速い講演)
    英語: 130 words/min (standard) ~ 150 words/min (native fast)

    引数:
        text: 発表原稿全体
        wpm_ja: 1 分あたりの字数 (日本語)
        wpm_en: 1 分あたりの語数 (英語)
    """
    cjk = sum(1 for c in text if "぀" <= c <= "ヿ" or "一" <= c <= "鿿")
    ja_ratio = cjk / max(1, len(text))
    is_japanese = ja_ratio > 0.3

    if is_japanese:
        n = len(text.replace(" ", "").replace("\n", ""))
        minutes = n / wpm_ja
        unit = f"{n} 字 / {wpm_ja} 字/分"
    else:
        words = len(text.split())
        minutes = words / wpm_en
        unit = f"{words} words / {wpm_en} wpm"

    return {
        "detected_language": "ja" if is_japanese else "en",
        "count_or_words": (len(text) if is_japanese else len(text.split())),
        "estimated_minutes": round(minutes, 2),
        "detail": unit,
        "hint": (
            "20-min talk = 300 字×20 = 6000 字 (ja) / 130×20 = 2600 words (en)。"
            " 質疑 5 分を別枠で確保する前提。"
        ),
    }


def presentation_check_slide_density(text: str,
                                      max_chars_per_slide: int = 80) -> dict:
    """1 スライドあたりの文字密度チェック (テキストを直接渡す)。

    目標: <= 80 字 / slide (標準) or <= 20 字 / slide (高橋メソッド)
    警告: >= 150 字 / slide (読まれない)

    引数:
        text: 1 スライドの内容 (短いブロックを想定)
        max_chars_per_slide: 閾値 (既定 80)
    """
    n = len(text.replace(" ", "").replace("\n", ""))
    over = n - max_chars_per_slide
    return {
        "char_count": n,
        "threshold": max_chars_per_slide,
        "within_limit": n <= max_chars_per_slide,
        "chars_over": max(0, over),
        "recommendation": (
            "高橋メソッド: 20 字以下 (大字・1 スライド 1 主張)。 "
            "standard: 80 字以下 (本文 + bullet)。 "
            "150 字以上は「読まれないスライド」。"
        ),
    }


def presentation_extract_pptx_text(pptx_path: str) -> dict:
    """pptx の各 slide のテキストを抽出。密度チェックや文字起こしに。

    note: python-pptx 依存。画像内の文字は拾えない (OCR 不要時のみ)。
    """
    try:
        import pptx  # type: ignore
    except ImportError:
        return {"error": "python-pptx not installed. `pip install python-pptx`"}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = pptx.Presentation(str(p))
    slides_text = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        # Bug H7 fix: walk grouped shapes recursively AND extract text
        # from tables (``has_table``) and SmartArt / diagram XML
        # (``a:t`` nodes). The old implementation silently dropped all
        # of these categories, which covers a large fraction of
        # figure-heavy pptx decks.
        for shape in _walk_shapes(slide.shapes):
            texts.extend(_extract_shape_text_items(shape))
        joined = "\n".join(texts)
        slides_text.append({
            "slide": i,
            "layout": slide.slide_layout.name,
            "char_count": len(joined.replace(" ", "").replace("\n", "")),
            "text_preview": joined[:200] + ("..." if len(joined) > 200 else ""),
        })
    return {
        "file": str(p),
        "total_slides": len(prs.slides),
        "slides": slides_text,
    }


# ------------------------------------------------------------------
# Deep-learning diagnostic tools (from 作図力学 + 木下 講演章)
# ------------------------------------------------------------------

def presentation_check_slide_line_count(pptx_path: str,
                                         max_horiz_lines: int = 8,
                                         max_vert_lines: int = 12,
                                         ) -> dict:
    """pptx の各 slide で text 行数が木下推奨の範囲内か検証。

    木下『理科系の作文技術』p.227: 横向き slide なら 8 行、縦向き slide
    なら 12 行を上限にせよ。これを超えると文字が小さくなりすぎる。

    Args:
        pptx_path: .pptx
        max_horiz_lines: 横 slide の行数上限 (既定 8)
        max_vert_lines: 縦 slide の行数上限 (既定 12)
    """
    try:
        import pptx  # type: ignore
    except ImportError:
        return {"error": "python-pptx not installed. `pip install python-pptx`"}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = pptx.Presentation(str(p))
    # Determine orientation from slide size
    w = prs.slide_width
    h = prs.slide_height
    is_horiz = w >= h
    limit = max_horiz_lines if is_horiz else max_vert_lines
    violations: list[dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        lines += 1
        if lines > limit:
            violations.append({"slide": i, "lines": lines, "limit": limit})
    return {
        "file": str(p),
        "orientation": "horizontal" if is_horiz else "vertical",
        "line_limit": limit,
        "total_slides": len(prs.slides),
        "violations": violations[:20],
        "violation_count": len(violations),
        "source": "木下『理科系の作文技術』p.227",
    }


def presentation_check_hedge_on_key_slides(pptx_path: str) -> dict:
    """pptx で Result / Conclusion / Summary スライドに弱気修飾語が
    含まれていないかチェック。木下 p.235 のズバリ話法。

    key slide = title が Result/Conclusion/Summary/まとめ/結論/結果
    を含むもの。

    弱気リスト: ようです、らしい、かもしれません、ではなかろうか、
    と思われる、と考えられる、だろう、でしょう、のようだ、と見られる
    """
    try:
        import pptx  # type: ignore
    except ImportError:
        return {"error": "python-pptx not installed. `pip install python-pptx`"}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = pptx.Presentation(str(p))

    key_title_pat = re.compile(
        r"(result|conclusion|summary|take.?home|まとめ|結論|結果|考察)",
        re.IGNORECASE)
    hedge_patterns = [
        (r"ようです", "ようです"),
        (r"らしい", "らしい"),
        (r"かもしれ(ない|ません)", "かもしれない"),
        (r"ではな(かろうか|いでしょうか)", "ではなかろうか"),
        (r"と(思われ|考えられ|見られ)(る|ます)", "〜と考えられる"),
        (r"だろう|でしょう", "だろう"),
        (r"ようだ[。、]", "ようだ"),
        (r"と見てよい", "と見てよい"),
    ]

    issues: list[dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        # Bug H3 fix: prefer the title placeholder (idx 0) and fall back
        # to shape-order only if no title placeholder exists. The shared
        # ``_slide_title`` helper encapsulates this.
        title_text = _slide_title(slide)
        body_text_parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            body_text_parts.append(shape.text_frame.text)
        body = "\n".join(body_text_parts)
        if not key_title_pat.search(title_text):
            continue
        # Scan body for hedges
        for pat, label in hedge_patterns:
            ms = re.findall(pat, body)
            if ms:
                issues.append({
                    "slide": i,
                    "title": title_text[:60],
                    "hedge": label,
                    "count": len(ms),
                })
    return {
        "file": str(p),
        "total_slides": len(prs.slides),
        "key_slide_hedge_issues": issues[:30],
        "issue_count": len(issues),
        "target": "0 hedges on Result/Conclusion slides",
        "source": "木下『理科系の作文技術』p.235",
    }


def presentation_check_script_paragraph_length(script: str,
                                                 min_chars: int = 150,
                                                 max_chars: int = 350,
                                                 ) -> dict:
    """発表原稿の 1 パラグラフが 200-300 字目安から大きく外れていないか。

    木下 p.225 + まんが p.158: 「区切りの明確なパラグラフを 200-300 字で
    十分考えて配列した原稿ができあがれば、その講演は半ば成功」

    Args:
        script: 発表原稿 (空行で段落区切り)
        min_chars: 下限 (既定 150)
        max_chars: 上限 (既定 350)
    """
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", script) if p.strip()]
    out: list[dict] = []
    for i, para in enumerate(paragraphs, start=1):
        n = len(para.replace(" ", "").replace("\n", ""))
        within = min_chars <= n <= max_chars
        if not within:
            out.append({"paragraph": i, "chars": n,
                        "head": para[:40] + ("..." if len(para) > 40 else "")})
    return {
        "total_paragraphs": len(paragraphs),
        "target_range": f"{min_chars}-{max_chars} chars",
        "violation_count": len(out),
        "violations": out[:15],
        "allow_ratio": f"up to 10% of paragraphs outside range",
        "source": "木下 p.225 / まんが p.158",
    }


def presentation_check_over_politeness(script: str) -> dict:
    """学会発表で過剰に丁寧な言い回しを検出。木下 p.235。

    学会は平等な場。「報告させていただきます」「であります」等は
    緊張の signal として聴衆に取られる。
    """
    patterns = {
        "報告させていただきます": r"報告させていただきま(す|した)",
        "発表させていただきます": r"発表させていただきま(す|した)",
        "であります": r"であります",
        "と思わせていただきます": r"と思わせていただき",
        "お話させていただきます": r"お話させていただきま",
        "〜させていただきたい": r"させていただ(きたい|こう)",
        "申し上げます": r"申し上げ(ます|たい)",
    }
    by_pat = {}
    total = 0
    for name, pat in patterns.items():
        n = len(re.findall(pat, script))
        if n:
            by_pat[name] = n
            total += n
    return {
        "total_over_polite": total,
        "by_pattern": by_pat,
        "target": "0 (学会 は 平等)",
        "recommendation": (
            "「報告させていただきます」→「報告します」、"
            "「であります」→「です」。"
            "学会は立場を対等と扱う場。過剰な丁寧は緊張の signal。"
        ),
        "source": "木下『理科系の作文技術』p.235",
    }


def presentation_check_time_13_rule(script: str,
                                     slot_min: float = 20.0,
                                     wpm_ja: int = 240,
                                     ) -> dict:
    """木下 1/3 則 — 前半で全員わかる話、中盤で大半が分かった気、
    後半で専門家向け、の layering が script に表れているか。

    簡易判定: script を 3 等分し、各 block 内の「専門用語密度」を比較。
    前 1/3 の専門用語密度 < 中 1/3 < 後 1/3 であれば layering 済。

    Args:
        script: 発表原稿
        slot_min: 持ち時間 (1/3 則の spacing 確認用)
        wpm_ja: 字数/分 (木下 学会は 240)
    """
    chars = list(script.replace("\n", " "))
    if not chars:
        return {"error": "script is empty"}
    third = len(chars) // 3
    blocks = [
        "".join(chars[:third]),
        "".join(chars[third:2 * third]),
        "".join(chars[2 * third:]),
    ]
    # Simplified specialist-marker density: count katakana runs >=4 chars
    # (concept names), digits + units ("kHz", "mT"), equations.
    # Bug H1 fix: the previous equation regex `[\w]_[\w]` over-fired on
    # filenames like "slide_1", "file_name", "result_v2". We now require
    # either LaTeX-style $...$ math delimiters, or a single-letter symbol
    # followed by a physics-style subscript index (i, j, k, l, m, n, x,
    # y, z) — the typical tensor / index notation in research talks.
    def _density(s):
        katakana = len(re.findall(r"[ァ-ヴー]{4,}", s))
        units = len(re.findall(r"\d+[.]?\d*\s*(kHz|mT|mm|V|A|Hz|μm|Ω|T|s)",
                                s))
        eq_latex = len(re.findall(r"\$[^$\n]+\$", s))
        eq_subscript = len(re.findall(r"\b[A-Za-z]_[ijklmnxyz]\b", s))
        eq_rel = len(re.findall(r"[=<>≤≥][ =<>≤≥]?", s))
        equations = eq_latex + eq_subscript + eq_rel
        return (katakana + units + equations) / max(1, len(s)) * 1000

    densities = [round(_density(b), 2) for b in blocks]
    layered = densities[0] <= densities[1] <= densities[2]

    return {
        "slot_min": slot_min,
        "wpm_ja": wpm_ja,
        "script_chars": len(script),
        "estimated_minutes": round(len(script) / wpm_ja, 1),
        "block_specialist_densities": densities,
        "layered_correctly": layered,
        "target": "density[前] <= density[中] <= density[後]",
        "source": "木下 1/3 則 p.224",
    }


# ------------------------------------------------------------------
# v0.4.0: 即戦力 5 tools (1/4 rule, takehome, title-verb, font, bullet-count)
# ------------------------------------------------------------------

def presentation_check_time_14_rule(script: str,
                                     slot_min: float = 10.0,
                                     wpm_ja: int = 240) -> dict:
    """木下 1/4 則 — 10 分講演を 4 等分 (intro/method/result/discussion) した
    予算に対して、原稿各 block の字数が逸脱していないか。
    """
    total_budget_chars = slot_min * wpm_ja
    quarter_budget = total_budget_chars / 4
    markers = {
        "intro":      r"(?:^#+\s*)?(?:序論|はじめに|背景|Introduction|Background)",
        "method":     r"(?:^#+\s*)?(?:手法|方法|Method|Approach)",
        "result":     r"(?:^#+\s*)?(?:結果|Result|Outcome)",
        "discussion": r"(?:^#+\s*)?(?:考察|議論|Discussion|まとめ|Summary|Conclusion)",
    }
    positions = {}
    for name, pat in markers.items():
        m = re.search(pat, script, re.MULTILINE | re.IGNORECASE)
        if m:
            positions[name] = m.start()
    if len(positions) < 3:
        return {"error": "Not enough section markers (need >=3)",
                "found": list(positions.keys())}
    ordered = sorted(positions.items(), key=lambda kv: kv[1])
    section_lengths = {}
    for i, (name, start) in enumerate(ordered):
        end = ordered[i+1][1] if i+1 < len(ordered) else len(script)
        section_lengths[name] = end - start
    issues = []
    for name, chars in section_lengths.items():
        if chars > quarter_budget * 1.25:
            issues.append({"section": name, "chars": chars,
                           "budget": round(quarter_budget, 0),
                           "verdict": "over budget"})
        elif chars < quarter_budget * 0.5:
            issues.append({"section": name, "chars": chars,
                           "budget": round(quarter_budget, 0),
                           "verdict": "under-developed"})
    return {
        "slot_min": slot_min,
        "budget_chars_per_quarter": round(quarter_budget, 0),
        "section_lengths": section_lengths,
        "issues": issues,
        "source": "木下『理科系の作文技術』第 10 章 1/4 則",
    }


def presentation_check_takehome_slide(pptx_path: str) -> dict:
    """pptx 最終 3 枚以内に Take-home / Summary / まとめ slide があるか確認。"""
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))
    n = len(prs.slides)
    last_slides = list(prs.slides)[-3:]
    takehome_pat = re.compile(
        r"(take.?home|summary|conclusion|まとめ|結論|結語|今後の|future\s*work)",
        re.IGNORECASE)
    found = []
    for i, s in enumerate(last_slides):
        idx = n - len(last_slides) + i + 1
        # Bug H3 fix: use shared `_slide_title` helper (title placeholder
        # preferred, shape order as fallback).
        title_text = _slide_title(s)
        if takehome_pat.search(title_text):
            found.append({"slide": idx, "title": title_text[:80]})
    return {
        "total_slides": n,
        "takehome_slides_in_last_3": found,
        "has_takehome": bool(found),
        "hint": "Take-home を最終 1-3 枚目に配置。",
        "source": "木下『理科系の作文技術』講演構成",
    }


def presentation_check_slide_title_specificity(
        pptx_path: str,
        max_title_chars: int = 28,
        min_title_chars: int = 5) -> dict:
    """各 slide title が短く具体的な「対象＋観点」になっているか点検。

    title は結果を言い切る場所ではない。「結果」「数値計算結果」のような
    汎用見出しではなく、「モデル1の計算精度評価」のように、何について
    何を示す slide かを少ない文字数で特定する。結果から分かったことは
    slide 最下部に置く。
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))
    if min_title_chars < 1 or max_title_chars < min_title_chars:
        return {"error": "Require 1 <= min_title_chars <= max_title_chars."}
    generic_titles = {
        "results", "result", "methodology", "method", "methods",
        "introduction", "background", "discussion", "conclusion",
        "overview", "summary", "analysis", "simulation results",
        "結果", "手法", "方法", "背景", "はじめに", "序論", "結論", "考察",
        "計算結果", "解析結果", "数値計算結果", "実験結果", "精度評価",
        "計算精度評価", "検証", "評価", "解析", "実験", "まとめ",
    }
    sentence_ending = re.compile(
        r"(?:した|する|なる|できる|可能|必要|示す|残る|進む|選ぶ|測る|"
        r"低減|削減|改善|向上|一致|整合|達成|実現)[。.!！]?$|"
        r"\b(?:is|are|was|were|shows?|achieves?|reduces?|increases?|"
        r"enables?|improves?|outperforms?|demonstrates?)\b[^.]*[.]?$",
        re.IGNORECASE,
    )
    numeric_result = re.compile(
        r"\d+(?:\.\d+)?\s*(?:%|％|倍|分の|GiB|MiB|GB|MB|ms|s|秒|(?:万)?自由度)",
        re.IGNORECASE,
    )
    viewpoint_terms = re.compile(
        r"(?:評価|検証|比較|条件|モデル|依存|分解|対応|性能|精度|閉包|"
        r"感度|頑健|高速化|適用|展開|限界|課題|定義|設計|測定|同定|"
        r"解析|原理|機構|構造|実装|停止則|モード|"
        r"evaluation|validation|comparison|condition|model|dependence|"
        r"decomposition|performance|accuracy|closure|sensitivity|robustness|"
        r"acceleration|application|limitation|design|measurement|identification|"
        r"analysis|principle|mechanism|structure|implementation|loop[- ]?free)",
        re.IGNORECASE,
    )
    structural_title = re.compile(
        r"^(?:title|agenda|contents?|references?|acknowledg(?:e)?ments?|"
        r"appendix|q\s*&\s*a|目次|参考文献|謝辞|付録|質疑)\s*$",
        re.IGNORECASE,
    )
    reports = []
    for i, s in enumerate(prs.slides, 1):
        title = _slide_title(s).strip()
        if not title or i == 1 or structural_title.match(title):
            reports.append({
                "slide": i,
                "title": title,
                "skipped": True,
                "reason": "title/structural slide",
            })
            continue
        compact = re.sub(r"\s+", "", title)
        normalized = title.lower().strip(" ：:。.!?！\t\r\n")
        issues = []
        has_target = len(compact) >= min_title_chars and normalized not in generic_titles
        has_viewpoint = bool(viewpoint_terms.search(title))
        is_concise = len(compact) <= max_title_chars and "\n" not in title
        numeric_values = numeric_result.findall(title)
        numeric_result_marker = re.search(
            r"(?:最大|最小|以内|以上|以下|誤差|差|へ|→|で解析|で求解|"
            r"低減|削減|改善|向上|一致|整合)",
            title,
        )
        is_result_sentence = bool(
            sentence_ending.search(title)
            or len(numeric_values) >= 2
            or (numeric_values and numeric_result_marker)
        )
        if not has_target:
            issues.append("title_not_specific")
        if not has_viewpoint:
            issues.append("viewpoint_not_explicit")
        if not is_concise:
            issues.append("title_too_long")
        if is_result_sentence:
            issues.append("title_is_result_sentence")
        reports.append({
            "slide": i,
            "title": title,
            "skipped": False,
            "char_count": len(compact),
            "normalized_title": normalized,
            "criteria": {
                "target_is_specific": has_target,
                "viewpoint_is_explicit": has_viewpoint,
                "concise_one_line": is_concise,
                "not_a_result_sentence": not is_result_sentence,
                "unique_in_deck": True,
            },
            "issues": issues,
        })
    title_counts: dict[str, int] = {}
    for item in reports:
        if item.get("skipped"):
            continue
        key = item["normalized_title"]
        title_counts[key] = title_counts.get(key, 0) + 1
    for item in reports:
        if item.get("skipped"):
            continue
        unique = title_counts[item["normalized_title"]] == 1
        item["criteria"]["unique_in_deck"] = unique
        if not unique:
            item["issues"].append("duplicate_title")
        item["score"] = 2 * sum(item["criteria"].values())
        item["score_max"] = 10
        item["passed"] = all(item["criteria"].values())
    evaluated = [item for item in reports if not item.get("skipped")]
    failed = [item for item in evaluated if not item["passed"]]
    return {
        "passed": not failed,
        "slides_checked": len(evaluated),
        "slides_passing": len(evaluated) - len(failed),
        "issue_count": len(failed),
        "issues": failed[:20],
        "slides": reports,
        "rule": (
            "Title = 少ない文字数で具体化した対象＋観点。"
            "Bottom = 図表・式・比較から分かったこと。"
        ),
        "acceptance_criteria": {
            "target_is_specific": "何について扱うかが特定できる。",
            "viewpoint_is_explicit": "評価・比較・条件・構造など、何を見るかが分かる。",
            "concise_one_line": f"一行かつ {max_title_chars} 文字以内を目安とする。",
            "not_a_result_sentence": "結果の数値や結論を言い切らない。",
            "unique_in_deck": "他 slide の title と区別できる。",
            "pass_condition": "5項目すべてを満たす。各2点、10点満点。",
        },
        "examples": {
            "bad": ["結果", "数値計算結果", "精度評価"],
            "good": ["モデル1の計算精度評価", "C型鉄心のメッシュ依存性"],
        },
    }


def presentation_check_slide_title_verb(pptx_path: str) -> dict:
    """後方互換 alias。title の動詞化ではなく具体性・簡潔性を点検。"""
    result = presentation_check_slide_title_specificity(pptx_path)
    if "error" not in result:
        result["deprecated_name"] = (
            "presentation_check_slide_title_verb is retained for compatibility; "
            "use presentation_check_slide_title_specificity."
        )
    return result


def presentation_check_slide_message_hierarchy(
        pptx_path: str,
        bottom_start_fraction: float = 0.68,
        footer_start_fraction: float = 0.94,
        min_takeaway_chars: int = 8,
        max_title_chars: int = 28) -> dict:
    """各 content slide の伝達意図 title と下端の知見を位置ベースで点検。

    title は「このスライドで何を伝えるか」を、短い「対象＋観点」で示す。
    「結果」「数値計算結果」のような汎用語や、結果の長い言い切りは避ける。
    takeaway は中央の図表・式・比較から「何が分かったか」を示す。
    footer、page number、URL、citation、所属は takeaway として数えない。
    """
    try:
        import pptx as _pptx
        from pptx.enum.shapes import PP_PLACEHOLDER
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    if not 0.0 < bottom_start_fraction < footer_start_fraction < 1.0:
        return {
            "error": (
                "Require 0 < bottom_start_fraction < "
                "footer_start_fraction < 1."
            )
        }
    if min_takeaway_chars < 1:
        return {"error": "min_takeaway_chars must be >= 1."}
    if max_title_chars < 5:
        return {"error": "max_title_chars must be >= 5."}

    prs = _pptx.Presentation(str(p))
    slide_h = float(prs.slide_height)
    structural_title = re.compile(
        r"^(?:title|agenda|contents?|references?|acknowledg(?:e)?ments?|"
        r"appendix|q\s*&\s*a|目次|参考文献|謝辞|付録|質疑|補足)\s*$",
        re.IGNORECASE,
    )
    generic_title = {
        "results", "result", "methodology", "method", "methods",
        "introduction", "background", "discussion", "conclusion",
        "overview", "summary", "analysis", "future work",
        "結果", "手法", "方法", "背景", "はじめに", "序論", "結論",
        "考察", "概要", "まとめ", "今後の課題", "解析", "実験",
        "計算結果", "解析結果", "数値計算結果", "実験結果", "精度評価",
        "計算精度評価", "検証", "評価",
    }
    claim_terms = re.compile(
        r"(?:\b(?:is|are|was|were|shows?|achieves?|reduces?|increases?|"
        r"enables?|improves?|outperforms?|solves?|demonstrates?|supports?|"
        r"requires?|remains?|preserves?)\b|"
        r"できる|可能|示す|示した|なる|保つ|維持|残る|減らす|低減|削減|"
        r"向上|改善|実現|達成|収束|一致|対応|整合|選べる|選択|主役|中核|"
        r"有効|高速|強い|優れる|必要|支配|保証|継承|適用|進む)",
        re.IGNORECASE,
    )
    result_sentence = re.compile(
        r"(?:した|する|なる|できる|可能|必要|示す|残る|進む|選ぶ|測る|"
        r"低減|削減|改善|向上|一致|整合|達成|実現)[。.!！]?$|"
        r"\b(?:is|are|was|were|shows?|achieves?|reduces?|increases?|"
        r"enables?|improves?|outperforms?|demonstrates?)\b[^.]*[.]?$",
        re.IGNORECASE,
    )
    numeric_result = re.compile(
        r"\d+(?:\.\d+)?\s*(?:%|％|倍|分の|GiB|MiB|GB|MB|ms|s|秒|(?:万)?自由度)",
        re.IGNORECASE,
    )
    viewpoint_terms = re.compile(
        r"(?:評価|検証|比較|条件|モデル|依存|分解|対応|性能|精度|閉包|"
        r"感度|頑健|高速化|適用|展開|限界|課題|定義|設計|測定|同定|"
        r"解析|原理|機構|構造|実装|停止則|モード|"
        r"evaluation|validation|comparison|condition|model|dependence|"
        r"decomposition|performance|accuracy|closure|sensitivity|robustness|"
        r"acceleration|application|limitation|design|measurement|identification|"
        r"analysis|principle|mechanism|structure|implementation|loop[- ]?free)",
        re.IGNORECASE,
    )
    takeaway_cue = re.compile(
        r"^(?:分かったこと|結論|示唆|要点|まとめ|以上より|したがって|"
        r"take.?away)\s*[：:]",
        re.IGNORECASE,
    )
    chrome_types = {
        getattr(PP_PLACEHOLDER, "FOOTER", None),
        getattr(PP_PLACEHOLDER, "DATE", None),
        getattr(PP_PLACEHOLDER, "SLIDE_NUMBER", None),
        getattr(PP_PLACEHOLDER, "HEADER", None),
    }

    def _placeholder_type(shape):
        try:
            if not shape.is_placeholder:
                return None
            return shape.placeholder_format.type
        except Exception:
            return None

    def _is_title_shape(shape) -> bool:
        try:
            return bool(shape.is_placeholder and shape.placeholder_format.idx == 0)
        except Exception:
            return False

    def _is_footer_like(text: str) -> bool:
        t = " ".join(text.split()).strip()
        if not t:
            return True
        if re.fullmatch(r"[\d\s/.-]+", t):
            return True
        if re.search(r"(?:https?://|www\.|doi\s*:|^\[?sources?\]?|^出典)",
                     t, re.IGNORECASE):
            return True
        if re.search(r"(?:大学|研究所|株式会社|University|Institute|©)", t,
                     re.IGNORECASE) and not claim_terms.search(t):
            return True
        return False

    reports = []
    checked = 0
    passed = 0
    title_failures = 0
    takeaway_failures = 0
    for index, slide in enumerate(prs.slides, 1):
        title = _slide_title(slide).strip()
        if index == 1 or structural_title.match(title):
            reports.append({
                "slide": index,
                "title": title,
                "skipped": True,
                "reason": "title/structural slide",
            })
            continue

        checked += 1
        title_norm = title.lower().strip(" ：:。.!?\t\r\n")
        title_content_chars = len(re.sub(r"\s+", "", title))
        title_is_specific = bool(
            title and title_norm not in generic_title and title_content_chars >= 5
        )
        title_is_concise = title_content_chars <= max_title_chars
        numeric_values = numeric_result.findall(title)
        numeric_result_marker = re.search(
            r"(?:最大|最小|以内|以上|以下|誤差|差|へ|→|で解析|で求解|"
            r"低減|削減|改善|向上|一致|整合)",
            title,
        )
        title_is_result_sentence = bool(
            result_sentence.search(title)
            or len(numeric_values) >= 2
            or (numeric_values and numeric_result_marker)
        )
        title_is_message = bool(
            title_is_specific and title_is_concise and not title_is_result_sentence
        )

        candidates = []
        for shape in _walk_shapes(slide.shapes):
            try:
                if not shape.has_text_frame or _is_title_shape(shape):
                    continue
                if _placeholder_type(shape) in chrome_types:
                    continue
                text_value = " ".join(shape.text_frame.text.split()).strip()
                if len(text_value) < min_takeaway_chars or _is_footer_like(text_value):
                    continue
                top_fraction = float(shape.top) / slide_h
                bottom_fraction = float(shape.top + shape.height) / slide_h
                if bottom_start_fraction <= top_fraction < footer_start_fraction:
                    candidates.append({
                        "text": text_value[:180],
                        "top_fraction": round(top_fraction, 3),
                    })
                elif bottom_fraction >= bottom_start_fraction:
                    lines = [line.strip() for line in shape.text_frame.text.splitlines()
                             if line.strip()]
                    if lines and takeaway_cue.search(lines[-1]):
                        candidates.append({
                            "text": lines[-1][:180],
                            "top_fraction": "last-line-in-bottom-spanning-box",
                        })
            except Exception:
                continue

        has_takeaway = bool(candidates)
        issues = []
        if not title_is_message:
            title_failures += 1
            if not title_is_specific:
                issues.append("title_not_specific")
            if not title_is_concise:
                issues.append("title_too_long")
            if title_is_result_sentence:
                issues.append("title_is_result_sentence")
        if not has_takeaway:
            takeaway_failures += 1
            issues.append("bottom_takeaway_missing")
        if not issues:
            passed += 1
        reports.append({
            "slide": index,
            "title": title,
            "skipped": False,
            "title_is_message": title_is_message,
            "title_char_count": title_content_chars,
            "title_is_specific": title_is_specific,
            "title_is_concise": title_is_concise,
            "title_is_result_sentence": title_is_result_sentence,
            "bottom_takeaway": candidates[0]["text"] if candidates else "",
            "has_bottom_takeaway": has_takeaway,
            "issues": issues,
        })

    checks = checked * 2
    failed_checks = title_failures + takeaway_failures
    score = 10.0 if checks == 0 else 10.0 * (checks - failed_checks) / checks
    return {
        "score": round(score, 1),
        "slides_checked": checked,
        "slides_passing": passed,
        "title_message_failures": title_failures,
        "bottom_takeaway_failures": takeaway_failures,
        "slides": reports,
        "rule": {
            "title": (
                "最上部に、そのスライドで伝えたい対象＋観点を短い具体語句で置く。"
                "『結果』『数値計算結果』や、結果の長い言い切りは避ける。"
            ),
            "bottom": "最下部に、中央の図表・式・比較から分かったことを一文で置く。",
        },
        "hint": (
            "title は短い具体的な伝達項目、下端文は得られた知見。下端文は title の"
            "同語反復ではなく、中央の証拠の解釈または含意にする。"
        ),
    }


def presentation_check_pptx_font_size(pptx_path: str,
                                        min_body_pt: int = 24,
                                        min_title_pt: int = 32,
                                        min_figure_pt: int = 20,
                                        exclude_chrome: bool = True) -> dict:
    """pptx audience-facing font size < 下限を検出。

    footer、page number、date、source/citation chrome は既定で除外する。
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))
    violations = []
    unresolved = 0
    excluded_chrome_runs = 0
    for i, slide in enumerate(prs.slides, 1):
        # Bug H7 fix: descend into grouped shapes while iterating.
        for shape in _walk_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            if exclude_chrome and _is_footer_or_page_chrome(shape, slide):
                excluded_chrome_runs += sum(
                    1 for para in shape.text_frame.paragraphs
                    for run in para.runs if run.text.strip()
                )
                continue
            is_title = _shape_is_slide_title(shape, slide)
            is_figure_text = getattr(shape, "name", "").startswith(
                "FIGURE_TEXT::"
            )
            limit = (
                min_title_pt if is_title
                else min_figure_pt if is_figure_text
                else min_body_pt
            )
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    # Bug H4 fix: walk the run -> paragraph -> layout
                    # placeholder -> slide master inheritance chain so
                    # we don't silently skip 100% of default-template
                    # presentations.
                    size_pt = _resolve_font_size(run, para, shape, slide)
                    if size_pt is None:
                        unresolved += 1
                        continue
                    if size_pt < limit:
                        violations.append({
                            "slide": i, "is_title": is_title,
                            "size_pt": size_pt, "limit_pt": limit,
                            "text": run.text[:40],
                        })
    return {
        "total_violations": len(violations),
        "min_body_pt": min_body_pt, "min_title_pt": min_title_pt,
        "min_figure_pt": min_figure_pt,
        "violations": violations[:20],
        "unresolved_runs": unresolved,
        "excluded_chrome_runs": excluded_chrome_runs,
        "exclude_chrome": exclude_chrome,
        "hint": ("Audience-facing body/caption/annotation/table/chart text "
                 ">= 24pt, title >= 32pt. Reconstructed figure text shapes "
                 "named FIGURE_TEXT:: use the pasted-figure floor >= 20pt. "
                 "Footer, page number, date, and "
                 "source/citation chrome are excluded by default. Text baked "
                 "into raster images must be audited separately with "
                 "presentation_check_embedded_figure_text_size; an unresolved "
                 "picture is not a pass. "
                 "unresolved_runs = number of text runs whose effective "
                 "font size could not be determined via run / paragraph "
                 "/ layout / master inheritance."),
    }


def presentation_check_bullet_count_per_slide(pptx_path: str,
                                                 max_bullets: int = 5) -> dict:
    """1 slide の bullet 数が上限超過を検出 (Miller 7±2)."""
    try:
        import pptx as _pptx
        from pptx.enum.shapes import PP_PLACEHOLDER
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))
    over_count = 0
    over_slides = []
    # Bug H6 fix: exclude title (placeholder idx 0) and
    # footer / date / slide-number placeholders so a slide with a
    # 2-line title + 3 real bullets does not get counted as 5 bullets.
    chrome_types = {
        getattr(PP_PLACEHOLDER, "FOOTER", None),
        getattr(PP_PLACEHOLDER, "DATE", None),
        getattr(PP_PLACEHOLDER, "SLIDE_NUMBER", None),
        getattr(PP_PLACEHOLDER, "HEADER", None),
    }
    for i, slide in enumerate(prs.slides, 1):
        bullets = 0
        for shape in _walk_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            # Skip title and chrome placeholders.
            try:
                pf = shape.placeholder_format
            except Exception:
                pf = None
            if pf is not None:
                try:
                    if pf.idx == 0:
                        continue
                except Exception:
                    pass
                try:
                    if pf.type in chrome_types:
                        continue
                except Exception:
                    pass
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    bullets += 1
        if bullets > max_bullets:
            over_count += 1
            over_slides.append({"slide": i, "bullets": bullets})
    return {
        "total_slides": len(prs.slides),
        "over_bullet_limit_count": over_count,
        "max_bullets_per_slide": max_bullets,
        "violations": over_slides[:15],
        "hint": ("Miller 7±2 下限 5 推奨。超過時は 2 枚に分割。"
                 "タイトル・フッター・日付・ページ番号は集計から除外。"),
    }


# ------------------------------------------------------------------
# v0.4.0: 中期 5 tools (qa/image-ratio/color/bullet-end/per-slide-time)
# ------------------------------------------------------------------

def presentation_check_qa_backup_slides(pptx_path: str,
                                          min_backup: int = 3) -> dict:
    """pptx に Q&A backup slide (hidden or named) が N 枚以上あるか確認."""
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))
    hidden_count = 0
    backup_named = 0
    backup_slides = []
    for i, s in enumerate(prs.slides, 1):
        try:
            is_hidden = bool(s._element.get("show") == "0")
        except Exception:
            is_hidden = False
        # Bug H3 fix: use shared `_slide_title` helper.
        title = _slide_title(s)
        is_backup_named = any(k in title.lower() for k in
                               ("backup", "q&a", "q & a", "qa", "予備",
                                "質問対策", "補足", "appendix"))
        if is_hidden:
            hidden_count += 1
            backup_slides.append({"slide": i, "title": title[:60], "reason": "hidden"})
        elif is_backup_named:
            backup_named += 1
            backup_slides.append({"slide": i, "title": title[:60], "reason": "named"})
    total_backup = hidden_count + backup_named
    return {
        "total_slides": len(prs.slides),
        "hidden_slides": hidden_count,
        "backup_named_slides": backup_named,
        "total_backup_candidates": total_backup,
        "min_recommended": min_backup,
        "verdict": "OK" if total_backup >= min_backup else "INSUFFICIENT",
        "backup_slides": backup_slides,
        "hint": f"Q&A backup {min_backup}-5 枚: method detail / failure / cost / scalability / future.",
    }


def presentation_check_final_deck_directory(
        directory_path: str,
        figure_text_ocr_backend: str = "none",
        figure_text_ocr_manifest_path: str = "",
        confirmed_textless_shapes: list[str] | None = None) -> dict:
    """最終発表資料のディレクトリ衛生と本体内重複を検査する。

    最終化後は presentation directory 直下の PPTX を正本1本に限定する。
    旧版にしかない有用内容は同じPPTX末尾の質問対策・補足スライドへ
    移したうえで、旧版PPTXとその検査ログは削除する。別のbackup deckや
    ``final2`` のような改訂コピーを残す運用は認めない。
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}

    root = pathlib.Path(directory_path)
    if not root.exists():
        return {"error": f"directory not found: {directory_path}"}
    if not root.is_dir():
        return {"error": f"not a directory: {directory_path}"}

    top_level_pptx = sorted(
        (path for path in root.glob("*.pptx") if path.is_file()),
        key=lambda path: path.name.casefold(),
    )
    revision_dir_markers = {
        "_archive", "archive", "archives", "old", "old_versions",
        "versions", "旧版", "改訂履歴",
    }
    archived_revision_pptx = []
    for path in root.rglob("*.pptx"):
        if path.parent == root:
            continue
        parent_parts = {
            part.casefold() for part in path.relative_to(root).parts[:-1]
        }
        if parent_parts & revision_dir_markers:
            archived_revision_pptx.append(str(path.relative_to(root)))

    duplicate_slide_groups = []
    backup_slide_count = 0
    backup_slides = []
    figure_text_audit = {}
    final_candidate = ""
    open_error = ""
    if len(top_level_pptx) == 1:
        final_path = top_level_pptx[0]
        final_candidate = final_path.name
        try:
            prs = _pptx.Presentation(str(final_path))
            signatures = {}
            for slide_index, slide in enumerate(prs.slides, 1):
                text_parts = []
                for shape in _walk_shapes(slide.shapes):
                    try:
                        if (not shape.has_text_frame
                                or _is_footer_or_page_chrome(shape, slide)):
                            continue
                        text_value = shape.text_frame.text.strip()
                    except Exception:
                        continue
                    if text_value:
                        text_parts.append(text_value)
                signature = re.sub(
                    r"[\W_]+", "", "".join(text_parts).casefold()
                )
                if len(signature) >= 20:
                    signatures.setdefault(signature, []).append(slide_index)
            duplicate_slide_groups = [
                slide_numbers for slide_numbers in signatures.values()
                if len(slide_numbers) > 1
            ]
            backup_report = presentation_check_qa_backup_slides(
                str(final_path), min_backup=0
            )
            backup_slide_count = backup_report.get(
                "total_backup_candidates", 0
            )
            backup_slides = backup_report.get("backup_slides", [])
            figure_text_audit = presentation_check_embedded_figure_text_size(
                str(final_path),
                ocr_backend=figure_text_ocr_backend,
                ocr_manifest_path=figure_text_ocr_manifest_path,
                confirmed_textless_shapes=confirmed_textless_shapes,
            )
        except Exception as exc:
            open_error = str(exc)

    issues = []
    if len(top_level_pptx) != 1:
        issues.append("top_level_pptx_count_must_be_one")
    if archived_revision_pptx:
        issues.append("archived_revision_pptx_must_be_deleted")
    if duplicate_slide_groups:
        issues.append("duplicate_slides_must_be_consolidated")
    if figure_text_audit and not figure_text_audit.get("passed", False):
        issues.append("embedded_figure_text_must_be_verified_at_20pt")
    if open_error:
        issues.append("final_deck_could_not_be_opened")

    return {
        "passed": not issues,
        "directory": str(root),
        "final_candidate": final_candidate,
        "top_level_pptx_count": len(top_level_pptx),
        "top_level_pptx": [path.name for path in top_level_pptx],
        "archived_revision_pptx_count": len(archived_revision_pptx),
        "archived_revision_pptx": archived_revision_pptx,
        "duplicate_slide_groups": duplicate_slide_groups,
        "backup_slide_count": backup_slide_count,
        "backup_slides": backup_slides,
        "figure_text_audit": figure_text_audit,
        "open_error": open_error,
        "issues": issues,
        "rule": (
            "最終版は正本1本だけを残す。旧版固有の有用内容は同じPPTX末尾の"
            "質問対策・補足スライドへ移し、重複は品質の高い一枚へ統合してから"
            "旧版PPTXと検査ログを削除する。貼付図中文字は表示寸法20 pt以上を"
            "検証し、未確認画像を残さない。"
        ),
    }


def presentation_check_image_aspect_ratio(
        pptx_path: str,
        relative_tolerance: float = 0.001) -> dict:
    """埋め込み元画像に対する非等方な拡大・縮小を検出する。

    PowerPoint の crop 値を考慮し、表示枠の縦横比と、切り抜き後に
    期待される画像比を照合する。crop / contain は許容するが、幅と高さを
    独立に変更して画像を引き伸ばす、または押しつぶす操作は違反とする。

    ``relative_tolerance`` は OOXML の丸め誤差だけを吸収するための値で、
    既定値 0.001 は 0.1% に相当する。
    """
    try:
        import pptx as _pptx
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    if relative_tolerance < 0:
        return {"error": "relative_tolerance must be >= 0."}

    try:
        prs = _pptx.Presentation(str(p))
    except Exception as exc:
        return {"error": f"failed to open: {exc}"}

    violations = []
    unresolved = []
    images_checked = 0
    for slide_index, slide in enumerate(prs.slides, 1):
        for shape in _walk_shapes(slide.shapes):
            try:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
            except Exception:
                continue

            images_checked += 1
            shape_name = getattr(shape, "name", "")
            try:
                image_width, image_height = shape.image.size
                frame_width = int(shape.width)
                frame_height = int(shape.height)
                crop_left = float(shape.crop_left or 0.0)
                crop_right = float(shape.crop_right or 0.0)
                crop_top = float(shape.crop_top or 0.0)
                crop_bottom = float(shape.crop_bottom or 0.0)
            except Exception as exc:
                unresolved.append({
                    "slide": slide_index,
                    "shape": shape_name,
                    "reason": str(exc),
                })
                continue

            visible_width = 1.0 - crop_left - crop_right
            visible_height = 1.0 - crop_top - crop_bottom
            if (image_width <= 0 or image_height <= 0
                    or frame_width <= 0 or frame_height <= 0
                    or visible_width <= 0 or visible_height <= 0):
                unresolved.append({
                    "slide": slide_index,
                    "shape": shape_name,
                    "reason": "invalid image, frame, or crop dimensions",
                })
                continue

            source_ratio = image_width / image_height
            expected_ratio = source_ratio * visible_width / visible_height
            frame_ratio = frame_width / frame_height
            relative_error = abs(frame_ratio / expected_ratio - 1.0)
            if relative_error > relative_tolerance:
                violations.append({
                    "slide": slide_index,
                    "shape": shape_name,
                    "issue": "image_aspect_ratio_changed",
                    "source_ratio": round(source_ratio, 6),
                    "expected_ratio_after_crop": round(expected_ratio, 6),
                    "frame_ratio": round(frame_ratio, 6),
                    "relative_error": round(relative_error, 6),
                    "relative_error_percent": round(relative_error * 100, 3),
                    "crop": {
                        "left": crop_left,
                        "right": crop_right,
                        "top": crop_top,
                        "bottom": crop_bottom,
                    },
                })

    passed = not violations and not unresolved
    return {
        "passed": passed,
        "images_checked": images_checked,
        "violation_count": len(violations),
        "unresolved_count": len(unresolved),
        "relative_tolerance": relative_tolerance,
        "violations": violations,
        "unresolved": unresolved,
        "rule": (
            "埋め込み元画像の縦横比を変更しない。枠に合わせる場合は、"
            "縦横比を固定した等方拡大・縮小と crop / contain を用いる。"
        ),
        "hint": (
            "LockAspectRatio は既に歪んだ枠比を固定する場合がある。"
            "元画像寸法へリセットしてから縦横比を固定し、一辺だけを変更する。"
        ),
    }


def presentation_check_embedded_figure_text_size(
        pptx_path: str,
        min_font_pt: float = 20.0,
        ocr_backend: str = "none",
        ocr_manifest_path: str = "",
        confirmed_textless_shapes: list[str] | None = None,
        min_confidence: float = 0.30,
        glyph_to_font_ratio: float = 0.72) -> dict:
    """画像へ焼き込まれた図中文字のスライド上換算サイズを検査する。

    Picture shape の表示高さと埋め込み画像の pixel 高さから OCR bounding
    box を point へ換算する。OCR box は字面の高さなので、既定では Latin
    capital-height に近い 0.72 em を用いて font size を推定する。20 pt は
    合格下限であり、OCR誤差と縮小余裕を考えて元図は24 pt以上を推奨する。

    ``ocr_backend``:
    - ``none``: 外部送信しない。全画像を未確認として返す。ただし
      ``confirmed_textless_shapes`` の ``"slide:shape name"`` は除外できる。
    - ``manifest``: ``ocr_manifest_path`` の再現可能な word boxes を用いる。
    - ``gcv``: Google Cloud Visionへ画像を送信してword boxesを取得する。
      未公開資料では明示的な許可なしに使用しない。

    manifest schema (OCR boxes or deterministic source-size evidence)::

        {"pictures": [{"slide": 3, "shape": "Picture 2", "words": [
          {"text": "HACApK", "confidence": 0.99,
           "bbox": [120, 40, 260, 72]}]},
          {"slide": 4, "shape": "Picture 3", "source_evidence": {
           "minimum_source_font_pt": 24, "source_width_cm": 16.5}}]}

    ``bbox`` は image pixel の ``[x0, y0, x1, y1]`` または4頂点とする。
    ``source_evidence`` は図生成コードで保証された最小フォントと物理幅を
    用い、PowerPoint上の表示幅（cropを含む）から最終ptへ換算する。
    OCR不能・manifest欠落・不完全なsource evidenceは合格扱いにしない。
    """
    try:
        import pptx as _pptx
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return {"error": "python-pptx not installed."}

    path = pathlib.Path(pptx_path)
    if not path.exists():
        return {"error": f"file not found: {pptx_path}"}
    if min_font_pt <= 0:
        return {"error": "min_font_pt must be > 0."}
    if not 0.0 <= min_confidence <= 1.0:
        return {"error": "min_confidence must be between 0 and 1."}
    if not 0.4 <= glyph_to_font_ratio <= 1.0:
        return {"error": "glyph_to_font_ratio must be between 0.4 and 1.0."}
    if ocr_backend not in {"none", "manifest", "gcv"}:
        return {"error": "ocr_backend must be none, manifest, or gcv."}

    confirmed_textless = set(confirmed_textless_shapes or [])
    manifest = {}
    if ocr_backend == "manifest":
        manifest_path = pathlib.Path(ocr_manifest_path)
        if not manifest_path.exists():
            return {"error": f"OCR manifest not found: {ocr_manifest_path}"}
        try:
            payload = json.loads(manifest_path.read_text(encoding="utf-8"))
            for item in payload.get("pictures", []):
                key = (int(item["slide"]), str(item["shape"]))
                manifest[key] = item
        except Exception as exc:
            return {"error": f"invalid OCR manifest: {exc}"}

    gcv_client = None
    gcv_vision = None
    if ocr_backend == "gcv":
        try:
            from google.cloud import vision as gcv_vision
            gcv_client = gcv_vision.ImageAnnotatorClient()
        except Exception as exc:
            return {"error": f"Google Cloud Vision is unavailable: {exc}"}

    def _bbox_height_px(bbox) -> float:
        if (isinstance(bbox, list) and len(bbox) == 4
                and all(isinstance(value, (int, float)) for value in bbox)):
            return abs(float(bbox[3]) - float(bbox[1]))
        if isinstance(bbox, list) and len(bbox) >= 4:
            points = [(float(point[0]), float(point[1])) for point in bbox[:4]]
            left = ((points[3][0] - points[0][0]) ** 2
                    + (points[3][1] - points[0][1]) ** 2) ** 0.5
            right = ((points[2][0] - points[1][0]) ** 2
                     + (points[2][1] - points[1][1]) ** 2) ** 0.5
            return max(left, right)
        return 0.0

    def _bbox_center_y(bbox) -> float:
        if (isinstance(bbox, list) and len(bbox) == 4
                and all(isinstance(value, (int, float)) for value in bbox)):
            return (float(bbox[1]) + float(bbox[3])) / 2.0
        if isinstance(bbox, list) and bbox:
            return sum(float(point[1]) for point in bbox) / len(bbox)
        return -1.0

    def _gcv_words(blob: bytes) -> list[dict]:
        image = gcv_vision.Image(content=blob)
        context = gcv_vision.ImageContext(language_hints=["ja", "en"])
        response = gcv_client.document_text_detection(
            image=image, image_context=context
        )
        if response.error.message:
            raise RuntimeError(response.error.message)
        words = []
        for page_result in response.full_text_annotation.pages:
            for block in page_result.blocks:
                for paragraph in block.paragraphs:
                    for word in paragraph.words:
                        text_value = "".join(
                            symbol.text for symbol in word.symbols
                        )
                        bbox = [
                            [vertex.x or 0, vertex.y or 0]
                            for vertex in word.bounding_box.vertices
                        ]
                        words.append({
                            "text": text_value,
                            "confidence": float(word.confidence or 0.0),
                            "bbox": bbox,
                        })
        return words

    try:
        prs = _pptx.Presentation(str(path))
    except Exception as exc:
        return {"error": f"failed to open: {exc}"}

    picture_reports = []
    violations = []
    unresolved = []
    pictures_checked = 0
    words_checked = 0
    emu_per_point = 12700.0
    for slide_index, slide in enumerate(prs.slides, 1):
        for shape in _walk_shapes(slide.shapes):
            try:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
            except Exception:
                continue

            pictures_checked += 1
            shape_name = getattr(shape, "name", "")
            shape_key = f"{slide_index}:{shape_name}"
            if shape_key in confirmed_textless:
                picture_reports.append({
                    "slide": slide_index,
                    "shape": shape_name,
                    "status": "confirmed_textless",
                    "minimum_estimated_font_pt": None,
                })
                continue

            try:
                image_width, image_height = shape.image.size
                if image_width <= 0 or image_height <= 0:
                    raise ValueError("invalid embedded image dimensions")
                if ocr_backend == "manifest":
                    manifest_item = manifest.get((slide_index, shape_name))
                    if manifest_item is None:
                        raise ValueError("picture is missing from OCR manifest")
                    if manifest_item.get("confirmed_textless"):
                        picture_reports.append({
                            "slide": slide_index,
                            "shape": shape_name,
                            "status": "confirmed_textless",
                            "minimum_estimated_font_pt": None,
                        })
                        continue
                    source_evidence = manifest_item.get("source_evidence")
                    if source_evidence is not None:
                        source_font_pt = float(
                            source_evidence["minimum_source_font_pt"]
                        )
                        if source_font_pt <= 0:
                            raise ValueError(
                                "minimum_source_font_pt must be positive"
                            )
                        if "source_width_pt" in source_evidence:
                            source_width_pt = float(
                                source_evidence["source_width_pt"]
                            )
                        elif "source_width_cm" in source_evidence:
                            source_width_pt = (
                                float(source_evidence["source_width_cm"])
                                / 2.54 * 72.0
                            )
                        elif "source_width_in" in source_evidence:
                            source_width_pt = (
                                float(source_evidence["source_width_in"])
                                * 72.0
                            )
                        else:
                            raise ValueError(
                                "source_evidence requires source_width_pt, "
                                "source_width_cm, or source_width_in"
                            )
                        if source_width_pt <= 0:
                            raise ValueError("source width must be positive")
                        visible_width_fraction = (
                            1.0
                            - float(shape.crop_left or 0.0)
                            - float(shape.crop_right or 0.0)
                        )
                        if visible_width_fraction <= 0:
                            raise ValueError("invalid horizontal crop dimensions")
                        displayed_width_pt = (
                            float(shape.width) / emu_per_point
                        )
                        scale = displayed_width_pt / (
                            source_width_pt * visible_width_fraction
                        )
                        minimum = source_font_pt * scale
                        report = {
                            "slide": slide_index,
                            "shape": shape_name,
                            "status": (
                                "violation"
                                if minimum < min_font_pt else "pass"
                            ),
                            "evidence_type": "source_size",
                            "minimum_source_font_pt": round(
                                source_font_pt, 2
                            ),
                            "source_width_pt": round(source_width_pt, 2),
                            "displayed_width_pt": round(
                                displayed_width_pt, 2
                            ),
                            "embed_scale": round(scale, 4),
                            "minimum_estimated_font_pt": round(minimum, 2),
                            "words_checked": 0,
                            "small_text": [],
                        }
                        picture_reports.append(report)
                        if minimum < min_font_pt:
                            violations.append(report)
                        continue
                    words = list(manifest_item.get("words", []))
                elif ocr_backend == "gcv":
                    words = _gcv_words(shape.image.blob)
                else:
                    raise ValueError("OCR or source-size evidence is required")

                crop_top = float(shape.crop_top or 0.0)
                crop_bottom = float(shape.crop_bottom or 0.0)
                visible_top_px = crop_top * image_height
                visible_bottom_px = (1.0 - crop_bottom) * image_height
                visible_height_px = visible_bottom_px - visible_top_px
                if visible_height_px <= 0:
                    raise ValueError("invalid vertical crop dimensions")
                displayed_height_pt = float(shape.height) / emu_per_point

                estimates = []
                low_confidence_words = 0
                for word in words:
                    bbox = word.get("bbox", [])
                    center_y = _bbox_center_y(bbox)
                    if not visible_top_px <= center_y <= visible_bottom_px:
                        continue
                    confidence = float(word.get("confidence", 1.0))
                    if confidence < min_confidence:
                        low_confidence_words += 1
                        continue
                    glyph_height_px = _bbox_height_px(bbox)
                    if glyph_height_px <= 0:
                        continue
                    glyph_height_pt = (
                        glyph_height_px * displayed_height_pt / visible_height_px
                    )
                    estimated_font_pt = glyph_height_pt / glyph_to_font_ratio
                    estimates.append({
                        "text": str(word.get("text", ""))[:40],
                        "confidence": round(confidence, 3),
                        "estimated_font_pt": round(estimated_font_pt, 2),
                    })

                if not estimates:
                    reason = "no reliable text boxes detected"
                    if low_confidence_words:
                        reason += f" ({low_confidence_words} below confidence)"
                    raise ValueError(reason)

                words_checked += len(estimates)
                minimum = min(
                    item["estimated_font_pt"] for item in estimates
                )
                small_words = [
                    item for item in estimates
                    if item["estimated_font_pt"] < min_font_pt
                ]
                report = {
                    "slide": slide_index,
                    "shape": shape_name,
                    "status": "violation" if small_words else "pass",
                    "minimum_estimated_font_pt": round(minimum, 2),
                    "words_checked": len(estimates),
                    "small_text": small_words[:20],
                }
                picture_reports.append(report)
                if small_words:
                    violations.append(report)
            except Exception as exc:
                item = {
                    "slide": slide_index,
                    "shape": shape_name,
                    "reason": str(exc),
                }
                unresolved.append(item)
                picture_reports.append({
                    **item,
                    "status": "unresolved",
                    "minimum_estimated_font_pt": None,
                })

    return {
        "passed": not violations and not unresolved,
        "min_font_pt": min_font_pt,
        "source_figure_target_pt": 24.0,
        "ocr_backend": ocr_backend,
        "pictures_checked": pictures_checked,
        "words_checked": words_checked,
        "violation_count": len(violations),
        "unresolved_count": len(unresolved),
        "violations": violations,
        "unresolved": unresolved,
        "pictures": picture_reports,
        "rule": (
            "貼り付けた図中文字は実際のスライド表示寸法で20 pt以上。"
            "元図は24 pt以上を標準とし、縮小後20 pt未満または未確認の図は"
            "最終版として合格させない。"
        ),
        "hint": (
            "小さい文字は元図を再生成するか、図から削除してPowerPointの"
            "ネイティブ文字で24 pt以上として置き直す。"
        ),
    }


def presentation_check_image_text_ratio(pptx_path: str,
                                          min_image_ratio: float = 0.3) -> dict:
    """1 slide の image 面積比が min 未満を検出 (Zen style)."""
    try:
        import pptx as _pptx
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))
    total_area = prs.slide_width * prs.slide_height
    low_image = []
    for i, s in enumerate(prs.slides, 1):
        img_area = 0
        txt_area = 0
        for shape in s.shapes:
            area = shape.width * shape.height
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_area += area
                elif shape.has_text_frame and shape.text_frame.text.strip():
                    txt_area += area
            except Exception:
                pass
        img_ratio = img_area / max(1, total_area)
        txt_ratio = txt_area / max(1, total_area)
        if img_ratio < min_image_ratio and txt_ratio > 0.15:
            low_image.append({"slide": i,
                              "image_ratio": round(img_ratio, 2),
                              "text_ratio": round(txt_ratio, 2)})
    return {
        "total_slides": len(prs.slides),
        "low_image_slide_count": len(low_image),
        "min_image_ratio": min_image_ratio,
        "low_image_slides": low_image[:15],
        "hint": "Zen: 画像 >= 50% 推奨。text を削り image を追加。",
    }


def presentation_check_color_count_per_slide(pptx_path: str,
                                                 max_colors: int = 5,
                                                 include_theme: bool = True) -> dict:
    """宮野『研究発表のためのスライドデザイン』S12: 3 色使い原則の検査 (v0.9.0)。

    **原則** (宮野 第 2 部): ベース色 + メイン色 + アクセント色 の **3 系統** を
    基本とし、計 3-5 色以内で統一する。無意味なグラデーションや多色使いを避ける。

    **検査対象**: 各 slide 内で text runs / shape fills / shape lines に使われた
    ユニークな RGB 色を列挙。`max_colors` を超えた slide を flag。

    Args:
        pptx_path: .pptx path。
        max_colors: 1 slide あたりの許容色数 (既定 5、宮野基準なら 3)。
        include_theme: True ならテーマ色も個別色として数える (既定 True)。

    Returns:
        {total_slides, violating_slides: [{slide, color_count, palette}],
         global_palette: [...], hint, source}
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed. "
                         "pip install mcp-server-document[pptx]"}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    try:
        prs = _pptx.Presentation(str(p))
    except Exception as exc:
        return {"error": f"failed to open: {exc}"}

    def _rgb_of(color_fmt):
        """Extract a canonical color token.

        Two namespaces share the returned set:

        * RGB colors: ``(int_r, int_g, int_b)`` (three ints).
        * Theme colors (when ``include_theme``): ``("theme", int_idx)``.

        These coexist inside the same ``palette`` set (``len(palette)``
        is the "distinct color" count). The display formatter below
        prints RGB as ``#rrggbb`` and theme colors as ``theme-<idx>``.
        """
        try:
            if color_fmt is None:
                return None
            # Try .rgb first
            rgb = getattr(color_fmt, "rgb", None)
            if rgb is not None:
                return (int(rgb[0]), int(rgb[1]), int(rgb[2]))
            # Theme colors: still track as theme_X if include_theme
            if include_theme:
                theme = getattr(color_fmt, "theme_color", None)
                if theme is not None:
                    return ("theme", int(theme))
        except Exception:
            return None
        return None

    def _fmt_color(c) -> str:
        """Pretty-print either a 3-int RGB tuple or a ('theme', idx)."""
        # Bug H5 fix: theme-color tuples start with the string "theme".
        # The original ``isinstance(c[0], int)`` branch fell through to
        # ``str(c)``, producing ugly ``('theme', 5)`` output. We now
        # format them as a readable ``theme-5`` token.
        if isinstance(c, tuple) and len(c) == 2 and c[0] == "theme":
            return f"theme-{c[1]}"
        if isinstance(c, tuple) and len(c) == 3 and all(isinstance(x, int) for x in c):
            return f"#{c[0]:02x}{c[1]:02x}{c[2]:02x}"
        return str(c)

    violating = []
    global_palette = set()

    for i, s in enumerate(prs.slides, 1):
        palette = set()
        for shape in s.shapes:
            # Text colors
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            c = _rgb_of(run.font.color)
                            if c is not None:
                                palette.add(c)
                        except Exception:
                            pass
            # Fill color
            try:
                fill = getattr(shape, "fill", None)
                if fill is not None:
                    fc = getattr(fill, "fore_color", None)
                    c = _rgb_of(fc)
                    if c is not None:
                        palette.add(c)
            except Exception:
                pass
            # Line color
            try:
                line = getattr(shape, "line", None)
                if line is not None:
                    lc = getattr(line, "color", None)
                    c = _rgb_of(lc)
                    if c is not None:
                        palette.add(c)
            except Exception:
                pass
        global_palette.update(palette)
        if len(palette) > max_colors:
            violating.append({
                "slide": i,
                "color_count": len(palette),
                "palette_sample": [_fmt_color(c) for c in list(palette)[:10]],
            })

    return {
        "total_slides": len(prs.slides),
        "max_colors_allowed": max_colors,
        "violating_slides": violating[:20],
        "violating_count": len(violating),
        "global_unique_colors": len(global_palette),
        "hint": (
            "宮野『研究発表のためのスライドデザイン』: ベース色 + メイン色 + "
            "アクセント色 の 3 系統を基本、max_colors=3 で厳格適用。**hint** — "
            "テーマ色は PowerPoint 内部で自動解決される場合があり、検出が不完全な "
            "こともある。最終確認は目視。"
        ),
        "source": "宮野公樹『研究発表のためのスライドデザイン』講談社 2010, S12 配色原則",
    }


def presentation_check_color_accessibility(pptx_path: str) -> dict:
    """R+G 近接色ペアを検出 (protanopia/deuteranopia で区別困難)."""
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))
    issues = []
    for i, s in enumerate(prs.slides, 1):
        colors = []
        for shape in s.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color and run.font.color.rgb:
                                rgb = run.font.color.rgb
                                colors.append((int(rgb[0]), int(rgb[1]), int(rgb[2])))
                        except Exception:
                            pass
        red_like = [c for c in colors if c[0] > 150 and c[1] < 100 and c[2] < 100]
        green_like = [c for c in colors if c[1] > 150 and c[0] < 100 and c[2] < 100]
        if red_like and green_like:
            issues.append({"slide": i, "red_colors": len(red_like),
                           "green_colors": len(green_like)})
    return {
        "total_slides": len(prs.slides),
        "accessibility_issues": issues[:10],
        "hint": "ColorBrewer 色盲対応パレット推奨。",
    }


def presentation_check_bullet_ending_style(pptx_path: str) -> dict:
    """bullet 末尾の「。」有無が統一されているか."""
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))
    with_period = 0
    without_period = 0
    examples = []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if not t:
                    continue
                if t.endswith(("。", ".")):
                    with_period += 1
                    if len(examples) < 5:
                        examples.append({"slide": i, "text": t[:60]})
                else:
                    without_period += 1
    return {
        "bullets_with_period": with_period,
        "bullets_without_period": without_period,
        "mixed": with_period > 0 and without_period > 0,
        "examples": examples,
        "recommendation": "bullet 末尾句点なしに統一推奨。",
    }


def presentation_check_japanese_copy_style(
        pptx_path: str,
        nominal_ratio_target: float = 0.60,
        min_prose_chars: int = 10) -> dict:
    """和文スライド本文の文節改行と体言止めを点検。

    宮野『研究発表のためのスライドデザイン』の「文脈を優先し、
    違和感のない位置で改行」と「長い説明文は体言止めで簡潔化」を
    PPTX 上で補助診断する。タイトル、footer/page chrome、数式、URL、
    疑問文は体言止め比率から除外する。日本語形態素解析を伴わない
    heuristic であり、自動修正ではなく読み上げ前の候補抽出に用いる。
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    if not 0.0 <= nominal_ratio_target <= 1.0:
        return {"error": "nominal_ratio_target must be between 0 and 1."}
    if min_prose_chars < 1:
        return {"error": "min_prose_chars must be >= 1."}

    prs = _pptx.Presentation(str(p))
    structural_title = re.compile(
        r"^(?:title|references?|acknowledg(?:e)?ments?|appendix|q\s*&\s*a|"
        r"参考文献|謝辞|付録|質疑|補足)\s*$",
        re.IGNORECASE,
    )
    dependent_line_end = re.compile(
        r"(?:について|によって|により|による|として|"
        r"および|及び|ならびに|または|又は|"
        r"から|まで|より|ため|の|を|に|が|は|へ|と|で|や|も)$"
    )
    dependent_line_start = re.compile(
        r"^(?:を|が|は|に|へ|と|で|も|の)(?=[^\s、。，．])"
    )
    verbal_ending = re.compile(
        r"(?:です|ます|ました|ません|でした|である|であった|"
        r"となる|になる|となった|になった|している|"
        r"される|された|できる|行う|示す|示した|確認した|"
        r"比較した|評価した|検討した|用いる|用いた|"
        r"得られた|考えられる|認められる|必要がある|"
        r"可能である|重要である|望ましい|難しい|新しい|"
        r"高い|低い|大きい|小さい|等しい|異なる|"
        r"一致する|対応する|維持する|低減する|向上する|"
        r"改善する|実現する|保証する|抑制する|支配する|"
        r"働く|含む|持つ|保つ|残る|進む|優れる|分かる|わかる|"
        r"ある|ない|する|した)$"
    )
    screen_only_lead_in = re.compile(
        r"^(?:そこで本研究では|本研究では|ここでは|次に)[、，,\s]*"
    )

    awkward_breaks = []
    verbal_examples = []
    lead_in_examples = []
    paragraphs_checked = 0
    nominal_like_count = 0
    verbal_sentence_count = 0
    skipped_paragraphs = 0
    slides_checked = 0

    def _is_prose_candidate(text_value: str) -> bool:
        compact = re.sub(r"\s+", "", text_value)
        if len(compact) < min_prose_chars:
            return False
        if text_value.rstrip().endswith(("?", "？")):
            return False
        if re.search(r"(?:https?://|www\.|doi\s*:)", text_value, re.IGNORECASE):
            return False
        if re.fullmatch(r"[\d\s.,:+\-*/=()\[\]{}<>%％×·・〜～]+", text_value):
            return False
        japanese_chars = len(re.findall(r"[ぁ-ゟ゠-ヿ㐀-鿿]", text_value))
        return japanese_chars >= 3

    for slide_index, slide in enumerate(prs.slides, 1):
        title = _slide_title(slide).strip()
        if slide_index == 1 or structural_title.match(title):
            continue
        slides_checked += 1
        for shape in _walk_shapes(slide.shapes):
            try:
                if (not shape.has_text_frame
                        or _shape_is_slide_title(shape, slide)
                        or _is_footer_or_page_chrome(shape, slide)):
                    continue
            except Exception:
                continue

            shape_name = getattr(shape, "name", "")
            for paragraph in shape.text_frame.paragraphs:
                raw = paragraph.text.replace("\r\n", "\n").replace("\r", "\n")
                lines = [line.strip() for line in re.split(r"[\n\v]", raw)
                         if line.strip()]
                for line_index in range(len(lines) - 1):
                    before = lines[line_index]
                    after = lines[line_index + 1]
                    reason = ""
                    if dependent_line_end.search(before):
                        reason = "助詞または連体修飾の直後で改行"
                    elif dependent_line_start.search(after):
                        reason = "改行後の行が従属的な助詞で開始"
                    elif (re.search(r"[A-Za-z0-9-]$", before)
                          and re.match(r"^[a-z0-9]", after)):
                        reason = "英単語または識別子の途中で改行した可能性"
                    elif re.search(r"[\(（「『【]$", before):
                        reason = "開き括弧が行末に孤立"
                    elif re.match(r"^[\)）」』】、。，．]", after):
                        reason = "閉じ括弧または句読点が行頭に孤立"
                    if reason:
                        awkward_breaks.append({
                            "slide": slide_index,
                            "shape": shape_name,
                            "before": before[:100],
                            "after": after[:100],
                            "reason": reason,
                        })

                text_value = "".join(lines).strip()
                if not _is_prose_candidate(text_value):
                    skipped_paragraphs += 1
                    continue
                paragraphs_checked += 1
                if screen_only_lead_in.search(text_value):
                    lead_in_examples.append({
                        "slide": slide_index,
                        "shape": shape_name,
                        "text": text_value[:140],
                    })
                ending = text_value.rstrip("　 \t、，,;:；：。．.!\uff01")
                if verbal_ending.search(ending):
                    verbal_sentence_count += 1
                    verbal_examples.append({
                        "slide": slide_index,
                        "shape": shape_name,
                        "text": text_value[:140],
                    })
                else:
                    nominal_like_count += 1

    ratio = (nominal_like_count / paragraphs_checked
             if paragraphs_checked else 1.0)
    ratio_passed = ratio >= nominal_ratio_target
    return {
        "passed": ratio_passed and not awkward_breaks and not lead_in_examples,
        "slides_checked": slides_checked,
        "paragraphs_checked": paragraphs_checked,
        "skipped_paragraphs": skipped_paragraphs,
        "nominal_like_count": nominal_like_count,
        "verbal_sentence_count": verbal_sentence_count,
        "nominal_like_ratio": round(ratio, 3),
        "nominal_ratio_target": nominal_ratio_target,
        "nominal_ratio_passed": ratio_passed,
        "awkward_line_break_count": len(awkward_breaks),
        "awkward_line_breaks": awkward_breaks[:50],
        "screen_only_lead_in_count": len(lead_in_examples),
        "screen_only_lead_ins": lead_in_examples[:30],
        "verbal_sentence_examples": verbal_examples[:30],
        "rule": {
            "line_break": "文字数でなく文脈・文節を優先し、違和感のない位置で改行。",
            "ending": "長い本文は体言止めを主とし、不要な接続句を画面から削除。",
        },
        "source": (
            "宮野公樹『研究発表のためのスライドデザイン』"
            "p.50-52, p.102"
        ),
        "hint": (
            "形態素解析なしの補助診断。指摘候補は声に出して確認し、"
            "タイトルの主張文や自然な説明文まで無理に体言止めにしない。"
        ),
    }


def presentation_estimate_per_slide_time(script: str,
                                           slide_boundaries: str = "",
                                           wpm_ja: int = 240) -> dict:
    """原稿を slide 境界で分割し、各 slide の発表時間を推定."""
    if slide_boundaries:
        markers = [m.strip() for m in slide_boundaries.split(",") if m.strip()]
    else:
        markers = ["[[SLIDE]]"]
    pattern = "|".join(re.escape(m) for m in markers)
    parts = re.split(pattern, script)
    parts = [p for p in parts if p.strip()]
    per_slide = []
    for i, part in enumerate(parts, 1):
        chars = len(part.replace(" ", "").replace("\n", ""))
        minutes = chars / wpm_ja
        per_slide.append({"slide": i, "chars": chars,
                          "minutes": round(minutes, 2)})
    total_min = sum(s["minutes"] for s in per_slide)
    if len(per_slide) >= 2:
        avg = total_min / len(per_slide)
        max_dev = max(abs(s["minutes"] - avg) for s in per_slide)
    else:
        avg = total_min
        max_dev = 0
    return {
        "slide_count": len(per_slide),
        "total_minutes": round(total_min, 2),
        "avg_min_per_slide": round(avg, 2),
        "max_deviation_min": round(max_dev, 2),
        "per_slide": per_slide[:20],
        "hint": "slide 間時間のばらつきが大きい場合、長い slide を 2 枚に分割。",
    }


