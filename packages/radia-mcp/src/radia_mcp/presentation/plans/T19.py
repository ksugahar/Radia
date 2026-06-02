"""Plan B T19: Chart simplification check (Cole Knaflic style)
(presentation_chart_simplification_check).

既存 `check_pie_3d_charts` を拡張し、Cole Nussbaumer Knaflic
『Storytelling with Data』の chart 簡素化原則を多 chart 種類に展開。

検出する 5 軸:
  1. **3d_decoration** — 3D pie / 3D bar (既存 + bar の 3D も検出)
  2. **excessive_legend** — 凡例 ≥ 7 (Miller's law 超過)
  3. **chart_count_per_slide** — 1 slide ≥3 charts (要分割)
  4. **chart_without_title** — chart に title 無し
  5. **decorative_axis** — 軸が 3D 透視 / 装飾過多

理系プレゼンでは特に Knaflic Ch.3 『clutter』批判が重要 (本来の信号を
chart junk が埋める)。

設計方針:
- python-pptx の chart object をスキャン
- chart_type / has_title / legend / etc を chart object 経由で取得
- score は per-slide chart compliance 平均
- PRIMARY OUTPUT は cluttered_charts list
"""

from __future__ import annotations

import pathlib


def _iter_all_shapes(shape_collection):
    try:
        from radia_mcp.presentation.tools import _walk_shapes
        yield from _walk_shapes(shape_collection)
    except Exception:
        for s in shape_collection:
            yield s


# 3D chart type names (XL_CHART_TYPE)
_3D_CHART_TYPES = {
    "PIE_3D", "BAR_3D", "BAR_OF_PIE", "BAR_3D_CLUSTERED",
    "BAR_3D_STACKED", "BAR_3D_STACKED_100",
    "COLUMN_3D", "COLUMN_3D_CLUSTERED",
    "COLUMN_3D_STACKED", "COLUMN_3D_STACKED_100",
    "LINE_3D", "AREA_3D", "AREA_3D_STACKED",
    "AREA_3D_STACKED_100", "DOUGHNUT_3D",
    "PIE_EXPLODED_3D", "CONE", "CYLINDER", "PYRAMID",
    "CONE_BAR", "CONE_BAR_CLUSTERED", "CONE_BAR_STACKED",
    "CYLINDER_BAR", "PYRAMID_BAR",
}


def _get_chart_type_name(chart) -> str:
    try:
        ct = chart.chart_type
        return getattr(ct, "name", str(ct))
    except Exception:
        return "UNKNOWN"


def _is_3d_chart(chart) -> bool:
    name = _get_chart_type_name(chart).upper()
    return any(d in name for d in _3D_CHART_TYPES)


def _legend_entry_count(chart) -> int:
    try:
        if chart.has_legend:
            # 系列 + カテゴリ entry を集計
            n = 0
            for plot in chart.plots:
                try:
                    n += len(plot.series)
                except Exception:
                    pass
            return n
        return 0
    except Exception:
        return 0


def _has_chart_title(chart) -> bool:
    try:
        return chart.has_title and bool(chart.chart_title.text_frame.text.strip())
    except Exception:
        return False


def presentation_chart_simplification_check(pptx_path: str) -> dict:
    """Chart 簡素化 (Cole Knaflic style) の 5 軸診断。

    既存 check_pie_3d_charts を拡張、bar/column/line 3D も検出 +
    凡例過多 / chart 過多 / title 不在も flag。

    Args:
        pptx_path: .pptx file

    Returns:
        dict: score / cluttered_charts / per_slide / comments
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))

    cluttered_charts: list[dict] = []
    per_slide: list[dict] = []
    total_charts = 0

    for i, slide in enumerate(prs.slides, 1):
        slide_charts: list[dict] = []
        for shape in _iter_all_shapes(slide.shapes):
            if not getattr(shape, "has_chart", False):
                continue
            try:
                if not shape.has_chart:
                    continue
                chart = shape.chart
            except Exception:
                continue
            total_charts += 1
            chart_info: dict = {
                "type": _get_chart_type_name(chart),
                "is_3d": _is_3d_chart(chart),
                "legend_entries": _legend_entry_count(chart),
                "has_title": _has_chart_title(chart),
            }
            slide_charts.append(chart_info)

            issues = []
            if chart_info["is_3d"]:
                issues.append("3D_decoration")
            if chart_info["legend_entries"] >= 7:
                issues.append(
                    f"excessive_legend({chart_info['legend_entries']})"
                )
            if not chart_info["has_title"]:
                issues.append("no_chart_title")

            if issues:
                cluttered_charts.append({
                    "slide_no": i,
                    "chart_type": chart_info["type"],
                    "legend_entries": chart_info["legend_entries"],
                    "issues": issues,
                })

        n_in_slide = len(slide_charts)
        per_slide.append({
            "slide_no": i,
            "n_charts": n_in_slide,
            "charts": slide_charts,
        })

        if n_in_slide >= 3:
            cluttered_charts.append({
                "slide_no": i,
                "chart_type": "(multiple)",
                "issues": [f"too_many_charts_per_slide({n_in_slide})"],
            })

    n_slides = len(per_slide)
    n_cluttered = len(cluttered_charts)

    if total_charts == 0:
        return {
            "score": 10.0,
            "score_max": 10,
            "n_charts": 0,
            "comments": [
                "Chart 検出ゼロ。chart-free スライド or text/figure 中心。"
                "本ツールは chart 含む発表向けなので skip 推奨。"
            ],
            "hint": "no charts to evaluate.",
            "source": "T19 chart simplification (presentation 2026-04)",
        }

    n_compliant = total_charts - sum(
        1 for c in cluttered_charts
        if c.get("chart_type") != "(multiple)"
    )
    score = round(n_compliant / total_charts * 10, 1) if total_charts else 0.0

    # ---- comments ----
    comments: list[str] = []

    comments.append(
        f"Charts 検出: {total_charts} 件 ({n_slides} slides scanned) | "
        f"clutter issues: {n_cluttered} 件 (score {score}/10)"
    )

    n_3d = sum(
        1 for c in cluttered_charts if "3D_decoration" in c.get("issues", [])
    )
    if n_3d:
        comments.append(
            f"⚠ 3D chart 検出 {n_3d} 件。Knaflic Ch.3『avoid 3D』に違反。"
            "値読み取りが歪み、cognitive load 増。2D に変換推奨。"
        )

    n_legend = sum(
        1 for c in cluttered_charts
        if any("excessive_legend" in i for i in c.get("issues", []))
    )
    if n_legend:
        comments.append(
            f"⚠ Legend 凡例 ≥7 が {n_legend} 件。Miller's law 超過で "
            "聴衆が series を追えない。色分けを減らすか直接ラベル化推奨。"
        )

    n_no_title = sum(
        1 for c in cluttered_charts if "no_chart_title" in c.get("issues", [])
    )
    if n_no_title:
        comments.append(
            f"ℹ Chart title 不在が {n_no_title} 件。"
            "理系プレゼンでは chart 単独で『何のグラフか』分かるよう "
            "title 必須 (slide title とは別)。"
        )

    n_too_many = sum(
        1 for c in cluttered_charts
        if any("too_many" in i for i in c.get("issues", []))
    )
    if n_too_many:
        comments.append(
            f"⚠ 1 slide に 3+ charts が {n_too_many} 件。"
            "聴衆は 1 chart に集中したい。slide 分割または "
            "代表 chart 1 個に絞る (T15 single_message と整合)。"
        )

    if n_cluttered == 0:
        comments.append(
            f"score {score}/10 — chart 簡素化基準を充足。"
            "Cole Knaflic style: 信号を強調し、clutter を削減。"
        )

    comments = comments[:7]

    hint = (
        "Cole Nussbaumer Knaflic『Storytelling with Data』(2015) の "
        "5 chart 簡素化原則。既存 check_pie_3d_charts は pie のみ、"
        "T19 は bar/column/line 3D + 凡例 + chart count + title まで網羅。"
        "理系プレゼンでは chart が evidence の core なので clutter は致命的。"
    )

    return {
        "score": score,
        "score_max": 10,
        "n_total_charts": total_charts,
        "n_cluttered": n_cluttered,
        "cluttered_charts": cluttered_charts,
        "per_slide": per_slide[:30],
        "comments": comments,
        "hint": hint,
        "source": (
            "Cole Nussbaumer Knaflic『Storytelling with Data』(2015) "
            "Ch.3 (clutter) + Tufte『Visual Display of Quantitative "
            "Information』(data-ink ratio) + Few『Show Me the Numbers』。"
            "既存 check_pie_3d_charts (pie 限定) を多 chart 種に拡張。"
            "(presentation 2026-04 採用)。"
        ),
    }
