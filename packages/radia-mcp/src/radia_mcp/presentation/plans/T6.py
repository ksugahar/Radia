"""Plan B T6: per-slide visual/text ratio 診断 (presentation_visual_text_ratio_score).

skill.md 宮野 S3 (KISS) / Reynolds『3 ideas / Presentation Zen』に基づく粗い
温度計。各スライドで「画像が占める面積」と「テキストを載せている shape の
面積」を EMU 単位で集計し、visual_ratio = image / (image + text + 1) の
分布から text-crowded slide を上位 3 件特定する。

設計方針:
- score は 0-10 の粗い温度計。final judgement は comments と per_slide_stats
  を見て人間判断。
- PRIMARY OUTPUT は comments (人間レビュアー風の助言文)。
- observations は raw measurement。
- 面積は EMU 単位 (1 EMU = 1/914400 inch)。slide master 背景画像は含まない。
"""

from __future__ import annotations

import pathlib


def _iter_all_shapes(shape_collection):
    """slide.shapes を group 対応で再帰的に yield する helper。

    presentation/tools.py の _walk_shapes が取り込めればそれを使う。
    取り込めない場合はフラット走査にフォールバック。
    """
    try:
        from radia_mcp.presentation.tools import _walk_shapes
    except Exception:
        _walk_shapes = None  # type: ignore

    if _walk_shapes is not None:
        try:
            yield from _walk_shapes(shape_collection)
            return
        except Exception:
            pass
    for s in shape_collection:
        yield s


def _shape_area_emu(shape) -> int:
    """(width * height) in EMU^2。None は 0 とみなす。"""
    w = getattr(shape, "width", 0) or 0
    h = getattr(shape, "height", 0) or 0
    try:
        return int(w) * int(h)
    except Exception:
        return 0


def _shape_text_len(shape) -> int:
    """text_frame の visible text 長 (改行込み)。"""
    try:
        if not getattr(shape, "has_text_frame", False):
            return 0
        txt = shape.text_frame.text or ""
        return len(txt)
    except Exception:
        return 0


def presentation_visual_text_ratio_score(pptx_path: str) -> dict:
    """per-slide visual/text ratio の distribution を score 化。

    Reynolds Zen: 画像面積 >= text 面積が理想。text-crowded slides 上位 3 件
    を特定する。

    引数:
        pptx_path: .pptx のパス

    返り値:
        Plan B 標準 dict: score / score_max / comments / observations
        / hint / source
    """
    hint = (
        "面積は EMU 単位、char_count は visible text 長。"
        "スライド master 背景画像は含まれない。"
    )
    source = "宮野 S3 (KISS) / Reynolds '3 ideas / Presentation Zen'"

    try:
        import pptx as _pptx  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except ImportError:
        return {
            "score": 10.0,
            "score_max": 10,
            "comments": [
                "python-pptx が未導入のため解析不能。`pip install python-pptx` 後に再実行。",
                "score は既定値 (10) を返す。判定には使わない。",
            ],
            "observations": {
                "total_slides": 0,
                "per_slide_stats": [],
                "crowded_slides": [],
                "median_visual_ratio": 0.0,
                "text_heavy_count": 0,
                "image_heavy_count": 0,
            },
            "hint": hint,
            "source": source,
        }

    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {
            "score": 10.0,
            "score_max": 10,
            "comments": [
                f"file not found: {pptx_path}",
                "score は既定値 (10) を返す。判定には使わない。",
            ],
            "observations": {
                "total_slides": 0,
                "per_slide_stats": [],
                "crowded_slides": [],
                "median_visual_ratio": 0.0,
                "text_heavy_count": 0,
                "image_heavy_count": 0,
            },
            "hint": hint,
            "source": source,
        }

    prs = _pptx.Presentation(str(p))
    total_slides = len(prs.slides)

    per_slide_stats: list[dict] = []

    for idx, slide in enumerate(prs.slides):
        image_area_emu = 0
        text_area_emu = 0
        char_count = 0

        for shape in _iter_all_shapes(slide.shapes):
            # picture?
            try:
                is_picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            except Exception:
                is_picture = False
            if is_picture:
                image_area_emu += _shape_area_emu(shape)
                continue

            # text shape?
            tlen = _shape_text_len(shape)
            if tlen > 0:
                text_area_emu += _shape_area_emu(shape)
                char_count += tlen

        denom = image_area_emu + text_area_emu + 1
        visual_ratio = image_area_emu / denom
        per_slide_stats.append({
            "idx": idx,
            "char_count": char_count,
            "image_area_emu": image_area_emu,
            "text_area_emu": text_area_emu,
            "visual_ratio": round(visual_ratio, 4),
        })

    # ---- aggregates ----
    ratios = [s["visual_ratio"] for s in per_slide_stats]
    if ratios:
        sorted_ratios = sorted(ratios)
        n = len(sorted_ratios)
        if n % 2 == 1:
            median_visual_ratio = float(sorted_ratios[n // 2])
        else:
            median_visual_ratio = float(
                (sorted_ratios[n // 2 - 1] + sorted_ratios[n // 2]) / 2
            )
    else:
        median_visual_ratio = 0.0

    median_visual_ratio = round(median_visual_ratio, 4)

    text_heavy_count = sum(1 for r in ratios if r < 0.3)
    image_heavy_count = sum(1 for r in ratios if r > 0.7)

    # crowded_slides: lowest visual_ratio の上位 3 件
    crowded_slides = sorted(
        per_slide_stats, key=lambda s: s["visual_ratio"]
    )[:3]

    observations = {
        "total_slides": total_slides,
        "per_slide_stats": per_slide_stats,
        "crowded_slides": crowded_slides,
        "median_visual_ratio": median_visual_ratio,
        "text_heavy_count": text_heavy_count,
        "image_heavy_count": image_heavy_count,
    }

    # ---- score ----
    if total_slides == 0:
        score: float = 10.0
    else:
        score = round(median_visual_ratio * 10, 1)
        # text_heavy の割合が過半なら penalty -3
        if text_heavy_count / total_slides > 0.5:
            score = max(0.0, round(score - 3.0, 1))

    # ---- comments ----
    comments: list[str] = []

    if total_slides == 0:
        comments.append("空デッキ (slide 数 0)。宮野 S3 / Reynolds Zen の評価対象なし。")
        comments.append("score は既定値 10 を返す。判定には使わない。")
    else:
        if text_heavy_count > 0:
            text_heavy_indices = [
                s["idx"] for s in per_slide_stats if s["visual_ratio"] < 0.3
            ]
            # 代表例: 最も crowded な slide
            worst = crowded_slides[0] if crowded_slides else None
            if worst is not None:
                comments.append(
                    f"text-crowded slides {text_heavy_count} 件 "
                    f"(visual_ratio < 0.3): indices {text_heavy_indices}。"
                    "Reynolds Zen 原則から: 1 slide = 1 visual + 数語 key "
                    f"message に近づける。例: slide {worst['idx']} は "
                    f"{worst['char_count']} 字 — bullets を図に変換を検討。"
                )
            else:
                comments.append(
                    f"text-crowded slides {text_heavy_count} 件 "
                    f"(visual_ratio < 0.3): indices {text_heavy_indices}。"
                    "Reynolds Zen 原則から: 1 slide = 1 visual + 数語 key "
                    "message に近づける。"
                )

        if image_heavy_count > total_slides * 0.5:
            comments.append(
                f"image-heavy すぎる可能性 ({image_heavy_count}/{total_slides})。"
                "Research presentation は数値・方程式の text weight を要する "
                "— image で置換できない slide も残す。"
            )

        # balanced: crowded が 0 かつ image-heavy も過半でない
        if text_heavy_count == 0 and image_heavy_count <= total_slides * 0.5:
            comments.append(
                f"text/image ratio は median {median_visual_ratio}、"
                "crowded slide なし。Zen 原則に近い。"
            )

        # 過半数が text-heavy の penalty 警告
        if text_heavy_count > 0 and text_heavy_count / total_slides > 0.5:
            comments.append(
                f"過半数 ({text_heavy_count}/{total_slides}) が text-heavy。"
                "penalty -3 を score に適用。章ごとに 1 slide だけでも図解に "
                "置き換えるとインパクトが大きい (宮野 S3 KISS)。"
            )

        # median の評価
        if (
            text_heavy_count == 0
            and image_heavy_count == 0
            and total_slides > 0
        ):
            comments.append(
                f"median visual_ratio = {median_visual_ratio}。"
                "全 slide が中庸 (0.3 <= ratio <= 0.7) で、text と画像の "
                "balance は取れているが、key slide では画像支配を検討。"
            )

    # 2-5 個に制限
    comments = comments[:5]
    while len(comments) < 2:
        comments.append(
            "score は粗い温度計。最終判断は per_slide_stats を確認。"
        )

    return {
        "score": float(score),
        "score_max": 10,
        "comments": comments,
        "observations": observations,
        "hint": hint,
        "source": source,
    }
