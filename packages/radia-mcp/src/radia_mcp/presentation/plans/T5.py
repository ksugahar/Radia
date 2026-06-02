"""Plan B T5: 進行位置 indicator 診断 (presentation_check_progress_indicator).

skill.md S15 (進行位置の明示) / 宮野『研究発表のためのスライドデザイン』に
基づく粗い温度計。section header / outline slide が無いまま content slide が
長く続く区間を検出し、章間 cue を挿入すべきかを助言する。

設計方針:
- score は 0-10 の粗い温度計。9 以上は誤差範囲、6 以下は改善サイン。
- PRIMARY OUTPUT は comments (人間レビュアー風の助言文)。
- observations は raw measurement。numeric な正誤判定には使わない。
- 検出 heuristic は 2 段: title keyword match + slide layout 名 match。
  タイトルなしの純粋な視覚的 divider slide は現状の heuristic では拾えない。
"""

from __future__ import annotations

import pathlib
import re


# ------------------------------------------------------------------
# Section header / outline 判定パターン
# ------------------------------------------------------------------
_TITLE_OUTLINE_PATTERNS = [
    r"目次",
    r"Outline",
    r"Agenda",
    r"Contents",
    r"Section",
    r"章",
    r"Part",
    r"Overview",
    r"Roadmap",
    r"進行",
    r"概要",
]

_TITLE_OUTLINE_RE = re.compile(
    "|".join(_TITLE_OUTLINE_PATTERNS), re.IGNORECASE
)

_LAYOUT_SECTION_MARKERS = ("section", "header", "divider")


# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------
def _slide_title_text(slide) -> str:
    """Return a best-effort title string for the slide.

    Prefers the title placeholder (idx == 0); falls back to the first
    non-empty text frame in shape order. Returns "" if nothing found.
    """
    # 1) Title placeholder
    try:
        for shape in slide.shapes:
            try:
                pf = shape.placeholder_format
            except Exception:
                pf = None
            if pf is None:
                continue
            try:
                if pf.idx == 0 and shape.has_text_frame:
                    txt = shape.text_frame.text or ""
                    if txt.strip():
                        return txt.splitlines()[0]
            except Exception:
                continue
    except Exception:
        pass
    # 2) Fallback: first non-empty text frame in shape order
    try:
        for shape in slide.shapes:
            try:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    return shape.text_frame.text.splitlines()[0]
            except Exception:
                continue
    except Exception:
        pass
    return ""


def _body_bullet_count(slide) -> int:
    """Rough count of body (non-title) bullet-like lines on the slide."""
    bullets = 0
    try:
        for shape in slide.shapes:
            try:
                pf = shape.placeholder_format
                is_title = (pf is not None and pf.idx == 0)
            except Exception:
                is_title = False
            if is_title:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            try:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        bullets += 1
            except Exception:
                continue
    except Exception:
        pass
    return bullets


def _is_section_header(slide) -> tuple[bool, str, str]:
    """Decide whether the slide is a section header / outline.

    Returns (is_section, method, title) where:
      - method in {"title_match", "body_list", "layout_match", ""}
      - title is the slide title (may be empty)
    """
    title = _slide_title_text(slide)

    # 1) Title keyword match (most reliable)
    if title and _TITLE_OUTLINE_RE.search(title):
        return True, "title_match", title

    # 2) Layout name match
    layout_name = ""
    try:
        layout_name = (slide.slide_layout.name or "").lower()
    except Exception:
        layout_name = ""
    if any(m in layout_name for m in _LAYOUT_SECTION_MARKERS):
        return True, "layout_match", title

    # 3) Body is mostly a bullet list AND title contains 概要 / 章
    if title and re.search(r"概要|章", title):
        if _body_bullet_count(slide) >= 3:
            return True, "title_match", title

    return False, "", title


def _longest_run(total: int, header_indices: set[int]) -> tuple[int, int, int]:
    """Return (longest_run_len, run_start, run_end) of consecutive slide
    indices NOT in header_indices.

    start/end are 0-indexed inclusive. If no run (total == 0), return
    (0, -1, -1).
    """
    best_len = 0
    best_start = -1
    best_end = -1
    cur_start = -1
    cur_len = 0
    for i in range(total):
        if i in header_indices:
            cur_start = -1
            cur_len = 0
            continue
        if cur_len == 0:
            cur_start = i
        cur_len += 1
        if cur_len > best_len:
            best_len = cur_len
            best_start = cur_start
            best_end = i
    return best_len, best_start, best_end


# ------------------------------------------------------------------
# Main
# ------------------------------------------------------------------
def presentation_check_progress_indicator(
    pptx_path: str,
    min_slides_without_cue: int = 10,
) -> dict:
    """outline / section-header slides for progress indication を検出。

    content slides が長く続く (>= min_slides_without_cue 枚) 区間を flag。
    宮野 S15 基準。

    引数:
        pptx_path: .pptx のパス
        min_slides_without_cue: 連続 content slide 許容上限 (既定 10)

    返り値:
        dict: score / score_max / comments / observations / hint / source
    """
    try:
        import pptx  # type: ignore
    except ImportError:
        return {"error": "python-pptx not installed. `pip install python-pptx`"}

    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}

    prs = pptx.Presentation(str(p))
    total_slides = len(prs.slides)

    section_header_indices: list[int] = []
    section_header_titles: list[str] = []
    methods_fired: set[str] = set()

    for i, slide in enumerate(prs.slides):
        is_section, method, title = _is_section_header(slide)
        if is_section:
            section_header_indices.append(i)
            section_header_titles.append(title)
            if method:
                methods_fired.add(method)

    header_set = set(section_header_indices)
    longest_run_without_cue, run_start, run_end = _longest_run(
        total_slides, header_set
    )
    is_long_run_present = longest_run_without_cue >= min_slides_without_cue

    # section_header_method summary
    if len(methods_fired) == 0:
        section_header_method = ""
    elif len(methods_fired) == 1:
        section_header_method = next(iter(methods_fired))
    else:
        section_header_method = "mixed"

    observations = {
        "total_slides": total_slides,
        "section_header_indices": section_header_indices,
        "section_header_titles": section_header_titles,
        "longest_run_without_cue": longest_run_without_cue,
        "is_long_run_present": bool(is_long_run_present),
        "section_header_method": section_header_method,
    }

    # ---- score ----
    if total_slides == 0:
        score: float = 10.0
    elif is_long_run_present and longest_run_without_cue > 0:
        raw = (min_slides_without_cue / longest_run_without_cue) * 10 - 1
        score = max(0.0, round(raw, 1))
    else:
        score = 10.0

    # ---- comments ----
    comments: list[str] = []

    if total_slides == 0:
        comments.append("empty deck — スライドが 1 枚も無い。")
        comments.append("最低 1 枚は冒頭に目次 (Outline) を置き、進行位置の cue を用意する (宮野 S15)。")
    else:
        if is_long_run_present:
            # 1-indexed for comment readability
            start_1 = run_start + 1
            end_1 = run_end + 1
            comments.append(
                f"スライド {start_1}-{end_1} ({longest_run_without_cue} 枚) に "
                "section header / outline slide が無い。audience は進行位置を失う。"
                "章間に 1 枚「次は … について」の切れ目を挿入 (宮野 S15)。"
            )

        if total_slides >= 15 and len(section_header_indices) == 0:
            comments.append(
                "進行位置の cue (Outline/Agenda/Section) が全く無い。"
                "最低 1 枚は冒頭に目次を置き、章ごとに現在位置を色で示す (宮野 S15)。"
            )

        if section_header_indices:
            # Human-readable 1-indexed position list
            idx_1 = [i + 1 for i in section_header_indices]
            comments.append(
                f"section header {len(section_header_indices)} 枚検出 "
                f"(slides: {idx_1})。進行位置 cue が機能している。"
            )

        if (not is_long_run_present) and total_slides >= min_slides_without_cue:
            comments.append(
                f"最大連続 content slide は {longest_run_without_cue} 枚で "
                f"閾値 {min_slides_without_cue} 枚を下回る。"
                "章ごとの切れ目は概ね配置されている。"
            )

        if section_header_method == "layout_match":
            comments.append(
                "section header の検出は layout 名のみに依存している。"
                "タイトルに「目次」「Outline」等の明示キーワードを載せると"
                "読み落とされにくい。"
            )

    # フォールバック: 最低 2 コメント
    if len(comments) < 2:
        if total_slides > 0 and total_slides < min_slides_without_cue:
            comments.append(
                f"スライド総数 {total_slides} 枚は閾値 {min_slides_without_cue} 枚"
                "未満。進行位置 cue の必要性は低いが、章構造が生まれた段階で"
                "目次 slide の導入を検討 (宮野 S15)。"
            )
        if len(comments) < 2:
            comments.append(
                "score は粗い指標。skill.md S15 の判定基準で最終確認を推奨。"
            )

    # 最大 5 個
    comments = comments[:5]

    hint = (
        "title keyword + layout 名の 2 段検出。"
        "タイトルなし slide (視覚的 section divider) は現状の heuristic では拾わない。"
    )

    source = "宮野『研究発表のためのスライドデザイン』S15 (進行位置を明示)"

    return {
        "score": float(score),
        "score_max": 10,
        "comments": comments,
        "observations": observations,
        "hint": hint,
        "source": source,
    }
