"""Plan B T15: Single message per slide (semantic)
(presentation_single_message_per_slide_semantic).

既存 `check_bullet_count_per_slide` は bullet 数のみで「1 主張」を判定するが、
**意味的に複数テーマ** (1 figure + 1 段落 + 1 bullet 群 = 3 主張) は
拾えない。本ツールは **shape 種別の多様性** で 1 主張違反を診断する。

検出ロジック (per slide):
  - text-frame 群 (≠ title)
  - image / picture
  - chart / graph
  - table
  - SmartArt / diagram

これらが **3 種以上 並存** + それぞれが意味的単位 (text 数 ≥ 30 字 / 図表
1 個以上) なら「複数テーマ slide」と判定。

設計方針:
- score は 「1 主張 slide 比率」(0-10)。
- 高密度 dashboard 風 slide を flag。
- PRIMARY OUTPUT は multi_theme_slides list (修正候補)。
"""

from __future__ import annotations

import pathlib


def _iter_all_shapes(shape_collection):
    try:
        from radia_mcp.presentation.tools import _walk_shapes
        yield from _walk_shapes(shape_collection)
    except Exception:
        for s in shape_collection:
            yield s


def _classify_shape(shape) -> str:
    """Return one of: title / text / image / chart / table / smartart / other."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        MSO_SHAPE_TYPE = None  # type: ignore

    # Title placeholder?
    try:
        pf = shape.placeholder_format
        if pf is not None and pf.idx == 0:
            return "title"
    except Exception:
        pass

    # Picture
    try:
        if MSO_SHAPE_TYPE is not None and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "image"
    except Exception:
        pass

    # Chart
    if getattr(shape, "has_chart", False):
        try:
            if shape.has_chart:
                return "chart"
        except Exception:
            pass

    # Table
    if getattr(shape, "has_table", False):
        try:
            if shape.has_table:
                return "table"
        except Exception:
            pass

    # Text
    if getattr(shape, "has_text_frame", False):
        try:
            if shape.has_text_frame and shape.text_frame.text.strip():
                return "text"
        except Exception:
            pass

    return "other"


def _text_unit_size(shape) -> int:
    try:
        if shape.has_text_frame:
            return len(shape.text_frame.text or "")
    except Exception:
        pass
    return 0


def presentation_single_message_per_slide_semantic(
    pptx_path: str, theme_threshold: int = 3
) -> dict:
    """1 主張 vs 複数テーマ slide を意味単位で診断。

    bullet 数だけでなく shape 種別の多様性で多テーマを検出。

    Args:
        pptx_path: .pptx file
        theme_threshold: 何種類以上を「多テーマ」とするか (default 3)

    Returns:
        dict: score / per_slide_classification / multi_theme_slides / comments
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))

    per_slide: list[dict] = []
    multi_theme_slides: list[dict] = []

    for i, slide in enumerate(prs.slides, 1):
        kind_counts = {
            "text_units": 0,  # text frames with >= 30 chars (excl title)
            "image": 0,
            "chart": 0,
            "table": 0,
            "smartart_or_other": 0,
        }
        for shape in _iter_all_shapes(slide.shapes):
            kind = _classify_shape(shape)
            if kind == "title":
                continue
            elif kind == "text":
                if _text_unit_size(shape) >= 30:
                    kind_counts["text_units"] += 1
            elif kind == "image":
                kind_counts["image"] += 1
            elif kind == "chart":
                kind_counts["chart"] += 1
            elif kind == "table":
                kind_counts["table"] += 1
            elif kind == "other":
                # SmartArt / freeforms — count as "other" theme
                kind_counts["smartart_or_other"] += 1

        # 意味的単位の数 (1 unit = 1 theme proxy)
        theme_units = (
            kind_counts["text_units"]
            + kind_counts["image"]
            + kind_counts["chart"]
            + kind_counts["table"]
        )
        # 種類の多様性 (3 種以上 = 多テーマの強い signal)
        n_kinds_used = sum(
            1 for k, v in kind_counts.items()
            if v >= 1 and k != "smartart_or_other"
        )

        is_multi_theme = (
            theme_units >= theme_threshold or n_kinds_used >= 3
        )

        per_slide.append({
            "slide_no": i,
            "theme_units": theme_units,
            "n_kinds_used": n_kinds_used,
            "kinds": {k: v for k, v in kind_counts.items() if v >= 1},
            "is_multi_theme": is_multi_theme,
        })

        if is_multi_theme:
            multi_theme_slides.append({
                "slide_no": i,
                "theme_units": theme_units,
                "n_kinds_used": n_kinds_used,
                "kinds": {k: v for k, v in kind_counts.items() if v >= 1},
                "advice": (
                    "複数テーマが並存。1 主張 slide に分割推奨。"
                    f"({n_kinds_used} 種 × {theme_units} 単位)"
                ),
            })

    n_slides = len(per_slide)
    n_multi = len(multi_theme_slides)
    n_single_or_low = n_slides - n_multi
    single_ratio = n_single_or_low / n_slides if n_slides else 0
    score = round(single_ratio * 10, 1)

    # ---- comments ----
    comments: list[str] = []

    comments.append(
        f"スライド {n_slides} 枚 | multi-theme {n_multi} 枚 "
        f"| single-message ratio {single_ratio*100:.0f}% (score {score}/10)"
    )

    if n_multi >= 1:
        comments.append(
            f"⚠ {n_multi} 枚で複数テーマ並存。"
            "宮野『研究発表のためのスライドデザイン』S4 の『1 主張 1 slide』"
            "原則違反。1 slide を 2-3 枚に分割を検討。"
        )
        for m in multi_theme_slides[:5]:
            comments.append(
                f"  • slide {m['slide_no']}: {m['kinds']} → {m['advice']}"
            )

    if n_multi == 0:
        comments.append(
            "全 slide で 1 主張 OK。理系プレゼンの core 原則充足。"
        )

    if n_multi >= n_slides * 0.3:
        comments.append(
            f"⚠ Multi-theme 率 {n_multi/n_slides*100:.0f}% は高い。"
            "全体的に 1 slide 詰め込みすぎ傾向。"
            "「目次として並べたら筋が通るか」(T13) 視点で再構成推奨。"
        )

    comments = comments[:7]

    hint = (
        "shape 種別 (text/image/chart/table/smartart) の多様性で"
        "1 主張違反を検出。"
        "既存 check_bullet_count_per_slide が拾えない『figure + paragraph + bullet』"
        "型の dashboard slide を flag。"
        "理系プレゼンでは figure 1 + 補足 1-2 行が標準で、それを超えると "
        "聴衆の認知負荷が急上昇。"
    )

    return {
        "score": score,
        "score_max": 10,
        "n_slides": n_slides,
        "n_multi_theme": n_multi,
        "single_message_ratio": round(single_ratio, 3),
        "multi_theme_slides": multi_theme_slides[:15],
        "per_slide": per_slide[:30],
        "comments": comments,
        "hint": hint,
        "source": (
            "宮野『研究発表のためのスライドデザイン』S4 (1 主張 1 slide) + "
            "高橋メソッド + Tufte『Cognitive Style of PowerPoint』"
            "(dashboard 過密批判)。既存 check_bullet_count_per_slide と相補。"
            "(presentation 2026-04 採用)。"
        ),
    }
