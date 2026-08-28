"""Presentation T28: Speaking pace estimate (verbal side 強化)
(presentation_speaking_pace_estimate).

既存 estimate_speaking_time は slide 数 × 分/slide の rough 推定。本ツールは
**speaker_note の文字数から WPM (words/characters per minute)** で正確に推定、
理系日本語発表の適正 range (330-390 字/分 = 中庸) と比較。

理系プレゼン適正 pace:
- 遅い (≤300 字/分): 聴衆 退屈リスク
- 中庸 (330-390 字/分): 標準、日本語学会発表の経験則
- 速い (≥420 字/分): 聴衆 理解困難リスク
- (英語の場合 130-160 WPM が中庸、本 tool は日本語前提 character count)
"""

from __future__ import annotations

import pathlib
import re


_SOURCES_BLOCK_RE = re.compile(r"(?im)^\s*\[sources\]\s*$")


def _spoken_note_text(note_text: str) -> tuple[str, bool]:
    """Return only the part of speaker notes intended to be spoken."""
    match = _SOURCES_BLOCK_RE.search(note_text)
    if not match:
        return note_text, False
    return note_text[:match.start()].rstrip(), True


def _iter_shapes(shape_collection):
    try:
        from radia_mcp.presentation.tools import _walk_shapes
        yield from _walk_shapes(shape_collection)
    except Exception:
        for s in shape_collection:
            yield s


def presentation_speaking_pace_estimate(
    pptx_path: str,
    target_duration_min: int = 0,
    pace_target_min: int = 330,
    pace_target_max: int = 390,
    include_backup: bool = False,
) -> dict:
    """speaker_note の文字数から WPM (日本語は文字数/分) で発表時間を推定。

    Args:
        pptx_path: .pptx file
        target_duration_min: 予定発表時間 (分)、0 なら target との比較 skip
        pace_target_min / pace_target_max: 適正 pace range (文字/分)、
            default 330-390 は日本語理系学会発表の経験則
        include_backup: True のときだけ補足・質疑用スライドを発表時間へ
            算入する。通常の本編時間監査では False。

    Returns:
        dict: per_slide / total_chars / required_pace /
              estimated_duration / recommendation / comments
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))

    from .._deck_sections import classify_deck_sections
    from ..tools import _slide_title

    slides = list(prs.slides)
    sections = classify_deck_sections(slides, _slide_title)
    backup_slide_numbers = set(sections["backup_slide_numbers"])

    per_slide: list[dict] = []
    total_chars = 0
    total_english_words = 0
    source_blocks_excluded = 0
    for i, slide in enumerate(slides, 1):
        excluded_backup = i in backup_slide_numbers and not include_backup
        if excluded_backup:
            per_slide.append({
                "slide_no": i,
                "note_chars": 0,
                "note_english_words": 0,
                "note_empty": False,
                "excluded_reason": "backup_slide",
            })
            continue
        note_text = ""
        try:
            if slide.has_notes_slide:
                note_text = slide.notes_slide.notes_text_frame.text or ""
        except Exception:
            pass
        note_text, source_block_excluded = _spoken_note_text(note_text)
        source_blocks_excluded += int(source_block_excluded)
        chars = len(note_text)
        # Count English words (space-delimited) as fallback if 日本語 chars 少ない
        eng_words = len(re.findall(r"\b[A-Za-z]+\b", note_text))
        total_chars += chars
        total_english_words += eng_words
        per_slide.append({
            "slide_no": i,
            "note_chars": chars,
            "note_english_words": eng_words,
            "note_empty": chars == 0,
            "source_block_excluded": source_block_excluded,
        })

    timed_slides = [
        s for s in per_slide if s.get("excluded_reason") != "backup_slide"
    ]
    n_slides = len(timed_slides)
    n_with_notes = sum(1 for s in timed_slides if not s["note_empty"])
    n_empty_notes = n_slides - n_with_notes

    # Pace estimation
    # 日本語想定: 330-390 chars/min が中庸、target=360 (median)
    # 英語併用なら 1 word ≈ 5 chars として fold back
    effective_chars = total_chars + total_english_words * 5
    median_pace = (pace_target_min + pace_target_max) / 2
    est_minutes_at_median = effective_chars / median_pace
    est_minutes_slow = effective_chars / pace_target_min  # slow pace → longer
    est_minutes_fast = effective_chars / pace_target_max  # fast pace → shorter

    # Target comparison
    recommendation = ""
    delta = 0.0
    if target_duration_min > 0:
        delta = est_minutes_at_median - target_duration_min
        if abs(delta) <= 1.0:
            recommendation = f"target {target_duration_min} 分に概ね適合 (差 {delta:+.1f} 分)"
        elif delta > 0:
            required_pace = effective_chars / target_duration_min
            recommendation = (
                f"target {target_duration_min} 分に対し {delta:+.1f} 分 超過。"
                f"注を削るか、pace を {required_pace:.0f} 字/分 に上げる "
                f"(現 median {median_pace:.0f})。"
            )
        else:
            required_pace = effective_chars / target_duration_min
            recommendation = (
                f"target {target_duration_min} 分に対し {-delta:.1f} 分 不足。"
                f"追加情報 or backup slide 追加、または pace を "
                f"{required_pace:.0f} 字/分 に下げる。"
            )

    # ---- comments ----
    comments: list[str] = []

    if n_empty_notes > n_slides * 0.3:
        comments.append(
            f"⚠ {n_empty_notes}/{n_slides} slides で speaker_note 未記入。"
            "音読時間の推定精度が落ちる。全 slide に note を入れると "
            "time estimate が正確になる + 発表リハーサル台本になる。"
        )

    comments.append(
        f"Speaker note 合計: {total_chars} 日本語文字 + "
        f"{total_english_words} 英単語 = {effective_chars} 実質文字相当"
    )
    comments.append(
        f"推定発表時間 (pace {pace_target_min}-{pace_target_max} 字/分):\n"
        f"  - slow ({pace_target_min}): {est_minutes_slow:.1f} 分\n"
        f"  - median ({median_pace:.0f}): {est_minutes_at_median:.1f} 分\n"
        f"  - fast ({pace_target_max}): {est_minutes_fast:.1f} 分"
    )

    if recommendation:
        comments.append(f"📌 Target 比較: {recommendation}")

    if effective_chars < 500 and n_with_notes >= 3:
        comments.append(
            "ℹ speaker_note 総量が少ない (<500 chars)。"
            "15+ 分発表なら台本が薄い可能性。"
        )

    # Score
    score = 10.0
    if n_empty_notes > n_slides * 0.3:
        score -= 3
    if target_duration_min > 0 and abs(delta) > 3:
        score -= 3
    if effective_chars < 500 and n_with_notes >= 3:
        score -= 2
    score = max(0.0, score)

    return {
        "score": round(score, 1),
        "score_max": 10,
        "n_slides": n_slides,
        "total_deck_slides": len(slides),
        "include_backup": include_backup,
        "backup_slide_count": sections["backup_slide_count"],
        "backup_slide_numbers": sections["backup_slide_numbers"],
        "n_with_notes": n_with_notes,
        "n_empty_notes": n_empty_notes,
        "total_chars_in_notes": total_chars,
        "total_english_words_in_notes": total_english_words,
        "effective_chars": effective_chars,
        "source_blocks_excluded": source_blocks_excluded,
        "pace_range_target": [pace_target_min, pace_target_max],
        "estimated_minutes": {
            "slow": round(est_minutes_slow, 1),
            "median": round(est_minutes_at_median, 1),
            "fast": round(est_minutes_fast, 1),
        },
        "target_duration_min": target_duration_min if target_duration_min else None,
        "delta_to_target_min": round(delta, 1) if target_duration_min else None,
        "recommendation": recommendation,
        "per_slide": per_slide[:30],
        "comments": comments,
        "hint": (
            "日本語理系発表の適正 pace 330-390 字/分 (median 360)。"
            "英語なら 130-160 WPM (fold back 済)。"
            "既存 estimate_speaking_time より精密で target 比較付き。"
        ),
        "source": (
            "日本音声学会 発話速度研究 (日本語平均 350 字/分) + "
            "英語 TED talks 統計 (平均 150 WPM) の rough adapt。"
            "理系学会発表 10-20 分向け。"
            "(presentation 2026-04 採用、v0.26.0)。"
        ),
    }
