"""Plan B T18: 理系プレゼン minimalism score
(presentation_rikei_minimalism_score).

Reynolds『Presentation Zen』の minimalism は corporate / TED 寄りで
「1 visual + ≤7 word」が標準。理系プレゼンでは少し違って:

  - **figure 1 + 数値 + 出典** が core (visual だけでは domain peer に薄い)
  - **caption / 軸ラベル / 単位** が必須 (corporate プレゼンでは省略可)
  - **数値 1+ per slide** (理系 evidence-based culture)
  - **bullet ≤ 4 per slide** (高橋メソッド より緩和: 理系は補足必要)
  - **font ≥ 24pt (本文・注釈・図表 label) / 32pt (title)**

設計方針:
- 理系 minimalism の 5 軸を per slide で評価
- score は 5 軸 coverage 率 (0-10)
- non-compliant slides を flag
"""

from __future__ import annotations

import pathlib
import re


def _iter_all_shapes(shape_collection):
    try:
        from radia_mcp.presentation.tools import _walk_shapes
        yield from _walk_shapes(shape_collection)
    except Exception:
        for s in shape_collection:
            yield s


def _resolve_size(run, para, shape, slide):
    try:
        from radia_mcp.presentation.tools import _resolve_font_size
        return _resolve_font_size(run, para, shape, slide)
    except Exception:
        try:
            return int(run.font.size.pt) if run.font.size else None
        except Exception:
            return None


def _get_title(slide) -> str:
    try:
        from radia_mcp.presentation.tools import _slide_title
        return _slide_title(slide)
    except Exception:
        return ""


def _has_figure_or_chart(slide) -> bool:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        MSO_SHAPE_TYPE = None
    for shape in _iter_all_shapes(slide.shapes):
        try:
            if MSO_SHAPE_TYPE is not None and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True
        except Exception:
            pass
        if getattr(shape, "has_chart", False):
            try:
                if shape.has_chart:
                    return True
            except Exception:
                pass
    return False


def _has_number_or_unit(text: str) -> bool:
    return bool(re.search(
        r"\d+(?:\.\d+)?\s*(?:%|％|x|×|倍|kHz|MHz|GHz|Hz|GB|MB|nm|μm|µm|mm|cm|m|s|sec|V|A|W)",
        text,
    ))


def _has_citation_marker(text: str) -> bool:
    return bool(re.search(
        r"\[\d+\]|\(\w+\s*(?:\d{4}|et\s+al)|出典|reference|cite",
        text, re.IGNORECASE,
    ))


def _has_axis_label_or_unit(slide) -> bool:
    """簡易: figure を持つ slide に軸ラベル / 単位 を示唆する文字列があるか。"""
    if not _has_figure_or_chart(slide):
        return True  # not applicable, treat as compliant
    full_text = ""
    for shape in _iter_all_shapes(slide.shapes):
        try:
            if shape.has_text_frame:
                full_text += " " + shape.text_frame.text
        except Exception:
            continue
    return bool(re.search(
        r"(?:横軸|縦軸|x軸|y軸|x-axis|y-axis|axis|unit|単位|"
        r"\[.*?(?:Hz|m|s|V|A|W|%|kg|mol|K)\]|"
        r"\([A-Za-z]{1,3}\)|\(.*?(?:Hz|m|s|V|A|W|%|kg|mol|K)\))",
        full_text, re.IGNORECASE,
    ))


def _bullet_count(slide) -> int:
    n = 0
    for shape in _iter_all_shapes(slide.shapes):
        try:
            pf = shape.placeholder_format
            if pf is not None and pf.idx == 0:
                continue
        except Exception:
            pass
        try:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        n += 1
        except Exception:
            continue
    return n


def _check_font_size(slide, min_body_pt=24, min_title_pt=32) -> bool:
    """全 run が下限を充たすか (理系学会標準)。"""
    try:
        from radia_mcp.presentation.tools import (
            _is_footer_or_page_chrome,
            _shape_is_slide_title,
        )
    except Exception:
        _is_footer_or_page_chrome = None
        _shape_is_slide_title = None
    for shape in _iter_all_shapes(slide.shapes):
        try:
            if not shape.has_text_frame:
                continue
            if (_is_footer_or_page_chrome is not None
                    and _is_footer_or_page_chrome(shape, slide)):
                continue
            if _shape_is_slide_title is not None:
                is_title = _shape_is_slide_title(shape, slide)
            else:
                is_title = False
                try:
                    pf = shape.placeholder_format
                    if pf is not None and pf.idx == 0:
                        is_title = True
                except Exception:
                    pass
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    sz = _resolve_size(run, para, shape, slide)
                    if sz is None:
                        continue
                    limit = min_title_pt if is_title else min_body_pt
                    if sz < limit:
                        return False
        except Exception:
            continue
    return True


def presentation_rikei_minimalism_score(pptx_path: str) -> dict:
    """理系プレゼンの minimalism 5 軸を per slide 診断。

    理系特化 (TED / corporate とは違う):
      1. figure_or_chart_present
      2. axis_label_unit_present (figure 持つ slide のみ評価)
      3. has_quantitative_evidence (数値 / 単位)
      4. has_source_or_citation (出典)
      5. bullet_within_4 + font_size_ok

    Args:
        pptx_path: .pptx file

    Returns:
        dict: score / per_slide / non_compliant_slides / comments
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))

    per_slide: list[dict] = []
    non_compliant: list[dict] = []

    for i, slide in enumerate(prs.slides, 1):
        all_text = ""
        for shape in _iter_all_shapes(slide.shapes):
            try:
                if shape.has_text_frame:
                    all_text += " " + (shape.text_frame.text or "")
            except Exception:
                continue

        figure_present = _has_figure_or_chart(slide)
        axis_ok = _has_axis_label_or_unit(slide)
        has_number = _has_number_or_unit(all_text)
        has_citation = _has_citation_marker(all_text)
        n_bullets = _bullet_count(slide)
        bullet_ok = n_bullets <= 4
        font_ok = _check_font_size(slide)

        # 5 軸 (figure と axis は figure-bearing slide のみ)
        if figure_present:
            n_pass = sum([
                figure_present, axis_ok, has_number,
                has_citation, bullet_ok and font_ok,
            ])
            n_total_axes = 5
        else:
            # text-only slide: 軸 / figure を除外して 3 軸評価
            n_pass = sum([
                has_number,
                has_citation,
                bullet_ok and font_ok,
            ])
            n_total_axes = 3

        slide_score = round(n_pass / n_total_axes * 10, 1)

        per_slide.append({
            "slide_no": i,
            "title": _get_title(slide)[:60],
            "figure_present": figure_present,
            "axis_label_ok": axis_ok,
            "has_quantitative": has_number,
            "has_citation": has_citation,
            "bullet_count": n_bullets,
            "bullet_within_4": bullet_ok,
            "font_size_ok": font_ok,
            "n_pass": n_pass,
            "n_total_axes": n_total_axes,
            "slide_score": slide_score,
        })

        if slide_score < 6:
            missing = []
            if figure_present and not axis_ok:
                missing.append("axis_label/unit")
            if not has_number:
                missing.append("quantitative_evidence")
            if not has_citation:
                missing.append("source/citation")
            if not bullet_ok:
                missing.append(f"bullet_overflow({n_bullets}>4)")
            if not font_ok:
                missing.append("font_too_small")
            non_compliant.append({
                "slide_no": i,
                "title": _get_title(slide)[:60],
                "score": slide_score,
                "missing": missing,
            })

    n_slides = len(per_slide)
    valid = [s for s in per_slide if s["slide_score"] is not None]
    overall = round(sum(s["slide_score"] for s in valid) / len(valid), 1) if valid else 0.0

    # ---- comments ----
    comments: list[str] = []

    n_non = len(non_compliant)
    comments.append(
        f"理系 minimalism score: {overall}/10 ({n_slides} slides, "
        f"non-compliant {n_non} slides)"
    )

    if n_non >= 1:
        comments.append(
            f"⚠ {n_non} 枚で理系 minimalism 不充足。"
            "理系プレゼンの evidence-based 文化として: "
            "数値 1+ / 出典明記 / figure 持つ slide は軸ラベル+単位 / "
            "bullet ≤4 / font ≥24pt が core 基準。"
        )
        for nc in non_compliant[:5]:
            comments.append(
                f"  • slide {nc['slide_no']} [{nc['score']}/10]: "
                f"{nc['title']} — missing {nc['missing']}"
            )

    n_no_evidence = sum(1 for s in per_slide if not s["has_quantitative"])
    if n_no_evidence > n_slides * 0.5:
        comments.append(
            f"⚠ {n_no_evidence} slides で数値/定量 evidence なし "
            "({n_no_evidence/n_slides*100:.0f}%)。"
            "理系プレゼンでは evidence-based slide が標準で、"
            "数値ゼロは『感想プレゼン』印象。"
        )

    n_no_citation = sum(1 for s in per_slide if not s["has_citation"])
    if n_no_citation > n_slides * 0.7:
        comments.append(
            f"ℹ {n_no_citation} slides で citation なし。"
            "introduction / related work / results 部分には "
            "出典 [1] / (Author 2023) を 1 件以上添える。"
        )

    if n_non == 0 and overall >= 8:
        comments.append(
            f"理系 minimalism {overall}/10 — 全 slide で核基準充足。"
            "evidence-based + 軸/単位 明記 + 簡潔 が両立。"
        )

    comments = comments[:7]

    hint = (
        "Reynolds Presentation Zen (TED 風 1 visual + ≤7 word) ではなく、"
        "理系プレゼン特化の minimalism (figure + 数値 + 出典 + axis/unit)。"
        "domain peer 前提なので citation / 軸ラベル / 単位 が必須要素。"
        "text-only slide では axis/figure 軸を除外して評価。"
    )

    return {
        "score": overall,
        "score_max": 10,
        "n_slides": n_slides,
        "n_non_compliant": n_non,
        "non_compliant_slides": non_compliant[:15],
        "per_slide": per_slide[:30],
        "comments": comments,
        "hint": hint,
        "source": (
            "宮野『研究発表のためのスライドデザイン』S5 (figure 中心 / "
            "evidence-based) + 高橋メソッド (大字 / 簡潔) + Reynolds "
            "『Presentation Zen』(general minimalism) の理系 adaptation。"
            "TED 用ではなく学会発表 / 研究室セミナー想定。"
            "(presentation 2026-04 採用)。"
        ),
    }
