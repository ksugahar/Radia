"""Plan B T3: 円グラフ / 3D グラフ NG 検出 (presentation_check_pie_3d_charts).

presentation/skill.md S8 (グラフの選択) / 宮野『研究発表のためのスライドデザイン』
に基づき、.pptx 内の chart shape を走査して pie / doughnut / 3D chart を
NG として検出する。

設計方針:
- score は 0-10 の粗い温度計。ok_count / total_charts 比で算出。
- PRIMARY OUTPUT は comments (人間レビュアー風の助言文)。
- observations は raw measurement。chart 種別の enum 名を並べるだけ。
"""

from __future__ import annotations

from typing import Any


# ------------------------------------------------------------------
# NG 判定: chart_type の name ベース (XL_CHART_TYPE enum に依存しない
# 実行時判定を行い、test 時に python-pptx の enum 変更にも耐える)
# ------------------------------------------------------------------
def _chart_type_name(chart_type: Any) -> str:
    """chart.chart_type から enum name (str) を取り出す。

    python-pptx の `XL_CHART_TYPE` enum は `.name` 属性を持つ。
    念のため属性がない場合は str() フォールバック。
    """
    name = getattr(chart_type, "name", None)
    if name is not None:
        return str(name)
    return str(chart_type)


def _is_ng_chart_type(type_name: str) -> bool:
    """chart_type 名 (例: "PIE", "THREE_D_BAR_CLUSTERED") が NG か判定。

    宮野 S8 の NG 条件:
    - PIE 系 (PIE, PIE_EXPLODED, PIE_OF_PIE, BAR_OF_PIE)
    - DOUGHNUT 系 (DOUGHNUT, DOUGHNUT_EXPLODED)
    - 3D 系 (THREE_D_* 全般)
    """
    if type_name.startswith("PIE"):
        return True
    if type_name.startswith("DOUGHNUT"):
        return True
    if type_name.startswith("BAR_OF_PIE"):
        return True
    if "THREE_D" in type_name:
        return True
    return False


# ------------------------------------------------------------------
# Main
# ------------------------------------------------------------------
def presentation_check_pie_3d_charts(pptx_path: str) -> dict:
    """pptx 内の chart shape を走査し、pie / doughnut / 3D chart を NG 検出。

    宮野『研究発表のためのスライドデザイン』S8 (グラフ選択) に基づく。
    円グラフは角度比較が難しく、3D グラフは奥行歪みで定量比較を阻害する
    ため、bar / horizontal bar / line / slopegraph への置換を推奨する。

    引数:
        pptx_path: 解析対象 .pptx ファイルパス

    返り値:
        dict: score / score_max / comments / observations / hint / source
    """
    # python-pptx lazy import (ない環境では例外を返す)
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:  # pragma: no cover - 環境依存
        raise ImportError(
            "python-pptx is required for presentation_check_pie_3d_charts. "
            "Install via `pip install python-pptx`."
        ) from exc

    prs = Presentation(pptx_path)

    ng_charts: list[dict] = []
    ok_charts: list[dict] = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            has_chart = getattr(shape, "has_chart", False)
            if not has_chart:
                continue
            chart = shape.chart
            type_name = _chart_type_name(chart.chart_type)
            entry = {
                "slide_idx": slide_idx,
                "chart_type": type_name,
                "shape_name": getattr(shape, "name", ""),
            }
            if _is_ng_chart_type(type_name):
                ng_charts.append(entry)
            else:
                ok_charts.append(entry)

    total_charts = len(ng_charts) + len(ok_charts)
    ng_count = len(ng_charts)
    ok_count = len(ok_charts)

    observations = {
        "total_charts": total_charts,
        "ng_charts": ng_charts,
        "ok_charts": ok_charts,
        "ng_count": ng_count,
        "ok_count": ok_count,
    }

    # ---- score ----
    if total_charts == 0:
        score: float = 10.0
    else:
        score = round((ok_count / total_charts) * 10, 1)

    # ---- comments (PRIMARY OUTPUT) ----
    comments: list[str] = []

    if total_charts == 0:
        comments.append("chart 未使用。表・画像のみか、言語情報中心のデッキ。")
        comments.append(
            "数値比較を載せる場合は bar / horizontal bar / line / slopegraph を"
            "選ぶ。pie / 3D は宮野 S8 により NG。"
        )
    elif ng_count == 0:
        comments.append(
            f"全 {total_charts} chart が bar / line / scatter 系。"
            "宮野 S8 基準を満たす。"
        )
        comments.append(
            "chart の強調 (S9) と軸目盛り密度 (S8) も別途確認することを推奨。"
        )
    else:
        # NG chart ごとに per-slide comment
        for ng in ng_charts:
            comments.append(
                f"スライド {ng['slide_idx'] + 1} に {ng['chart_type']} 検出。"
                "宮野 S8 により pie / 3D chart は NG。"
                "bar chart / horizontal bar / slopegraph に置換。"
            )
        # サマリ行
        if ok_count > 0:
            comments.append(
                f"OK chart {ok_count} 件 / NG chart {ng_count} 件。"
                "NG を修正すれば score は上がる。"
            )
        else:
            comments.append(
                f"NG chart {ng_count} 件すべてが pie / 3D 系。"
                "角度比較・奥行歪みで定量読解を阻害する。全て置換推奨。"
            )

    # 最大 6 個に制限 (長大 deck で NG が多い場合)
    if len(comments) > 6:
        comments = comments[:5] + [
            f"... 他 {len(ng_charts) - 4} 件の NG chart を省略。observations.ng_charts を参照。"
        ]

    # 少なくとも 2 個を担保
    if len(comments) < 2:
        comments.append(
            "score は粗い指標。skill.md S8 の判定基準で最終確認を推奨。"
        )

    hint = (
        "chart_type 列挙ベースの粗い判定。stacked bar の 3D 版も NG 扱い。"
        "横棒 (bar) / 線グラフ / slopegraph を推奨。"
    )

    source = (
        "宮野『研究発表のためのスライドデザイン』S8 (グラフ選択: pie / 3D NG) / Tufte"
    )

    return {
        "score": float(score),
        "score_max": 10,
        "comments": comments,
        "observations": observations,
        "hint": hint,
        "source": source,
    }
