"""Plan B T10: pptx 下線密度診断 (presentation_check_underline_in_pptx).

pptx の全 run を走査し `font.underline` の数を数える粗い温度計。

skill.md T3 (count_underlines の pptx 版) / 宮野『研究発表のためのスライド
デザイン』S14 (下線密度) に基づく。下線を多用すると視覚的な壁紙と化して
emphasis としての機能を失う — 宮野 S14 の指摘を slide deck 全体 +
slide 単位で計測する。

設計方針:
- score は 0-10 の粗い温度計。9 以上は誤差範囲、6 以下は改善サイン。
- PRIMARY OUTPUT は comments (人間レビュアー風の助言文)。
- observations は raw measurement。numeric な正誤判定には使わない。
- `run.font.underline` の判定は python-pptx の仕様に従う:
    - `None`: layout / master から継承 → 下線なし扱い
    - `False`: 明示的に下線なし
    - `True`: 下線あり (既定のシングル)
    - MSO_TEXT_UNDERLINE_TYPE enum: NONE 以外はすべて下線扱い
"""

from __future__ import annotations

import pathlib


# ------------------------------------------------------------------
# Underline 判定
# ------------------------------------------------------------------
def _is_underlined(run) -> bool:
    """Return True iff the run is explicitly underlined.

    python-pptx 的 semantics:
      - None   : layout / master inherit (→ 既定は下線なし扱い)
      - False  : 明示的に off
      - True   : 明示的に on (single line)
      - MSO_TEXT_UNDERLINE_TYPE: NONE 以外は下線扱い。MIXED (-2) は一部だけ
        下線が掛かっている状態 → conservatively "下線あり" とカウント。
    """
    try:
        u = run.font.underline
    except Exception:
        return False
    if u is None or u is False:
        return False
    if u is True:
        return True
    # Enum-like. Compare by name to avoid version-specific import path.
    try:
        name = getattr(u, "name", None) or str(u)
    except Exception:
        return True  # conservative: treat unknown truthy value as underlined
    name = str(name).upper()
    if "NONE" in name and "DASH" not in name and "DOT" not in name:
        # "NONE" as exact enum, not a DASH/DOT variant
        return False
    return True


def _iter_runs_in_shape(shape):
    """Yield every text run nested in ``shape`` (text_frame + table cells)."""
    # 1) text frames
    try:
        if getattr(shape, "has_text_frame", False):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    yield run
    except Exception:
        pass

    # 2) tables
    try:
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    try:
                        tf = cell.text_frame
                    except Exception:
                        continue
                    for para in tf.paragraphs:
                        for run in para.runs:
                            yield run
    except Exception:
        pass


# ------------------------------------------------------------------
# Main
# ------------------------------------------------------------------
def presentation_check_underline_in_pptx(pptx_path: str) -> dict:
    """pptx runs の font.underline を走査して下線密度を診断する。

    宮野 S14: 下線を多用すると壁紙化して読みにくくなる。目安は deck 全体
    で 3-6 箇所以内、>= 10 箇所で warning。slide 単位の過剰 (>= 4) も別
    集計する。

    引数:
        pptx_path: .pptx のパス

    返り値:
        dict: score / score_max / comments / observations / hint / source
    """
    hint = (
        "MSO_TEXT_UNDERLINE_TYPE の None = 継承 default "
        "(→ underline なし扱い)。layout 経由で一律下線にしている場合は拾えない。"
    )
    source = (
        "宮野『研究発表のためのスライドデザイン』S14 / "
        "skill.md T3 (count_underlines の pptx 版)"
    )

    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {
            "score": 0,
            "score_max": 10,
            "comments": [
                f"pptx が見つからない: {pptx_path}",
                "ファイル path を確認して再実行すること。",
            ],
            "observations": {"error": f"file not found: {pptx_path}"},
            "hint": hint,
            "source": source,
        }

    try:
        import pptx as _pptx  # type: ignore
    except ImportError:
        return {
            "score": 0,
            "score_max": 10,
            "comments": [
                "python-pptx 未インストール。`pip install python-pptx` 後に再実行。",
            ],
            "observations": {"error": "python-pptx not installed"},
            "hint": hint,
            "source": source,
        }

    try:
        prs = _pptx.Presentation(str(p))
    except Exception as exc:
        return {
            "score": 0,
            "score_max": 10,
            "comments": [f"pptx 読み込み失敗: {exc}"],
            "observations": {"error": f"failed to open: {exc}"},
            "hint": hint,
            "source": source,
        }

    # `_walk_shapes` は tools.py 側に実装済み。group shape 対応で一貫性を保つ。
    try:
        from radia_mcp.presentation.tools import _walk_shapes  # type: ignore
    except Exception:
        _walk_shapes = None  # type: ignore

    slides = list(prs.slides)
    total_slides = len(slides)

    per_slide_underline_count: list[int] = []
    for slide in slides:
        count = 0
        if _walk_shapes is not None:
            try:
                shape_iter = _walk_shapes(slide.shapes)
            except Exception:
                shape_iter = slide.shapes
        else:
            shape_iter = slide.shapes
        try:
            for shape in shape_iter:
                for run in _iter_runs_in_shape(shape):
                    if _is_underlined(run):
                        count += 1
        except Exception:
            pass
        per_slide_underline_count.append(count)

    total_underlined_runs = sum(per_slide_underline_count)
    # 1-indexed slide numbers for over-threshold slides (>= 4 underlines)
    over_threshold_slides: list[int] = [
        i + 1 for i, c in enumerate(per_slide_underline_count) if c >= 4
    ]
    total_underline_density: float = (
        round(total_underlined_runs / total_slides, 3)
        if total_slides > 0 else 0.0
    )

    observations = {
        "total_slides": total_slides,
        "total_underlined_runs": total_underlined_runs,
        "per_slide_underline_count": per_slide_underline_count,
        "over_threshold_slides": over_threshold_slides,
        "total_underline_density": total_underline_density,
    }

    # ---- score ----
    if total_slides == 0:
        score = 10.0
    else:
        score = 10.0
        if total_underlined_runs > total_slides * 3:
            score -= 2
        # -1 per over_threshold slide, capped at -4
        score -= min(4, len(over_threshold_slides))
        score = max(0.0, min(10.0, score))

    # ---- comments ----
    comments: list[str] = []

    if total_slides == 0:
        comments.append("empty deck — スライドが 1 枚も無い。下線計測は対象外。")
        comments.append(
            "宮野 S14: 本文に下線を多用すると壁紙化して emphasis が機能しない — "
            "3-6 箇所/deck 以内が目安。"
        )
    else:
        if total_underlined_runs > 10:
            comments.append(
                f"deck 全体で下線 {total_underlined_runs} 箇所 — 壁紙化リスク。"
                "宮野 S14: 下線は emphasis として機能しなくなる。"
                "**bold** or color emphasis に置換を検討。"
            )

        if over_threshold_slides:
            comments.append(
                f"1 slide に下線 4 個以上: slide {over_threshold_slides}。"
                "審査員の目が麻痺する — 3 個以下に削る。"
            )

        if total_underlined_runs > total_slides * 3 and total_underlined_runs <= 10:
            comments.append(
                f"slide 平均 {total_underline_density} 箇所/枚 — 3 箇所を超えており"
                "密度過剰。宮野 S14 の壁紙化 risk、色/太字への置換を検討。"
            )

        if (total_underlined_runs <= 10
                and not over_threshold_slides
                and total_underlined_runs > 0
                and total_underlined_runs <= total_slides * 3):
            comments.append(
                f"下線 {total_underlined_runs} 箇所、deck 全体で適正。"
            )

        if total_underlined_runs == 0:
            comments.append(
                "下線 0 箇所、deck 全体で適正。"
                "key number/term は **bold** or color で emphasis する想定か確認。"
            )

    # フォールバック: 最低 2 コメント
    if len(comments) < 2:
        comments.append(
            "score は粗い指標。skill.md S14 / T3 の判定基準で最終確認を推奨。"
        )
    # cap at 5
    comments = comments[:5]

    return {
        "score": float(score),
        "score_max": 10,
        "comments": comments,
        "observations": observations,
        "hint": hint,
        "source": source,
    }
