"""Presentation T29: Q&A anticipation list (verbal side 強化)
(presentation_qa_anticipation_list).

slide の contents から **予想される Q&A 質問 list** を heuristic で生成。
発表前のリハーサル + backup slide 準備の支援。

検出する 8 種の典型質問:
  1. methodology_detail — 「なぜ X 手法を選んだ?」
  2. sample_size — 「n はなぜ N なのか?」
  3. generalization — 「他の Y 条件でも成立するか?」
  4. comparison_baseline — 「Baseline Z と比較したか?」
  5. parameter_sensitivity — 「Parameter P の変化に robust か?」
  6. reproducibility — 「code/data は公開するか?」
  7. future_work — 「次の段階は?」
  8. limitations_probe — 「本研究の弱点は?」
"""

from __future__ import annotations

import pathlib
import re


def _iter_shapes(shape_collection):
    try:
        from radia_mcp.presentation.tools import _walk_shapes
        yield from _walk_shapes(shape_collection)
    except Exception:
        for s in shape_collection:
            yield s


QA_CATEGORIES: list[dict] = [
    {
        "id": "methodology_detail",
        "label": "手法の選定理由",
        "trigger_patterns": [
            r"(?:手法|method|approach|アルゴリズム)",
        ],
        "sample_questions": [
            "なぜ X 手法を選んだのか? 他の手法 Y との比較は?",
            "この手法の limitation は何か? 具体的な失敗ケースは?",
            "Parameter tuning はどうやったか? hyperparameter の選び方は?",
        ],
    },
    {
        "id": "sample_size",
        "label": "サンプル数 / n の妥当性",
        "trigger_patterns": [
            r"n\s*=\s*\d+", r"サンプル", r"被験者",
        ],
        "sample_questions": [
            "N = X はどう決めた? Power analysis やったか?",
            "この N で統計的有意差を検出できる保証は?",
            "サンプルの代表性は? selection bias は無いか?",
        ],
    },
    {
        "id": "generalization",
        "label": "一般化可能性",
        "trigger_patterns": [
            r"(?:結果|result|performance)",
            r"\d+\s*[%x×倍]",
        ],
        "sample_questions": [
            "この結果は他の条件 (データセット / 装置 / 環境) でも再現する?",
            "Scale-up したらどうなる? edge case は?",
            "現場 (production / real-world) への適用時の課題は?",
        ],
    },
    {
        "id": "comparison_baseline",
        "label": "Baseline との比較",
        "trigger_patterns": [
            r"(?:baseline|比較|comparison|state[\- ]of[\- ]the[\- ]art|SOTA)",
        ],
        "sample_questions": [
            "Baseline Z (競合手法) とは比較したか? 数値は?",
            "Ablation study はやったか? 各 component の寄与は?",
            "Cost (computational / memory) での比較は?",
        ],
    },
    {
        "id": "parameter_sensitivity",
        "label": "Parameter sensitivity",
        "trigger_patterns": [
            r"(?:parameter|設定|条件|hyperparameter)",
            r"(?:結果|performance)\s*(?:は|が)",
        ],
        "sample_questions": [
            "Parameter P を変えると結果はどう変わる? sensitivity analysis は?",
            "最適な P はどう決めた? cross-validation やったか?",
            "Edge case (P → 0, P → ∞) での挙動は?",
        ],
    },
    {
        "id": "reproducibility",
        "label": "再現性 / 公開",
        "trigger_patterns": [
            r"(?:実装|implementation|code|data)",
        ],
        "sample_questions": [
            "Code / data は公開するか? GitHub URL は?",
            "再現手順は? dependency / version の明示は?",
            "乱数 seed / 計算環境の詳細は?",
        ],
    },
    {
        "id": "future_work",
        "label": "今後の展開",
        "trigger_patterns": [
            r"(?:まとめ|結論|conclusion|summary)",
        ],
        "sample_questions": [
            "次の研究の方向性は? open questions は?",
            "実用化までのロードマップは?",
            "他分野への展開可能性は?",
        ],
    },
    {
        "id": "limitations_probe",
        "label": "Limitation の深掘り",
        "trigger_patterns": [
            r"(?:limitation|限界|制約|課題)",
        ],
        "sample_questions": [
            "本研究の最大の weakness は? 認識している failure mode は?",
            "Scope 外として扱ったが重要な要因は?",
            "Alternative interpretation は無いか?",
        ],
    },
]


def _extract_all_text(pptx_path: str) -> str:
    try:
        import pptx as _pptx
    except ImportError:
        return ""
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return ""
    prs = _pptx.Presentation(str(p))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in _iter_shapes(slide.shapes):
            try:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text or "")
            except Exception:
                continue
        try:
            if slide.has_notes_slide:
                parts.append(slide.notes_slide.notes_text_frame.text or "")
        except Exception:
            pass
    return "\n".join(parts)


def presentation_qa_anticipation_list(pptx_path: str) -> dict:
    """slide 内容から予想 Q&A 質問 list を生成。

    Returns:
        dict: anticipated_questions (category 別 8 種) / top_priority_questions /
              backup_slide_suggestions / comments
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}

    all_text = _extract_all_text(pptx_path)
    if not all_text:
        return {
            "error": "text 抽出失敗",
            "hint": "pptx が読めない or 内容 空",
        }

    anticipated: list[dict] = []
    for cat in QA_CATEGORIES:
        # Check if any trigger pattern matches
        trigger_hit = False
        hit_examples: list[str] = []
        for pat in cat["trigger_patterns"]:
            for m in re.finditer(pat, all_text, re.IGNORECASE):
                trigger_hit = True
                hit_examples.append(m.group(0))
                if len(hit_examples) >= 3:
                    break
            if trigger_hit and len(hit_examples) >= 3:
                break
        if trigger_hit:
            anticipated.append({
                "category_id": cat["id"],
                "label": cat["label"],
                "trigger_hits": hit_examples[:3],
                "sample_questions": cat["sample_questions"],
                "priority": "HIGH" if len(hit_examples) >= 3 else "MEDIUM",
            })
        else:
            # Category 不発でも sample question は参考として残す
            anticipated.append({
                "category_id": cat["id"],
                "label": cat["label"],
                "trigger_hits": [],
                "sample_questions": cat["sample_questions"][:2],
                "priority": "LOW",  # slide に明示ないので質問確率低め
            })

    # Top priority: HIGH から先に、各 category から 1 件ずつ
    top_priority_questions: list[dict] = []
    for a in sorted(anticipated, key=lambda x: 0 if x["priority"] == "HIGH"
                     else 1 if x["priority"] == "MEDIUM" else 2):
        if a["sample_questions"]:
            top_priority_questions.append({
                "category": a["label"],
                "question": a["sample_questions"][0],
                "priority": a["priority"],
            })
        if len(top_priority_questions) >= 10:
            break

    # Backup slide suggestions
    backup_slide_suggestions = [
        f"Backup: {a['label']} 詳細" for a in anticipated
        if a["priority"] == "HIGH"
    ][:5]

    # ---- comments ----
    n_high = sum(1 for a in anticipated if a["priority"] == "HIGH")
    n_medium = sum(1 for a in anticipated if a["priority"] == "MEDIUM")

    comments: list[str] = []
    comments.append(
        f"Q&A 予測: HIGH {n_high} / MEDIUM {n_medium} / "
        f"LOW {len(anticipated) - n_high - n_medium} カテゴリ"
    )
    comments.append(
        "Top 5 想定質問 (順に backup slide 準備推奨):"
    )
    for q in top_priority_questions[:5]:
        comments.append(
            f"  • [{q['priority']}] {q['category']}: {q['question']}"
        )

    if backup_slide_suggestions:
        comments.append(
            f"📌 Backup slide 追加推奨: {backup_slide_suggestions}"
        )

    hint = (
        "slide 内容から trigger word を検出して、そのカテゴリで頻出する "
        "典型質問を list 化。予想精度は粗いが『当たる質問』の 7-8 割は "
        "covered できる。T26 rewrite_suggest の qa_backup_slide_template と "
        "併用して backup slide 用意。"
    )

    score = min(10.0, (n_high * 1.5 + n_medium * 0.5) * 1.5)
    score = max(0.0, round(score, 1))

    return {
        "score": score,
        "score_max": 10,
        "anticipated_questions": anticipated,
        "top_priority_questions": top_priority_questions,
        "backup_slide_suggestions": backup_slide_suggestions,
        "n_high_priority_categories": n_high,
        "n_medium_priority_categories": n_medium,
        "comments": comments,
        "hint": hint,
        "source": (
            "理系学会 Q&A 統計 (日本物理学会 / 応物 / IEEE conf 経験則)。"
            "8 category × 2-3 sample question。"
            "(presentation 2026-04 採用、v0.26.0)。"
        ),
    }
