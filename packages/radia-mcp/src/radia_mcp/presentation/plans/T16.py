"""Plan B T16: Western-style text density per slide
(presentation_text_density_per_slide_western_style).

欧米 PPT 文化は **text-heavy** (1 slide に段落型本文 + 多数 bullet) が
許容される。一方で日本理系プレゼンは **大字 / 簡潔** が標準で、欧米式
slide を直接使うと「読まされる感」で聴衆の集中が切れる。

本ツールは欧米式 text-heavy slide を検出:
  1. 1 slide あたり総文字数 (excl title) が overload
  2. 段落型 textbox の存在 (1 textbox に句点 ≥3、改行少)
  3. font size が小さい (細字 + 過密 = 読み手のみ意図)
  4. bullet level の深さ (level 3+ の入れ子 — 欧米軍用文書系)

設計方針:
- score は「1 主張で大字」slide 比率 (0-10、高いほど日本理系向き)
- 欧米寄り slide を flag、修正提案
- font size 取得には既存 _resolve_font_size を間接参照
- PRIMARY OUTPUT は western_heavy_slides list
"""

from __future__ import annotations

import pathlib
import re


def _iter_all_shapes(shape_collection):
    try:
        from radia_mcp.presentation.tools import _walk_shapes
        yield from _walk_shapes(shape_collection)
    except Exception:
        for s in shape_collection:
            yield s


def _resolve_size(run, para, shape, slide):
    try:
        from radia_mcp.presentation.tools import _resolve_font_size
        return _resolve_font_size(run, para, shape, slide)
    except Exception:
        try:
            return int(run.font.size.pt) if run.font.size else None
        except Exception:
            return None


def _shape_chars_excl_title(slide) -> tuple[int, int]:
    """Return (total_chars, paragraph_textboxes_count)."""
    total = 0
    para_textboxes = 0
    for shape in _iter_all_shapes(slide.shapes):
        try:
            pf = shape.placeholder_format
            if pf is not None and pf.idx == 0:
                continue  # skip title
        except Exception:
            pass
        try:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text or ""
            chars = len(txt)
            total += chars
            # Paragraph-style: ≥3 句点 in single textbox, low newline
            n_periods = txt.count("。") + txt.count(".") + txt.count("．")
            n_newlines = txt.count("\n")
            if n_periods >= 3 and n_newlines / max(n_periods, 1) < 1.0:
                para_textboxes += 1
        except Exception:
            continue
    return (total, para_textboxes)


def _max_bullet_level(slide) -> int:
    max_level = 0
    for shape in _iter_all_shapes(slide.shapes):
        try:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                lvl = getattr(para, "level", 0) or 0
                if lvl > max_level:
                    max_level = lvl
        except Exception:
            continue
    return max_level


def _avg_body_font_size(slide) -> float | None:
    sizes: list[int] = []
    for shape in _iter_all_shapes(slide.shapes):
        try:
            pf = shape.placeholder_format
            if pf is not None and pf.idx == 0:
                continue
        except Exception:
            pass
        try:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    sz = _resolve_size(run, para, shape, slide)
                    if sz:
                        sizes.append(sz)
        except Exception:
            continue
    if sizes:
        return sum(sizes) / len(sizes)
    return None


def presentation_text_density_per_slide_western_style(
    pptx_path: str,
    char_limit: int = 100,
    rikei_target_max: int = 50,
) -> dict:
    """欧米式 text-heavy slide を検出し、日本理系向けに修正提案。

    Args:
        pptx_path: .pptx file
        char_limit: 1 slide 文字数上限 (default 100、超えで warn)
        rikei_target_max: 日本理系 推奨上限 (default 50)

    Returns:
        dict: score / per_slide / western_heavy_slides / comments
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))

    per_slide: list[dict] = []
    western_heavy_slides: list[dict] = []

    for i, slide in enumerate(prs.slides, 1):
        chars, para_count = _shape_chars_excl_title(slide)
        max_level = _max_bullet_level(slide)
        avg_font = _avg_body_font_size(slide)

        # Western-heavy 判定 (3 つのうち 2+ が hit で western-heavy)
        signals = {
            "text_overload": chars > char_limit,
            "paragraph_style": para_count >= 1,
            "deep_nesting": max_level >= 3,
            "small_font": avg_font is not None and avg_font < 18,
        }
        n_signals = sum(1 for v in signals.values() if v)
        is_western_heavy = n_signals >= 2

        per_slide.append({
            "slide_no": i,
            "char_count": chars,
            "paragraph_textboxes": para_count,
            "max_bullet_level": max_level,
            "avg_body_font_pt": (
                round(avg_font, 1) if avg_font is not None else None
            ),
            "signals": signals,
            "is_western_heavy": is_western_heavy,
        })

        if is_western_heavy:
            western_heavy_slides.append({
                "slide_no": i,
                "char_count": chars,
                "signals_hit": [k for k, v in signals.items() if v],
                "advice": (
                    f"chars={chars} (target ≤{rikei_target_max}, "
                    f"limit ≤{char_limit}); "
                    "日本理系プレゼンとしては text-heavy。"
                    "1 主張に絞り、bullet level を 1-2 に抑え、"
                    "段落型は箇条書きに分解、font ≥24pt。"
                ),
            })

    n_slides = len(per_slide)
    n_western = len(western_heavy_slides)
    n_rikei = n_slides - n_western
    score = round(n_rikei / n_slides * 10, 1) if n_slides else 0.0

    # Aggregate stats
    total_chars = sum(s["char_count"] for s in per_slide)
    avg_chars = total_chars / n_slides if n_slides else 0
    n_over_target = sum(
        1 for s in per_slide if s["char_count"] > rikei_target_max
    )
    n_over_limit = sum(
        1 for s in per_slide if s["char_count"] > char_limit
    )

    # ---- comments ----
    comments: list[str] = []

    comments.append(
        f"スライド {n_slides} 枚 | avg {avg_chars:.0f} chars/slide | "
        f"target ≤{rikei_target_max} 超: {n_over_target} 枚 | "
        f"limit ≤{char_limit} 超: {n_over_limit} 枚 | "
        f"western-heavy: {n_western} 枚 (score {score}/10)"
    )

    if n_western >= 1:
        comments.append(
            f"⚠ 欧米式 text-heavy slide が {n_western} 枚。"
            "日本理系プレゼンの聴衆 (大学院生～教員) は『読まされる感』で "
            "集中が切れる。1 slide 1 主張 + 大字を徹底 (skill.md 高橋メソッド + 宮野)。"
        )
        for ws in western_heavy_slides[:5]:
            comments.append(
                f"  • slide {ws['slide_no']} ({ws['char_count']} chars, "
                f"signals: {ws['signals_hit']})"
            )

    if avg_chars > char_limit:
        comments.append(
            f"⚠ 全体平均 {avg_chars:.0f} chars/slide が limit を超過。"
            "プレゼン全体が text-heavy。slide 数を増やして 1 slide あたりを薄く。"
        )

    if n_western == 0:
        comments.append(
            "全 slide で日本理系プレゼン基準を充足 (text density OK)。"
        )

    comments = comments[:7]

    hint = (
        "理系プレゼンの聴衆は domain peer なので欧米式 text-heavy は不要。"
        "Tufte『Cognitive Style of PowerPoint』も同根の批判。"
        "char_limit と rikei_target_max は調整可 — 学会発表は 50-80 chars/slide、"
        "勉強会は 30-50 chars/slide が標準。"
    )

    return {
        "score": score,
        "score_max": 10,
        "n_slides": n_slides,
        "avg_chars_per_slide": round(avg_chars, 1),
        "n_over_target": n_over_target,
        "n_over_limit": n_over_limit,
        "n_western_heavy": n_western,
        "western_heavy_slides": western_heavy_slides[:15],
        "per_slide": per_slide[:30],
        "comments": comments,
        "hint": hint,
        "source": (
            "宮野『研究発表のためのスライドデザイン』S5 (1 slide 1 主張) + "
            "高橋メソッド + Tufte『Cognitive Style of PowerPoint』"
            "(2003) — bullet 過多 / 段落型 batchexpression slide 批判。"
            "(presentation 2026-04 採用)。"
        ),
    }
