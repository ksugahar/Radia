"""Plan B T13: Slide titles outline coherence
(presentation_slide_titles_outline_coherence).

skill.md L21「目標: スライドタイトルを並べたら目次にして筋が通る」を
実測化。**全 slide のタイトルだけを抜き出して並べた時に、論理が一直線に
通るか** を診断する。

検出する 5 軸:
  1. **title_density** — タイトル付き slide の比率 (no-title slide 過多 NG)
  2. **topic_specificity** — 対象＋観点が具体的なタイトルの比率
  3. **logical_flow_markers** — 「まず」「次に」「最後に」「結論」「対策」
     等の繋がり語の出現
  4. **topical_diversity** — 連続スライドの title token overlap
     (重複過多 = 同主張連発、過少 = 飛躍)
  5. **outline_completeness** — Hook / 背景 / 主張 / 証拠 / まとめ
     の典型構造を充たす keyword の検出

設計方針:
- pptx_path から全 title を抽出 → outline 配列を生成
- 5 軸独立に評価、score 化
- PRIMARY OUTPUT は extracted_outline (タイトル並び) + flow analysis
- 短い発表 (≤5 slides) は flow 評価 less strict
"""

from __future__ import annotations

import pathlib
import re


def _iter_all_shapes(shape_collection):
    try:
        from radia_mcp.presentation.tools import _walk_shapes
        yield from _walk_shapes(shape_collection)
    except Exception:
        for s in shape_collection:
            yield s


def _get_title(slide) -> str:
    try:
        from radia_mcp.presentation.tools import _slide_title
        return _slide_title(slide)
    except Exception:
        try:
            for shape in slide.shapes:
                if (getattr(shape, "has_text_frame", False)
                        and shape.text_frame.text.strip()):
                    return shape.text_frame.text.splitlines()[0]
        except Exception:
            pass
        return ""


_FLOW_MARKERS = {
    "opening": ["はじめに", "introduction", "背景", "background",
                 "motivation", "なぜ", "問題", "課題", "閉包条件", "従来法"],
    "middle": ["まず", "次に", "そして", "first", "second", "third",
                "next", "approach", "method", "手法", "アプローチ", "分解",
                "保存", "局所閉包", "検証条件", "実装"],
    "evidence": ["結果", "results", "実験", "experiment", "evidence",
                  "performance", "比較", "comparison", "偏差", "計算規模",
                  "実測性能", "整合性", "FEM"],
    "closing": ["まとめ", "結論", "conclusion", "summary", "takeaway",
                 "今後", "future", "次の段階", "謝辞", "acknowledg", "成果",
                 "限界", "展望", "HDiv"],
}

_CLAIM_VERBS = re.compile(
    r"\b(?:is|are|was|were|shows?|achieves?|reduces?|increases?|"
    r"enables?|improves?|outperforms?|solves?|demonstrates?|"
    r"proposes?|presents?|delivers?|cuts?)\b",
    re.IGNORECASE,
)

_CLAIM_JP_KEYWORDS = ["倍", "速", "削", "改善", "向上", "実現",
                       "達成", "減らす", "上回る", "可能に", "未達"]


def _is_claim_title(title: str) -> bool:
    if _CLAIM_VERBS.search(title):
        return True
    if any(k in title for k in _CLAIM_JP_KEYWORDS):
        return True
    # 数値 + 単位
    if re.search(r"\d+(?:\.\d+)?\s*(?:%|％|x|×|倍|kHz|MHz|GB)", title):
        return True
    return False


def _extract_meaningful_tokens(title: str) -> set[str]:
    """Title から noun-like tokens を集合で返す (overlap 判定用)。"""
    tokens: set[str] = set()
    # 漢字 2+
    for m in re.finditer(r"[一-鿿]{2,}", title):
        tokens.add(m.group(0))
    # カタカナ 3+
    for m in re.finditer(r"[ァ-ヴー]{3,}", title):
        tokens.add(m.group(0))
    # 英大文字始まり 4+
    for m in re.finditer(r"\b[A-Z][A-Za-z0-9]{3,}\b", title):
        tokens.add(m.group(0))
    return tokens


def presentation_slide_titles_outline_coherence(pptx_path: str) -> dict:
    """全 slide のタイトルだけ並べて outline 化し、論理整合を診断。

    skill.md L21 で言及されているが未実装だった「タイトル並べたら目次に
    して筋が通る」を 5 軸で評価。

    Args:
        pptx_path: .pptx ファイルパス

    Returns:
        dict: score / extracted_outline / per_axis / comments
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))

    from .._deck_sections import classify_deck_sections

    all_slides = list(prs.slides)
    sections = classify_deck_sections(all_slides, _get_title)
    backup_slide_numbers = set(sections["backup_slide_numbers"])

    titles: list[dict] = []
    for i, slide in enumerate(all_slides, 1):
        if i in backup_slide_numbers:
            continue
        title = _get_title(slide).strip()
        tokens = _extract_meaningful_tokens(title)
        titles.append({
            "slide_no": i,
            "title": title,
            "has_title": bool(title),
            "is_claim": _is_claim_title(title) if title else False,
            "tokens": tokens,
            "is_specific_topic": bool(tokens),
        })

    n_slides = len(titles)
    if n_slides == 0:
        return {
            "score": None,
            "comments": ["スライドが 1 枚も検出されず。"],
            "hint": "pptx 読み込み失敗の可能性。",
            "source": "T13 outline coherence (presentation 2026-04)",
        }

    # ---- 1. title_density ----
    n_with_title = sum(1 for t in titles if t["has_title"])
    title_density = n_with_title / n_slides

    # ---- 2. topic specificity / result-sentence diagnostic ----
    n_claims = sum(1 for t in titles if t["is_claim"])
    claim_density = n_claims / n_with_title if n_with_title else 0
    n_specific = sum(1 for t in titles if t["is_specific_topic"])
    topic_specificity = n_specific / n_with_title if n_with_title else 0

    # ---- 3. logical_flow_markers ----
    flow_hits: dict[str, list[int]] = {k: [] for k in _FLOW_MARKERS}
    for t in titles:
        if not t["has_title"]:
            continue
        title_lc = t["title"].lower()
        for phase, kws in _FLOW_MARKERS.items():
            if any(kw in title_lc or kw in t["title"] for kw in kws):
                flow_hits[phase].append(t["slide_no"])

    n_phases_present = sum(1 for v in flow_hits.values() if v)

    # ---- 4. topical_diversity (隣接 slide title overlap) ----
    overlaps: list[float] = []
    for i in range(len(titles) - 1):
        a, b = titles[i]["tokens"], titles[i + 1]["tokens"]
        if a or b:
            overlap = (
                len(a & b) / min(len(a), len(b)) if a and b else 0
            )
            overlaps.append(overlap)
    avg_overlap = sum(overlaps) / len(overlaps) if overlaps else 0

    # 0.0 = 完全飛躍、1.0 = 完全重複 (両方 NG、0.2-0.5 が target)
    if 0.15 <= avg_overlap <= 0.55:
        diversity_status = "BALANCED"
    elif avg_overlap < 0.15:
        diversity_status = "JUMPY"  # 飛躍多
    else:
        diversity_status = "REDUNDANT"  # 重複多

    # ---- 5. outline_completeness ----
    # opening / middle / evidence / closing が揃っているか
    completeness = n_phases_present / 4

    # ---- score ----
    # Title = concise target + viewpoint; bottom band = learned result.
    # Therefore claim density is diagnostic only and is not rewarded here.
    score = (
        title_density * 3
        + topic_specificity * 2
        + (n_phases_present / 4) * 3
        + (1 if diversity_status == "BALANCED" else 0.5
           if diversity_status != "REDUNDANT" else 0) * 2
    )
    score = round(min(10.0, score), 1)

    # ---- comments ----
    comments: list[str] = []

    comments.append(
        f"スライド {n_slides} 枚 | タイトル付き {n_with_title} ({title_density*100:.0f}%) "
        f"| 具体的な対象・観点 {n_specific} ({topic_specificity*100:.0f}%) "
        f"| 4 phase 内 {n_phases_present} 出現 | 隣接 overlap 平均 {avg_overlap:.2f} "
        f"({diversity_status})"
    )

    if title_density < 0.7:
        comments.append(
            f"⚠ Title density {title_density*100:.0f}% は低い。"
            "タイトル無し slide が多いと 目次として並べた時に筋が見えない。"
            "全 slide にタイトル必須。"
        )

    if claim_density > 0.4:
        comments.append(
            f"⚠ 結果言い切り型タイトルが {claim_density*100:.0f}%。"
            "タイトルはスライドで伝える具体的な対象・観点、下端文は図表から"
            "分かった結論として役割を分ける。"
        )

    if n_phases_present < 3:
        missing_phases = [k for k, v in flow_hits.items() if not v]
        comments.append(
            f"⚠ 4 phases 中 {n_phases_present} のみ出現。不在: {missing_phases}。"
            "opening (背景) / middle (アプローチ) / evidence (結果) / closing (まとめ) "
            "の流れが欠けている。"
        )

    if diversity_status == "REDUNDANT":
        comments.append(
            f"隣接 slide title 重複度 {avg_overlap:.2f} が高い。"
            "「Results 1」「Results 2」「Results 3」のように同じ "
            "テーマを連発している可能性。各 slide に異なる主張を。"
        )
    elif diversity_status == "JUMPY":
        comments.append(
            f"隣接 slide title 重複度 {avg_overlap:.2f} が低い。"
            "話題が飛躍している可能性。trans 句 (「次に〜」「これを踏まえ」) "
            "で繋ぎ言葉を加えるか、中間 slide を追加。"
        )

    if score >= 8:
        comments.append(
            f"Outline coherence score {score}/10 — タイトル並べて目次として "
            "筋が通っている。reviewer の斜め読みにも対応。"
        )

    # Outline 出力 (人間が読んで判定する用)
    outline_md = "\n".join(
        f"  {t['slide_no']:3d}. {t['title'][:60]} "
        f"{'[result]' if t['is_claim'] else '[topic]' if t['has_title'] else '[no title]'}"
        for t in titles[:30]
    )
    if len(titles) > 30:
        outline_md += f"\n  ... ({len(titles) - 30} more slides)"

    comments.append("📋 Extracted outline (titles only):\n" + outline_md)

    comments = comments[:8]

    hint = (
        "「タイトル並べたら目次として筋が通る」(skill.md L21) を実測化。"
        "score 高 = 具体的な対象・観点を示すタイトルが論理的に並んでいる。"
        "結果文は下端の takeaway に置き、タイトルとの役割を分ける。"
        "extracted_outline を音読して『目次として通るか』を最終判定。"
        "短い発表 (5 slides 以下) は phase 検出が thin になりがち。"
    )

    return {
        "score": score,
        "score_max": 10,
        "n_slides": n_slides,
        "total_deck_slides": len(all_slides),
        "backup_slide_count": sections["backup_slide_count"],
        "backup_slide_numbers": sections["backup_slide_numbers"],
        "title_density": round(title_density, 3),
        "claim_density": round(claim_density, 3),
        "topic_specificity": round(topic_specificity, 3),
        "n_phases_present": n_phases_present,
        "phases_detected": {k: v for k, v in flow_hits.items() if v},
        "avg_adjacent_overlap": round(avg_overlap, 3),
        "diversity_status": diversity_status,
        "extracted_outline": [
            {"slide": t["slide_no"], "title": t["title"],
             "is_claim": t["is_claim"]}
            for t in titles
        ],
        "comments": comments,
        "hint": hint,
        "source": (
            "skill.md L21 (タイトル並べたら目次として筋が通る) + "
            "宮野『研究発表のためのスライドデザイン』S2 (構成) + "
            "Reynolds『Presentation Zen』Ch.5 (storyboard)。"
            "(presentation 2026-04 採用)。"
        ),
    }
