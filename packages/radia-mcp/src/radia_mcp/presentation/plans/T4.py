"""Plan B T4: 全スライド ロゴ重複検出 (presentation_check_logo_on_every_slide).

skill.md 宮野 S13 「全スライドにロゴを入れない — 記憶容量を食うだけ」に基づく
粗い温度計。位置・サイズが揃った picture shape が多数スライドで繰り返されて
いれば、「ロゴの過剰配置」と推定する。

設計方針:
- score は 0-10 の粗い温度計。最終判断は comments と位置情報を読み、
  企業テンプレ運用のポリシーと合わせて人間判断で確定する。
- PRIMARY OUTPUT は comments (人間レビュアー風の助言文)。
- observations は raw measurement。
- ratio = 1.0 (全スライド) → score = 0。ratio が threshold 未満なら score = 10。
"""
from __future__ import annotations

import pathlib
from collections import Counter


# EMU rounding quantum — 100,000 EMU ≈ 0.1 cm = 1.06 mm の精度
_EMU_ROUND = -5


def _make_position_key(left, top, width, height) -> tuple:
    """EMU 座標を 100K 単位で丸めた位置キーを作る。

    python-pptx の shape.left 等は None になり得る (layout 継承ケース)。
    None は 0 とみなしてキー化する。
    """
    def _r(v):
        if v is None:
            return 0
        return int(round(int(v), _EMU_ROUND))
    return (_r(left), _r(top), _r(width), _r(height))


def presentation_check_logo_on_every_slide(
    pptx_path: str,
    threshold_ratio: float = 0.8,
) -> dict:
    """全スライドに同じロゴ画像が繰り返し配置されているかを検出。

    宮野『研究発表のためのスライドデザイン』S13「全スライドにロゴを入れるな
    — 記憶容量を食うだけ」に基づく粗い判定。全 slide の picture shape を
    収集し、(left, top, width, height) を 100K EMU 単位に丸めた位置キーで
    集計する。同じ位置キーが閾値以上のスライド数で観測されれば「logo
    推定」。

    引数:
        pptx_path: .pptx ファイルのパス
        threshold_ratio: 「全スライドに配置されている」と判定する比率
            (0.0-1.0)。既定 0.8。

    返り値:
        Plan B 標準 dict: score / score_max / comments / observations
        / hint / source
    """
    try:
        import pptx as _pptx
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except ImportError:
        return {
            "score": 10.0,
            "score_max": 10,
            "comments": [
                "python-pptx が未導入のため解析不能。pip install python-pptx 後に再実行。",
                "score は既定値 (10) を返す。判定には使わない。",
            ],
            "observations": {
                "total_slides": 0,
                "picture_shapes_per_slide": [],
                "repeated_picture_stats": {},
                "most_repeated_key": None,
                "most_repeated_ratio": 0.0,
                "is_logo_detected": False,
            },
            "hint": (
                "画像の (x, y, w, h) 位置マッチによる粗い判定。"
                "slide master 経由の継承も検出可能。企業 PPT template が"
                "慣例的にロゴを master に埋め込んでいる場合は本ツールでも"
                "拾えるが、template 運用のポリシーと調整。"
            ),
            "source": (
                "宮野『研究発表のためのスライドデザイン』S13 "
                "(photo / icon 一貫性、ロゴ過剰 NG)"
            ),
        }

    path = pathlib.Path(pptx_path)
    if not path.exists():
        return {
            "score": 10.0,
            "score_max": 10,
            "comments": [
                f"file not found: {pptx_path}",
                "score は既定値 (10) を返す。判定には使わない。",
            ],
            "observations": {
                "total_slides": 0,
                "picture_shapes_per_slide": [],
                "repeated_picture_stats": {},
                "most_repeated_key": None,
                "most_repeated_ratio": 0.0,
                "is_logo_detected": False,
            },
            "hint": (
                "画像の (x, y, w, h) 位置マッチによる粗い判定。"
                "slide master 経由の継承も検出可能。企業 PPT template が"
                "慣例的にロゴを master に埋め込んでいる場合は本ツールでも"
                "拾えるが、template 運用のポリシーと調整。"
            ),
            "source": (
                "宮野『研究発表のためのスライドデザイン』S13 "
                "(photo / icon 一貫性、ロゴ過剰 NG)"
            ),
        }

    prs = _pptx.Presentation(str(path))
    total_slides = len(prs.slides)

    picture_shapes_per_slide: list[list[dict]] = []
    # per slide: 該当 slide で使われた position_key の set
    # 同じ slide 内で同一 key が 2 枚ある場合でも、slide 単位では 1 回
    # カウントするため set を使う。
    slide_keys: list[set] = []

    for slide in prs.slides:
        picts: list[dict] = []
        keys_in_slide: set = set()
        for shape in slide.shapes:
            try:
                is_picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            except Exception:
                is_picture = False
            if not is_picture:
                continue
            pos_key = _make_position_key(
                getattr(shape, "left", 0),
                getattr(shape, "top", 0),
                getattr(shape, "width", 0),
                getattr(shape, "height", 0),
            )
            picts.append({
                "position_key": pos_key,
                "name": getattr(shape, "name", "") or "",
                "shape_id": getattr(shape, "shape_id", None),
            })
            keys_in_slide.add(pos_key)
        picture_shapes_per_slide.append(picts)
        slide_keys.append(keys_in_slide)

    # position_key ごとに「何枚のスライドに出現したか」を集計
    key_counter: Counter = Counter()
    for keys_in_slide in slide_keys:
        for k in keys_in_slide:
            key_counter[k] += 1

    repeated_picture_stats = dict(key_counter)

    if key_counter:
        most_repeated_key, most_repeated_count = key_counter.most_common(1)[0]
    else:
        most_repeated_key = None
        most_repeated_count = 0

    if total_slides > 0 and most_repeated_count > 0:
        most_repeated_ratio = most_repeated_count / total_slides
    else:
        most_repeated_ratio = 0.0

    is_logo_detected = (
        most_repeated_ratio >= threshold_ratio
        and most_repeated_count >= 3
    )

    observations = {
        "total_slides": total_slides,
        "picture_shapes_per_slide": picture_shapes_per_slide,
        "repeated_picture_stats": repeated_picture_stats,
        "most_repeated_key": most_repeated_key,
        "most_repeated_ratio": round(most_repeated_ratio, 3),
        "is_logo_detected": is_logo_detected,
    }

    # ---- score ----
    any_picture = any(picture_shapes_per_slide)
    if total_slides == 0 or not any_picture:
        score: float = 10.0
    elif is_logo_detected:
        # ratio が 1.0 → 0, 0.8 → 2.0, 0.9 → 1.0
        score = max(0.0, round((1.0 - most_repeated_ratio) * 10, 1))
    else:
        score = 10.0

    # ---- comments ----
    comments: list[str] = []

    if total_slides == 0:
        comments.append("空デッキ (slide 数 0)。宮野 S13 の評価対象なし。")
        comments.append("score は既定値 10 を返す。判定には使わない。")
    elif not any_picture:
        comments.append(
            "picture shape 不在。全スライドに画像が無いため、"
            "宮野 S13『全スライドにロゴを入れない』は自動満たされる。"
        )
        comments.append(
            "ただし図表・写真を一切使わない構成は別の問題 "
            "(image_text_ratio) を招く可能性があるため、並行チェック推奨。"
        )
    elif is_logo_detected:
        x, y, w, h = most_repeated_key  # type: ignore[misc]
        comments.append(
            f"スライド {most_repeated_count}/{total_slides} 枚に同じ位置 "
            f"(x={x}, y={y}) で同じサイズの画像 (logo 推定) が配置されている。"
            "宮野 S13 は『全スライドにロゴを入れない — 記憶容量を食うだけ』。"
            "タイトルスライド + 最終スライドのみに限定することを推奨。"
        )
        comments.append(
            f"ratio = {most_repeated_ratio:.2f} (threshold {threshold_ratio:.2f})。"
            f"位置キー (EMU 100K 単位): left={x}, top={y}, width={w}, height={h}。"
        )
        if most_repeated_ratio >= 0.99:
            comments.append(
                "全スライドに出現している。聴衆の目線が毎スライドでロゴに"
                "引かれ、本題への認知コストが増える (宮野 S13)。"
            )
        comments.append(
            "slide master / layout に埋め込まれているロゴは本ツールでも"
            "拾うため、template 運用のポリシー (学会指定 template など) "
            "と照合した上で除外判定すること。"
        )
    else:
        comments.append("画像の一律配置なし。宮野 S13 基準を満たす。")
        if most_repeated_count >= 2:
            comments.append(
                f"最頻出の位置キーは {most_repeated_count}/{total_slides} 枚 "
                f"(ratio {most_repeated_ratio:.2f}) に留まり、threshold "
                f"{threshold_ratio:.2f} 未満。logo 重複ではなく通常の図版使用と推定。"
            )
        else:
            comments.append(
                "同一位置・同一サイズの画像が 2 枚以上重複するスライド"
                "ペアなし。ロゴ過剰の兆候は検出されない。"
            )

    # 2-5 個に制限
    comments = comments[:5]
    while len(comments) < 2:
        comments.append(
            "score は粗い温度計。最終判断は comments と位置情報で確認。"
        )

    hint = (
        "画像の (x, y, w, h) 位置マッチによる粗い判定。slide master 経由の"
        "継承も検出可能。企業 PPT template が慣例的にロゴを master に"
        "埋め込んでいる場合は本ツールでも拾えるが、template 運用のポリシーと"
        "調整。"
    )
    source = (
        "宮野『研究発表のためのスライドデザイン』S13 "
        "(photo / icon 一貫性、ロゴ過剰 NG)"
    )

    return {
        "score": float(score),
        "score_max": 10,
        "comments": comments,
        "observations": observations,
        "hint": hint,
        "source": source,
    }
