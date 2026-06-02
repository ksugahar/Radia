"""Plan B T7: speaker note ratio 診断 (presentation_speaker_note_ratio).

木下『理科系の作文技術』の memo card 論 / 宮野『研究発表のためのスライドデザイン』
S1 storyboard 実践に基づく粗い温度計。

設計方針:
- score は 0-10 の粗い温度計。
- PRIMARY OUTPUT は comments (人間レビュアー風の助言文)。
- observations は raw measurement。numeric な正誤判定には使わない。
- speaker note が皆無 = 朗読型 risk (slide を読み上げるだけ) の疑い。
- speaker note が slide text の 5 倍超 = 台本朗読型 risk。
"""

from __future__ import annotations

import pathlib


# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------
def _slide_text_nontitle(slide) -> str:
    """Concatenate all non-title text from a slide."""
    parts: list[str] = []
    try:
        for shape in slide.shapes:
            try:
                pf = shape.placeholder_format
                is_title = (pf is not None and pf.idx == 0)
            except Exception:
                is_title = False
            if is_title:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            try:
                txt = shape.text_frame.text or ""
            except Exception:
                txt = ""
            if txt.strip():
                parts.append(txt)
    except Exception:
        pass
    return "\n".join(parts)


def _slide_note_text(slide) -> str:
    """Return stripped speaker-note text, or '' if absent / empty.

    Note: python-pptx creates a notes_slide object even for empty notes,
    so we must strip() to check for real content.
    """
    try:
        if not getattr(slide, "has_notes_slide", False):
            return ""
    except Exception:
        return ""
    try:
        ns = slide.notes_slide
    except Exception:
        return ""
    try:
        tf = ns.notes_text_frame
        txt = tf.text or ""
    except Exception:
        return ""
    return txt if txt.strip() else ""


# ------------------------------------------------------------------
# Main
# ------------------------------------------------------------------
def presentation_speaker_note_ratio(pptx_path: str) -> dict:
    """slide.notes_slide の speaker note 長 vs slide text 長の比率を検査。

    note >> text = 発表準備済み; note 不在 = 朗読型 risk。

    引数:
        pptx_path: .pptx のパス

    返り値:
        dict: score / score_max / comments / observations / hint / source
    """
    try:
        import pptx  # type: ignore  # noqa: F401
    except ImportError:
        return {"error": "python-pptx not installed. `pip install python-pptx`"}

    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}

    import pptx as _pptx
    prs = _pptx.Presentation(str(p))
    total_slides = len(prs.slides)

    per_slide_ratios: list[dict] = []
    slides_with_notes = 0
    slides_without_notes = 0
    ratios_for_mean: list[float] = []
    overly_note_heavy_slides: list[int] = []

    for i, slide in enumerate(prs.slides):
        slide_text = _slide_text_nontitle(slide)
        note_text = _slide_note_text(slide)
        slide_char_count = len(slide_text)
        note_char_count = len(note_text)

        if note_char_count > 0:
            slides_with_notes += 1
        else:
            slides_without_notes += 1

        # note_to_slide_ratio:
        # - slide_char_count == 0 and note > 0  -> treat as very high (note only)
        # - both zero -> 0.0
        if slide_char_count == 0 and note_char_count == 0:
            ratio = 0.0
        elif slide_char_count == 0:
            # note only — treat as "infinity" but use large finite number
            # for mean computation. We pick note_char_count itself as a
            # conservative magnitude.
            ratio = float(note_char_count)
        else:
            ratio = note_char_count / slide_char_count

        per_slide_ratios.append({
            "idx": i,
            "slide_char_count": slide_char_count,
            "note_char_count": note_char_count,
            "note_to_slide_ratio": round(ratio, 3),
        })

        # mean は note のある slide のみで計算 (note 無し slide の 0 で
        # 薄まるのを避ける)
        if note_char_count > 0:
            ratios_for_mean.append(ratio)

        if ratio > 5.0 and note_char_count > 0:
            overly_note_heavy_slides.append(i)

    if total_slides > 0:
        notes_ratio_coverage = round(slides_with_notes / total_slides, 3)
    else:
        notes_ratio_coverage = 0.0

    if ratios_for_mean:
        mean_note_to_slide_ratio = round(
            sum(ratios_for_mean) / len(ratios_for_mean), 3
        )
    else:
        mean_note_to_slide_ratio = 0.0

    observations = {
        "total_slides": total_slides,
        "slides_with_notes": slides_with_notes,
        "slides_without_notes": slides_without_notes,
        "notes_ratio_coverage": notes_ratio_coverage,
        "per_slide_ratios": per_slide_ratios,
        "mean_note_to_slide_ratio": mean_note_to_slide_ratio,
        "overly_note_heavy_slides": overly_note_heavy_slides,
    }

    # ---- score ----
    if total_slides == 0:
        score: float = 10.0
    else:
        score = 0.0
        if notes_ratio_coverage >= 0.3:
            score += 5.0
        if 1.0 <= mean_note_to_slide_ratio <= 5.0:
            score += 3.0
        if len(overly_note_heavy_slides) > total_slides * 0.3:
            score -= 2.0
        score = max(0.0, min(10.0, score))

    # ---- comments ----
    comments: list[str] = []

    if total_slides == 0:
        comments.append("empty deck — スライドが 1 枚も無い。")
        comments.append(
            "最低 1 枚は用意し、key slide には 2-3 文の speaker note を "
            "記す (木下 memo card 論)。"
        )
    else:
        if notes_ratio_coverage == 0:
            comments.append(
                "speaker note 全く無し。口頭発表の準備が追えない — "
                "Qiita 風に slide だけ読むと抑揚が出ない。"
                "key slides に 2-3 文で話すこと書く。"
            )

        n_heavy = len(overly_note_heavy_slides)
        if n_heavy > 0:
            comments.append(
                f"note が slide text の 5 倍超の slides {n_heavy} 件。"
                "台本朗読型の risk — audience の顔を見ず読み上げるプレゼン。bullet 化。"
            )

        if (
            notes_ratio_coverage >= 0.3
            and 1.0 <= mean_note_to_slide_ratio <= 5.0
        ):
            pct = round(notes_ratio_coverage * 100, 1)
            comments.append(
                f"note coverage {pct}%、平均比率 {mean_note_to_slide_ratio}。"
                "発表準備として適切。"
            )

        if (
            0 < notes_ratio_coverage < 0.3
            and "speaker note 全く無し" not in "\n".join(comments)
        ):
            pct = round(notes_ratio_coverage * 100, 1)
            comments.append(
                f"note coverage {pct}% — 準備 slide が 3 割未満。"
                "key slide (Intro / 結果 / Take-home) から note を埋める (宮野 S1)。"
            )

        if (
            notes_ratio_coverage >= 0.3
            and mean_note_to_slide_ratio < 1.0
            and n_heavy == 0
        ):
            comments.append(
                f"note coverage は {round(notes_ratio_coverage * 100, 1)}% と "
                f"あるが、平均比率 {mean_note_to_slide_ratio} は slide text より短い。"
                "要点のみの memo で良いが、intro slide は完全文を推奨 (木下)。"
            )

    # フォールバック: 最低 2 コメント
    if len(comments) < 2:
        comments.append(
            "score は粗い指標。木下 memo card 論 / 宮野 S1 で最終確認を推奨。"
        )

    # 最大 5 個
    comments = comments[:5]

    hint = (
        "note_text は notes_slide.notes_text_frame.text 経由。"
        "空 note でも object は存在することが多いので strip が必要。"
    )

    source = "木下 memo card 論 / 宮野 S1 storyboard 実践"

    return {
        "score": float(score),
        "score_max": 10,
        "comments": comments,
        "observations": observations,
        "hint": hint,
        "source": source,
    }
