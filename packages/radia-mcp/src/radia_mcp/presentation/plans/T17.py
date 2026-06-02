"""Plan B T17: Mini-IMRAD structure check (理系プレゼン特化)
(presentation_mini_imrad_structure_check).

理系プレゼンは論文の **mini IMRAD** 構造を踏襲するのが標準:

  Title → Background/Motivation → Problem/Question →
    Methods/Approach → Results/Evidence → Discussion → Conclusion

(Duarte の sparkline は corporate / TED 寄りなので採用せず、
理系プレゼン専用 mini-IMRAD を実装)。

検出する 7 phases:
  1. **title_card** — タイトル / 著者 (slide 1)
  2. **background** — 背景 / 動機 / 関連研究
  3. **problem** — 課題 / 仮説 / 問い
  4. **methods** — 手法 / アプローチ / 実装
  5. **results** — 結果 / 評価 / 比較
  6. **discussion** — 考察 / 限界 / 今後
  7. **conclusion** — まとめ / 結論 / 謝辞

設計方針:
- 全 slide title をスキャン → 各 phase keyword で分類
- 出現順序が標準 (title→bg→prob→method→results→discussion→conclusion) か
- 各 phase が複数 slide にまたがるとき密度バランスも check
- 短い発表 (< 6 slides) は phase 簡易判定
"""

from __future__ import annotations

import pathlib
import re


def _get_title(slide) -> str:
    try:
        from radia_mcp.presentation.tools import _slide_title
        return _slide_title(slide)
    except Exception:
        return ""


_PHASE_KEYWORDS = {
    "title_card": [
        # title slide is usually slide 1, no body
    ],
    "background": [
        "背景", "background", "motivation", "動機", "はじめに",
        "introduction", "概要", "overview", "現状", "context",
        "関連研究", "先行研究", "related\\s+work",
    ],
    "problem": [
        "課題", "problem", "challenge", "question", "問い", "目的",
        "objective", "目標", "仮説", "hypothesis", "未解決",
        "research\\s+question",
    ],
    "methods": [
        "手法", "method", "approach", "アプローチ", "アルゴリズム",
        "algorithm", "実装", "implementation", "構成", "system",
        "model", "モデル", "framework", "フレームワーク",
        "design", "設計", "プロトコル",
    ],
    "results": [
        "結果", "result", "実験", "experiment", "evaluation",
        "評価", "比較", "comparison", "performance", "性能",
        "ベンチマーク", "benchmark", "計測", "measurement",
        "データ", "data analysis",
    ],
    "discussion": [
        "考察", "discussion", "implication", "解釈", "interpretation",
        "限界", "limitation", "今後", "future", "発展", "次の段階",
        "課題と展望",
    ],
    "conclusion": [
        "まとめ", "conclusion", "結論", "summary", "key\\s+takeaway",
        "takeaway", "謝辞", "acknowledg", "Q&A", "ご清聴",
        "ありがとう", "thank\\s+you",
    ],
}


def _classify_slide(slide_no: int, title: str, body: str) -> str:
    """Return best matching phase, or 'unknown'."""
    if slide_no == 1 and len(title) <= 60 and not body.strip():
        return "title_card"

    title_lc = title.lower()
    body_lc = body.lower()
    full = title_lc + " " + body_lc

    # Score each phase by keyword hits in title (weight 3) + body (weight 1)
    best_phase = "unknown"
    best_score = 0
    for phase, kws in _PHASE_KEYWORDS.items():
        if not kws:
            continue
        score = 0
        for kw in kws:
            try:
                if re.search(kw, title_lc, re.IGNORECASE) or kw in title:
                    score += 3
                if re.search(kw, body_lc, re.IGNORECASE) or kw in body:
                    score += 1
            except re.error:
                continue
        if score > best_score:
            best_score = score
            best_phase = phase
    return best_phase if best_score > 0 else "unknown"


def _get_body_text(slide) -> str:
    parts: list[str] = []
    title = _get_title(slide).strip()
    try:
        from radia_mcp.presentation.tools import _walk_shapes
        shapes = _walk_shapes(slide.shapes)
    except Exception:
        shapes = slide.shapes
    for shape in shapes:
        try:
            pf = shape.placeholder_format
            if pf is not None and pf.idx == 0:
                continue
        except Exception:
            pass
        try:
            if shape.has_text_frame:
                txt = shape.text_frame.text or ""
                if txt.strip() != title:
                    parts.append(txt)
        except Exception:
            continue
    return "\n".join(parts)


def presentation_mini_imrad_structure_check(pptx_path: str) -> dict:
    """理系プレゼンの mini-IMRAD 構造 (7 phases) 充足度を診断。

    Args:
        pptx_path: .pptx file

    Returns:
        dict: score / phase_distribution / order_check / comments
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))

    classifications: list[dict] = []
    for i, slide in enumerate(prs.slides, 1):
        title = _get_title(slide).strip()
        body = _get_body_text(slide).strip()
        phase = _classify_slide(i, title, body)
        classifications.append({
            "slide_no": i,
            "title": title[:60],
            "phase": phase,
        })

    n_slides = len(classifications)
    phase_distribution: dict[str, list[int]] = {
        p: [] for p in _PHASE_KEYWORDS
    }
    phase_distribution["unknown"] = []
    for c in classifications:
        phase_distribution[c["phase"]].append(c["slide_no"])

    n_phases_present = sum(
        1 for k, v in phase_distribution.items()
        if v and k != "unknown"
    )

    # Order check: do phases appear in expected order?
    expected_order = [
        "title_card", "background", "problem", "methods",
        "results", "discussion", "conclusion",
    ]
    order_violations: list[str] = []
    last_seen_idx = -1
    for c in classifications:
        if c["phase"] in expected_order:
            curr_idx = expected_order.index(c["phase"])
            if curr_idx < last_seen_idx - 1:  # allow small forward jumps
                order_violations.append(
                    f"slide {c['slide_no']} ({c['phase']}) "
                    f"appears after later phase"
                )
            last_seen_idx = max(last_seen_idx, curr_idx)

    # Phase count balance
    n_methods = len(phase_distribution["methods"])
    n_results = len(phase_distribution["results"])
    if n_results > 0 and n_methods == 0:
        balance_issue = "Results あるが Methods 無し (再現不能 評価)"
    elif n_methods > n_slides * 0.5:
        balance_issue = "Methods 過剰 (時間不足リスク)"
    elif n_results > n_slides * 0.5:
        balance_issue = "Results 過剰 (考察薄リスク)"
    else:
        balance_issue = ""

    # ---- score ----
    score = 0.0
    score += n_phases_present / 7 * 5  # phase coverage (max 5)
    score += (1 - len(order_violations) / max(n_slides, 1)) * 3  # order (max 3)
    score += 2 if not balance_issue else (1 if balance_issue else 0)
    score = round(min(10.0, max(0.0, score)), 1)

    # ---- comments ----
    comments: list[str] = []

    phases_present = [
        k for k, v in phase_distribution.items()
        if v and k != "unknown"
    ]
    comments.append(
        f"Mini-IMRAD: {n_phases_present}/7 phases 出現 ({phases_present})。"
        f"score {score}/10。"
    )

    missing = [
        k for k in expected_order
        if not phase_distribution.get(k)
    ]
    if missing:
        comments.append(
            f"⚠ 不在 phases: {missing}。理系プレゼンは title→background→"
            "problem→methods→results→discussion→conclusion の "
            "mini-IMRAD が標準。"
        )

    if order_violations:
        comments.append(
            f"⚠ Phase 出現順違反 {len(order_violations)} 件。"
            "聴衆は IMRAD 順を期待しているので、Methods が Results の "
            "後に来る等は誤読の元。"
        )
        for v in order_violations[:3]:
            comments.append(f"  • {v}")

    if balance_issue:
        comments.append(f"⚠ Phase balance: {balance_issue}")

    n_unknown = len(phase_distribution["unknown"])
    if n_unknown > n_slides * 0.3:
        comments.append(
            f"ℹ {n_unknown} slides で phase 判定不能。"
            "title が phase keyword を含まない可能性 (figure-only slides 等)。"
            "section divider slide を入れると視認性向上。"
        )

    if score >= 8 and not balance_issue and not missing:
        comments.append(
            f"score {score}/10 — mini-IMRAD 構造を完全充足。"
            "聴衆は『どこに居るか』を常に把握できる。"
        )

    comments = comments[:7]

    hint = (
        "Duarte sparkline (TED 風) ではなく、理系プレゼン専用の "
        "mini-IMRAD (論文 IMRAD 縮小版) で評価。"
        "T13 (titles outline coherence) と相補的: T13 は title 並び論理、"
        "T17 は phase 構造充足度。"
        "section divider slide があると phase 検出が安定。"
    )

    return {
        "score": score,
        "score_max": 10,
        "n_slides": n_slides,
        "n_phases_present": n_phases_present,
        "phase_distribution": phase_distribution,
        "order_violations": order_violations,
        "balance_issue": balance_issue,
        "expected_order": expected_order,
        "classifications": classifications[:30],
        "comments": comments,
        "hint": hint,
        "source": (
            "Day & Gastel『How to Write & Publish a Scientific Paper』(7th ed) "
            "Ch.10 (IMRAD) を理系プレゼンに縮小適用 + 宮野『研究発表の "
            "ためのスライドデザイン』S2 (構成)。Duarte sparkline (corporate) "
            "ではなく学会発表向けの構造軸。"
            "(presentation 2026-04 採用)。"
        ),
    }
