"""Plan B T22: Results slide statistical evidence (理系プレゼン特化)
(presentation_results_slide_statistical_evidence).

理系プレゼンの Results slide では論文と同じく **統計報告の compliance**
が必要。paper T12 (statistical_reporting_compliance) と同根だが、
プレゼン slide では「視覚的に伝わる」必要があり、別軸の評価が必要。

Results slide で要求される 4 要素:
  1. **error_bars** — figure に error bar / 信頼区間 / 標準偏差の視覚化
  2. **sample_size** — n = X / N=X の明示
  3. **statistical_test** — p 値 or test 名 (t-test / ANOVA / etc)
  4. **effect_size** — magnitude (倍率 / %改善 / Cohen's d 等)

設計方針:
- Results / 結果 / 評価 keyword を含む slide を抽出
- 各 slide で 4 要素 keyword detection
- score は per-Results-slide 平均
- PRIMARY OUTPUT は statistically_weak_slides list
"""

from __future__ import annotations

import pathlib
import re


def _iter_all_shapes(shape_collection):
    try:
        from radia_mcp.presentation.tools import _walk_shapes
        yield from _walk_shapes(shape_collection)
    except Exception:
        for s in shape_collection:
            yield s


def _get_title(slide) -> str:
    try:
        from radia_mcp.presentation.tools import _slide_title
        return _slide_title(slide)
    except Exception:
        return ""


_RESULTS_TITLE_PATTERNS = [
    r"(?:results?|結果|実験結果|評価|evaluation|performance|"
    r"性能|比較|comparison|benchmark|ベンチマーク)",
]


def _is_results_slide(slide) -> bool:
    title = _get_title(slide).lower()
    if any(re.search(p, title, re.IGNORECASE) for p in _RESULTS_TITLE_PATTERNS):
        return True
    # Body 内に強い結果 markers があれば results slide とみなす
    full_text = ""
    for shape in _iter_all_shapes(slide.shapes):
        try:
            if shape.has_text_frame:
                full_text += " " + (shape.text_frame.text or "")
        except Exception:
            continue
    # 数値 + improvement/reduction が複数あれば results slide
    n_results_signals = len(re.findall(
        r"\d+(?:\.\d+)?\s*(?:%|％|x|×|倍|fold).*?"
        r"(?:improve|reduc|enhanc|向上|削減|短縮|増加|減少)",
        full_text, re.IGNORECASE,
    ))
    return n_results_signals >= 2


_ERROR_BAR_PATTERNS = [
    r"(?:error\s*bar|エラーバー|信頼区間|confidence\s+interval|"
    r"95%\s*CI|standard\s+deviation|標準偏差|SD|SEM|±)",
    r"\d+(?:\.\d+)?\s*±\s*\d+",  # 5.2 ± 0.3
]

_SAMPLE_SIZE_PATTERNS = [
    r"\bn\s*=\s*\d+",
    r"\bN\s*=\s*\d+",
    r"被験者\s*\d+\s*名",
    r"サンプル\s*(?:数|サイズ)\s*[=:]?\s*\d+",
    r"\d+\s+(?:participants?|subjects?|trials?|runs?)",
]

_STAT_TEST_PATTERNS = [
    r"\bp\s*[<>=]\s*0?\.\d",
    r"(?:t-?test|ANOVA|chi-?square|Kruskal|Mann-?Whitney|Wilcoxon|"
    r"Tukey|Bonferroni|Fisher|χ²|F\s*\(\d+,\s*\d+\))",
    r"(?:t検定|分散分析|ノンパラメトリック|有意水準)",
    r"\bp\s*[-=]\s*value",
]

_EFFECT_SIZE_PATTERNS = [
    r"\d+(?:\.\d+)?\s*(?:%|％|x|×|倍|fold)",
    r"(?:Cohen'?s\s*d|η²|eta\s*squared|odds\s*ratio|OR|risk\s*ratio|RR)",
    r"(?:improvement|reduction|increase|decrease)\s+(?:of|by)?\s*\d",
    r"\d+(?:\.\d+)?-?(?:fold|倍)\s*(?:improvement|reduc|enhanc|向上|削減)",
]


def _check_pattern(text: str, patterns: list[str]) -> bool:
    return any(re.search(p, text, re.IGNORECASE) for p in patterns)


def presentation_results_slide_statistical_evidence(pptx_path: str) -> dict:
    """Results slide の統計報告 4 要素 compliance を診断 (paper T12 の plot 版)。

    Args:
        pptx_path: .pptx file

    Returns:
        dict: score / results_slides / statistically_weak_slides / comments
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))

    results_slides: list[dict] = []
    n_total = 0

    for i, slide in enumerate(prs.slides, 1):
        n_total += 1
        if not _is_results_slide(slide):
            continue

        full_text = ""
        for shape in _iter_all_shapes(slide.shapes):
            try:
                if shape.has_text_frame:
                    full_text += "\n" + (shape.text_frame.text or "")
            except Exception:
                continue

        has_error_bars = _check_pattern(full_text, _ERROR_BAR_PATTERNS)
        has_n = _check_pattern(full_text, _SAMPLE_SIZE_PATTERNS)
        has_test = _check_pattern(full_text, _STAT_TEST_PATTERNS)
        has_effect = _check_pattern(full_text, _EFFECT_SIZE_PATTERNS)

        n_pass = sum([has_error_bars, has_n, has_test, has_effect])
        slide_score = round(n_pass / 4 * 10, 1)

        missing: list[str] = []
        if not has_error_bars:
            missing.append("error_bars/CI")
        if not has_n:
            missing.append("sample_size_n")
        if not has_test:
            missing.append("statistical_test/p_value")
        if not has_effect:
            missing.append("effect_size/magnitude")

        results_slides.append({
            "slide_no": i,
            "title": _get_title(slide)[:60],
            "has_error_bars": has_error_bars,
            "has_sample_size": has_n,
            "has_statistical_test": has_test,
            "has_effect_size": has_effect,
            "compliance_score": slide_score,
            "missing": missing,
        })

    n_res = len(results_slides)
    if n_res == 0:
        return {
            "score": 10.0,
            "score_max": 10,
            "n_total_slides": n_total,
            "n_results_slides": 0,
            "comments": [
                "Results slide 検出ゼロ。理論系・proof of concept 系発表 か、"
                "title が『Results』『結果』を含まない可能性。"
                "本ツールは実験系 results slide 向け。"
            ],
            "hint": "no results slides to evaluate.",
            "source": "T22 results statistical evidence (presentation 2026-04)",
        }

    overall = round(
        sum(s["compliance_score"] for s in results_slides) / n_res, 1
    )
    statistically_weak = [
        s for s in results_slides if s["compliance_score"] < 5
    ]

    # ---- comments ----
    comments: list[str] = []

    n_weak = len(statistically_weak)
    comments.append(
        f"Results slide: {n_res}/{n_total} 枚 | weak {n_weak} | "
        f"avg compliance {overall}/10"
    )

    if n_weak >= 1:
        comments.append(
            f"⚠ {n_weak} 枚で統計報告 weak。"
            "理系プレゼンの Results slide は paper T12 と同じ 4 要素 "
            "(error bar / n / test / effect size) を articulate。"
            "Q&A で『有意差は?』『N は?』を予防。"
        )
        for s in statistically_weak[:5]:
            comments.append(
                f"  • slide {s['slide_no']} [{s['compliance_score']}/10]: "
                f"{s['title']} — missing {s['missing']}"
            )

    n_no_eb = sum(
        1 for s in results_slides if not s["has_error_bars"]
    )
    n_no_n = sum(
        1 for s in results_slides if not s["has_sample_size"]
    )
    n_no_test = sum(
        1 for s in results_slides if not s["has_statistical_test"]
    )

    if n_no_eb > n_res * 0.5:
        comments.append(
            f"⚠ {n_no_eb} Results slides で error bar / 信頼区間言及なし。"
            "実験系では『誤差表示無し棒グラフ』は信頼性 ⇣ 印象。"
        )
    if n_no_n > n_res * 0.5:
        comments.append(
            f"⚠ {n_no_n} Results slides で sample size n の明示なし。"
            "「n=10 で統計検定?」と reviewer に疑念。"
        )
    if n_no_test > n_res * 0.5:
        comments.append(
            f"ℹ {n_no_test} Results slides で統計検定言及なし。"
            "差を主張するなら test 名 (t-test 等) と p 値を 1 行で。"
        )

    if n_weak == 0:
        comments.append(
            f"全 {n_res} Results slide で 4 要素充足。"
            "Reviewer の統計関連 Q&A はほぼ予防済。"
        )

    comments = comments[:7]

    hint = (
        "paper T12 (statistical_reporting_compliance) の slide 版。"
        "Results slide は『見て分かる』必要があるので error bar の "
        "視覚化が特に重要。理論系 / proof-of-concept 系では適用外。"
        "Image-only 結果図 (image 化された棒グラフ) は本 tool では拾えない。"
    )

    return {
        "score": overall,
        "score_max": 10,
        "n_total_slides": n_total,
        "n_results_slides": n_res,
        "n_statistically_weak": n_weak,
        "results_slides": results_slides,
        "statistically_weak_slides": statistically_weak,
        "comments": comments,
        "hint": hint,
        "source": (
            "APA Publication Manual 7th ed. (2020) + CONSORT 2010 + "
            "STROBE Statement (2007) を理系プレゼン Results slide に "
            "適用。paper T12 の slide 版。"
            "(presentation 2026-04 採用、理系特化)。"
        ),
    }
