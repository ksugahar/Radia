"""Plan B T11: slide 内 text 密度バランス診断 (presentation_slide_density_balance).

student anti-pattern (レビューでしばしば指摘) と
横徹流『発表 50 枚 vs 20 枚』議論に基づく粗い温度計。
deck 内の char-count 分布の不均衡を Gini-like 指標で評価し、
2-3 slide に全 deck text の 50% が集中するような偏りを flag する。

設計方針:
- score は 0-10 の粗い温度計。9 以上は誤差範囲、6 以下は改善サイン。
- PRIMARY OUTPUT は comments (人間レビュアー風の助言文)。
- observations は raw measurement。numeric な正誤判定には使わない。
- 偏り指標は top-3 集中比率と Gini coefficient の 2 本立て。
- Research talk では結果 slide に text が集中しがちで OK の場合もある。
  method slide で text が多いなら structuring 弱のサイン。
"""
from __future__ import annotations

import pathlib


# ------------------------------------------------------------------
# Gini coefficient (standard discrete formula)
# ------------------------------------------------------------------
def _gini(values: list[int]) -> float:
    """Return the Gini coefficient of *values* in [0, 1].

    0 = perfectly equal distribution.
    1 = all mass concentrated in a single element (in the limit n -> inf).

    Uses the sorted formula:
        G = (2 * sum_{i=1..n} i * x_i) / (n * sum x) - (n + 1) / n
    """
    if not values:
        return 0.0
    vals = sorted(values)
    n = len(vals)
    total = sum(vals)
    if total <= 0:
        return 0.0
    cum = 0
    for i, v in enumerate(vals, 1):
        cum += i * v
    return (2 * cum) / (n * total) - (n + 1) / n


def _median(values: list[int]) -> int:
    if not values:
        return 0
    s = sorted(values)
    n = len(s)
    mid = n // 2
    if n % 2 == 1:
        return int(s[mid])
    return int((s[mid - 1] + s[mid]) / 2)


def _percentile(values: list[int], p: float) -> int:
    """Nearest-rank percentile. p in [0, 1]."""
    if not values:
        return 0
    s = sorted(values)
    n = len(s)
    # nearest-rank: ceil(p * n), clamped to [1, n]
    import math
    k = max(1, min(n, int(math.ceil(p * n))))
    return int(s[k - 1])


def _slide_char_count(slide) -> int:
    """Sum of visible text characters across all shapes on *slide*.

    Uses presentation.tools._walk_shapes / _extract_shape_text_items when
    available so group shapes, tables, and SmartArt text are counted.
    Falls back to a local text_frame walk otherwise.
    """
    try:
        from radia_mcp.presentation.tools import (
            _extract_shape_text_items,
            _walk_shapes,
        )
    except Exception:
        _walk_shapes = None  # type: ignore
        _extract_shape_text_items = None  # type: ignore

    total = 0

    if _walk_shapes is not None and _extract_shape_text_items is not None:
        try:
            for shape in _walk_shapes(slide.shapes):
                try:
                    for frag in _extract_shape_text_items(shape):
                        total += len(frag)
                except Exception:
                    continue
            return total
        except Exception:
            pass

    # Fallback walk
    try:
        for shape in slide.shapes:
            try:
                if getattr(shape, "has_text_frame", False):
                    total += len(shape.text_frame.text or "")
            except Exception:
                continue
    except Exception:
        pass
    return total


# ------------------------------------------------------------------
# Main
# ------------------------------------------------------------------
def presentation_slide_density_balance(pptx_path: str) -> dict:
    """deck 内の char-count 分布の不均衡を Gini-like 指標で評価。

    2-3 slide に全 deck text の 50% が集中するような偏りを flag する。

    引数:
        pptx_path: .pptx のパス

    返り値:
        Plan B 標準 dict: score / score_max / comments / observations
        / hint / source
    """
    try:
        import pptx as _pptx  # type: ignore
    except ImportError:
        return {
            "score": 10.0,
            "score_max": 10,
            "comments": [
                "python-pptx が未導入のため解析不能。pip install python-pptx 後に再実行。",
                "score は既定値 (10) を返す。判定には使わない。",
            ],
            "observations": {
                "total_slides": 0,
                "total_text_chars": 0,
                "per_slide_chars": [],
                "median_chars": 0,
                "p90_chars": 0,
                "top_3_ratio": 0.0,
                "gini_coefficient": 0.0,
            },
            "hint": (
                "Gini の通念: > 0.5 = 顕著な不均衡。Research talk では"
                "結果 slide に text が集中しがち (OK)、method slide で "
                "text が多いなら structuring 弱。"
            ),
            "source": (
                "student anti-pattern (レビューでしばしば指摘) / "
                "横徹流『発表 50 枚 vs 20 枚』議論"
            ),
        }

    path = pathlib.Path(pptx_path)
    if not path.exists():
        return {"error": f"file not found: {pptx_path}"}

    prs = _pptx.Presentation(str(path))
    total_slides = len(prs.slides)

    per_slide_chars: list[int] = []
    for slide in prs.slides:
        per_slide_chars.append(_slide_char_count(slide))

    total_text_chars = sum(per_slide_chars)
    median_chars = _median(per_slide_chars)
    p90_chars = _percentile(per_slide_chars, 0.9)

    # top-3 ratio
    if total_text_chars > 0:
        top3 = sorted(per_slide_chars, reverse=True)[:3]
        top_3_ratio = sum(top3) / total_text_chars
    else:
        top_3_ratio = 0.0

    gini_coefficient = _gini(per_slide_chars)

    observations = {
        "total_slides": total_slides,
        "total_text_chars": total_text_chars,
        "per_slide_chars": per_slide_chars,
        "median_chars": median_chars,
        "p90_chars": p90_chars,
        "top_3_ratio": round(top_3_ratio, 3),
        "gini_coefficient": round(gini_coefficient, 3),
    }

    # ---- score ----
    if total_slides < 3 or total_text_chars == 0:
        score: float = 10.0
    else:
        clipped_g = min(max(gini_coefficient, 0.0), 1.0)
        score = round((1.0 - clipped_g) * 10, 1)
        if top_3_ratio > 0.5:
            score = max(0.0, round(score - 2.0, 1))

    # ---- comments ----
    comments: list[str] = []

    if total_slides == 0:
        comments.append("空デッキ (slide 数 0)。density balance 評価対象なし。")
        comments.append("score は既定値 10 を返す。判定には使わない。")
    elif total_slides < 3:
        comments.append(
            f"slide 総数 {total_slides} 枚は density balance 評価には少なすぎる "
            "(最低 3 枚)。score は既定値 10 を返す。"
        )
        comments.append(
            "章構造が生まれた段階で再実行して偏りを確認すること。"
        )
    elif total_text_chars == 0:
        comments.append(
            "visible text 0 文字。図版中心 deck の可能性はあるが、"
            "image_text_ratio 等の併用推奨。"
        )
        comments.append("score は既定値 10 を返す。判定には使わない。")
    else:
        pct_top3 = round(top_3_ratio * 100, 1)
        g_round = round(gini_coefficient, 2)

        if top_3_ratio > 0.5:
            comments.append(
                f"上位 3 slides に全 text の {pct_top3}% が集中。"
                "密度偏りが大きすぎる — 詰まっている slides を 2-3 枚に分割。"
            )
        if gini_coefficient > 0.5:
            comments.append(
                f"char 分布 Gini {g_round}。slide ごとの text 量にばらつきあり。"
            )
        if top_3_ratio <= 0.5 and gini_coefficient <= 0.5:
            comments.append(
                f"slide 間の text 量バランス良好 (Gini {g_round}, "
                f"top-3 比率 {pct_top3}%)。"
            )

        # Diagnostic context: median vs p90 gap
        if median_chars > 0 and p90_chars >= median_chars * 3:
            comments.append(
                f"p90 ({p90_chars} 字) が median ({median_chars} 字) の "
                f"{round(p90_chars / max(median_chars, 1), 1)} 倍。"
                "text-heavy slide が偏在。結果 slide なら許容、"
                "method slide なら structuring 弱。"
            )

        # Positive / neutral fallback
        if gini_coefficient <= 0.3 and top_3_ratio <= 0.4:
            comments.append(
                f"Gini {g_round} は低く、text 密度は全 slide に概ね均等配分。"
            )

    # 2-5 個に制限
    comments = comments[:5]
    while len(comments) < 2:
        comments.append(
            "score は粗い温度計。最終判断は per_slide_chars の内訳を確認。"
        )

    hint = (
        "Gini の通念: > 0.5 = 顕著な不均衡。Research talk では"
        "結果 slide に text が集中しがち (OK)、method slide で "
        "text が多いなら structuring 弱。"
    )
    source = (
        "student anti-pattern (レビューでしばしば指摘) / "
        "横徹流『発表 50 枚 vs 20 枚』議論"
    )

    return {
        "score": float(score),
        "score_max": 10,
        "comments": comments,
        "observations": observations,
        "hint": hint,
        "source": source,
    }
