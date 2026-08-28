"""Plan B T9: 矢印使用量の粗い温度計 (presentation_arrow_usage).

skill.md / 宮野『研究発表のためのスライドデザイン』S10「矢印は多用しない」に
基づく粗い温度計。line connector や arrow-end を持つ shape が 1 枚のスライド
に過剰配置されていないか、太線・塗りつぶしの「騒がしい」矢印が混入していない
かを検出する。

設計方針:
- score は 0-10 の粗い温度計。comments (PRIMARY) と observations を読んで
  人間レビュアーが最終判断する。
- 矢印 0 本は「問題なし」で満点 (宮野 S10 は『矢印使用自体』ではなく
  『矢印の過剰使用・目立ちすぎる矢印』を戒めるため)。
- python-pptx は arrow-end を完全に API 化していない (line.end_arrow_type
  が存在しない)。よって XML を直接 parse して `<a:tailEnd>` / `<a:headEnd>`
  の type 属性から arrow-end を判定する。
"""

from __future__ import annotations

import pathlib

# DrawingML namespace
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_LN = f"{{{_A_NS}}}ln"
_TAIL_END = f"{{{_A_NS}}}tailEnd"
_HEAD_END = f"{{{_A_NS}}}headEnd"
_SOLID_FILL = f"{{{_A_NS}}}solidFill"
_GRAD_FILL = f"{{{_A_NS}}}gradFill"
_PATT_FILL = f"{{{_A_NS}}}pattFill"
_PRST_GEOM = f"{{{_A_NS}}}prstGeom"

# 2pt ≈ 25400 EMU (1 pt = 12700 EMU). 太い矢印 threshold。
_HEAVY_LINE_EMU = 25400

# arrow-end 判定対象となる type 属性値。"none" は除外。
_ARROW_TYPES = {"triangle", "arrow", "stealth", "oval", "diamond"}

# Filled AutoShape arrows.  These do not carry an <a:headEnd> marker, so the
# preset geometry must also be inspected when checking semantic arrow layers.
_ARROW_PRESETS = {
    "rightArrow", "leftArrow", "upArrow", "downArrow",
    "leftRightArrow", "upDownArrow", "quadArrow",
    "bentArrow", "uturnArrow", "leftUpArrow", "bentUpArrow",
    "curvedRightArrow", "curvedLeftArrow", "curvedUpArrow",
    "curvedDownArrow", "stripedRightArrow", "notchedRightArrow",
    "circularArrow", "leftCircularArrow", "leftRightCircularArrow",
    "swooshArrow",
}


def _iter_shapes(shapes):
    """Yield all shapes including those inside groups (flat walk)."""
    for shape in shapes:
        # group: recurse into its .shapes
        try:
            if _is_group_shape(shape):
                yield from _iter_shapes(shape.shapes)
                continue
        except Exception:
            pass
        yield shape


def _is_group_shape(shape) -> bool:
    """Return whether *shape* is a group, comparing against the enum member.

    Matching ``str(shape.shape_type)`` against the literal ``"GROUP (6)"``
    depends on python-pptx's enum repr; if that repr ever changes, groups stop
    being recognised and their contents get audited as one opaque shape.
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except Exception:
        MSO_SHAPE_TYPE = None
    try:
        shape_type = shape.shape_type
    except Exception:
        return False
    if shape_type is None:
        return False
    if MSO_SHAPE_TYPE is not None:
        return shape_type == MSO_SHAPE_TYPE.GROUP
    return str(shape_type) == "GROUP (6)"


def _find_ln_element(shape):
    """Return the <a:ln> element for a shape, or None."""
    try:
        elem = shape._element  # type: ignore[attr-defined]
    except Exception:
        return None
    if elem is None:
        return None
    # <a:ln> lives inside <p:spPr> (or equivalent). Use a deep search to be
    # robust to cxnSp vs sp element differences.
    try:
        return elem.find(f".//{_LN}")
    except Exception:
        return None


def _has_arrow_end(ln_elem) -> bool:
    """True if the <a:ln> element has a tailEnd / headEnd with an
    arrow-like type."""
    if ln_elem is None:
        return False
    for tag in (_TAIL_END, _HEAD_END):
        end = ln_elem.find(tag)
        if end is None:
            continue
        t = end.get("type") or ""
        if t.lower() in _ARROW_TYPES:
            return True
    return False


def _line_width_emu(ln_elem) -> int:
    """Return the line width in EMU (0 if unspecified)."""
    if ln_elem is None:
        return 0
    w = ln_elem.get("w")
    if w is None:
        return 0
    try:
        return int(w)
    except (TypeError, ValueError):
        return 0


def _has_line_fill(ln_elem) -> bool:
    """True if the <a:ln> element carries an explicit solid/grad/pattern
    fill child (indicating a custom coloured / filled stroke)."""
    if ln_elem is None:
        return False
    for tag in (_SOLID_FILL, _GRAD_FILL, _PATT_FILL):
        if ln_elem.find(tag) is not None:
            return True
    return False


def _is_line_like_shape(shape) -> bool:
    """True if the shape is a connector / line-like shape.

    Uses python-pptx shape_type == LINE as the primary signal, with a
    fallback on XML tag (p:cxnSp = connector) for robustness.
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except Exception:
        MSO_SHAPE_TYPE = None

    try:
        st = shape.shape_type
    except Exception:
        st = None

    if MSO_SHAPE_TYPE is not None and st == MSO_SHAPE_TYPE.LINE:
        return True

    # Fallback: XML tag-based detection for connectors.
    try:
        tag = shape._element.tag  # type: ignore[attr-defined]
    except Exception:
        tag = ""
    if tag.endswith("}cxnSp"):
        return True

    # Artifact Tool exports zero-width/zero-height straight lines as an
    # ordinary <p:sp> AutoShape with preset geometry "line" rather than as a
    # <p:cxnSp> connector.  Treat that representation as line-like as well.
    if _preset_geometry(shape) == "line":
        return True

    return False


def _preset_geometry(shape) -> str:
    """Return DrawingML preset geometry, or an empty string."""
    try:
        elem = shape._element  # type: ignore[attr-defined]
        geom = elem.find(f".//{_PRST_GEOM}")
    except Exception:
        return ""
    return "" if geom is None else (geom.get("prst") or "")


def _is_semantic_arrow(shape) -> bool:
    """True for an arrow-ended line or a filled arrow AutoShape."""
    if _has_arrow_end(_find_ln_element(shape)):
        return True
    return _preset_geometry(shape) in _ARROW_PRESETS


def _shape_bbox_pt(shape) -> tuple[float, float, float, float] | None:
    """Return (left, top, width, height) in points, or None if unresolvable.

    A shape without an ``<a:xfrm>`` that python-pptx cannot resolve from a
    layout reports ``None`` geometry.  Such a shape has no comparable bounding
    box, so it is reported as unresolved rather than crashing the audit.
    """
    emu_per_pt = 12700.0
    try:
        values = (shape.left, shape.top, shape.width, shape.height)
    except Exception:
        return None
    if any(value is None for value in values):
        return None
    try:
        return tuple(float(value) / emu_per_pt for value in values)
    except (TypeError, ValueError):
        return None


def _bbox_intersects(
    a: tuple[float, float, float, float],
    b: tuple[float, float, float, float],
    tolerance_pt: float,
) -> bool:
    """Return True when two boxes meet, including zero-width divider lines."""
    ax0, ay0, aw, ah = a
    bx0, by0, bw, bh = b
    ax1, ay1 = ax0 + aw, ay0 + ah
    bx1, by1 = bx0 + bw, by0 + bh
    return not (
        ax1 < bx0 - tolerance_pt
        or bx1 < ax0 - tolerance_pt
        or ay1 < by0 - tolerance_pt
        or by1 < ay0 - tolerance_pt
    )


def presentation_check_arrow_layering(
    pptx_path: str,
    tolerance_pt: float = 1.0,
) -> dict:
    """Check that semantic arrows stay in front of intersecting straight lines.

    A directional arrow carries more visual meaning than a neutral divider or
    reference line.  If both must intersect, the line is background furniture
    and the arrow must have the larger z-order.  The check covers top-level
    arrow AutoShapes and explicit arrow-ended lines; grouped geometry is left
    for visual review because group-local and slide-global z-order differ.
    """
    try:
        import pptx  # type: ignore
    except ImportError:
        return {"error": "python-pptx not installed. `pip install python-pptx`"}

    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    if tolerance_pt < 0:
        return {"error": "tolerance_pt must be non-negative"}

    prs = pptx.Presentation(str(p))
    violations: list[dict] = []
    arrow_count = 0
    straight_line_count = 0
    intersection_count = 0
    skipped_group_count = 0
    unresolved: list[dict] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        arrows: list[dict] = []
        straight_lines: list[dict] = []
        for z_order, shape in enumerate(slide.shapes, start=1):
            if _is_group_shape(shape):
                skipped_group_count += 1
                continue

            is_arrow = _is_semantic_arrow(shape)
            is_line = False if is_arrow else _is_line_like_shape(shape)
            if not is_arrow and not is_line:
                continue

            bbox_pt = _shape_bbox_pt(shape)
            if bbox_pt is None:
                unresolved.append({
                    "slide": slide_index,
                    "shape": getattr(shape, "name", ""),
                    "kind": "arrow" if is_arrow else "line",
                    "reason": "shape has no resolvable position/size",
                })
                continue

            record = {
                "name": getattr(shape, "name", ""),
                "z_order": z_order,
                "bbox_pt": bbox_pt,
            }
            if is_arrow:
                arrows.append(record)
                arrow_count += 1
            else:
                straight_lines.append(record)
                straight_line_count += 1

        for arrow in arrows:
            for line in straight_lines:
                if not _bbox_intersects(
                    arrow["bbox_pt"], line["bbox_pt"], tolerance_pt
                ):
                    continue
                intersection_count += 1
                if arrow["z_order"] > line["z_order"]:
                    continue
                violations.append({
                    "slide": slide_index,
                    "issue": "arrow_behind_straight_line",
                    "arrow_name": arrow["name"],
                    "arrow_z_order": arrow["z_order"],
                    "line_name": line["name"],
                    "line_z_order": line["z_order"],
                    "arrow_bbox_pt": [round(v, 3) for v in arrow["bbox_pt"]],
                    "line_bbox_pt": [round(v, 3) for v in line["bbox_pt"]],
                })

    return {
        "file": str(p),
        "passed": len(violations) == 0,
        "slides_checked": len(prs.slides),
        "arrow_count": arrow_count,
        "straight_line_count": straight_line_count,
        "intersection_count": intersection_count,
        "violation_count": len(violations),
        "violations": violations,
        "skipped_group_count": skipped_group_count,
        "unresolved_count": len(unresolved),
        "unresolved": unresolved,
        "rule": (
            "An intersecting directional arrow must be in front of a neutral "
            "straight divider/reference line."
        ),
    }


def presentation_arrow_usage(pptx_path: str) -> dict:
    """矢印 shape (line connector with arrow) の過剰使用検出。

    宮野 S10 基準: 1 スライドに矢印 3 本以下、太線 / 塗りつぶしの矢印は NG。
    画面が騒がしくなり、位置関係だけで流れを示す方が軽い。

    Args:
        pptx_path: .pptx のパス

    Returns:
        Plan B 標準 dict: score / score_max / comments / observations /
        hint / source
    """
    try:
        import pptx  # type: ignore
    except ImportError:
        return {"error": "python-pptx not installed. `pip install python-pptx`"}

    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}

    prs = pptx.Presentation(str(p))
    total_slides = len(prs.slides)

    per_slide_arrow_count: list[int] = []
    heavy_arrows: list[dict] = []

    for slide_idx, slide in enumerate(prs.slides):
        count = 0
        for shape in _iter_shapes(slide.shapes):
            if not _is_line_like_shape(shape):
                continue

            ln_elem = _find_ln_element(shape)
            # arrow-end があるものだけカウント。連結線でも端点装飾がなければ
            # 「矢印」とは見なさない。
            # ただし cxnSp (connector) は styleRef から arrow-end を継承する
            # ことがあり、XML には現れない。保守的に: arrow-end が明示され
            # ていない connector も「矢印扱い」でカウントする
            # (テンプレート経由の矢印を見逃すよりマシ)。
            is_arrow = _has_arrow_end(ln_elem)
            try:
                tag = shape._element.tag  # type: ignore[attr-defined]
            except Exception:
                tag = ""
            if not is_arrow and tag.endswith("}cxnSp"):
                is_arrow = True

            if not is_arrow:
                continue

            count += 1

            width_emu = _line_width_emu(ln_elem)
            has_fill = _has_line_fill(ln_elem)
            if width_emu > _HEAVY_LINE_EMU or has_fill:
                heavy_arrows.append({
                    "slide_idx": slide_idx,
                    "line_weight": width_emu,
                    "has_fill": bool(has_fill),
                })

        per_slide_arrow_count.append(count)

    total_arrow_count = sum(per_slide_arrow_count)
    max_per_slide = max(per_slide_arrow_count) if per_slide_arrow_count else 0
    overused_slides = [
        i for i, c in enumerate(per_slide_arrow_count) if c > 3
    ]

    observations = {
        "total_slides": total_slides,
        "per_slide_arrow_count": per_slide_arrow_count,
        "overused_slides": overused_slides,
        "heavy_arrows": heavy_arrows,
        "total_arrow_count": total_arrow_count,
        "max_per_slide": max_per_slide,
    }

    # ---- score ----
    if total_arrow_count == 0:
        score: float = 10.0
    else:
        penalty = 0.0
        # -2 per overused slide, capped at -6
        penalty += min(len(overused_slides) * 2, 6)
        # -1 per heavy arrow, capped at -3
        penalty += min(len(heavy_arrows) * 1, 3)
        score = max(0.0, round(10.0 - penalty, 1))

    # ---- comments (PRIMARY) ----
    comments: list[str] = []

    if total_arrow_count == 0:
        comments.append(
            "矢印 shape 0 本。宮野 S10『矢印は多用しない』基準は自動満たされる。"
        )
        comments.append(
            "ただし流れ説明で矢印が必要な場合は細線 1pt matte 黒 1-2 本に抑える。"
        )
    else:
        if overused_slides:
            # 1-indexed for readability
            idx_1 = [i + 1 for i in overused_slides]
            comments.append(
                f"矢印過多 slide: {idx_1}。宮野 S10: 1 slide に 3 個以下に削る。"
                "矢印なしで位置・隣接だけで流れを示す方が軽い。"
            )

        if heavy_arrows:
            comments.append(
                f"太線 or 塗りつぶし矢印 {len(heavy_arrows)} 件。"
                "画面が騒々しい。細線 1pt matte 黒に統一。"
            )

        if not overused_slides and not heavy_arrows:
            comments.append(
                f"矢印使用は {total_arrow_count} 本で適正範囲。"
                f"最大 slide 当たり {max_per_slide} 本、太線・塗りつぶしなし。"
            )

        if max_per_slide > 3 and len(overused_slides) >= 2:
            comments.append(
                "複数 slide で 3 本超過。矢印は流れの『補助』であり"
                "『主役』ではない (宮野 S10)。表・レイアウトでの代替検討を。"
            )

    # 2-5 個に制限 & 不足時 padding
    while len(comments) < 2:
        comments.append(
            "score は粗い温度計。最終判断は comments と observations で確認。"
        )
    comments = comments[:5]

    hint = (
        "connector の arrow-end 判定は XML 経由で粗い "
        "(python-pptx は完全 API 未提供)。line width は EMU 単位、"
        "2pt ≈ 25400 EMU。"
    )

    source = "宮野『研究発表のためのスライドデザイン』S10 (矢印過剰 NG)"

    return {
        "score": float(score),
        "score_max": 10,
        "comments": comments,
        "observations": observations,
        "hint": hint,
        "source": source,
    }
