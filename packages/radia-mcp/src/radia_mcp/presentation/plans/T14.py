"""Plan B T14: Title-body alignment check
(presentation_title_body_alignment_check).

各 slide で **タイトルの対象・観点が本文内容と対応しているか** を診断。
`check_slide_title_specificity` は title 単体の具体性・簡潔性を判定し、
本ツールは **title vs body の内容対応** を見る。

検出する 3 軸 per slide:
  1. **title_body_overlap** — title token と body token の overlap
  2. **title_summarizes_body** — title の対象・観点が body と一致するか
  3. **title_alone_understandable** — title 単独で意味が取れるか (5W1H)

設計方針:
- 各 slide で title と body を抽出
- token overlap + 数値/固有名 一致 + 5W1H signals
- score は per-slide 平均
- PRIMARY OUTPUT は misaligned_slides list (修正対象)
"""

from __future__ import annotations

import pathlib
import re


def _iter_all_shapes(shape_collection):
    try:
        from radia_mcp.presentation.tools import _walk_shapes
        yield from _walk_shapes(shape_collection)
    except Exception:
        for s in shape_collection:
            yield s


def _get_title(slide) -> str:
    try:
        from radia_mcp.presentation.tools import _slide_title
        return _slide_title(slide)
    except Exception:
        return ""


def _get_body_text(slide) -> str:
    """Return visible body plus notes that explain image-only evidence."""
    parts: list[str] = []
    title = _get_title(slide).strip()
    for shape in _iter_all_shapes(slide.shapes):
        try:
            if not shape.has_text_frame:
                continue
            # Skip title placeholder
            try:
                pf = shape.placeholder_format
                if pf is not None and pf.idx == 0:
                    continue
            except Exception:
                pass
            txt = shape.text_frame.text or ""
            # Skip if it IS the title (fallback case)
            if txt.strip() == title:
                continue
            parts.append(txt)
        except Exception:
            continue
    try:
        if slide.has_notes_slide:
            note_text = slide.notes_slide.notes_text_frame.text or ""
            note_text = re.split(r"\n\s*\[Sources\]\s*", note_text, maxsplit=1)[0]
            if note_text.strip():
                parts.append(note_text)
    except Exception:
        pass
    return "\n".join(parts)


def _extract_tokens(text: str) -> set[str]:
    tokens: set[str] = set()
    for m in re.finditer(r"[一-鿿]{2,}", text):
        tokens.add(m.group(0))
    for m in re.finditer(r"[ァ-ヴー]{3,}", text):
        tokens.add(m.group(0))
    for m in re.finditer(r"\b[A-Z][A-Za-z0-9]{2,}\b", text):
        tokens.add(m.group(0))
    return tokens


def _extract_numbers_with_units(text: str) -> set[str]:
    matches = re.findall(
        r"\d+(?:\.\d+)?\s*(?:%|％|x|×|倍|kHz|MHz|GHz|Hz|GB|MB|nm|μm|mm|cm|m|s)",
        text,
    )
    return set(matches)


def _has_5w1h_signal(title: str) -> dict:
    """Title が 5W1H 要素 (Who/What/When/Where/Why/How) を含むか rough 判定。"""
    return {
        "who": bool(re.search(r"(?:we|our|本|申請者|研究室|チーム|システム|手法|アルゴリズム)",
                                title, re.IGNORECASE)),
        "what": len(_extract_tokens(title)) >= 1,
        "why": bool(re.search(r"(?:なぜ|why|理由|目的|because|to\s+\w)",
                                title, re.IGNORECASE)),
        "how": bool(re.search(r"(?:による|を用いた|with|using|via|method|approach)",
                                title, re.IGNORECASE)),
        "quantified": bool(re.search(r"\d", title)),
    }


def _matched_title_tokens(title_tokens: set[str], body_tokens: set[str]) -> set[str]:
    """Match Japanese compound variants without requiring exact token identity."""
    matched: set[str] = set()
    for title_token in title_tokens:
        title_folded = title_token.casefold()
        for body_token in body_tokens:
            body_folded = body_token.casefold()
            if title_folded == body_folded:
                matched.add(title_token)
                break
            if min(len(title_folded), len(body_folded)) >= 3 and (
                title_folded in body_folded or body_folded in title_folded
            ):
                matched.add(title_token)
                break
    return matched


def presentation_title_body_alignment_check(pptx_path: str) -> dict:
    """Title の対象・観点が body の内容と対応しているかを診断。

    Args:
        pptx_path: .pptx ファイルパス

    Returns:
        dict: score / per_slide / misaligned_slides / comments
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))

    from .._deck_sections import classify_deck_sections, is_closing_title

    all_slides = list(prs.slides)
    sections = classify_deck_sections(all_slides, _get_title)
    backup_slide_numbers = set(sections["backup_slide_numbers"])

    per_slide: list[dict] = []
    misaligned: list[dict] = []

    for i, slide in enumerate(all_slides, 1):
        title = _get_title(slide).strip()
        body = _get_body_text(slide).strip()

        if i == 1 or i in backup_slide_numbers or is_closing_title(title):
            per_slide.append({
                "slide_no": i,
                "title": title,
                "body_chars": len(body),
                "skip_reason": (
                    "backup_slide" if i in backup_slide_numbers
                    else "structural_slide"
                ),
                "alignment_score": None,
            })
            continue

        # Skip if no title or no body (titlecard slides etc.)
        if not title or not body:
            per_slide.append({
                "slide_no": i,
                "title": title,
                "body_chars": len(body),
                "skip_reason": "no_title" if not title else "no_body",
                "alignment_score": None,
            })
            continue

        title_tokens = _extract_tokens(title)
        body_tokens = _extract_tokens(body)
        title_nums = _extract_numbers_with_units(title)
        body_nums = _extract_numbers_with_units(body)

        # 1. title_body_overlap (Jaccard-like)
        if title_tokens and body_tokens:
            common = _matched_title_tokens(title_tokens, body_tokens)
            overlap = len(common) / len(title_tokens)
        else:
            overlap = 0
            common = set()

        # 2. title_summarizes_body
        # Title の number が body にも出る → strong alignment
        # Title token の 50% 以上が body にもある → moderate
        nums_in_body = title_nums & body_nums
        if title_nums and not nums_in_body:
            summarize_status = "TITLE_HAS_UNCLAIMED_NUMBER"
        elif overlap >= 0.5:
            summarize_status = "STRONG"
        elif overlap >= 0.25:
            summarize_status = "MODERATE"
        else:
            summarize_status = "WEAK"

        # 3. title_alone_understandable (5W1H)
        wh = _has_5w1h_signal(title)
        n_5w1h = sum(1 for v in wh.values() if v)
        title_understandable = n_5w1h >= 2  # at least 2 of 5W1H

        # Per-slide alignment score (0-10)
        slide_score = (
            overlap * 4
            + (3 if summarize_status == "STRONG"
               else 2 if summarize_status == "MODERATE" else 0)
            + (3 if title_understandable else 1.5)
        )
        slide_score = round(min(10.0, slide_score), 1)

        per_slide.append({
            "slide_no": i,
            "title": title[:80],
            "body_chars": len(body),
            "title_tokens": len(title_tokens),
            "body_tokens": len(body_tokens),
            "common_tokens": sorted(common)[:5],
            "overlap_ratio": round(overlap, 3),
            "title_numbers": sorted(title_nums),
            "numbers_in_body": sorted(nums_in_body),
            "summarize_status": summarize_status,
            "title_understandable_alone": title_understandable,
            "5w1h_satisfied": n_5w1h,
            "alignment_score": slide_score,
        })

        if slide_score < 5:
            misaligned.append({
                "slide_no": i,
                "title": title[:80],
                "alignment_score": slide_score,
                "reason": (
                    "title の対象・観点が body の内容と対応していない"
                    if summarize_status == "WEAK"
                    else "title 単独で意味が取れない"
                    if not title_understandable
                    else "title に未裏付けの数値"
                    if summarize_status == "TITLE_HAS_UNCLAIMED_NUMBER"
                    else "整合性 weak"
                ),
            })

    # Overall score (averaging non-skipped)
    valid_scores = [s["alignment_score"] for s in per_slide
                     if s["alignment_score"] is not None]
    if valid_scores:
        score = round(sum(valid_scores) / len(valid_scores), 1)
    else:
        score = 0.0

    # ---- comments ----
    comments: list[str] = []

    n_total = len(per_slide)
    n_evaluated = len(valid_scores)
    n_skip = n_total - n_evaluated
    n_misaligned = len(misaligned)

    comments.append(
        f"スライド {n_total} 枚 | 評価対象 {n_evaluated} 枚 (no title/body skip {n_skip}) "
        f"| misaligned {n_misaligned} 枚 (score < 5) | overall {score}/10"
    )

    if n_misaligned >= 1:
        comments.append(
            f"⚠ {n_misaligned} 枚で title-body alignment が weak。"
            "聴衆は title から slide の対象・観点を取るので、"
            "title が body と一致しないと混乱を生む。"
        )
        for m in misaligned[:5]:
            comments.append(
                f"  • slide {m['slide_no']} [{m['alignment_score']}/10]: "
                f"{m['title']} — {m['reason']}"
            )

    # 数値だけ title にあって body に裏付け無いケース
    unclaimed_num_slides = [
        s for s in per_slide
        if s.get("summarize_status") == "TITLE_HAS_UNCLAIMED_NUMBER"
    ]
    if unclaimed_num_slides:
        comments.append(
            f"⚠ {len(unclaimed_num_slides)} 枚で title に数値あるが "
            "body に裏付けが無い。title 過大広告 (paper T10 三角形と同根)。"
        )

    if n_skip > n_total * 0.3:
        comments.append(
            f"ℹ {n_skip} 枚で title または body が無い。"
            "title-only / image-only slides ならば OK だが、本文が "
            "あるべき slide で body 検出失敗ならば再確認。"
        )

    if score >= 8:
        comments.append(
            f"Title-body alignment {score}/10 — title の対象・観点が body と "
            "良く対応している。聴衆の認知負荷低い。"
        )

    comments = comments[:8]

    hint = (
        "check_slide_title_specificity は title 単体の具体性・簡潔性を check、"
        "T14 は title vs body 内容対応を見る点で相補。"
        "alignment_score 5+ 推奨、低い slide は title 書き直し。"
        "title だけ見て slide の核が分かるかが reviewer 視点。"
    )

    return {
        "score": score,
        "score_max": 10,
        "n_total_slides": n_total,
        "main_slide_count": sections["main_slide_count"],
        "backup_slide_count": sections["backup_slide_count"],
        "backup_slide_numbers": sections["backup_slide_numbers"],
        "n_evaluated": n_evaluated,
        "n_misaligned": n_misaligned,
        "misaligned_slides": misaligned,
        "per_slide": per_slide[:30],
        "comments": comments,
        "hint": hint,
        "source": (
            "Reynolds『Presentation Zen』Ch.6 (one slide = one message) + "
            "宮野『研究発表のためのスライドデザイン』S4 (slide structure) + "
            "Tufte『Cognitive Style of PowerPoint』(2003) — title-body coupling "
            "の重要性。check_slide_title_specificity と相補。"
            "(presentation 2026-04 採用)。"
        ),
    }
