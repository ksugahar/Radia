"""Plan B T1: 冒頭 Hook 強度診断 (presentation_opening_hook_strength).

pptx の最初 2 枚のテキストから Hook 強度を測る粗い温度計。

skill.md T1 (冒頭 Hook 強度) / 宮野『研究発表のためのスライドデザイン』S1
(作成前に全体構成を検討) / Reynolds "Presentation Zen" の opening 原則 /
木下『理科系の作文技術』原則 2 (重点先行) に基づく。

設計方針:
- score は 0-10 の粗い温度計。9 以上は誤差範囲、6 以下は改善サイン。
- PRIMARY OUTPUT は comments (人間レビュアー風の助言文)。
- observations は raw measurement。自動化した numeric な正誤判定には使わない。
- 発話 (speaker note) は含まない — 口頭の Hook は別枠で評価する。
"""

from __future__ import annotations

import pathlib
import re


# ------------------------------------------------------------------
# 冒頭 Hook 判定用パターン
# ------------------------------------------------------------------
_GENERIC_OPENER_PATTERN = re.compile(
    r"^\s*(?:本日は|本発表は|この発表|I\s+will|Today|We\s+will\s+present)",
    re.IGNORECASE,
)

_DIGIT_UNIT_PATTERN = re.compile(
    r"\d+\s*(?:kHz|MHz|V|A|年|%|倍|兆円|億円)"
)
_DIGIT_STANDALONE_PATTERN = re.compile(r"\d{3,}")

_CONTRAST_PATTERNS = [
    r"しかし",
    r"だが",
    r"従来",
    r"limit",
    r"versus",
    r"対して",
    r"超える",
    r"以前",
    r"現状",
    r"unlike",
]

# 現場・組織・人物語 (locus)
_SCENE_NOUN_PATTERNS = [
    r"研究室",
    r"大学",
    r"工場",
    r"現場",
    r"案件",
    r"[社教先][長授生]",  # 社長/教授/先生 等
    r"patient",
    r"user",
]

_QUESTION_OPENER_BODY_PATTERNS = [
    r"Why\b",
    r"なぜ",
    r"How do we",
    r"ご存じだろうか",
    r"did you know",
]


# ------------------------------------------------------------------
# pptx ヘルパー (tools.py の _walk_shapes / _slide_title と互換)
# ------------------------------------------------------------------
def _extract_slide_title_and_body(slide) -> tuple[str, str]:
    """Return (title, body) text for a single slide.

    presentation/tools.py の _walk_shapes / _slide_title が存在すればそれを
    流用し、同じ group-shape 走査ロジックで一貫性を確保する。
    """
    try:
        from radia_mcp.presentation.tools import (
            _slide_title,
            _walk_shapes,
        )
    except Exception:
        _slide_title = None  # type: ignore
        _walk_shapes = None  # type: ignore

    # ---- title ----
    if _slide_title is not None:
        try:
            title = _slide_title(slide) or ""
        except Exception:
            title = ""
    else:
        title = ""
        try:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                    title = shape.text_frame.text.splitlines()[0]
                    break
        except Exception:
            pass

    # ---- body: concat all non-title shapes' text_frame paragraphs ----
    body_parts: list[str] = []
    shape_iter = None
    if _walk_shapes is not None:
        try:
            shape_iter = list(_walk_shapes(slide.shapes))
        except Exception:
            shape_iter = None
    if shape_iter is None:
        try:
            shape_iter = list(slide.shapes)
        except Exception:
            shape_iter = []

    for shape in shape_iter:
        try:
            if not getattr(shape, "has_text_frame", False):
                continue
            # Skip the title placeholder (idx == 0)
            try:
                pf = shape.placeholder_format
                if pf is not None and pf.idx == 0:
                    continue
            except Exception:
                pass
            for para in shape.text_frame.paragraphs:
                t = para.text
                if t and t.strip():
                    body_parts.append(t)
        except Exception:
            continue
    body = "\n".join(body_parts)
    return title, body


# ------------------------------------------------------------------
# Main
# ------------------------------------------------------------------
def presentation_opening_hook_strength(pptx_path: str) -> dict:
    """First 2 slides' text hook 強度診断。

    Detects: "発表します" opener (弱) / concrete numbers / comparison /
    before-after / question opening。宮野 S1 + Reynolds Zen + 木下 重点先行。

    引数:
        pptx_path: .pptx のパス

    返り値:
        dict: score / score_max / comments / observations / hint / source
    """
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {
            "score": 0,
            "score_max": 10,
            "comments": [
                f"pptx が見つからない: {pptx_path}",
                "ファイル path を確認して再実行すること。",
            ],
            "observations": {"error": f"file not found: {pptx_path}"},
            "hint": (
                "pptx の最初 2 枚のテキストから粗い Hook 強度を測る。"
                "発話 (speaker note) は含まない — 口頭の Hook は別。"
                "skill.md 宮野 S1 + 演習で判定。"
            ),
            "source": (
                "宮野『研究発表のためのスライドデザイン』S1 / "
                "Reynolds Zen / 木下 原則 2 (重点先行)"
            ),
        }

    try:
        import pptx as _pptx  # type: ignore
    except ImportError:
        return {
            "score": 0,
            "score_max": 10,
            "comments": [
                "python-pptx 未インストール。`pip install python-pptx` 後に再実行。",
            ],
            "observations": {"error": "python-pptx not installed"},
            "hint": (
                "pptx の最初 2 枚のテキストから粗い Hook 強度を測る。"
                "発話 (speaker note) は含まない — 口頭の Hook は別。"
                "skill.md 宮野 S1 + 演習で判定。"
            ),
            "source": (
                "宮野『研究発表のためのスライドデザイン』S1 / "
                "Reynolds Zen / 木下 原則 2 (重点先行)"
            ),
        }

    try:
        prs = _pptx.Presentation(str(p))
    except Exception as exc:
        return {
            "score": 0,
            "score_max": 10,
            "comments": [f"pptx 読み込み失敗: {exc}"],
            "observations": {"error": f"failed to open: {exc}"},
            "hint": (
                "pptx の最初 2 枚のテキストから粗い Hook 強度を測る。"
                "発話 (speaker note) は含まない — 口頭の Hook は別。"
                "skill.md 宮野 S1 + 演習で判定。"
            ),
            "source": (
                "宮野『研究発表のためのスライドデザイン』S1 / "
                "Reynolds Zen / 木下 原則 2 (重点先行)"
            ),
        }

    slides = list(prs.slides)
    n_total = len(slides)

    first_title, first_body = "", ""
    second_title, second_body = "", ""
    if n_total >= 1:
        first_title, first_body = _extract_slide_title_and_body(slides[0])
    if n_total >= 2:
        second_title, second_body = _extract_slide_title_and_body(slides[1])

    slides_analyzed = min(2, n_total)
    combined_first_two = "\n".join(
        [first_title, first_body, second_title, second_body]
    )

    # ---- observation flags ----
    has_generic_opener = bool(_GENERIC_OPENER_PATTERN.search(first_body))

    has_digit_in_opening = bool(
        _DIGIT_UNIT_PATTERN.search(combined_first_two)
        or _DIGIT_STANDALONE_PATTERN.search(combined_first_two)
    )

    has_contrast_word = any(
        re.search(p, combined_first_two, re.IGNORECASE)
        for p in _CONTRAST_PATTERNS
    )

    has_scene_noun = any(
        re.search(p, combined_first_two, re.IGNORECASE)
        for p in _SCENE_NOUN_PATTERNS
    )

    # question opener: first slide title ends with ? or ? OR body has marker
    title_stripped = first_title.rstrip()
    title_ends_q = title_stripped.endswith("?") or title_stripped.endswith("？")
    body_has_q = any(
        re.search(p, first_body, re.IGNORECASE)
        for p in _QUESTION_OPENER_BODY_PATTERNS
    )
    has_question_opener = bool(title_ends_q or body_has_q)

    observations = {
        "slides_analyzed": slides_analyzed,
        "first_slide_title": first_title,
        "first_slide_body": first_body,
        "second_slide_title": second_title,
        "second_slide_body": second_body,
        "has_generic_opener": has_generic_opener,
        "has_digit_in_opening": has_digit_in_opening,
        "has_contrast_word": has_contrast_word,
        "has_scene_noun": has_scene_noun,
        "has_question_opener": has_question_opener,
    }

    # ---- score ----
    score = 0
    if has_digit_in_opening:
        score += 3
    if has_contrast_word:
        score += 2
    if has_scene_noun:
        score += 2
    if has_question_opener:
        score += 2
    if has_generic_opener:
        score -= 3
    # baseline: 悪い flag がひとつも無ければ +1
    if not has_generic_opener:
        if not (has_digit_in_opening or has_contrast_word
                or has_scene_noun or has_question_opener):
            score += 1
    score = max(0, min(10, score))

    # ---- comments ----
    comments: list[str] = []

    if has_generic_opener:
        comments.append(
            "「本日は〜発表します」は教科書的な弱い opening。Hook は 3-5 秒で"
            "掴む勝負。最初のスライドを「年間 120 億円の商用 CAE に依存して"
            "きた。本研究で置換する」式の事実 + 対立で再構成 (Reynolds)。"
        )
    if not has_digit_in_opening:
        comments.append(
            "冒頭 2 枚に数値が無い。「500 kHz」「14 年」「年間 120 億円」等"
            "の具体数値で audience の脳を掴む (宮野 S1)。"
        )
    if not has_contrast_word:
        comments.append(
            "冒頭 2 枚に対立構造がない。「従来 X → 本研究 Y」「しかし Z」で "
            "before/after を視覚化。"
        )
    if not has_scene_noun:
        comments.append(
            "冒頭 2 枚に現場・組織の語が無い。「○○研究室」「工場の現場で」"
            "「患者が」と locus を置くと audience が情景を描きやすい。"
        )

    # all good path
    if (has_digit_in_opening and has_contrast_word and has_scene_noun
            and not has_generic_opener):
        comments.append(
            "冒頭 hook 強度十分。数値・対立・現場のいずれも揃う。"
        )

    # fallback — ensure at least 2 comments
    if len(comments) < 2:
        comments.append(
            "score は粗い指標。skill.md T1 の判定基準で最終確認を推奨。"
        )
    # cap at 5
    comments = comments[:5]

    hint = (
        "pptx の最初 2 枚のテキストから粗い Hook 強度を測る。"
        "発話 (speaker note) は含まない — 口頭の Hook は別。"
        "skill.md 宮野 S1 + 演習で判定。"
    )

    source = (
        "宮野『研究発表のためのスライドデザイン』S1 / "
        "Reynolds Zen / 木下 原則 2 (重点先行)"
    )

    return {
        "score": float(score),
        "score_max": 10,
        "comments": comments,
        "observations": observations,
        "hint": hint,
        "source": source,
    }
