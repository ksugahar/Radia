"""Plan B T8: deck 全体のフォント統一度診断
(presentation_font_consistency).

宮野『研究発表のためのスライドデザイン』S11 に基づく粗い温度計。
deck 内で使用されているフォントファミリー (run.font.name) の数を集計し、
1-2 種 (理想: 日本語用 + 英数字用) を超えて混在していないか判定する。

設計方針:
- score は 0-10 の粗い温度計。9 以上は誤差範囲、6 以下は改善サイン。
- PRIMARY OUTPUT は comments (人間レビュアー風の助言文)。
- observations は raw measurement。numeric な正誤判定には使わない。
- run.font.name が None の場合は "(default)" 扱い。
  layout / master 経由の継承は追わない (= 見た目上の実フォントとは乖離する
  可能性がある — それは hint で明示)。
"""

from __future__ import annotations

import pathlib
from collections import Counter


# ------------------------------------------------------------------
# Constants
# ------------------------------------------------------------------
_DEFAULT_LABEL = "(default)"


# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------
def _iter_shapes(slide):
    """Yield all non-group shapes from a slide, recursing into groups.

    tools.py の _walk_shapes があれば再利用し、group-shape 走査の一貫性を
    確保する。無ければ素朴な非再帰 iteration にフォールバック。
    """
    try:
        from radia_mcp.presentation.tools import _walk_shapes
    except Exception:
        _walk_shapes = None  # type: ignore

    if _walk_shapes is not None:
        try:
            yield from _walk_shapes(slide.shapes)
            return
        except Exception:
            pass
    # Fallback
    try:
        for s in slide.shapes:
            yield s
    except Exception:
        return


def _collect_font_names(slide) -> list[str]:
    """Return a list of font-name strings per run on this slide.

    None / 空文字列は "(default)" として記録する。
    """
    names: list[str] = []
    for shape in _iter_shapes(slide):
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            tf = shape.text_frame
        except Exception:
            continue
        try:
            paragraphs = list(tf.paragraphs)
        except Exception:
            continue
        for para in paragraphs:
            try:
                runs = list(para.runs)
            except Exception:
                continue
            for run in runs:
                try:
                    name = run.font.name
                except Exception:
                    name = None
                if name is None or str(name).strip() == "":
                    names.append(_DEFAULT_LABEL)
                else:
                    names.append(str(name))
    return names


# ------------------------------------------------------------------
# Main
# ------------------------------------------------------------------
def presentation_font_consistency(pptx_path: str) -> dict:
    """deck 内で使用されているフォントファミリーの数を集計。

    宮野 S11: フォントは 1-2 種に限る (日本語用 + 英数字用) が理想。
    4 種以上はコピペ起源の divergence のサイン。

    引数:
        pptx_path: .pptx のパス

    返り値:
        dict: score / score_max / comments / observations / hint / source
    """
    try:
        import pptx  # type: ignore
    except ImportError:
        return {"error": "python-pptx not installed. `pip install python-pptx`"}

    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}

    prs = pptx.Presentation(str(p))
    total_slides = len(prs.slides)

    counter: Counter[str] = Counter()
    for slide in prs.slides:
        for name in _collect_font_names(slide):
            counter[name] += 1

    font_family_counts: dict[str, int] = dict(counter)
    distinct_font_count = len(font_family_counts)

    # dominant = top 3 by count (stable ordering: count desc, then name asc)
    ranked = sorted(
        font_family_counts.items(),
        key=lambda kv: (-kv[1], kv[0]),
    )
    dominant_fonts: list[str] = [name for name, _c in ranked[:3]]

    # rare = count < 3
    rare_fonts: list[str] = sorted(
        [name for name, c in font_family_counts.items() if c < 3]
    )

    observations = {
        "total_slides": total_slides,
        "font_family_counts": font_family_counts,
        "distinct_font_count": distinct_font_count,
        "dominant_fonts": dominant_fonts,
        "rare_fonts": rare_fonts,
    }

    # ---- score ----
    if distinct_font_count <= 2:
        score: float = 10.0
    elif distinct_font_count == 3:
        score = 7.0
    elif distinct_font_count == 4:
        score = 4.0
    else:  # >= 5
        score = 0.0

    # -1 penalty if rare_fonts は dominant と重ならない余剰 font (accent じゃなく divergence)
    # 指示: "If rare_fonts are a subset of dominant (e.g., accent use),
    #       -1 penalty but not below 0"
    if rare_fonts:
        dominant_set = set(dominant_fonts)
        if set(rare_fonts).issubset(dominant_set):
            score = max(0.0, score - 1.0)

    # ---- comments ----
    comments: list[str] = []

    families_str = ", ".join(dominant_fonts) if dominant_fonts else "(none)"

    if distinct_font_count >= 4:
        comments.append(
            f"font {distinct_font_count} 種類混在 ({families_str})。"
            "宮野 S11 原則 (1-2 種に統一) から逸脱。"
            "全 slide master で日本語 + 英数字 1 組にする設定で対応可能。"
        )

    if rare_fonts:
        rare_str = ", ".join(rare_fonts)
        comments.append(
            f"1-2 回しか出ない font が {len(rare_fonts)} 件 ({rare_str})。"
            "コピペ起源の可能性が高い — paste-from-text コマンドでデフォルト統一。"
        )

    if distinct_font_count <= 2 and distinct_font_count > 0:
        comments.append(
            f"font 統一済み ({families_str})。宮野 S11 基準を満たす。"
        )

    if distinct_font_count == 3:
        comments.append(
            f"font 3 種 ({families_str})。宮野 S11 の理想 (1-2 種) からは"
            "1 種超過。日本語用 + 英数字用 + 強調用 の構成なら許容範囲だが、"
            "master で統一すれば余剰 font を排除できる。"
        )

    if total_slides == 0:
        comments.append("empty deck — スライドが 1 枚も無いためフォント診断は無意味。")

    if distinct_font_count == 0 and total_slides > 0:
        comments.append(
            "テキスト run が 1 件も検出できず font 集計不能。"
            "画像のみの deck もしくは text_frame が空の可能性。"
        )

    # フォールバック: 最低 2 コメント
    if len(comments) < 2:
        comments.append(
            "run.font.name は layout / master 継承を辿らないため、"
            "見た目の実フォントと乖離する場合がある (hint 参照)。"
        )
        if len(comments) < 2:
            comments.append(
                "score は粗い指標。skill.md S11 の判定基準で最終確認を推奨。"
            )

    # 最大 5 個
    comments = comments[:5]

    hint = (
        "run.font.name が None の場合は (default) 扱い — "
        "layout 経由の継承は追わない。"
    )

    source = "宮野『研究発表のためのスライドデザイン』S11 (フォント統一)"

    return {
        "score": float(score),
        "score_max": 10,
        "comments": comments,
        "observations": observations,
        "hint": hint,
        "source": source,
    }
