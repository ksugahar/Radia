"""Plan B T21: Figure slide compliance (理系プレゼン特化)
(presentation_figure_slide_compliance).

理系プレゼンの figure slide には論文と同じ厳格性が必要:
  軸ラベル / 単位 / 凡例 / caption / 出典 / scale。

corporate プレゼンでは省略可能だが、理系学会では これらが無いと
reviewer の質問が confounded になる。

検出する 5 軸 per figure slide:
  1. **axis_labels** — x/y 軸ラベル文字列の存在
  2. **units_specified** — 単位 [m/s], [Hz] 等の明記
  3. **legend_or_label** — 凡例 or 直接ラベル (T19 と相補)
  4. **caption_or_title** — figure title (slide title とは別)
  5. **source_citation** — 出典 [N] / (Author 2023) / 自作明示

設計方針:
- figure / chart / image を含む slide を対象
- 各 slide で text を集めて 5 軸 keyword detection
- chart 内部の axis title は python-pptx で取得試行
- score は per-figure-slide 平均
"""

from __future__ import annotations

import pathlib
import re


def _iter_all_shapes(shape_collection):
    try:
        from radia_mcp.presentation.tools import _walk_shapes
        yield from _walk_shapes(shape_collection)
    except Exception:
        for s in shape_collection:
            yield s


def _get_title(slide) -> str:
    try:
        from radia_mcp.presentation.tools import _slide_title
        return _slide_title(slide)
    except Exception:
        return ""


def _slide_has_figure(slide) -> bool:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        MSO_SHAPE_TYPE = None
    for shape in _iter_all_shapes(slide.shapes):
        try:
            if MSO_SHAPE_TYPE is not None and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True
        except Exception:
            pass
        if getattr(shape, "has_chart", False):
            try:
                if shape.has_chart:
                    return True
            except Exception:
                pass
    return False


def _chart_axis_titles(slide) -> list[str]:
    titles: list[str] = []
    for shape in _iter_all_shapes(slide.shapes):
        if not getattr(shape, "has_chart", False):
            continue
        try:
            if not shape.has_chart:
                continue
            chart = shape.chart
            for axis_attr in ("category_axis", "value_axis"):
                try:
                    ax = getattr(chart, axis_attr, None)
                    if ax is not None and ax.has_title:
                        t = ax.axis_title.text_frame.text
                        if t.strip():
                            titles.append(t.strip())
                except Exception:
                    continue
        except Exception:
            continue
    return titles


_AXIS_LABEL_PATTERNS = [
    r"(?:横軸|縦軸|x軸|y軸|x-axis|y-axis|abscissa|ordinate)",
    r"\([A-Za-zα-ωΑ-Ω][^()]{0,12}\)\s*(?:の|を|axis)",
    r"^\s*[A-Za-zα-ωΑ-Ω][^\n]{0,30}\[\w",  # Label [unit]
]

_UNIT_PATTERNS = [
    r"\[\s*[A-Za-zα-ω/·\d\^\-]{1,15}\s*\]",  # [m/s] [Hz]
    r"\b(?:Hz|kHz|MHz|GHz|m/s|kg|N|J|W|Pa|V|A|°C|K|mol|°|deg|rad)\b",
    r"単位\s*[:：]?\s*\S",
]

_LEGEND_OR_LABEL_PATTERNS = [
    r"(?:legend|凡例)",
    r"(?:line|symbol)\s*[:：]\s*\w",
    r":\s*[a-zA-Z]+(?:\s+vs\.?\s+[a-zA-Z])",  # A vs B
]

_FIGURE_CAPTION_PATTERNS = [
    r"(?:Fig(?:ure)?|図)\s*\.?\s*\d+",
    r"(?:caption|キャプション)",
]

_SOURCE_PATTERNS = [
    r"\[\d+\]",
    r"\(.*?\d{4}\)",
    r"出典", r"source", r"adapted\s+from",
    r"(?:本研究|自作|originally)",
]


def _check_axis_with_chart_titles(text: str, chart_titles: list[str]) -> bool:
    if any(t.strip() for t in chart_titles):
        return True
    for pat in _AXIS_LABEL_PATTERNS:
        if re.search(pat, text, re.IGNORECASE):
            return True
    return False


def _check_pattern(text: str, patterns: list[str]) -> bool:
    return any(re.search(p, text, re.IGNORECASE) for p in patterns)


def presentation_figure_slide_compliance(pptx_path: str) -> dict:
    """Figure / chart slide の理系プレゼン compliance を診断。

    Args:
        pptx_path: .pptx file

    Returns:
        dict: score / figure_slides / per_slide / comments
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))

    figure_slides: list[dict] = []
    n_total = 0

    for i, slide in enumerate(prs.slides, 1):
        n_total += 1
        if not _slide_has_figure(slide):
            continue

        full_text = ""
        for shape in _iter_all_shapes(slide.shapes):
            try:
                if shape.has_text_frame:
                    full_text += "\n" + (shape.text_frame.text or "")
            except Exception:
                continue

        chart_axis_titles = _chart_axis_titles(slide)

        has_axis = _check_axis_with_chart_titles(
            full_text, chart_axis_titles
        )
        has_unit = _check_pattern(full_text, _UNIT_PATTERNS)
        has_legend = _check_pattern(full_text, _LEGEND_OR_LABEL_PATTERNS)
        has_caption = _check_pattern(full_text, _FIGURE_CAPTION_PATTERNS)
        has_source = _check_pattern(full_text, _SOURCE_PATTERNS)

        n_pass = sum([has_axis, has_unit, has_legend, has_caption, has_source])
        slide_score = round(n_pass / 5 * 10, 1)

        missing: list[str] = []
        if not has_axis:
            missing.append("axis_labels")
        if not has_unit:
            missing.append("units")
        if not has_legend:
            missing.append("legend/label")
        if not has_caption:
            missing.append("figure_caption")
        if not has_source:
            missing.append("source/citation")

        figure_slides.append({
            "slide_no": i,
            "title": _get_title(slide)[:60],
            "chart_axis_titles_detected": chart_axis_titles,
            "has_axis_labels": has_axis,
            "has_units": has_unit,
            "has_legend_or_label": has_legend,
            "has_figure_caption": has_caption,
            "has_source": has_source,
            "compliance_score": slide_score,
            "missing": missing,
        })

    n_fig = len(figure_slides)
    if n_fig == 0:
        return {
            "score": 10.0,
            "score_max": 10,
            "n_total_slides": n_total,
            "n_figure_slides": 0,
            "comments": [
                "Figure / chart slide 検出ゼロ。本ツール対象外。"
                "テキスト中心プレゼンなら他 tool 使用推奨。"
            ],
            "hint": "no figures to evaluate.",
            "source": "T21 figure slide compliance (presentation 2026-04)",
        }

    overall = round(
        sum(s["compliance_score"] for s in figure_slides) / n_fig, 1
    )

    # ---- comments ----
    comments: list[str] = []

    n_non = sum(1 for s in figure_slides if s["compliance_score"] < 6)
    comments.append(
        f"Figure slide: {n_fig}/{n_total} 枚 | non-compliant {n_non} | "
        f"avg compliance {overall}/10"
    )

    if n_non >= 1:
        comments.append(
            f"⚠ {n_non} 枚の figure slide で compliance 不充足。"
            "理系プレゼンの figure は『軸ラベル + 単位 + 凡例 + caption + 出典』"
            "が core 5 要素 (論文 figure と同じ厳格性)。"
        )
        for s in figure_slides[:5]:
            if s["compliance_score"] < 6:
                comments.append(
                    f"  • slide {s['slide_no']} [{s['compliance_score']}/10]: "
                    f"{s['title']} — missing {s['missing']}"
                )

    n_no_unit = sum(
        1 for s in figure_slides if not s["has_units"]
    )
    n_no_axis = sum(
        1 for s in figure_slides if not s["has_axis_labels"]
    )
    n_no_source = sum(
        1 for s in figure_slides if not s["has_source"]
    )

    if n_no_axis > n_fig * 0.4:
        comments.append(
            f"⚠ {n_no_axis} figure slides で軸ラベル不在。"
            "Q&A で『何を plot してる?』の質問必至。"
        )
    if n_no_unit > n_fig * 0.4:
        comments.append(
            f"⚠ {n_no_unit} figure slides で単位明記なし。"
            "理系では SI 単位 / dimensional 明示が標準。"
        )
    if n_no_source > n_fig * 0.5:
        comments.append(
            f"ℹ {n_no_source} figure slides で出典 / 自作明示なし。"
            "「自作 (本研究)」「[Smith 2023] より」のような 1 文を添える。"
        )

    if n_non == 0:
        comments.append(
            f"全 {n_fig} figure slide で 5 要素充足。"
            "理系学会 publication-grade 基準を満たす figure。"
        )

    comments = comments[:7]

    hint = (
        "理系プレゼン特化 — corporate プレゼンでは省略可な要素を厳格 check。"
        "chart 内部の axis title は python-pptx で取得 (image figure には適用不能)。"
        "T19 (chart simplification) と相補: T19 は 3D / 凡例過多、T21 は不足要素。"
    )

    return {
        "score": overall,
        "score_max": 10,
        "n_total_slides": n_total,
        "n_figure_slides": n_fig,
        "figure_slides": figure_slides,
        "comments": comments,
        "hint": hint,
        "source": (
            "Tufte『Visual Display of Quantitative Information』(1983) + "
            "宮野『研究発表のためのスライドデザイン』S6 (figure 標準) + "
            "IEEE/IEEJ 論文 figure caption 規定 の理系プレゼン版。"
            "(presentation 2026-04 採用、理系特化)。"
        ),
    }
