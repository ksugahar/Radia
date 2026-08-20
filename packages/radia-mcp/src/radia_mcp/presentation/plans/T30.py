"""Presentation T30: Script vs slide coverage (verbal side 強化)
(presentation_script_vs_slide_coverage).

speaker_note (台本) が slide 内容を網羅しているかを診断。
台本に言及されない slide content は「発表中言い忘れる」リスク。
逆に、slide に無いが台本にある内容は補足説明として OK。

検出ロジック:
  - 各 slide の body text tokens と speaker_note tokens を抽出
  - per slide: coverage = speaker_note でカバーされた body token の比率
  - missing_from_script: slide にあるが note で言及されない tokens
  - extra_in_script: note にあるが slide に無い tokens (補足、基本 OK)
"""

from __future__ import annotations

import pathlib
import re


def _iter_shapes(shape_collection):
    try:
        from radia_mcp.presentation.tools import _walk_shapes
        yield from _walk_shapes(shape_collection)
    except Exception:
        for s in shape_collection:
            yield s


# A Latin token bounded by ASCII, NOT by \b.  Python's \w matches CJK, so
# "のHACApK" has no word boundary before the H and \b[A-Z] never fires --
# which made every acronym spoken inside a Japanese sentence invisible.
_LATIN = re.compile(r"(?<![A-Za-z0-9])[A-Z][A-Za-z0-9]{2,}(?![A-Za-z0-9])")
_KANJI = re.compile(r"[一-鿿]{2,}")
_KATAKANA = re.compile(r"[ァ-ヴー]{3,}")


def _extract_tokens(text: str) -> set[str]:
    tokens: set[str] = set()
    for pattern in (_KANJI, _KATAKANA, _LATIN):
        for m in pattern.finditer(text):
            tokens.add(m.group(0))
    return tokens


def _picture_rects(slide) -> list[tuple[int, int, int, int]]:
    """Bounding rectangles of the pictures on this slide."""
    rects: list[tuple[int, int, int, int]] = []
    for shape in _iter_shapes(slide.shapes):
        try:
            if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
                if None not in (shape.left, shape.top, shape.width, shape.height):
                    rects.append((shape.left, shape.top,
                                  shape.left + shape.width,
                                  shape.top + shape.height))
        except Exception:
            continue
    return rects


def _inside_any(rects, shape) -> bool:
    """Does this shape's centre sit inside one of the rectangles?"""
    try:
        if None in (shape.left, shape.top, shape.width, shape.height):
            return False
        cx = shape.left + shape.width // 2
        cy = shape.top + shape.height // 2
    except Exception:
        return False
    return any(x0 <= cx <= x1 and y0 <= cy <= y1 for x0, y0, x1, y1 in rects)


def classify_slide_text(slide, slide_height: int = 6858000) -> dict:
    """Split a slide's text into the parts that play different roles.

    Not everything printed on a slide is something a speaker says.  A label
    sitting on top of a figure ("Mesh A", "165,600 DoF") is read by the eye
    while the speaker talks about something else, and a table cell is a
    reference the audience scans.  Scoring those against the script marks a
    deck down for text that was never meant to be spoken.

    The take-home line along the bottom is the opposite case: it IS the
    slide's claim, so it gets its own check rather than being averaged in
    with everything else.
    """
    try:
        from radia_mcp.presentation.tools import _slide_title
        title = _slide_title(slide)
    except Exception:
        title = ""

    rects = _picture_rects(slide)
    body: list[str] = []
    figure_labels: list[str] = []
    table_text: list[str] = []
    banner = ""

    for shape in _iter_shapes(slide.shapes):
        try:
            if getattr(shape, "has_table", False):
                tbl = shape.table
                table_text.extend(c.text for r in tbl.rows for c in r.cells)
                continue
            if not shape.has_text_frame:
                continue
            txt = (shape.text_frame.text or "")
            if not txt.strip() or txt.strip() == title:
                continue
            if _inside_any(rects, shape):
                figure_labels.append(txt)
                continue
            top = shape.top
            if top is not None and 0.85 * slide_height < top < 0.94 * slide_height:
                banner = txt.replace("\n", "")
                continue
            body.append(txt)
        except Exception:
            continue

    return {
        "body": "\n".join(body),
        "figure_labels": "\n".join(figure_labels),
        "table_text": "\n".join(table_text),
        "banner": banner,
    }


def _get_slide_body(slide) -> str:
    return classify_slide_text(slide)["body"]


def _is_spoken(token: str, notes: str, notes_tokens: set[str]) -> bool:
    """Is this slide token said in the script?

    Three ways, loosest last:

    1. the same token,
    2. the token appearing anywhere in the script text -- saying
       "Gram二重積分" is saying "二重積分",
    3. most of the token's character pairs appearing in the script.  A slide
       compresses ("独立参照", "補償点法") where a speaker inflects ("独立な
       参照解", "補償点磁荷法"), and a maximal-kanji-run tokenizer makes those
       different words.  Requiring two thirds of the pairs, and only for
       tokens of three characters or more, accepts the inflected form without
       accepting a token whose pieces merely scatter across the script.
    """
    if token in notes_tokens or token in notes:
        return True
    if len(token) < 3 or not any("一" <= c <= "鿿" for c in token):
        return False
    pairs = [token[j:j + 2] for j in range(len(token) - 1)]
    hits = sum(1 for p in pairs if p in notes)
    return hits * 3 >= len(pairs) * 2


def _get_notes(slide) -> str:
    try:
        if slide.has_notes_slide:
            return slide.notes_slide.notes_text_frame.text or ""
    except Exception:
        pass
    return ""


def presentation_script_vs_slide_coverage(pptx_path: str) -> dict:
    """台本 (speaker_note) が slide 内容を網羅しているかを per slide 診断。

    Returns:
        dict: per_slide / avg_coverage / low_coverage_slides /
              uncovered_tokens / comments
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))

    per_slide: list[dict] = []
    low_coverage_slides: list[dict] = []

    # Deck furniture -- the affiliation in the corner, a running section label
    # -- appears on nearly every slide and is not spoken.  Scoring it against
    # the script marks every slide down for the same non-omission.
    slides = list(prs.slides)
    parts = [classify_slide_text(s) for s in slides]
    seen: dict[str, int] = {}
    for part in parts:
        for token in _extract_tokens(part["body"]):
            seen[token] = seen.get(token, 0) + 1
    boilerplate = {t for t, n in seen.items()
                   if len(slides) >= 4 and n >= 0.8 * len(slides)}

    # The take-home line along the bottom is the slide's claim, so it gets its
    # own check.  Averaged in with the rest it disappears: it is a couple of
    # tokens against a whole slide, and a deck can score well while the
    # speaker never lands a single one of its claims out loud.
    banners: list[dict] = []
    for i, (slide, part) in enumerate(zip(slides, parts), 1):
        banner = part["banner"].strip()
        if i == 1 or not banner:
            continue
        notes = _get_notes(slide).strip()
        tokens = _extract_tokens(banner)
        if not tokens:
            continue
        ntok = _extract_tokens(notes)
        said = {t for t in tokens if _is_spoken(t, notes, ntok)}
        banners.append({
            "slide_no": i,
            "banner": banner,
            "spoken_ratio": round(len(said) / len(tokens), 3),
            "not_spoken": sorted(tokens - said)[:6],
        })
    banners_unspoken = [b for b in banners if b["spoken_ratio"] < 0.5]

    for i, slide in enumerate(slides, 1):
        if i == 1:
            # The cover carries the formal title, the English subtitle and the
            # author list.  skill.md says the speaker need not read the formal
            # title aloud, so scoring it as unspoken content is a false alarm --
            # the title checkers skip slide 1 for the same reason.
            per_slide.append({
                "slide_no": i, "coverage": None,
                "skip_reason": "cover_slide",
            })
            continue
        body = _get_slide_body(slide).strip()
        notes = _get_notes(slide).strip()
        body_tokens = _extract_tokens(body)
        notes_tokens = _extract_tokens(notes)

        # Skip if either is empty
        if not body_tokens and not notes_tokens:
            per_slide.append({
                "slide_no": i, "coverage": None,
                "skip_reason": "empty_both",
            })
            continue
        if not body_tokens:
            per_slide.append({
                "slide_no": i, "coverage": None,
                "skip_reason": "title_only_slide",
                "notes_chars": len(notes),
            })
            continue
        if not notes_tokens:
            per_slide.append({
                "slide_no": i,
                "coverage": 0.0,
                "body_tokens": len(body_tokens),
                "notes_tokens": 0,
                "missing_from_script": sorted(body_tokens)[:10],
                "warning": "speaker_note 空で slide に内容あり",
            })
            low_coverage_slides.append({
                "slide_no": i,
                "coverage": 0.0,
                "reason": "speaker_note 未記入",
                "missing_tokens": sorted(body_tokens)[:5],
            })
            continue

        body_tokens = body_tokens - boilerplate
        if not body_tokens:
            per_slide.append({
                "slide_no": i, "coverage": None,
                "skip_reason": "only_deck_furniture",
            })
            continue

        covered = {t for t in body_tokens if _is_spoken(t, notes, notes_tokens)}
        missing = body_tokens - covered
        extra = notes_tokens - body_tokens
        coverage = len(covered) / len(body_tokens)

        per_slide.append({
            "slide_no": i,
            "coverage": round(coverage, 3),
            "body_tokens": len(body_tokens),
            "notes_tokens": len(notes_tokens),
            "covered_tokens": sorted(covered)[:5],
            "missing_from_script": sorted(missing)[:5],
            "extra_in_script": sorted(extra)[:5],
        })

        if coverage < 0.5:
            low_coverage_slides.append({
                "slide_no": i,
                "coverage": round(coverage, 3),
                "reason": f"slide 内容の {coverage*100:.0f}% しか言及されない",
                "missing_tokens": sorted(missing)[:5],
            })

    # Overall stats
    valid_coverages = [s["coverage"] for s in per_slide
                       if s.get("coverage") is not None]
    avg_coverage = (
        sum(valid_coverages) / len(valid_coverages) if valid_coverages else 0
    )
    n_low = len(low_coverage_slides)

    # Score
    score = round(avg_coverage * 10, 1)
    if n_low > len(valid_coverages) * 0.3 and valid_coverages:
        score -= 2
    score = max(0.0, round(score, 1))

    # Comments
    comments: list[str] = []
    comments.append(
        f"平均 coverage: {avg_coverage*100:.0f}% "
        f"({len(valid_coverages)} slides 評価対象) | "
        f"low-coverage (<50%): {n_low} slides"
    )

    if n_low >= 1:
        comments.append(
            f"⚠ {n_low} slides で台本 coverage <50%。"
            "発表中言い忘れるリスク。speaker_note に言及追加推奨。"
        )
        for lc in low_coverage_slides[:5]:
            comments.append(
                f"  • slide {lc['slide_no']} ({lc['coverage']*100:.0f}%): "
                f"missing {lc['missing_tokens']}"
            )

    if avg_coverage >= 0.7 and n_low == 0:
        comments.append(
            f"✓ 平均 coverage {avg_coverage*100:.0f}%、全 slide で balance OK"
        )

    hint = (
        "body token と note token の overlap で coverage を推定。"
        "missing は『slide に書いてあるが台本に無い → 発表中飛ばすリスク』。"
        "extra は『台本にあるが slide に無い → 補足説明、基本 OK』。"
        "coverage 50% 未満は警戒、70%+ で良好。"
    )

    return {
        "score": score,
        "score_max": 10,
        "avg_coverage_pct": round(avg_coverage * 100, 1),
        "n_slides_evaluated": len(valid_coverages),
        "n_low_coverage_slides": n_low,
        "deck_boilerplate_ignored": sorted(boilerplate),
        "banner_checked": len(banners),
        "banners_not_spoken": banners_unspoken,
        "banner_detail": banners,
        "low_coverage_slides": low_coverage_slides,
        "per_slide": per_slide[:30],
        "comments": comments,
        "hint": hint,
        "source": (
            "presentation リハーサル支援。既存 speaker_note_ratio (T7) は "
            "note の有無、本 tool は note と slide の内容一致度を見る。"
            "(presentation 2026-04 採用、v0.26.0)。"
        ),
    }
