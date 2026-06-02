"""Plan B T20: Equation slide compliance (理系プレゼン特化)
(presentation_equation_slide_compliance).

理系プレゼンで数式を出す場合、論文と同じく「定義した記号は説明する」
原則が必須。LaTeX 論文では \\nomenclature {} があるが、プレゼンでは
slide 内で記号定義 + 物理意味 + 出典を articulate する必要がある。

検出する 4 軸 per equation slide:
  1. **symbol_definitions** — 各 symbol に定義文があるか
  2. **physical_meaning** — 数式が何を表すかの一文 (「これは X を意味」)
  3. **dimension_check** — 物理量に単位 / 次元の言及
  4. **source_or_derivation** — 出典 [N] or 「式 X から導出」

検出ロジック:
  - LaTeX 数式マーカー ($, $$, \\begin{equation}, etc) を持つ slide
  - PowerPoint 数式 object (a:oMath) も検出
  - Image-only equation (画像化された式) は OCR 不可のため対象外と明記

設計方針:
- equation 不在 slide は score=10 (該当なし)
- equation あり slide のみ評価
- PRIMARY OUTPUT は per_equation_slide findings
"""

from __future__ import annotations

import pathlib
import re


def _iter_all_shapes(shape_collection):
    try:
        from radia_mcp.presentation.tools import _walk_shapes
        yield from _walk_shapes(shape_collection)
    except Exception:
        for s in shape_collection:
            yield s


def _get_title(slide) -> str:
    try:
        from radia_mcp.presentation.tools import _slide_title
        return _slide_title(slide)
    except Exception:
        return ""


# Inline / display math markers
_LATEX_MATH_PATTERNS = [
    re.compile(r"\$[^$]+\$"),
    re.compile(r"\\\(.*?\\\)"),
    re.compile(r"\\\[.*?\\\]"),
    re.compile(r"\\begin\{(?:equation|align|gather|eqnarray)\*?\}"),
]

# Common equation symbols that suggest formula presence
_FORMULA_SYMBOLS = re.compile(
    r"(?:[α-ωΑ-Ω]|"  # Greek letters
    r"\\(?:alpha|beta|gamma|delta|epsilon|theta|lambda|mu|nu|pi|rho|"
    r"sigma|tau|phi|chi|psi|omega|nabla|partial|infty|sum|int|prod)|"
    r"[=±∞≈≡≠≤≥]\s*[a-zA-Z]|"  # LHS = RHS
    r"\b\\?(?:sin|cos|tan|log|ln|exp|sqrt)\b|"
    r"\^[\d\{]|_\{)"
)

_DEFINITION_PATTERNS = [
    r"(?:where|ここで)[^\n。]{5,40}(?:は|は\s*[a-zα-ω])",
    r"(?:[a-zA-Zα-ωΑ-Ω]|\\\w+)\s*(?:は|=)\s*[^\n。]{5,40}(?:を表|意味|示)",
    r"(?:記号|symbol|notation)\s*(?:定義|table)",
    r":\s*[A-Za-zα-ωΑ-Ω][^\n,，。]{3,30}\s*\(",  # x: description (unit)
]

_PHYSICAL_MEANING_PATTERNS = [
    r"(?:この式|本式|式\s*\(\d+\))\s*(?:は|では)\s*[^\n。]{5,60}",
    r"(?:this\s+equation|the\s+equation)\s+(?:represents|shows|describes)",
    r"(?:意味|表す|示す|describes|expresses|denotes)",
    r"(?:物理的に|直感的に|意味として)",
]

_DIMENSION_PATTERNS = [
    r"\[\s*[A-Za-zα-ωΑ-Ω/·\d\^\-]{1,15}\s*\]",  # [m/s] [kg·m^2]
    r"(?:単位|unit)\s*[:：]?\s*\S",
    r"(?:dimensional|次元)\s*(?:check|分析|analysis)",
    r"\d+(?:\.\d+)?\s*(?:m/s|kg|Pa|N|J|W|Hz|V|A|°C|K|mol)",
]

_SOURCE_PATTERNS = [
    r"\[\d+\]",  # citation
    r"\(.*?\d{4}\)",  # (Author 2023)
    r"(?:式\s*\(\d+\))\s*(?:より|から導出|を変形)",
    r"(?:see|参照|cf\.)\s+(?:eq|式)",
]


def _has_equation(slide) -> tuple[bool, str]:
    """Detect if slide contains an equation. Returns (has_eq, evidence)."""
    full_text = ""
    for shape in _iter_all_shapes(slide.shapes):
        try:
            if shape.has_text_frame:
                full_text += " " + (shape.text_frame.text or "")
        except Exception:
            continue

    # 1. LaTeX markers
    for pat in _LATEX_MATH_PATTERNS:
        m = pat.search(full_text)
        if m:
            return (True, f"latex_marker: {m.group(0)[:30]}")

    # 2. Formula symbols (multiple greek + operators)
    sym_matches = _FORMULA_SYMBOLS.findall(full_text)
    if len(sym_matches) >= 3:
        return (True, f"formula_symbols ({len(sym_matches)} signals)")

    # 3. Title-based hint
    title = _get_title(slide).lower()
    if any(kw in title for kw in [
        "equation", "式", "formula", "derivation", "導出",
    ]):
        return (True, f"equation_title_hint: {title[:30]}")

    # 4. PowerPoint OMath element via XML inspection
    try:
        for shape in _iter_all_shapes(slide.shapes):
            try:
                xml = shape.element.xml if hasattr(shape, "element") else ""
                if "<m:oMath" in xml or 'xmlns:m="' in xml and "oMath" in xml:
                    return (True, "pptx_omath")
            except Exception:
                continue
    except Exception:
        pass

    return (False, "")


def _check_axis(text: str, patterns: list) -> bool:
    if isinstance(patterns[0], str):
        return any(
            re.search(p, text, re.IGNORECASE) for p in patterns
        )
    # already compiled
    return any(p.search(text) for p in patterns)


def presentation_equation_slide_compliance(pptx_path: str) -> dict:
    """数式 slide の理系プレゼン compliance を診断。

    Args:
        pptx_path: .pptx file

    Returns:
        dict: score / equation_slides / per_slide / comments
    """
    try:
        import pptx as _pptx
    except ImportError:
        return {"error": "python-pptx not installed."}
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}
    prs = _pptx.Presentation(str(p))

    equation_slides: list[dict] = []
    n_total = 0

    for i, slide in enumerate(prs.slides, 1):
        n_total += 1
        has_eq, evidence = _has_equation(slide)
        if not has_eq:
            continue

        full_text = ""
        for shape in _iter_all_shapes(slide.shapes):
            try:
                if shape.has_text_frame:
                    full_text += "\n" + (shape.text_frame.text or "")
            except Exception:
                continue

        has_definitions = _check_axis(full_text, _DEFINITION_PATTERNS)
        has_meaning = _check_axis(full_text, _PHYSICAL_MEANING_PATTERNS)
        has_dimension = _check_axis(full_text, _DIMENSION_PATTERNS)
        has_source = _check_axis(full_text, _SOURCE_PATTERNS)

        n_pass = sum([has_definitions, has_meaning, has_dimension, has_source])
        slide_score = round(n_pass / 4 * 10, 1)

        missing: list[str] = []
        if not has_definitions:
            missing.append("symbol_definitions")
        if not has_meaning:
            missing.append("physical_meaning")
        if not has_dimension:
            missing.append("dimension/unit")
        if not has_source:
            missing.append("source/derivation")

        equation_slides.append({
            "slide_no": i,
            "title": _get_title(slide)[:60],
            "evidence": evidence,
            "has_symbol_definitions": has_definitions,
            "has_physical_meaning": has_meaning,
            "has_dimension_unit": has_dimension,
            "has_source": has_source,
            "compliance_score": slide_score,
            "missing": missing,
        })

    n_eq = len(equation_slides)
    if n_eq == 0:
        return {
            "score": 10.0,
            "score_max": 10,
            "n_total_slides": n_total,
            "n_equation_slides": 0,
            "comments": [
                "数式 slide 検出ゼロ。本ツールは数式を含む発表向け。"
                "数式が image 化されていれば本 tool では検出不能 — "
                "OCR で text 化してから再診断。"
            ],
            "hint": "no equations to evaluate.",
            "source": "T20 equation slide compliance (presentation 2026-04)",
        }

    overall = round(
        sum(s["compliance_score"] for s in equation_slides) / n_eq, 1
    )

    # ---- comments ----
    comments: list[str] = []

    n_non = sum(1 for s in equation_slides if s["compliance_score"] < 6)
    comments.append(
        f"数式 slide: {n_eq}/{n_total} 枚 | non-compliant {n_non} | "
        f"avg compliance {overall}/10"
    )

    if n_non >= 1:
        comments.append(
            f"⚠ {n_non} 枚の数式 slide で compliance 不充足。"
            "理系プレゼンの数式は『記号定義 + 物理意味 + 単位 + 出典』が "
            "core 4 要素。論文と同じ厳格性が求められる。"
        )
        for s in equation_slides[:5]:
            if s["compliance_score"] < 6:
                comments.append(
                    f"  • slide {s['slide_no']} [{s['compliance_score']}/10]: "
                    f"{s['title']} — missing {s['missing']}"
                )

    # 各要素の global 不足統計
    n_no_def = sum(
        1 for s in equation_slides if not s["has_symbol_definitions"]
    )
    n_no_dim = sum(
        1 for s in equation_slides if not s["has_dimension_unit"]
    )
    if n_no_def > n_eq * 0.5:
        comments.append(
            f"⚠ {n_no_def} 数式 slides で symbol 定義なし。"
            "「ここで X は ..., Y は ...」形式を必須として徹底。"
        )
    if n_no_dim > n_eq * 0.5:
        comments.append(
            f"⚠ {n_no_dim} 数式 slides で 次元 / 単位言及なし。"
            "理系では物理量の dimension check が標準 (SI 単位明示)。"
        )

    if n_non == 0:
        comments.append(
            f"全 {n_eq} 数式 slide で 4 要素充足。"
            "Equation literacy 高、reviewer の質問を予防。"
        )

    comments = comments[:7]

    hint = (
        "理系プレゼン特化 (TED にはない概念)。Image 化された式は "
        "本 tool では検出不能なので OCR 必須。"
        "純情報学系 (式が無い) では適用外で score=10 を返す。"
        "Q&A で『この記号は?』『単位は?』を予防する効果。"
    )

    return {
        "score": overall,
        "score_max": 10,
        "n_total_slides": n_total,
        "n_equation_slides": n_eq,
        "equation_slides": equation_slides,
        "comments": comments,
        "hint": hint,
        "source": (
            "理系プレゼン慣行 + IEEE/IEEJ 学会発表ガイドライン "
            "(数式 slide compliance) + 宮野『研究発表のためのスライド "
            "デザイン』S6 (数式記述) の rough adaptation。"
            "(presentation 2026-04 採用、理系特化)。"
        ),
    }
