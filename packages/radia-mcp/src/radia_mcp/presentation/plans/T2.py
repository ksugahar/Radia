"""Plan B T2: Take-home スライド強度診断 (presentation_takehome_strength).

skill.md 木下『理科系の作文技術』p.230 / 宮野公樹『研究発表のためのスライド
デザイン』S14 (1 枚 1 メッセージ) に基づく粗い温度計。

設計方針:
- score は 0-10 の粗い温度計。9 以上は誤差範囲、6 以下は改善サイン。
- PRIMARY OUTPUT は comments (人間レビュアー風の助言文)。
- observations は raw measurement。自動化した numeric な正誤判定には使わない。
- 最終スライド群から take-home を推定する粗い判定。複数 take-home を
  並べる多メッセージデッキは想定外 — 単一 message 原則が前提。
"""

from __future__ import annotations

import pathlib
import re


# ------------------------------------------------------------------
# Patterns
# ------------------------------------------------------------------
# Q&A / Thank you / 謝辞 / 参考文献 など「締めの後ろ飾り」スライド。
# これらは take-home 候補から除外する (最後の content slide まで遡る)。
_QA_TITLE_PATTERN = re.compile(
    r"(q\s*&\s*a|q\s*and\s*a|thank\s*you|thanks|ご清聴|ご質問|"
    r"質疑|謝辞|acknowledg|参考文献|references?|bibliography|"
    r"付録|appendix|補遺|おわり|end\s*of\s*presentation)",
    re.IGNORECASE,
)

# 最終スライドが Q&A 系であれば deck の構造的準備として positive signal。
_QA_EXPLICIT_PATTERN = re.compile(
    r"(q\s*&\s*a|q\s*and\s*a|thank\s*you|ご清聴|ご質問)",
    re.IGNORECASE,
)

# Key number: 数値 + 単位 (kHz / MHz / % / 倍 / × / 年 / 件)。
_KEY_NUMBER_PATTERN = re.compile(
    r"\d+(?:\.\d+)?\s*(?:kHz|MHz|GHz|Hz|%|％|倍|×|x|年|件|枚|名|回|ms|μs|ns|"
    r"dB|mT|T|V|A|W|kW|MW)",
    re.IGNORECASE,
)

# Action verb: take-home で audience に記憶すべき動作を指示する動詞。
# title または body が該当 verb で始まる/含む かを判定。
_ACTION_VERB_PATTERN = re.compile(
    r"(提案|示す|示した|実現|確立|達成|実証|開発|構築|"
    r"propose[sd]?|show[sn]?|demonstrate[sd]?|enable[sd]?|"
    r"achieve[sd]?|establish(?:ed|es)?|build[st]?|built)",
    re.IGNORECASE,
)


# ------------------------------------------------------------------
# Main
# ------------------------------------------------------------------
def presentation_takehome_strength(pptx_path: str) -> dict:
    """最終スライド or last-3-slides の Take-home 品質診断。

    skill.md 木下 p.230 (「Take-home: 提案手法で解析を 4× 加速」1 行) /
    宮野 S14 (1 枚 1 メッセージ) に基づく。

    観測項目:
    - 最終 content slide (Q&A/参考文献除く) の title と body
    - body の行数 (1-3 行が標準)
    - key number の有無 (「12%」「3 倍」「500 kHz」等)
    - action verb の有無 (「提案」「実現」「propose」等)
    - deck に Q&A/ご清聴 slide が明示的にあるか

    Args:
        pptx_path: .pptx path。

    Returns:
        dict: score / score_max / comments / observations / hint / source
    """
    try:
        import pptx  # type: ignore
    except ImportError:
        return {
            "error": "python-pptx not installed. "
                     "pip install mcp-server-document[pptx]"
        }
    p = pathlib.Path(pptx_path)
    if not p.exists():
        return {"error": f"file not found: {pptx_path}"}

    # Import shared helpers from tools.py (H3/H7 fixes).
    from ..tools import _slide_title, _walk_shapes

    try:
        prs = pptx.Presentation(str(p))
    except Exception as exc:  # pragma: no cover - corrupt file
        return {"error": f"failed to open: {exc}"}

    slides = list(prs.slides)
    slides_scanned = len(slides)

    from .._deck_sections import classify_deck_sections

    sections = classify_deck_sections(slides, _slide_title)
    backup_slide_numbers = set(sections["backup_slide_numbers"])

    # A closing slide may correctly precede a trailing appendix.  Looking only
    # at the physical last slide used to miss that common research-talk layout.
    has_qa_or_questions_title = any(
        _QA_EXPLICIT_PATTERN.search(_slide_title(slide))
        for slide_no, slide in enumerate(slides, 1)
        if slide_no not in backup_slide_numbers
    )

    # Walk backward from the last slide, skipping Q&A/References/Thanks.
    takehome_slide_idx = -1
    takehome_title = ""
    takehome_body = ""

    for i in range(len(slides) - 1, -1, -1):
        slide = slides[i]
        if i + 1 in backup_slide_numbers:
            continue
        title = _slide_title(slide)
        if title and _QA_TITLE_PATTERN.search(title):
            continue
        # Collect non-title body text via _walk_shapes (groups + tables).
        body_parts: list[str] = []
        for shape in _walk_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                # Try tables anyway (has_text_frame = False for tables)
                try:
                    if getattr(shape, "has_table", False) and shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                try:
                                    tf = cell.text_frame
                                except Exception:
                                    continue
                                for para in tf.paragraphs:
                                    t = para.text.strip()
                                    if t and t != title.strip():
                                        body_parts.append(t)
                except Exception:
                    pass
                continue
            # Skip the title shape itself.
            try:
                pf = shape.placeholder_format
            except Exception:
                pf = None
            is_title_shape = False
            try:
                if pf is not None and pf.idx == 0:
                    is_title_shape = True
            except Exception:
                pass
            if is_title_shape:
                continue
            try:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if not t:
                        continue
                    # Defensive: also skip the exact title string if it
                    # leaked into a non-placeholder shape.
                    if title and t == title.strip():
                        continue
                    body_parts.append(t)
            except Exception:
                continue

        body = "\n".join(body_parts)
        # Accept this slide as take-home candidate if it has any title
        # or body content (i.e., it is a content slide).
        if title.strip() or body.strip():
            takehome_slide_idx = i
            takehome_title = title
            takehome_body = body
            break
        # Otherwise keep walking back (empty slide).

    body_lines = [
        ln for ln in takehome_body.split("\n") if ln.strip()
    ]
    body_line_count = len(body_lines)
    body_char_count = len(takehome_body.replace(" ", "").replace("\n", ""))

    # has_key_number: body (primary) or title の数値 + 単位。
    scan_text_for_number = takehome_body + "\n" + takehome_title
    has_key_number = bool(_KEY_NUMBER_PATTERN.search(scan_text_for_number))

    # has_action_verb: title または body 内に action verb がある
    # (「starts with」は緩和: body は bullets が並ぶので含まれていれば OK)。
    scan_text_for_verb = takehome_title + "\n" + takehome_body
    has_action_verb = bool(_ACTION_VERB_PATTERN.search(scan_text_for_verb))

    observations = {
        "slides_scanned": slides_scanned,
        "main_slide_count": sections["main_slide_count"],
        "backup_slide_count": sections["backup_slide_count"],
        "backup_slide_numbers": sections["backup_slide_numbers"],
        "takehome_slide_idx": takehome_slide_idx,
        "takehome_title": takehome_title[:120],
        "takehome_body": takehome_body[:400],
        "body_line_count": body_line_count,
        "body_char_count": body_char_count,
        "has_key_number": has_key_number,
        "has_action_verb": has_action_verb,
        "has_qa_or_questions_title": has_qa_or_questions_title,
    }

    # ---- score (粗い温度計) ----
    score = 0.0
    if body_line_count <= 3 and takehome_slide_idx >= 0:
        score += 3
    if has_key_number:
        score += 2
    if has_action_verb:
        score += 2
    if takehome_slide_idx >= 0:
        score += 2
    if body_line_count > 6:
        score -= 2
    if takehome_slide_idx == -1:
        score -= 3
    if has_qa_or_questions_title:
        score += 1
    score = max(0.0, min(10.0, score))

    # ---- comments (PRIMARY OUTPUT) ----
    comments: list[str] = []

    if takehome_slide_idx == -1:
        comments.append(
            "Take-home スライドが見当たらない。最終スライドの 1 行で"
            "「本日の話を 1 文で」と締める (木下 p.230)。"
        )
    else:
        if body_line_count > 6:
            comments.append(
                f"Take-home の body が {body_line_count} 行。"
                "「1 枚 1 メッセージ」原則から大きく外れる。3 行以内に圧縮。"
            )
        if not has_action_verb:
            comments.append(
                "Take-home タイトル/body に action verb がない。"
                "「提案する」「実現する」「示す」「propose」「demonstrate」"
                "等で audience に記憶すべき動作を指示。"
            )
        if not has_key_number:
            comments.append(
                "Take-home に 1 つ key number が欲しい。"
                "「12 % 改善」「3 倍高速」「500 kHz 動作実証」等、"
                "数値 1 つで記憶に残る。"
            )
        # 全 OK: 肯定コメント
        if (body_line_count <= 3 and has_action_verb
                and has_key_number and body_line_count > 0):
            comments.append(
                "Take-home slide detected, concise, with number + "
                "action verb. 宮野 S2 要件を満たす。"
            )

    if has_qa_or_questions_title:
        comments.append(
            "最終スライドに Q&A / ご清聴 系タイトルを配置 — "
            "audience の question-handling 準備あり、丁寧な構成。"
        )

    # 最低 2 個のコメントを保証 (単純な肯定だけの場合に補う)
    if len(comments) < 2:
        comments.append(
            "score は粗い指標。skill.md 木下 p.230 / 宮野 S14 "
            "(1 枚 1 メッセージ) の判定基準で最終確認を推奨。"
        )

    # 最大 5 個に制限
    comments = comments[:5]

    hint = (
        "最終スライド群から take-home を推定する粗い判定。"
        "複数の take-home を並べる多メッセージデッキは想定外 — "
        "単一 message 原則が前提。score は粗い温度計であり、"
        "最終判断は comments を読み skill.md の判定基準と照合すること。"
    )

    return {
        "score": float(score),
        "score_max": 10,
        "comments": comments,
        "observations": observations,
        "hint": hint,
        "source": ("木下 p.230 / 宮野『研究発表のためのスライドデザイン』 "
                   "S14 (1 枚 1 メッセージ)"),
    }
