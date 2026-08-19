"""Put LaTeX equations into Word and PowerPoint as native Office math.

The equation arrives as an equation: Office's own tools edit it, it follows the
theme font and colour, and the reader needs nothing installed.  Not a picture,
not an OLE object, no MTEF.

Word wraps math as bare ``<m:oMath>`` inside ``<w:p>``; PowerPoint wraps it in
``<a14:m>`` inside ``<a:p>``.  That is the only structural difference.

python-docx / python-pptx / lxml are imported where they are used, so importing
this module costs nothing when only the C++ side is wanted.
"""

from __future__ import annotations

import re
import zipfile
from typing import Iterable, Sequence

from radia._equation import (MarkdownDoc, MdBlock, MdSegment, md_blocks,
                             tex_to_mathml, tex_to_omml, tex_to_rtf)

# The raster the clipboard carries.  Google Slides has no equation object and
# rejects SVG on upload, so for a slide the picture IS the equation and stays
# one -- it has to be print quality, not screen quality.  Equation Editor's era
# could not have afforded this; here it is a few tens of kilobytes.
PASTE_DPI = 600.0

W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"
M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"


class Piece:
    """One run of a block: literal text, or an equation."""

    __slots__ = ("text", "latex", "display")

    def __init__(self, text: str = "", latex: str = "", display: bool = False):
        self.text = text
        self.latex = latex
        self.display = display

    @property
    def is_math(self) -> bool:
        return bool(self.latex)

    def __repr__(self) -> str:
        if self.is_math:
            kind = "display" if self.display else "inline"
            return f"<Piece {kind} {self.latex!r}>"
        return f"<Piece text {self.text!r}>"


def split_math(markdown: str) -> list[Piece]:
    """Split text into literal runs and equations.

    Uses the same scanner the editor uses, so what becomes an equation here is
    exactly what the editor lets you edit: a ``$`` inside code stays a dollar
    sign in both.
    """
    doc = MarkdownDoc()
    doc.load(markdown)
    pieces: list[Piece] = []
    for seg in doc.segments():
        if seg.is_math:
            pieces.append(Piece(latex=seg.body.strip(),
                                display=(seg.kind == MdSegment.DisplayMath)))
        elif seg.kind == MdSegment.CodeSpan:
            pieces.append(Piece(text=seg.body))
        else:
            pieces.append(Piece(text=seg.source))
    return pieces


def _omml_element(latex: str, display: bool):
    from lxml import etree

    xml = tex_to_omml(latex)
    if display:
        xml = xml.replace("<m:oMath", "<m:oMathPara><m:oMath", 1)
        xml = xml.replace("</m:oMath>", "</m:oMath></m:oMathPara>", 1)
        xml = xml.replace("<m:oMathPara>", f'<m:oMathPara xmlns:m="{M_NS}">', 1)
    elif "xmlns:m=" not in xml:
        xml = xml.replace("<m:oMath", f'<m:oMath xmlns:m="{M_NS}"', 1)
    return etree.fromstring(xml.encode("utf-8"))


def omml_paragraph(paragraph, pieces: Iterable[Piece], powerpoint: bool = False):
    """Fill a Word or PowerPoint paragraph with text and equations.

    Mixing the two is what puts an equation inside a sentence rather than on a
    line of its own.
    """
    from lxml import etree

    ns = A_NS if powerpoint else W_NS
    p = paragraph._p
    end_props = p.find(f"{{{A_NS}}}endParaRPr") if powerpoint else None

    for piece in pieces:
        if piece.is_math:
            math = _omml_element(piece.latex, piece.display)
            if powerpoint:
                node = etree.SubElement(p, f"{{{A14_NS}}}m", nsmap={"a14": A14_NS})
                node.append(math)
            else:
                p.append(math)
                node = math
        elif piece.text:
            node = etree.SubElement(p, f"{{{ns}}}r")
            t = etree.SubElement(node, f"{{{ns}}}t")
            if not powerpoint:
                t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
            t.text = piece.text
        else:
            continue
        if end_props is not None:
            end_props.addprevious(node)
    return paragraph


def markdown_to_docx(markdown: str, out_path: str) -> str:
    """Write Markdown to a .docx, equations included as native Office math.

    Headings, paragraphs, fenced code, bullet and numbered lists are handled --
    what a technical note needs.  Tables, footnotes, images and inline emphasis
    are not, and pass through as literal text rather than being dropped
    silently.
    """
    from docx import Document
    from docx.shared import Pt
    from lxml import etree

    doc = Document()
    # The same block scanner the viewer uses.  A second copy here would drift
    # from it, and the two would disagree about what a file says.
    for block in md_blocks(markdown):
        if block.kind == MdBlock.Blank:
            continue

        if block.kind == MdBlock.Code:
            para = doc.add_paragraph()
            run = para.add_run(block.text)
            run.font.name = "Consolas"
            run.font.size = Pt(9)
            continue

        if block.kind == MdBlock.Heading:
            para = doc.add_heading(level=min(block.level, 4))
            # add_heading seeds a run with the text; start from an empty one
            for r in list(para._p):
                if etree.QName(r).localname == "r":
                    para._p.remove(r)
        elif block.kind == MdBlock.Bullet:
            para = doc.add_paragraph(style="List Bullet")
        elif block.kind == MdBlock.Numbered:
            para = doc.add_paragraph(style="List Number")
        else:
            para = doc.add_paragraph()

        omml_paragraph(para, split_math(block.text))

    doc.save(out_path)
    return out_path


def _script_text(block):
    """The spoken line in this block, or None if the block is not one.

    A blockquote is the script.  The scanner has no blockquote kind -- the
    marker arrives at the front of the text -- so it is read here rather than
    by teaching the shared scanner a rule only slides use.
    """
    if block.kind not in (MdBlock.Paragraph, MdBlock.Bullet):
        return None
    lines = block.text.split("\n")
    if not lines or not lines[0].lstrip().startswith(">"):
        return None
    out = []
    for line in lines:
        stripped = line.lstrip()
        out.append(stripped[1:].lstrip() if stripped.startswith(">") else line)
    return "\n".join(out).strip()


def _no_bullet(paragraph, centre: bool = False):
    """Turn the bullet off for this paragraph, and optionally centre it.

    A display equation is its own line, not an item in a list -- PowerPoint
    puts a bullet in front of every paragraph until told otherwise.
    """
    from lxml import etree

    pPr = paragraph._p.find(f"{{{A_NS}}}pPr")
    if pPr is None:
        pPr = etree.SubElement(paragraph._p, f"{{{A_NS}}}pPr")
        paragraph._p.insert(0, pPr)
    if centre:
        pPr.set("algn", "ctr")
    for tag in ("buChar", "buAutoNum", "buNone"):
        for e in pPr.findall(f"{{{A_NS}}}{tag}"):
            pPr.remove(e)
    etree.SubElement(pPr, f"{{{A_NS}}}buNone")


def _shrink_to_fit(text_frame):
    """Ask PowerPoint to shrink the text rather than run off the slide."""
    from lxml import etree

    bodyPr = text_frame._txBody.find(f"{{{A_NS}}}bodyPr")
    if bodyPr is None:
        return
    for tag in ("normAutofit", "spAutoFit", "noAutofit"):
        for e in bodyPr.findall(f"{{{A_NS}}}{tag}"):
            bodyPr.remove(e)
    etree.SubElement(bodyPr, f"{{{A_NS}}}normAutofit")


def markdown_to_pptx(markdown: str, out_path: str, title: str = "",
                     font_pt: int = 20) -> str:
    """Write Markdown to a .pptx, equations included as native Office math.

    A heading starts a slide and becomes its title.  Under it:

        > a blockquote is the SCRIPT -- what you say -- and goes into the
          speaker notes, not onto the slide
        everything else is what the audience SEES

    That is the whole layout rule: a deck built from a note follows the note's
    own headings rather than a second structure invented here, and script and
    slide live in one file so they can be compared.

    Equations become real PowerPoint equations, so they follow the theme font,
    scale with the text, and can be edited in PowerPoint's own editor.  A
    picture would be none of those things.
    """
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    body_layout = prs.slide_layouts[1]

    slide = None
    body = None
    script: list[str] = []

    def flush_script():
        """Put what has been said into this slide's notes."""
        nonlocal script
        if slide is not None and script:
            notes = slide.notes_slide.notes_text_frame
            notes.clear()
            first = True
            for para_text in script:
                para = notes.paragraphs[0] if first else notes.add_paragraph()
                first = False
                omml_paragraph(para, split_math(para_text), powerpoint=True)
        script = []

    def new_slide(heading):
        nonlocal slide, body
        flush_script()
        slide = prs.slides.add_slide(body_layout)
        slide.shapes.title.text = heading
        body = slide.placeholders[1].text_frame
        body.clear()
        body.word_wrap = True
        _shrink_to_fit(body)
        return body

    if title:
        cover = prs.slides.add_slide(title_layout)
        cover.shapes.title.text = title
        if len(cover.placeholders) > 1:
            cover.placeholders[1].text = ""

    first_para = True
    for block in md_blocks(markdown):
        if block.kind == MdBlock.Blank:
            continue

        if block.kind == MdBlock.Heading:
            new_slide(block.text.strip())
            first_para = True
            continue

        spoken = _script_text(block)
        if spoken is not None:
            script.append(spoken)
            continue

        if body is None:
            body = new_slide(title or "")
            first_para = True

        para = body.paragraphs[0] if first_para else body.add_paragraph()
        first_para = False

        if block.kind == MdBlock.Code:
            run = para.add_run()
            run.text = block.text
            run.font.name = "Consolas"
            run.font.size = Pt(14)
            continue

        if block.kind == MdBlock.Bullet:
            para.level = 1
        elif block.kind == MdBlock.Numbered:
            para.level = 1

        pieces = split_math(block.text)
        display_only = (len(pieces) == 1 and pieces[0].is_math
                        and pieces[0].display)
        if display_only or block.kind == MdBlock.Paragraph:
            _no_bullet(para, centre=display_only)
        para.font.size = Pt(font_pt)

        omml_paragraph(para, pieces, powerpoint=True)

    flush_script()
    prs.save(out_path)
    return out_path


def count_equations(path: str) -> int:
    """Count native Office equations in a .docx or .pptx.

    "<m:oMath" is a prefix of "<m:oMathPara", so the element itself is matched
    rather than the string.
    """
    n = 0
    with zipfile.ZipFile(path) as z:
        names: Sequence[str] = z.namelist()
        parts = [p for p in names
                 if p == "word/document.xml"
                 or re.match(r"ppt/slides/slide\d+\.xml$", p)]
        for part in parts:
            n += len(re.findall(r"<m:oMath[ >]", z.read(part).decode("utf-8")))
    return n


def copy_to_clipboard(latex: str, display: bool = False,
                      pictures: bool = True) -> list[str]:
    """Put one equation on the Windows clipboard for every target at once.

    The clipboard holds many formats and each application takes the one it
    understands best, so a single Copy serves every target with no mode switch.
    Which format each one needs was measured, not assumed:

      Rich Text Format   Word reads it as maths; PowerPoint reads it as text
      MathML             PowerPoint reads it as maths; Word refuses it alone
      CF_ENHMETAFILE     a vector picture, for anywhere without an equation
                         object -- Excel keeps the metafile beside the bitmap,
                         and it is how a vector reaches Google Slides via
                         Google Drawings
      PNG                the picture that always works; Google Slides rejects
                         SVG on upload and takes raster only
      CF_UNICODETEXT     the LaTeX itself, for Markdown, Jupyter, any editor

    The pictures are offered *after* the equation formats deliberately: Word and
    PowerPoint were re-checked with them present and still produce native
    equations, but an application that prefers a picture would silently
    downgrade, so `pictures=False` is there to take them away.

    Returns the format names actually placed.  Windows only.
    """
    import ctypes

    import win32clipboard as cb          # pywin32, Windows only

    from radia._equation import (MathMLOptions, RtfOptions, SvgStyle,
                                 tex_to_dib, tex_to_emf, tex_to_png)

    rtf_opt = RtfOptions()
    rtf_opt.display = display
    mml_opt = MathMLOptions()
    mml_opt.display = display

    payload = [
        ("Rich Text Format",
         tex_to_rtf(latex, rtf_opt).encode("latin-1", "replace")),
        ("MathML", tex_to_mathml(latex, mml_opt).encode("utf-8")),
    ]

    style = SvgStyle()
    emf_bytes = png_bytes = dib_bytes = b""
    if pictures:
        emf_bytes = tex_to_emf(latex, style)
        png_bytes = tex_to_png(latex, style, PASTE_DPI / 72.0)
        dib_bytes = tex_to_dib(latex, style, PASTE_DPI / 72.0)

    cb.OpenClipboard()
    try:
        cb.EmptyClipboard()
        placed = []
        for name, data in payload:
            cb.SetClipboardData(cb.RegisterClipboardFormat(name), data)
            placed.append(name)

        if emf_bytes:
            # CF_ENHMETAFILE wants a metafile handle, not bytes, and the system
            # owns it once it is on the clipboard -- it must not be deleted.
            gdi = ctypes.windll.gdi32
            handle = gdi.SetEnhMetaFileBits(len(emf_bytes), emf_bytes)
            if handle:
                cb.SetClipboardData(14, handle)      # CF_ENHMETAFILE
                placed.append("CF_ENHMETAFILE")
        if png_bytes:
            cb.SetClipboardData(cb.RegisterClipboardFormat("PNG"), png_bytes)
            placed.append("PNG")
        if dib_bytes:
            # CF_DIB is what an application takes when it pastes a picture, and
            # the only image format a browser finds here: Windows synthesises
            # it from a bitmap but not from a metafile.  Without it an equation
            # dropped into Google Slides -- which has no equation object and
            # takes raster only -- arrived as its LaTeX.
            cb.SetClipboardData(8, dib_bytes)          # CF_DIB
            placed.append("CF_DIB")

        cb.SetClipboardData(cb.CF_UNICODETEXT, latex)
        placed.append("CF_UNICODETEXT")
    finally:
        cb.CloseClipboard()
    return placed
